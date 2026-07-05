"""
Blue Robot Document Tools
=========================
Document management, indexing, and search.
"""

from __future__ import annotations

import base64
import hashlib
import json
import os
import subprocess
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional

import requests
from werkzeug.utils import secure_filename

# ================================================================================
# CONFIGURATION
# ================================================================================

UPLOAD_FOLDER = Path(os.environ.get("UPLOAD_FOLDER", "uploads"))
DOCUMENTS_FOLDER = Path(os.environ.get("DOCUMENTS_FOLDER", "documents"))
CAMERA_FOLDER = Path(os.environ.get("CAMERA_DIR", "camera_captures"))
MAX_FILE_SIZE = 16 * 1024 * 1024  # 16MB
ALLOWED_EXTENSIONS = {
    'bmp', 'csv', 'doc', 'docx', 'gif', 'html', 'jpeg', 'jpg',
    'json', 'md', 'pdf', 'png', 'pptx', 'rtf', 'tiff', 'txt',
    'webp', 'xlsx', 'xml'
}
DOCUMENT_INDEX_FILE = "document_index.json"
LM_STUDIO_RAG_URL = os.environ.get("LM_STUDIO_RAG_URL", "http://127.0.0.1:1234/v1/rag")


# ================================================================================
# INDEX MANAGEMENT
# ================================================================================

def load_document_index() -> Dict:
    """Load the document index from disk."""
    if os.path.exists(DOCUMENT_INDEX_FILE):
        try:
            with open(DOCUMENT_INDEX_FILE, 'r') as f:
                return json.load(f)
        except Exception:
            pass
    return {"documents": []}


def save_document_index(index: Dict):
    """Save the document index to disk."""
    with open(DOCUMENT_INDEX_FILE, 'w') as f:
        json.dump(index, f, indent=2)


def allowed_file(filename: str) -> bool:
    """Check if file extension is allowed."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def ensure_unique_path(directory: str, filename: str) -> str:
    """Ensure a unique file path by adding numbers if file already exists."""
    base_path = os.path.join(directory, secure_filename(filename))

    if not os.path.exists(base_path):
        return base_path

    name, ext = os.path.splitext(filename)
    counter = 1

    while True:
        new_filename = f"{name}_{counter}{ext}"
        new_path = os.path.join(directory, secure_filename(new_filename))

        if not os.path.exists(new_path):
            return new_path

        counter += 1

        if counter > 9999:
            import time
            timestamp = int(time.time())
            new_filename = f"{name}_{timestamp}{ext}"
            return os.path.join(directory, secure_filename(new_filename))


def get_file_hash(filepath: str) -> str:
    """Get MD5 hash of file for deduplication."""
    hash_md5 = hashlib.md5()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hash_md5.update(chunk)
    return hash_md5.hexdigest()


# ================================================================================
# IMAGE ENCODING
# ================================================================================

def encode_image_to_base64(filepath: str) -> Optional[Dict[str, Any]]:
    """Encode an image file to base64 for vision model viewing."""
    ext = filepath.rsplit('.', 1)[1].lower()
    if ext not in ['png', 'jpg', 'jpeg', 'tiff', 'bmp', 'gif', 'webp']:
        return None

    try:
        with open(filepath, 'rb') as f:
            image_data = base64.b64encode(f.read()).decode('utf-8')

        mime_types = {
            'jpg': 'image/jpeg',
            'jpeg': 'image/jpeg',
            'png': 'image/png',
            'gif': 'image/gif',
            'webp': 'image/webp',
            'bmp': 'image/bmp',
            'tiff': 'image/tiff'
        }

        mime_type = mime_types.get(ext, 'image/jpeg')

        return {
            "type": "image_url",
            "image_url": {
                "url": f"data:{mime_type};base64,{image_data}"
            }
        }
    except Exception as e:
        print(f"   [ERROR] Failed to encode image {filepath}: {e}")
        return None


# ================================================================================
# TEXT EXTRACTION
# ================================================================================

def extract_text_isolated(filepath: str, timeout: int = 600) -> str:
    """Extract text in a throwaway subprocess.

    pypdf killed the server twice on 2026-07-04 with a native access
    violation mid-extraction — an uncatchable crash, and on a different PDF
    each run. Batch indexing (reindex, folder watcher) calls this instead of
    extract_text_from_file so a dying parser costs one document, not the
    whole process.
    """
    worker = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                          "extract_worker.py")
    try:
        proc = subprocess.run(
            [sys.executable, worker, str(filepath)],
            capture_output=True, timeout=timeout)
    except subprocess.TimeoutExpired:
        return f"Error: text extraction timed out after {timeout}s"
    except OSError as e:
        return f"Error: could not start extraction worker: {e}"
    if proc.returncode != 0:
        tail = proc.stderr.decode("utf-8", "replace").strip().splitlines()
        detail = tail[-1] if tail else f"exit code {proc.returncode}"
        return f"Error: text extraction crashed ({detail})"
    return proc.stdout.decode("utf-8", "replace")


def extract_text_from_file(filepath: str) -> str:
    """Extract text from various file types."""
    ext = filepath.rsplit('.', 1)[1].lower() if '.' in filepath else ''

    # Files NUL-filled by an unclean shutdown make pypdf die with a native
    # access violation that no except clause can catch (it took the whole
    # server down on 2026-07-04), so refuse them before any parser runs.
    try:
        with open(filepath, 'rb') as f:
            head = f.read(4096)
        if head and not head.strip(b'\x00'):
            return "Error: file is NUL-corrupted (unclean shutdown?) and needs to be restored"
    except OSError as e:
        return f"Error reading file: {e}"

    if ext in ('txt', 'md'):
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    elif ext == 'pdf':
        try:
            try:
                from pypdf import PdfReader
            except ImportError:
                from PyPDF2 import PdfReader
            text = []
            with open(filepath, 'rb') as f:
                for page in PdfReader(f).pages:
                    text.append(page.extract_text() or '')
            return '\n'.join(text)
        except ImportError:
            return "Error: pypdf not installed. Install with: pip install pypdf"
        except Exception as e:
            return f"Error extracting PDF: {str(e)}"

    elif ext in ('doc', 'docx'):
        try:
            import docx
            doc = docx.Document(filepath)
            return '\n'.join(paragraph.text for paragraph in doc.paragraphs)
        except ImportError:
            return "Error: python-docx not installed. Install with: pip install python-docx"
        except Exception as e:
            return f"Error extracting Word doc: {str(e)}"

    elif ext in ('csv', 'tsv'):
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    elif ext in ('json', 'xml', 'html', 'htm'):
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            raw = f.read()
        if ext in ('html', 'htm'):
            try:
                from bs4 import BeautifulSoup
                return BeautifulSoup(raw, 'html.parser').get_text(separator='\n')
            except ImportError:
                return raw
        return raw

    elif ext == 'rtf':
        try:
            from striprtf.striprtf import rtf_to_text
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return rtf_to_text(f.read())
        except ImportError:
            import re
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                raw = f.read()
            return re.sub(r'\\[a-z]+\d* ?|[{}]', '', raw)

    elif ext == 'xlsx':
        try:
            from openpyxl import load_workbook
            wb = load_workbook(filepath, read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"# Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    cells = ['' if v is None else str(v) for v in row]
                    if any(cells):
                        lines.append('\t'.join(cells))
            return '\n'.join(lines)
        except ImportError:
            return "Error: openpyxl not installed. Install with: pip install openpyxl"
        except Exception as e:
            return f"Error extracting xlsx: {str(e)}"

    elif ext == 'pptx':
        try:
            from pptx import Presentation
            lines = []
            for i, slide in enumerate(Presentation(filepath).slides, 1):
                lines.append(f"# Slide {i}")
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        lines.append(shape.text)
            return '\n'.join(lines)
        except ImportError:
            return "Error: python-pptx not installed. Install with: pip install python-pptx"
        except Exception as e:
            return f"Error extracting pptx: {str(e)}"

    elif ext in ('png', 'jpg', 'jpeg', 'tiff', 'bmp', 'gif', 'webp'):
        try:
            from PIL import Image
            image = Image.open(filepath)
            width, height = image.size
            mode = image.mode
            file_size = os.path.getsize(filepath)

            return (f"[IMAGE FILE - Vision model will view directly]\n"
                    f"Dimensions: {width}x{height}\n"
                    f"Color mode: {mode}\n"
                    f"File size: {file_size} bytes\n"
                    f"Format: {ext.upper()}")
        except ImportError:
            return f"[IMAGE FILE] {os.path.basename(filepath)} - PIL required to read metadata"
        except Exception as e:
            return f"[IMAGE FILE] {os.path.basename(filepath)} - Error reading: {str(e)}"

    return f"Unsupported file type: .{ext}"


# ================================================================================
# RAG FUNCTIONS
# ================================================================================

def add_document_to_rag(filepath: str, filename: str) -> bool:
    """Add a document to LM Studio's RAG system."""
    try:
        text_content = extract_text_from_file(filepath)

        if text_content.startswith("Error"):
            print(f"   [ERROR] {text_content}")
            return False

        payload = {
            "file_path": str(filepath),
            "content": text_content,
            "metadata": {
                "filename": filename,
                "source": "blue_middleware"
            }
        }

        response = requests.post(
            f"{LM_STUDIO_RAG_URL}/documents",
            json=payload,
            timeout=30
        )

        if response.status_code in [200, 201]:
            print(f"   [OK] Document indexed in RAG system")
            return True
        else:
            print(f"   [WARN] RAG indexing returned: {response.status_code}")
            return True

    except requests.exceptions.RequestException as e:
        print(f"   [WARN] RAG system not available: {e}")
        return True
    except Exception as e:
        print(f"   [ERROR] Error adding to RAG: {e}")
        return False


def search_documents_rag(query: str, max_results: int = 3) -> str:
    """Search documents using LM Studio's RAG system."""
    print(f"   [FIND] Searching documents for: '{query}'")

    try:
        payload = {
            "query": query,
            "max_results": max_results
        }

        response = requests.post(
            f"{LM_STUDIO_RAG_URL}/search",
            json=payload,
            timeout=10
        )

        if response.status_code == 200:
            try:
                results = response.json()

                if isinstance(results, dict):
                    if 'results' in results:
                        results = results['results']
                    elif 'data' in results:
                        results = results['data']
                    elif 'documents' in results:
                        results = results['documents']

                if isinstance(results, list) and len(results) > 0:
                    formatted_results = []

                    for i, result in enumerate(results[:max_results], 1):
                        if isinstance(result, dict):
                            filename = result.get('metadata', {}).get('filename', 'Unknown')
                            content = result.get('content', result.get('text', ''))
                            score = result.get('score', 0)

                            content_preview = str(content)[:500] if content else "No content"

                            formatted_results.append(
                                f"[{i}] From: {filename} (relevance: {score:.2f})\n{content_preview}"
                            )
                        else:
                            formatted_results.append(f"[{i}] {str(result)[:500]}")

                    if formatted_results:
                        return "Here's what I found in your documents:\n\n" + "\n\n".join(formatted_results)

                return search_documents_local(query, max_results)

            except Exception as e:
                print(f"   [WARN] Error parsing RAG response: {e}")
                return search_documents_local(query, max_results)
        else:
            return search_documents_local(query, max_results)

    except requests.exceptions.RequestException as e:
        print(f"   [WARN] RAG connection error: {e}")
        return search_documents_local(query, max_results)
    except Exception as e:
        print(f"   [ERROR] Unexpected error: {e}")
        return search_documents_local(query, max_results)


def search_documents_local(query: str, max_results: int = 3) -> str:
    """Fallback: Simple local search through document index."""
    print(f"   [FOLDER] Using local document search...")
    index = load_document_index()
    documents = index.get("documents", [])

    if not documents:
        return (
            "I don't have any documents to search through yet! "
            "You can upload documents at http://127.0.0.1:5000/documents - "
            "I can read PDFs, Word docs, text files, markdown, and images."
        )

    query_lower = query.lower()
    matches = []

    if "all" in query_lower and ("document" in query_lower or "summarize" in query_lower):
        doc_list = []
        for i, doc in enumerate(documents[:10], 1):
            preview = doc.get('text_preview', 'No preview available')[:200]
            doc_list.append(f"{i}. {doc['filename']}\n   Preview: {preview}...")

        summary = "\n\n".join(doc_list)
        if len(documents) > 10:
            summary += f"\n\n...and {len(documents) - 10} more documents."

        return f"I have {len(documents)} document(s) uploaded:\n\n{summary}"

    for doc in documents:
        relevance = 0
        filename_lower = doc['filename'].lower()

        if query_lower in filename_lower:
            relevance += 3

        text_content = doc.get('text_preview', '').lower()
        if query_lower in text_content:
            relevance += 5

        for word in query_lower.split():
            if len(word) > 3:
                if word in filename_lower:
                    relevance += 1
                if word in text_content:
                    relevance += 2

        if relevance > 0:
            matches.append((doc, relevance))

    if not matches:
        return (
            f"I couldn't find any documents matching '{query}'. "
            f"I have {len(documents)} document(s) uploaded. "
            "Try using different keywords, or ask me to list all documents."
        )

    matches.sort(key=lambda x: x[1], reverse=True)

    text_results = []
    for doc, score in matches[:max_results]:
        filepath = doc.get('filepath', '')
        filename = doc['filename']

        if os.path.exists(filepath):
            try:
                full_text = extract_text_from_file(filepath)
                content = full_text[:2000] if len(full_text) > 2000 else full_text
                text_results.append(f"[FILE] **{filename}** (relevance: {score})\n\n{content}\n")
            except Exception:
                preview = doc.get('text_preview', 'No preview available')[:500]
                text_results.append(f"[FILE] **{filename}** (relevance: {score})\n\n{preview}...\n")
        else:
            preview = doc.get('text_preview', 'No preview available')[:500]
            text_results.append(f"[FILE] **{filename}** (relevance: {score})\n\n{preview}...\n")

    return "Here's what I found in your documents:\n\n" + "\n---\n\n".join(text_results)


# ================================================================================
# DOCUMENT CREATION
# ================================================================================

def create_document_file(filename: str, content: str, file_type: str = "txt") -> str:
    """Create a new document file."""
    try:
        DOCUMENTS_FOLDER.mkdir(parents=True, exist_ok=True)

        if not filename.endswith(f'.{file_type}'):
            filename = f"{filename}.{file_type}"

        filepath = ensure_unique_path(str(DOCUMENTS_FOLDER), filename)
        actual_filename = os.path.basename(filepath)

        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)

        index = load_document_index()
        index['documents'].append({
            'filename': actual_filename,
            'filepath': filepath,
            'text_preview': content[:500],
            'hash': get_file_hash(filepath)
        })
        save_document_index(index)

        add_document_to_rag(filepath, actual_filename)

        return json.dumps({
            "success": True,
            "filename": actual_filename,
            "filepath": filepath,
            "message": f"Created document: {actual_filename}",
            "size": len(content)
        })

    except Exception as e:
        return json.dumps({
            "success": False,
            "error": f"Failed to create document: {str(e)}"
        })


__all__ = [
    'UPLOAD_FOLDER', 'DOCUMENTS_FOLDER', 'MAX_FILE_SIZE',
    'ALLOWED_EXTENSIONS', 'DOCUMENT_INDEX_FILE',
    'load_document_index', 'save_document_index', 'allowed_file',
    'ensure_unique_path', 'get_file_hash',
    'encode_image_to_base64', 'extract_text_from_file',
    'add_document_to_rag', 'search_documents_rag', 'search_documents_local',
    'create_document_file',
]
