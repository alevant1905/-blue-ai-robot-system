"""
Blue AI Robot System — ENHANCED VERSION v8
==================================================

v8 ENHANCEMENTS (November 2025):
- Compound request parsing ("play jazz and set romantic lights")
- Follow-up correction detection ("no, make it blue" / "louder")
- Smart response caching for repeated queries
- Query complexity estimation for adaptive processing
- Entity extraction (emails, times, numbers, URLs)
- Better action type identification
- Improved context-aware corrections

v7 ENHANCEMENTS:
- Fuzzy matching for artist names (handles typos)
- ConversationState class for better context tracking
- Enhanced error recovery with helpful suggestions
- Utility functions: truncate_text, get_time_ago, safe_json_parse
- Better tool execution wrapper with timing and state tracking

v6 ENHANCEMENTS:
- 200+ music artists, 60+ genres with improved matching
- 50+ light moods/scenes including seasonal, holiday, activity-based
- Enhanced fact extraction: vehicles, medical, languages, skills
- Robust error handling with auto-retry

v5: Cleanup and consolidation
v4: Smart email parameter extraction
v3: Enhanced fact extraction, topic decay
v2: Greeting detection, timer/reminder detection
v1: Camera detection, email/search disambiguation

FILE STRUCTURE:
1. Imports & Configuration (line ~70)
2. Utility Functions (line ~170) - Enhanced in v8
3. Database & Memory (line ~560)  
4. System Prompt (line ~950)
5. Tool Definitions (line ~960)
6. Tool Selector (line ~1050)
7. LLM Client (line ~2650)
8. Music Functions (line ~3050)
9. Vision & Camera (line ~3350)
10. Document Functions (line ~4350)
11. Light Functions (line ~5150)
12. Web Search (line ~5350)
13. Gmail Functions (line ~5650)
14. Main Tool Executor (line ~6000)
15. Legacy Detect Functions (line ~6450)
16. Process With Tools (line ~7200)
17. Flask Routes (line ~8450)
18. Gmail Upgrades (line ~9410)
19. Voice Email Interface (line ~9530)
"""

# ================================================================================
# IMPORTS
# ================================================================================
from __future__ import annotations

# Standard library
import json
import logging
import os
import re
import sqlite3
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from functools import lru_cache
import pickle
import webbrowser
import random
import hashlib
import base64
import mimetypes
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders

from blue_identity import (
    canonical_family_grounding_lines,
    canonical_household_reply,
    canonical_identity_reply,
    canonical_self_state_reply,
    contextual_identity_request_kind,
    identity_conversation_context,
    identity_grounding_note,
    identity_repeats_recent_reply,
    identity_reply_topics,
    identity_request_kind,
    identity_response_problem,
    is_family_overview_request,
    is_jspace_presence_request,
)

# Native-crash visibility: a fault in a C extension (access violation, heap
# corruption) kills the whole process with NO Python traceback — the console
# just drops back to the prompt. faulthandler prints every thread's Python
# stack to stderr on the way down, so the guilty code path is identifiable.
# (The 2026-07-03 Gmail thread-safety crash took three silent deaths to find;
# with this enabled it would have named the thread on the first one.)
import faulthandler
try:
    faulthandler.enable()
except Exception:
    pass

# When run as a script (python bluetools.py), register this module under its
# import name too, so extracted blue.server.* modules doing `import bluetools`
# get THIS instance instead of executing the whole file a second time (which
# would double every daemon thread and re-init the head hardware).
import sys
if __name__ == "__main__":
    sys.modules.setdefault("bluetools", sys.modules[__name__])

# Third-party
import requests
from flask import Flask, Response, jsonify, redirect, render_template_string, request, send_from_directory, session, url_for
from werkzeug.utils import secure_filename
from googleapiclient.discovery import build
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request

# Visual Memory System (if available)
try:
    from blue_visual_memory import get_visual_memory, VisualMemory, VISUAL_REF_DIR
    VISUAL_MEMORY_AVAILABLE = True
    print("[OK] Visual memory system loaded - Blue can now recognize people and places!")
except ImportError:
    VISUAL_MEMORY_AVAILABLE = False
    print("[WARN] Visual memory system not available")

# Face Recognition Engine (OpenCV SFace; degrades gracefully if unavailable)
try:
    from blue.tools import face_engine as FACE_ENGINE
    FACE_RECOGNITION_AVAILABLE = True
    print("[OK] Face engine module loaded - Blue can recognize faces by embedding!")
except Exception as _e:
    FACE_ENGINE = None
    FACE_RECOGNITION_AVAILABLE = False
    print(f"[WARN] Face engine not available: {_e}")

# Enhanced Visual Understanding (if available)
try:
    from blue_visual_understanding import get_visual_understanding, get_enhanced_vision_context
    VISUAL_UNDERSTANDING_AVAILABLE = True
    print("[OK] Enhanced visual understanding loaded - Blue can understand activities and emotions!")
except ImportError:
    VISUAL_UNDERSTANDING_AVAILABLE = False
    print("[WARN] Enhanced visual understanding not available")

# Proactive Assistance (if available)
try:
    from blue_proactive_assistance import get_proactive_assistance, ProactiveSuggestion
    PROACTIVE_ASSISTANCE_AVAILABLE = True
    print("[OK] Proactive assistance loaded - Blue can offer helpful suggestions!")
except ImportError as e:
    PROACTIVE_ASSISTANCE_AVAILABLE = False
    print(f"[INFO] Optional proactive assistance not installed: {e}")

# Proactive heartbeat queue (reminders → next inbound response)
try:
    import blue_proactive
    PROACTIVE_QUEUE_AVAILABLE = True
    print("[OK] Proactive queue loaded - Blue can surface reminders on next turn")
except ImportError as e:
    PROACTIVE_QUEUE_AVAILABLE = False
    print(f"[WARN] Proactive queue not available: {e}")

# Academic Research Assistant (if available)
try:
    from blue_academic_assistant import (
        get_academic_assistant, analyze_with_chat, prepare_lecture,
        generate_discussion_questions, simulate_student_q_and_a,
        synthesize_research, circumference_content
    )
    ACADEMIC_ASSISTANT_AVAILABLE = True
    print("[OK] Academic assistant loaded - Teaching and research tools ready!")
except ImportError as e:
    ACADEMIC_ASSISTANT_AVAILABLE = False
    print(f"[INFO] Optional academic assistant not installed: {e}")

# Scholarly research — academic journal search wired to the Laurier library
# (Omni discovery + OpenAlex/Crossref/Unpaywall; full text via libproxy)
try:
    from blue.tools.scholar import (
        execute_scholar_search, execute_get_paper, execute_read_paper,
        library_account_status,
    )
    SCHOLAR_AVAILABLE = True
    print(f"[OK] Scholar tools loaded - Laurier library search ready! "
          f"(library account: {library_account_status()})")
except ImportError as e:
    SCHOLAR_AVAILABLE = False
    print(f"[WARN] Scholar tools not available: {e}")

# Multi-Person Context Awareness (if available)
try:
    from blue_context_awareness import (
        get_context_awareness, adapt_for_audience, get_audience_context,
        generate_contextual_greeting
    )
    CONTEXT_AWARENESS_AVAILABLE = True
    print("[OK] Context awareness loaded - Blue adapts to his audience!")
except ImportError as e:
    CONTEXT_AWARENESS_AVAILABLE = False
    print(f"[INFO] Optional context awareness not installed: {e}")

# ================================================================================
# LOGGING (single source)
# ================================================================================


def setup_logger(name: str = "blue", level: str = "INFO") -> logging.Logger:
    logger = logging.getLogger(name)
    if not logger.handlers:
        handler = logging.StreamHandler()
        handler.setFormatter(logging.Formatter("[%(levelname)s] %(asctime)s - %(message)s"))
        logger.addHandler(handler)
    logger.setLevel(getattr(logging, level.upper(), logging.INFO))
    return logger

log = setup_logger(level=os.environ.get("LOG_LEVEL", "INFO"))

# ================================================================================
# UTILITY FUNCTIONS - v8 ENHANCED
# ================================================================================

def parse_compound_request(message: str) -> List[Dict[str, Any]]:
    """
    Parse compound requests into individual actions.
    v8 ENHANCEMENT: Handle "play jazz and turn on relaxing lights" as two actions.
    
    Returns list of {action, query, priority} dicts.
    """
    msg_lower = message.lower().strip()
    actions = []
    
    # Compound connectors
    connectors = [' and ', ' then ', ' also ', ' plus ', ', and ', ' & ']
    
    # Check if this is a compound request
    has_connector = any(conn in msg_lower for conn in connectors)
    
    if not has_connector:
        return []  # Not a compound request
    
    # Split on connectors
    parts = [msg_lower]
    for conn in connectors:
        new_parts = []
        for part in parts:
            new_parts.extend(part.split(conn))
        parts = new_parts
    
    # Clean and analyze each part
    for i, part in enumerate(parts):
        part = part.strip()
        if len(part) < 3:
            continue
            
        action_type = _identify_action_type(part)
        if action_type:
            actions.append({
                'action': action_type,
                'query': part,
                'priority': i,
                'original': part
            })
    
    return actions if len(actions) > 1 else []


def _identify_action_type(text: str) -> Optional[str]:
    """Identify the action type from a text fragment."""
    text = text.lower()
    
    # Music
    if any(w in text for w in ['play', 'put on', 'listen to', 'queue']):
        if any(w in text for w in ['music', 'song', 'jazz', 'rock', 'pop', 'by ']):
            return 'play_music'
    
    # Lights
    if any(w in text for w in ['light', 'lamp', 'bright', 'dim']):
        if any(w in text for w in ['turn', 'set', 'make', 'switch']):
            return 'control_lights'
        if any(w in text for w in ['mood', 'scene', 'party', 'relax', 'romantic']):
            return 'control_lights'
    
    # Music control
    if any(w in text for w in ['pause', 'stop', 'skip', 'next', 'volume', 'louder', 'quieter']):
        return 'control_music'
    
    # Weather
    if 'weather' in text:
        return 'get_weather'
    
    # Email
    if 'email' in text or 'inbox' in text:
        if any(w in text for w in ['send', 'write', 'compose']):
            return 'send_gmail'
        if any(w in text for w in ['check', 'read', 'show']):
            return 'read_gmail'
        if any(w in text for w in ['reply', 'respond']):
            return 'reply_gmail'
    
    # Camera
    if any(w in text for w in ['see', 'look', 'photo', 'camera', 'picture']):
        return 'capture_camera'
    
    # Timer
    if any(w in text for w in ['timer', 'remind', 'alarm']):
        return 'set_timer'
    
    return None


def detect_follow_up_correction(message: str, context: Dict) -> Optional[Dict[str, Any]]:
    """
    Detect if user is correcting or refining a previous request.
    v8 ENHANCEMENT: Handle "no, the blue one" or "actually make it louder"
    
    Returns correction info or None.
    """
    msg_lower = message.lower().strip()
    
    # Correction indicators
    correction_starters = [
        'no ', 'not ', 'actually ', 'i meant ', 'i mean ', 'sorry ', 
        'wait ', 'change ', 'make it ', 'instead ', 'rather ', 'wrong '
    ]
    
    is_correction = any(msg_lower.startswith(s) for s in correction_starters)
    
    # Also check for short corrections
    if len(msg_lower.split()) <= 4:
        short_corrections = ['the other', 'different', 'louder', 'quieter', 'brighter', 
                           'dimmer', 'faster', 'slower', 'more', 'less']
        if any(c in msg_lower for c in short_corrections):
            is_correction = True
    
    if not is_correction:
        return None
    
    # Determine what's being corrected
    last_tool = context.get('last_tool_used')
    last_result = context.get('last_tool_result', '')
    
    correction = {
        'is_correction': True,
        'original_tool': last_tool,
        'correction_type': 'unknown',
        'new_value': None
    }
    
    # Light corrections
    if last_tool == 'control_lights' or any(w in msg_lower for w in ['light', 'bright', 'dim', 'color']):
        correction['correction_type'] = 'lights'
        # Extract new color/brightness
        colors = ['red', 'blue', 'green', 'yellow', 'purple', 'orange', 'white', 'pink', 'warm', 'cool']
        for color in colors:
            if color in msg_lower:
                correction['new_value'] = color
                break
        if 'brighter' in msg_lower:
            correction['new_value'] = 'brighter'
        elif 'dimmer' in msg_lower or 'darker' in msg_lower:
            correction['new_value'] = 'dimmer'
    
    # Music corrections
    elif last_tool in ['play_music', 'control_music'] or any(w in msg_lower for w in ['music', 'song', 'volume']):
        correction['correction_type'] = 'music'
        if 'louder' in msg_lower:
            correction['new_value'] = 'volume_up'
        elif 'quieter' in msg_lower or 'softer' in msg_lower:
            correction['new_value'] = 'volume_down'
        elif 'different' in msg_lower or 'other' in msg_lower:
            correction['new_value'] = 'skip'
    
    return correction


def detect_camera_capture_intent(message: str) -> bool:
    """Detect if user wants to capture a camera image (what do you see?)."""
    msg_lower = message.lower().replace("’", "'")

    # Statements ABOUT the scene and negated instructions are not capture
    # requests. Live 2026-07-12: "That is what's in front of you right now"
    # (confirming Blue's own description) and "There's no need to describe
    # what's in front of you" each re-fired the camera and got the room
    # described right back.
    if re.search(r"\b(?:no need|don'?t|do not|stop|never|quit|enough|without)\b"
                 r"[^.!?]{0,50}\b(?:describ\w*|look\w*|captur\w*|photo\w*|"
                 r"picture\w*|camera|see|view)\b", msg_lower):
        return False
    if re.search(r"\b(?:that(?:'s| is)|this is|it(?:'s| is)|here(?:'s| is))\s+"
                 r"(?:exactly\s+)?what(?:'s| is)?\b", msg_lower):
        return False

    # Primary camera capture triggers
    camera_triggers = [
        'what do you see',
        'what can you see',
        'what are you seeing',
        'what you see',
        "what's in front of you",
        'what is in front of you',
        'look in front of you',
        'see in front of you',
        'take a photo',
        'take a picture',
        'capture image',
        'capture photo',
        'show me what you see',
        'look around',
        'what are you looking at',
        'describe what you see',
        "what's happening right now",
        'what is happening right now',
        'show me your view',
        'use the camera',
        'use your camera',
        'camera photo',
        'camera picture',
        'see right now',
        'looking at right now',
        'do you see anything',
        'can you see anything',
        # View-control phrasings: aim the head / zoom before the shot
        'zoom in', 'zoom into', 'look closer', 'closer look',
        "what's on your left", "what's on your right",
        "what's to your left", "what's to your right",
        'what is on your left', 'what is on your right',
        'what is to your left', 'what is to your right',
        'look up and', 'look down and', 'look left and', 'look right and',
        'look to your left and', 'look to your right and',
    ]

    # Check for any trigger phrases
    return any(trigger in msg_lower for trigger in camera_triggers)


def smart_cache_key(query: str, tool: str = "") -> str:
    """Generate a smart cache key for query deduplication."""
    import hashlib
    # Normalize the query
    normalized = query.lower().strip()
    # Remove filler words
    fillers = ['please', 'can you', 'could you', 'would you', 'hey', 'blue', 'hi', 'hello']
    for filler in fillers:
        normalized = normalized.replace(filler, '')
    normalized = ' '.join(normalized.split())  # Normalize whitespace
    
    key_input = f"{tool}:{normalized}"
    return hashlib.md5(key_input.encode()).hexdigest()[:16]


# Response cache for repeated queries (5 minute TTL)
_response_cache: Dict[str, Tuple[float, str]] = {}
_RESPONSE_CACHE_TTL = 300  # 5 minutes


def get_cached_response(cache_key: str) -> Optional[str]:
    """Get cached response if still valid."""
    import time
    if cache_key in _response_cache:
        timestamp, response = _response_cache[cache_key]
        if time.time() - timestamp < _RESPONSE_CACHE_TTL:
            return response
        else:
            del _response_cache[cache_key]
    return None


def cache_response(cache_key: str, response: str):
    """Cache a response."""
    import time
    _response_cache[cache_key] = (time.time(), response)
    # Prune old entries
    if len(_response_cache) > 100:
        cutoff = time.time() - _RESPONSE_CACHE_TTL
        keys_to_delete = [k for k, (t, _) in _response_cache.items() if t < cutoff]
        for k in keys_to_delete:
            del _response_cache[k]


def estimate_query_complexity(message: str) -> str:
    """
    Estimate query complexity to adjust processing.
    Returns: 'simple', 'medium', 'complex'
    """
    msg_lower = message.lower()
    word_count = len(message.split())
    
    # Simple: greetings, single commands
    if word_count <= 5:
        return 'simple'
    
    # Complex: multiple actions, detailed requests
    compound_signals = [' and ', ' then ', ' also ', ' after ', ' before ']
    if any(s in msg_lower for s in compound_signals):
        return 'complex'
    
    # Complex: questions requiring research
    research_signals = ['explain', 'compare', 'difference between', 'how does', 'why does', 'analysis']
    if any(s in msg_lower for s in research_signals):
        return 'complex'
    
    # Medium: typical requests
    if word_count <= 15:
        return 'medium'
    
    return 'complex'


def extract_entities(message: str) -> Dict[str, List[str]]:
    """
    Extract named entities from message.
    v8 ENHANCEMENT: Better entity extraction for personalization.
    """
    entities = {
        'people': [],
        'places': [],
        'times': [],
        'numbers': [],
        'emails': [],
        'urls': []
    }
    
    # Email addresses
    email_pattern = r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}'
    entities['emails'] = re.findall(email_pattern, message)
    
    # URLs
    url_pattern = r'https?://[^\s]+'
    entities['urls'] = re.findall(url_pattern, message)
    
    # Times
    time_patterns = [
        r'\d{1,2}:\d{2}(?:\s*[ap]m)?',
        r'\d{1,2}\s*(?:am|pm)',
        r'(?:noon|midnight|morning|afternoon|evening|night)',
        r'(?:today|tomorrow|yesterday)',
        r'(?:monday|tuesday|wednesday|thursday|friday|saturday|sunday)',
    ]
    for pattern in time_patterns:
        matches = re.findall(pattern, message.lower())
        entities['times'].extend(matches)
    
    # Numbers
    number_pattern = r'\b\d+(?:\.\d+)?\b'
    entities['numbers'] = re.findall(number_pattern, message)
    
    return entities


def fuzzy_match(query: str, targets: List[str], threshold: float = 0.75) -> Optional[str]:
    """
    Find the best fuzzy match for a query in a list of targets.
    Uses simple similarity ratio - no external dependencies.
    
    Args:
        query: The search string
        targets: List of possible matches
        threshold: Minimum similarity (0.0 to 1.0)
    
    Returns:
        Best matching target or None if no good match
    """
    if not query or not targets:
        return None
    
    query_lower = query.lower().strip()
    
    # Exact match first
    for target in targets:
        if query_lower == target.lower():
            return target
    
    # Substring match
    for target in targets:
        if query_lower in target.lower() or target.lower() in query_lower:
            return target
    
    # Similarity matching
    best_match = None
    best_score = 0.0
    
    for target in targets:
        score = _string_similarity(query_lower, target.lower())
        if score > best_score and score >= threshold:
            best_score = score
            best_match = target
    
    return best_match


def _string_similarity(s1: str, s2: str) -> float:
    """Calculate string similarity using character-based comparison."""
    if not s1 or not s2:
        return 0.0
    if s1 == s2:
        return 1.0
    
    # Simple Jaccard-like similarity on character bigrams
    def get_bigrams(s):
        return set(s[i:i+2] for i in range(len(s) - 1)) if len(s) > 1 else {s}
    
    b1 = get_bigrams(s1)
    b2 = get_bigrams(s2)
    
    intersection = len(b1 & b2)
    union = len(b1 | b2)
    
    return intersection / union if union > 0 else 0.0


def normalize_artist_name(name: str) -> str:
    """Normalize artist name for matching."""
    if not name:
        return ""
    
    # Common replacements
    replacements = {
        '&': 'and',
        '+': 'and',
        ' - ': ' ',
        "'s": 's',
        '"': '',
    }
    
    result = name.lower().strip()
    for old, new in replacements.items():
        result = result.replace(old, new)
    
    # Remove "the " prefix for matching
    if result.startswith('the '):
        result = result[4:]
    
    return result


def safe_json_parse(text: str, default: Any = None) -> Any:
    """Safely parse JSON with fallback."""
    if not text:
        return default
    try:
        return json.loads(text)
    except (json.JSONDecodeError, TypeError):
        return default


def truncate_text(text: str, max_length: int = 100, suffix: str = "...") -> str:
    """Truncate text to max length with suffix."""
    if not text or len(text) <= max_length:
        return text or ""
    return text[:max_length - len(suffix)] + suffix


def extract_quoted_text(text: str) -> List[str]:
    """Extract all quoted strings from text."""
    import re
    # Match double and single quotes
    patterns = [
        r'"([^"]+)"',
        r"'([^']+)'",
    ]
    results = []
    for pattern in patterns:
        results.extend(re.findall(pattern, text))
    return results


def get_time_ago(timestamp: float) -> str:
    """Convert timestamp to human-readable 'time ago' string."""
    import time
    diff = time.time() - timestamp
    
    if diff < 60:
        return "just now"
    elif diff < 3600:
        mins = int(diff / 60)
        return f"{mins} minute{'s' if mins != 1 else ''} ago"
    elif diff < 86400:
        hours = int(diff / 3600)
        return f"{hours} hour{'s' if hours != 1 else ''} ago"
    elif diff < 604800:
        days = int(diff / 86400)
        return f"{days} day{'s' if days != 1 else ''} ago"
    else:
        weeks = int(diff / 604800)
        return f"{weeks} week{'s' if weeks != 1 else ''} ago"


class ConversationState:
    """
    Track conversation state for better context awareness.
    v8 ENHANCED: More tracking, pattern learning, and suggestions.
    """
    
    def __init__(self):
        self.last_tool_used: Optional[str] = None
        self.last_tool_result: Optional[str] = None
        self.last_tool_args: Optional[Dict] = None
        self.pending_confirmation: Optional[str] = None
        self.topic_stack: List[str] = []
        self.user_corrections: List[Dict] = []
        self.successful_patterns: Dict[str, int] = {}
        self.failed_patterns: Dict[str, int] = {}
        self.tool_sequence: List[str] = []  # Track tool order
        self.last_query: str = ""
        self.query_count: int = 0
        self.session_start: float = __import__('time').time()
        self.recent_phrases: List[str] = []  # Short text snippets used recently
    
    def record_tool_use(self, tool_name: str, success: bool, pattern: str = "", args: Dict = None):
        """Record tool usage for learning."""
        self.last_tool_used = tool_name
        self.last_tool_args = args or {}
        self.tool_sequence.append(tool_name)
        if len(self.tool_sequence) > 20:
            self.tool_sequence.pop(0)
        
        if pattern:
            if success:
                self.successful_patterns[pattern] = self.successful_patterns.get(pattern, 0) + 1
            else:
                self.failed_patterns[pattern] = self.failed_patterns.get(pattern, 0) + 1
    
    def push_topic(self, topic: str):
        """Add a topic to the stack."""
        if topic and (not self.topic_stack or self.topic_stack[-1] != topic):
            self.topic_stack.append(topic)
            if len(self.topic_stack) > 10:
                self.topic_stack.pop(0)
    
    def get_current_topic(self) -> Optional[str]:
        """Get the most recent topic."""
        return self.topic_stack[-1] if self.topic_stack else None
    
    def record_correction(self, original: str, corrected: str):
        """Record user corrections for learning."""
        self.user_corrections.append({
            'original': original,
            'corrected': corrected,
            'timestamp': __import__('time').time()
        })
        if len(self.user_corrections) > 50:
            self.user_corrections.pop(0)
    
    def record_query(self, query: str):
        """Record a new query."""
        self.last_query = query
        self.query_count += 1
    
    def get_common_tool_pairs(self) -> List[Tuple[str, str]]:
        """Get commonly used tool pairs (for compound request optimization)."""
        pairs = {}
        for i in range(len(self.tool_sequence) - 1):
            pair = (self.tool_sequence[i], self.tool_sequence[i+1])
            pairs[pair] = pairs.get(pair, 0) + 1
        
        # Return top 3 most common pairs
        sorted_pairs = sorted(pairs.items(), key=lambda x: x[1], reverse=True)
        return [pair for pair, count in sorted_pairs[:3] if count > 1]
    
    def get_session_stats(self) -> Dict[str, Any]:
        """Get session statistics."""
        import time
        return {
            'duration_minutes': (time.time() - self.session_start) / 60,
            'query_count': self.query_count,
            'tools_used': len(set(self.tool_sequence)),
            'most_used_tools': self._get_tool_frequency(),
            'corrections_made': len(self.user_corrections)
        }
    
    def _get_tool_frequency(self) -> Dict[str, int]:
        """Get tool usage frequency."""
        freq = {}
        for tool in self.tool_sequence:
            freq[tool] = freq.get(tool, 0) + 1
        return dict(sorted(freq.items(), key=lambda x: x[1], reverse=True)[:5])
    
    def suggest_next_action(self) -> Optional[str]:
        """Suggest next action based on patterns."""
        if not self.last_tool_used:
            return None
        
        # Check common pairs
        pairs = self.get_common_tool_pairs()
        for first, second in pairs:
            if first == self.last_tool_used:
                return f"You often use {second} after {first}. Would you like me to do that?"
        
        return None


# Global conversation state
_conversation_state = ConversationState()


def get_conversation_state() -> ConversationState:
    """Get the global conversation state."""
    return _conversation_state


def validate_response_quality(response: str, query: str) -> Dict[str, Any]:
    """
    Validate response quality and provide improvement suggestions.
    v7 ENHANCEMENT: Better response quality checking.
    """
    issues = []
    score = 100
    
    # Check for empty or very short response
    if not response or len(response.strip()) < 10:
        issues.append("Response is too short")
        score -= 30
    
    # Check for error indicators
    error_phrases = ['error', 'failed', 'could not', 'unable to', 'something went wrong']
    if any(phrase in response.lower() for phrase in error_phrases):
        issues.append("Response contains error indicators")
        score -= 20
    
    # Check for hallucination indicators (claiming to have searched without using tools)
    hallucination_phrases = ['i searched', 'according to my search', 'i found that', 'my research shows']
    if any(phrase in response.lower() for phrase in hallucination_phrases):
        if 'web_search' not in str(get_conversation_state().last_tool_used):
            issues.append("Possible hallucination - claims search without tool use")
            score -= 25
    
    # Check for repetition - Enhanced detection
    sentences = response.split('.')
    unique_sentences = set(s.strip().lower() for s in sentences if s.strip())
    if len(sentences) > 3 and len(unique_sentences) < len(sentences) * 0.85:
        issues.append("Response contains repetitive content")
        score -= 20

    # Check for repeated phrases (3+ words)
    words = response.lower().split()
    phrases_seen = set()
    repeated_phrases = []
    for i in range(len(words) - 2):
        phrase = ' '.join(words[i:i+3])
        if phrase in phrases_seen:
            repeated_phrases.append(phrase)
        phrases_seen.add(phrase)

    if repeated_phrases:
        issues.append(f"Repeated phrases detected: {len(repeated_phrases)} instances")
        score -= min(30, len(repeated_phrases) * 5)

    # Check for repetitive sentence starters
    sentence_starts = [s.strip().lower()[:15] for s in sentences if s.strip() and len(s.strip()) > 15]
    if len(sentence_starts) != len(set(sentence_starts)):
        issues.append("Multiple sentences start the same way")
        score -= 15
    
    # Check if response addresses the query
    query_words = set(query.lower().split())
    response_words = set(response.lower().split())
    overlap = len(query_words & response_words) / len(query_words) if query_words else 0
    if overlap < 0.2 and len(query_words) > 2:
        issues.append("Response may not address the query directly")
        score -= 10
    
    return {
        'score': max(0, score),
        'issues': issues,
        'is_valid': score >= 50
    }


def clean_response_text(text: str) -> str:
    """Clean up response text for better presentation."""
    if not text:
        return ""

    # Remove excessive whitespace
    import re
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' {2,}', ' ', text)

    # Remove common artifacts
    artifacts = ['```json', '```', '{"', '"}']
    for artifact in artifacts:
        if text.startswith(artifact) or text.endswith(artifact):
            text = text.strip(artifact)

    return text.strip()


def check_response_against_history(response: str, conversation_messages: List[Dict]) -> Dict:
    """
    Check if a response is too similar to recent assistant messages.
    Returns dict with 'is_duplicate' (bool) and 'similarity_score' (float).
    """
    from difflib import SequenceMatcher

    # Get last 5 assistant messages
    recent_responses = []
    for msg in reversed(conversation_messages):
        if msg.get('role') == 'assistant' and msg.get('content'):
            recent_responses.append(msg['content'])
            if len(recent_responses) >= 5:
                break

    if not recent_responses:
        return {'is_duplicate': False, 'similarity_score': 0.0, 'issues': []}

    # Normalize response for comparison — truncate to 300 chars to keep SequenceMatcher fast
    _normalize_re = re.compile(r'[^\w\s]')
    _whitespace_re = re.compile(r'\s+')

    def normalize(text):
        text = _normalize_re.sub('', text.lower())
        text = _whitespace_re.sub(' ', text)
        return text.strip()[:300]

    normalized_response = normalize(response)
    max_similarity = 0.0
    issues = []

    for prev_response in recent_responses:
        normalized_prev = normalize(prev_response)

        # Check for exact phrase matches first (fast)
        if len(normalized_response) > 20 and normalized_response in normalized_prev:
            issues.append("Exact duplicate of previous response")
            return {'is_duplicate': True, 'similarity_score': 1.0, 'issues': issues}

        # Calculate similarity using SequenceMatcher (on truncated text)
        similarity = SequenceMatcher(None, normalized_response, normalized_prev).ratio()
        max_similarity = max(max_similarity, similarity)

        if similarity > 0.85:
            issues.append(f"Very similar to recent response (similarity: {similarity:.2f})")

    is_duplicate = max_similarity > 0.85

    return {
        'is_duplicate': is_duplicate,
        'similarity_score': max_similarity,
        'issues': issues
    }


def extract_action_from_query(query: str) -> Dict[str, Any]:
    """
    Extract the intended action and parameters from a user query.
    v7 ENHANCEMENT: Better query understanding.
    """
    query_lower = query.lower().strip()
    
    # Action patterns with their tool mappings
    action_patterns = {
        'play_music': [r'^play\s+(.+)', r'^put on\s+(.+)', r'^listen to\s+(.+)'],
        'control_lights': [r'^turn (on|off)\s+(.+)?lights?', r'^set lights? to\s+(.+)', r'^(\w+) mode for lights'],
        'web_search': [r'^search (?:for\s+)?(.+)', r'^google\s+(.+)', r'^look up\s+(.+)'],
        'read_gmail': [r'^check (?:my )?email', r'^read (?:my )?email', r'^show (?:my )?inbox'],
        'send_gmail': [r'^send (?:an )?email to\s+(.+)', r'^email\s+(\S+@\S+)'],
        'get_weather': [r'^weather (?:in\s+)?(.+)?', r'^what\'?s the weather'],
        'capture_camera': [r'^what do you see', r'^take a photo', r'^look around'],
    }
    
    for tool, patterns in action_patterns.items():
        for pattern in patterns:
            match = re.search(pattern, query_lower)
            if match:
                return {
                    'action': tool,
                    'params': match.groups() if match.groups() else None,
                    'confidence': 0.9,
                    'raw_match': match.group(0)
                }
    
    return {
        'action': None,
        'params': None,
        'confidence': 0.0,
        'raw_match': None
    }


# ================================================================================

# --- Document storage configuration ---
DOCUMENTS_FOLDER = os.environ.get("DOCUMENTS_DIR", os.path.join(os.getcwd(), "documents"))
CAMERA_FOLDER = os.environ.get("CAMERA_DIR", os.path.join(os.getcwd(), "camera_captures"))
os.makedirs(DOCUMENTS_FOLDER, exist_ok=True)
os.makedirs(CAMERA_FOLDER, exist_ok=True)

# CONFIG (single source)
# ================================================================================

# Core facts DB
BLUE_FACTS_DB = os.environ.get("BLUE_FACTS_DB", "data/blue.db")
BLUE_FACTS: Dict[str, str] = {}

# Model/API (left as env-driven to match your runtime)
MAX_TOKENS = int(os.environ.get("MAX_TOKENS", "4096"))

# Conversation context (resolved the previous inconsistencies; default 40)
MAX_CONTEXT_MESSAGES = int(os.environ.get("MAX_CONTEXT_MESSAGES", "40"))

# Files
UPLOAD_FOLDER = Path(os.environ.get("UPLOAD_FOLDER", "uploads"))
MAX_FILE_SIZE = 16 * 1024 * 1024  # 16MB
ALLOWED_EXTENSIONS = {'bmp', 'csv', 'doc', 'docx', 'gif', 'html', 'jpeg', 'jpg', 'json', 'md', 'pdf', 'png', 'pptx', 'rtf', 'tiff', 'txt', 'webp', 'xlsx', 'xml'}

# DB (conversations)
CONVERSATION_DB = os.environ.get("CONVERSATION_DB", "data/conversations.db")

# Address book
ADDRESS_BOOK_PATH = Path(os.environ.get("BLUE_ADDRESS_BOOK", "/mnt/data/blue_address_book.json"))

# Gmail scopes — unified so read/label/send/compose all work
GMAIL_SCOPES = [
    "https://www.googleapis.com/auth/gmail.modify",
    "https://www.googleapis.com/auth/gmail.send",
    "https://www.googleapis.com/auth/gmail.compose",
]
# ================================================================================
# BLUE CORE MEMORY SYSTEM - IMPROVED VERSION
# ================================================================================

try:
    from blue_memory_improved import get_memory_system
    memory_system = get_memory_system()
    ENHANCED_MEMORY_AVAILABLE = True
    print("[OK] Enhanced memory system loaded - Blue will remember better!")
except ImportError as e:
    ENHANCED_MEMORY_AVAILABLE = False
    memory_system = None
    print(f"[WARN] Enhanced memory not available: {e}")
    print("[WARN] Using legacy memory system")


_facts_cache = {"data": None, "timestamp": 0}
_FACTS_CACHE_TTL = 30  # seconds - facts don't change that often

def load_blue_facts(db_path: str = BLUE_FACTS_DB) -> Dict[str, str]:
    """Load facts using improved memory system if available. Cached for speed."""
    import time as _time
    now = _time.time()
    if _facts_cache["data"] is not None and (now - _facts_cache["timestamp"]) < _FACTS_CACHE_TTL:
        return _facts_cache["data"]

    if ENHANCED_MEMORY_AVAILABLE and memory_system:
        result = memory_system.load_facts()
        _facts_cache["data"] = result
        _facts_cache["timestamp"] = now
        return result

    # Fallback to legacy system
    facts: Dict[str, str] = {}
    try:
        if not os.path.exists(db_path):
            log.warning(f"[MEM] facts DB not found: {db_path}")
            return facts
        conn = sqlite3.connect(db_path)
        conn.row_factory = sqlite3.Row
        rows = conn.execute("SELECT fact_key, values_concat FROM facts_top").fetchall()
        for r in rows:
            facts[r["fact_key"]] = r["values_concat"]
    except Exception as e:
        log.warning(f"[MEM] failed to load facts: {e}")
    finally:
        try:
            conn.close()
        except Exception:
            pass
    if facts:
        log.info(f"[MEM] loaded {len(facts)} core facts from {db_path}")
    _facts_cache["data"] = facts
    _facts_cache["timestamp"] = now
    return facts


def invalidate_facts_cache():
    """Call after saving new facts to refresh cache on next load."""
    _facts_cache["data"] = None
    _facts_cache["timestamp"] = 0

def _facts_block() -> str:
    items: List[str] = []
    # Get fresh facts if enhanced memory is available
    facts = load_blue_facts() if ENHANCED_MEMORY_AVAILABLE else BLUE_FACTS
    
    def add(label: str, key: str) -> None:
        v = facts.get(key)
        if v:
            items.append(f"{label}: {v}")
    for label, key in [
        ("Name", "name"),
        ("Identity", "identity"),
        ("Created by", "created_by"),
        ("Original form", "original_form"),
        ("Upgraded by", "upgraded_by"),
        ("Privacy", "privacy"),
        ("Physical features", "physical_features"),
        ("Tools", "tool"),
        ("Has memory", "has_memory"),
        ("Moods", "mood"),
    ]:
        add(label, key)
    return " | ".join(items)



def save_blue_facts(facts: Dict[str, str], db_path: str = None) -> bool:
    """Save facts using improved memory system if available."""
    global BLUE_FACTS
    BLUE_FACTS.update(facts)
    invalidate_facts_cache()
    
    if ENHANCED_MEMORY_AVAILABLE and memory_system:
        return memory_system.save_facts(facts)
    
    # Fallback to legacy system
    if db_path is None:
        db_path = BLUE_FACTS_DB
    
    try:
        os.makedirs(os.path.dirname(db_path), exist_ok=True)
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()
        
        # Create table if needed
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS facts_top (
                fact_key TEXT PRIMARY KEY,
                values_concat TEXT,
                last_updated TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        
        # Save facts
        saved_count = 0
        for fact_key, values_concat in facts.items():
            cursor.execute("""
                INSERT INTO facts_top (fact_key, values_concat, last_updated)
                VALUES (?, ?, CURRENT_TIMESTAMP)
                ON CONFLICT(fact_key) DO UPDATE SET
                    values_concat = excluded.values_concat,
                    last_updated = CURRENT_TIMESTAMP
            """, (fact_key, values_concat))
            saved_count += 1
        
        conn.commit()
        conn.close()
        log.info(f"[MEM] Saved {saved_count} facts to database")
        return True
    except Exception as e:
        log.error(f"[MEM] Failed to save facts: {e}")
        return False


def extract_and_save_facts(messages: list) -> bool:
    """Extract facts from conversation and save to database."""
    # Try enhanced system first
    if ENHANCED_MEMORY_AVAILABLE and memory_system:
        try:
            return memory_system.extract_and_save_facts(messages)
        except Exception as e:
            log.warning(f"[MEM] Enhanced extraction failed, using legacy: {e}")
    
    # Fallback to legacy extraction
    if not messages:
        return False
    
    import re
    facts_to_save = {}
    
    for msg in messages:
        if msg.get('role') not in ['user', 'assistant']:
            continue
        
        content = msg.get('content', '')
        content_lower = content.lower()
        
        if len(content) < 10 or content.strip().startswith(('{', '[', '```', 'import ')):
            continue
        
        # === NAME EXTRACTION ===
        name_patterns = [
            r"my name is ([A-Z][a-z]+(?:\s+[A-Z][a-z]+)*)",
            r"i'?m ([A-Z][a-z]+)(?:\s|,|\.|$)",
            r"call me ([A-Z][a-z]+)",
            r"this is ([A-Z][a-z]+) speaking",
            r"it'?s ([A-Z][a-z]+) here"
        ]
        for pattern in name_patterns:
            match = re.search(pattern, content)
            if match:
                name = match.group(1).strip()
                if 2 <= len(name) <= 30 and name.replace(' ', '').isalpha():
                    facts_to_save['user_name'] = name
                    log.info(f"[MEM] Learned name: {name}")
                    break
        
        # === LOCATION EXTRACTION ===
        location_patterns = [
            r"i live in ([A-Z][a-zA-Z\s,]+?)(?:\.|,|$|\sand\s|\swith\s)",
            r"i'?m (?:from|in|based in) ([A-Z][a-zA-Z\s]+?)(?:\.|,|$)",
            r"my (?:city|town|home) is ([A-Z][a-zA-Z\s]+)",
            r"we live in ([A-Z][a-zA-Z\s,]+?)(?:\.|,|$)"
        ]
        for pattern in location_patterns:
            match = re.search(pattern, content)
            if match:
                location = match.group(1).strip().rstrip('.,;')
                if 2 <= len(location) <= 50:
                    facts_to_save['location'] = location
                    log.info(f"[MEM] Learned location: {location}")
                    break
        
        # === WORK/EDUCATION ===
        work_patterns = [
            (r"i (?:work|teach) at ([A-Z][a-zA-Z\s&.]+?)(?:\.|,|$)", 'workplace'),
            (r"i work for ([A-Z][a-zA-Z\s&]+?)(?:\.|,|$)", 'workplace'),
            (r"i'?m (?:a|an) ([a-z][a-z\s]+(?:teacher|professor|engineer|developer|doctor|scientist|artist|writer|designer|manager|director))", 'occupation'),
            (r"i studied (?:at )?([A-Z][a-zA-Z\s&]+?)(?:\.|,|$)", 'education'),
            (r"i graduated from ([A-Z][a-zA-Z\s&]+?)(?:\.|,|$)", 'education'),
            (r"i run ([A-Z][a-zA-Z\s&]+?)(?:\.|,|$)", 'business'),
            (r"my company is ([A-Z][a-zA-Z\s&]+?)(?:\.|,|$)", 'business')
        ]
        for pattern, key in work_patterns:
            match = re.search(pattern, content)
            if match:
                value = match.group(1).strip().rstrip('.,;')
                if 2 <= len(value) <= 100:
                    facts_to_save[key] = value
                    log.info(f"[MEM] Learned {key}: {value}")
        
        # === FAMILY EXTRACTION ===
        family_relations = ['partner', 'wife', 'husband', 'spouse', 'daughter', 'son', 'child', 
                           'mother', 'father', 'mom', 'dad', 'brother', 'sister', 'girlfriend', 'boyfriend']
        for relation in family_relations:
            patterns = [
                rf"my {relation}(?:'s name)? is ([A-Z][a-z]+)",
                rf"my {relation},? ([A-Z][a-z]+)",
                rf"(?:this is |meet )?my {relation} ([A-Z][a-z]+)"
            ]
            for pattern in patterns:
                match = re.search(pattern, content)
                if match:
                    name = match.group(1).strip()
                    if 2 <= len(name) <= 30 and name.isalpha():
                        facts_to_save[f'{relation}_name'] = name
                        log.info(f"[MEM] Learned {relation}: {name}")
                        break
        
        # Multiple children (e.g., "my daughters are Emmy, Athena, and Vilda")
        if re.search(r"my (?:daughters?|sons?|children|kids) (?:are|named) ", content_lower):
            match = re.search(r"my (?:daughters?|sons?|children|kids) (?:are|named) ([A-Z][a-zA-Z,\s&]+?)(?:\.|$)", content)
            if match:
                names = match.group(1).strip()
                if 2 <= len(names) <= 100:
                    facts_to_save['children_names'] = names
                    log.info(f"[MEM] Learned children: {names}")
        
        # === PETS ===
        pet_types = ['dog', 'cat', 'pet', 'puppy', 'kitten', 'bird', 'fish', 'hamster', 'rabbit']
        for pet in pet_types:
            patterns = [
                rf"my {pet}(?:'s name)? is ([A-Z][a-z]+)",
                rf"my {pet},? ([A-Z][a-z]+)",
                rf"i have a {pet} (?:named |called )?([A-Z][a-z]+)",
                rf"(?:this is |meet )?my {pet} ([A-Z][a-z]+)"
            ]
            for pattern in patterns:
                match = re.search(pattern, content)
                if match:
                    name = match.group(1).strip()
                    if 2 <= len(name) <= 30 and name.isalpha():
                        facts_to_save[f'{pet}_name'] = name
                        log.info(f"[MEM] Learned {pet}: {name}")
                        break
        
        # === HOBBIES & INTERESTS ===
        hobby_patterns = [
            (r"i (?:love|enjoy|like) (?:to )?([a-z]+ing)", 'hobby'),
            (r"my hobbies? (?:is|are|include) ([a-zA-Z,\s&]+?)(?:\.|$)", 'hobbies'),
            (r"i'?m (?:really )?into ([a-zA-Z\s]+?)(?:\.|,|$)", 'interest'),
            (r"i collect ([a-zA-Z\s]+?)(?:\.|,|$)", 'collection')
        ]
        for pattern, key in hobby_patterns:
            match = re.search(pattern, content_lower)
            if match:
                value = match.group(1).strip().rstrip('.,;')
                if 3 <= len(value) <= 50:
                    facts_to_save[key] = value.title()
                    log.info(f"[MEM] Learned {key}: {value}")
        
        # === PREFERENCES ===
        if 'my favorite' in content_lower or 'i prefer' in content_lower:
            match = re.search(r"my favorite ([a-z\s]+) is ([a-zA-Z0-9\s]+?)(?:\.|,|$)", content_lower)
            if match:
                pref_type = match.group(1).strip().replace(' ', '_')
                pref_value = match.group(2).strip().title()
                if len(pref_type) <= 20 and len(pref_value) <= 50:
                    facts_to_save[f'favorite_{pref_type}'] = pref_value
                    log.info(f"[MEM] Learned {pref_type}: {pref_value}")
            
            match = re.search(r"i prefer ([a-zA-Z\s]+) (?:over|to) ([a-zA-Z\s]+)", content_lower)
            if match:
                preference = f"{match.group(1).strip()} over {match.group(2).strip()}"
                facts_to_save['preference'] = preference.title()
                log.info(f"[MEM] Learned preference: {preference}")
        
        # === BIRTHDAY/AGE ===
        if "i'm " in content_lower or "i am " in content_lower:
            match = re.search(r"i'?m (\d{1,2}) years old", content_lower)
            if match:
                age = match.group(1)
                if 5 <= int(age) <= 120:
                    facts_to_save['age'] = age
                    log.info(f"[MEM] Learned age: {age}")
        
        birthday_patterns = [
            r"my birthday is ([A-Za-z]+ \d{1,2})",
            r"i was born (?:on )?([A-Za-z]+ \d{1,2})",
            r"my birthday'?s? (?:on )?([A-Za-z]+ \d{1,2})"
        ]
        for pattern in birthday_patterns:
            match = re.search(pattern, content)
            if match:
                birthday = match.group(1).strip()
                facts_to_save['birthday'] = birthday
                log.info(f"[MEM] Learned birthday: {birthday}")
                break
        
        # === ALLERGIES/DIETARY ===
        allergy_patterns = [
            r"i'?m allergic to ([a-zA-Z\s,]+?)(?:\.|$)",
            r"i have (?:a |an )?([a-zA-Z\s]+) allergy",
            r"i can'?t eat ([a-zA-Z\s]+?)(?:\.|,|$)"
        ]
        for pattern in allergy_patterns:
            match = re.search(pattern, content_lower)
            if match:
                allergy = match.group(1).strip()
                if 2 <= len(allergy) <= 50:
                    facts_to_save['allergy'] = allergy.title()
                    log.info(f"[MEM] Learned allergy: {allergy}")
                    break
        
        dietary_patterns = [
            r"i'?m (?:a )?(vegetarian|vegan|pescatarian|gluten[- ]free|lactose[- ]intolerant|keto|paleo)",
            r"i (?:don'?t|do not) eat ([a-zA-Z\s]+?)(?:\.|,|$)"
        ]
        for pattern in dietary_patterns:
            match = re.search(pattern, content_lower)
            if match:
                diet = match.group(1).strip()
                facts_to_save['dietary'] = diet.title()
                log.info(f"[MEM] Learned dietary: {diet}")
                break
        
        # === VEHICLES (v6 enhancement) ===
        vehicle_patterns = [
            (r"i drive (?:a |an )?(\d{4} )?([A-Z][a-zA-Z]+ [A-Z]?[a-zA-Z]+)", 'vehicle'),
            (r"my car is (?:a |an )?(\d{4} )?([A-Z][a-zA-Z]+ [A-Z]?[a-zA-Z]+)", 'vehicle'),
            (r"i have (?:a |an )?(\d{4} )?([A-Z][a-zA-Z]+ [A-Z]?[a-zA-Z]+) (?:car|truck|suv|vehicle)", 'vehicle'),
        ]
        for pattern, key in vehicle_patterns:
            match = re.search(pattern, content)
            if match:
                year = match.group(1).strip() if match.group(1) else ""
                make_model = match.group(2).strip()
                vehicle = f"{year}{make_model}".strip()
                if 3 <= len(vehicle) <= 50:
                    facts_to_save['vehicle'] = vehicle
                    log.info(f"[MEM] Learned vehicle: {vehicle}")
                    break
        
        # === LANGUAGES (v6 enhancement) ===
        language_patterns = [
            r"i speak ([A-Z][a-z]+(?:,? (?:and )?[A-Z][a-z]+)*)",
            r"i'?m fluent in ([A-Z][a-z]+(?:,? (?:and )?[A-Z][a-z]+)*)",
            r"my (?:native|first) language is ([A-Z][a-z]+)",
            r"i'?m learning ([A-Z][a-z]+)"
        ]
        for pattern in language_patterns:
            match = re.search(pattern, content)
            if match:
                languages = match.group(1).strip()
                if 3 <= len(languages) <= 100:
                    facts_to_save['languages'] = languages
                    log.info(f"[MEM] Learned languages: {languages}")
                    break
        
        # === SKILLS/EXPERTISE (v6 enhancement) ===
        skill_patterns = [
            r"i'?m (?:good|great|skilled|experienced) at ([a-zA-Z\s,]+?)(?:\.|$)",
            r"i know (?:how to )?([a-zA-Z\s]+?)(?:\.|$)",
            r"i can ([a-zA-Z\s]+) (?:well|professionally|expertly)",
            r"my skills? (?:include|are|is) ([a-zA-Z\s,]+?)(?:\.|$)"
        ]
        for pattern in skill_patterns:
            match = re.search(pattern, content_lower)
            if match:
                skill = match.group(1).strip()
                if 3 <= len(skill) <= 100 and skill not in ['do', 'be', 'help']:
                    facts_to_save['skills'] = skill.title()
                    log.info(f"[MEM] Learned skill: {skill}")
                    break
        
        # === MEDICAL (v6 enhancement) ===
        medical_patterns = [
            r"i have ([a-zA-Z\s]+(?:diabetes|asthma|arthritis|condition|disease|disorder))",
            r"i'?m (?:on|taking) ([a-zA-Z\s]+) (?:medication|medicine|pills)",
            r"i wear ([a-zA-Z\s]+(?:glasses|contacts|hearing aid|braces))"
        ]
        for pattern in medical_patterns:
            match = re.search(pattern, content_lower)
            if match:
                medical = match.group(1).strip()
                if 3 <= len(medical) <= 50:
                    facts_to_save['medical'] = medical.title()
                    log.info(f"[MEM] Learned medical: {medical}")
                    break
        
        # === TIMEZONE/SCHEDULE (v6 enhancement) ===
        if 'timezone' in content_lower or 'time zone' in content_lower:
            match = re.search(r"(?:my )?time ?zone is ([A-Z]{2,4}|[A-Z][a-z]+/[A-Z][a-z]+)", content)
            if match:
                tz = match.group(1).strip()
                facts_to_save['timezone'] = tz
                log.info(f"[MEM] Learned timezone: {tz}")
        
        # === PHONE/CONTACT (v6 enhancement) ===
        phone_match = re.search(r"my (?:phone|number|cell) is ([0-9\-\(\)\s]{10,20})", content)
        if phone_match:
            phone = phone_match.group(1).strip()
            facts_to_save['phone'] = phone
            log.info(f"[MEM] Learned phone: {phone}")
        
        # === EMAIL (v6 enhancement) ===
        email_match = re.search(r"my email is ([a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,})", content)
        if email_match:
            email = email_match.group(1).strip()
            facts_to_save['email'] = email
            log.info(f"[MEM] Learned email: {email}")
    
    if facts_to_save:
        global BLUE_FACTS
        BLUE_FACTS.update(facts_to_save)
        return save_blue_facts(BLUE_FACTS)
    
    return False


def build_system_preamble(robot_name: str = "Blue") -> str:
    # No bot-identity facts here anymore. The facts table's identity rows
    # describe BLUE ("name: Blue", "identity: Friendly home robot assistant",
    # the privacy blurb) and this preamble stamped them on EVERY robot as
    # "ground-truth... Do not contradict them" — on Hexia's page the robot
    # introduced itself as Blue and denied being Hexia (2026-07-12). They
    # were also the source of Blue's own brochure recitals. Robot identity
    # lives in persona_line + the self-profile + the j_space now; the
    # <known_facts> block carries the USER's facts (and already filters bot
    # facts on its side).
    # Hardcoded user pronouns — Alex uses he/him, and the model has been
    # caught defaulting to "she" otherwise.
    return (
        f"You are {robot_name}. Alex uses he/him pronouns — always refer to "
        "Alex as he/him, never as she/her."
    )

# Load facts at import time (only if enhanced memory not available)
try:
    if not ENHANCED_MEMORY_AVAILABLE:
        BLUE_FACTS = load_blue_facts()
    else:
        # With enhanced memory, facts are loaded fresh each conversation
        BLUE_FACTS = {}
except Exception:
    BLUE_FACTS = {}

# ===== Enhanced Tools Import =====
try:
    from blue_tools_enhanced import (
        CalendarManager,
        ContactManager,
        TaskManager,
        NoteManager,
        SystemController,
        FileOperations,
        TimerManager,
        StorytellingTools,
        LocationServices,
        SmartHomeController,
        MusicController
    )
    ENHANCED_TOOLS_AVAILABLE = True
    print("[OK] Enhanced tools loaded successfully!")
except ImportError as e:
    ENHANCED_TOOLS_AVAILABLE = False
    print(f"[WARN] Enhanced tools not available: {e}")

# ===== Conversation Persistence Setup =====
try:
    from blue_database import create_database
    db = create_database()
    CONVERSATION_DB_AVAILABLE = True
    print("[OK] Conversation database connected - long-term memory enabled!")
except Exception as e:
    CONVERSATION_DB_AVAILABLE = False
    db = None
    if ENHANCED_MEMORY_AVAILABLE and memory_system:
        print(f"[INFO] Legacy conversation database not loaded: {e}")
        print("[OK] Enhanced memory is active for cross-session conversations")
    else:
        print(f"[WARN] Conversation database not available: {e}")
        print("[WARN] Blue will not remember conversations across sessions")


# ================================================================================
# IMPROVED TOOL SELECTION SYSTEM (Integrated Version - October 2025)
# ================================================================================
# This section contains the confidence-based tool selection system.
# When enabled (USE_IMPROVED_SELECTOR=True), it provides:
# - Confidence scoring (0.0-1.0) for each tool
# - Priority-based conflict resolution
# - Context awareness from conversation history
# - Disambiguation when confidence is low
# - Negative signals to prevent false positives
# ================================================================================

"""
Blue Robot Tool Selection - ENHANCED VERSION
=============================================

Key improvements over original:
1. Clear positive AND negative signals for each tool
2. Better email disambiguation (read vs send vs reply)
3. Better search disambiguation (web vs documents)
4. Context awareness from conversation history
5. Specific disambiguation questions when uncertain
"""

from dataclasses import dataclass, field
from typing import List, Dict, Optional, Tuple, Set
import re
from collections import defaultdict, Counter
import json


# ================================================================================
# TOOL SELECTOR - Now using modular package
# ================================================================================
#
# The tool selector has been refactored into a modular package: blue/tool_selector/
# This provides better maintainability, testability, and extensibility.
#
# Old implementation: 1,700+ lines in this file
# New implementation: Modular package with 19 files
#
# Benefits:
# - 90% smaller files (largest is 298 lines vs 3,028)
# - Each detector is independently testable
# - Easy to add new detectors (create new file)
# - 10-20% performance improvement
# - 100% backward compatible
#
# For details, see: README_REFACTORING.md
# ================================================================================

from blue.tool_selector import (
    ToolIntent,
    ToolSelectionResult,
    ImprovedToolSelector,
    integrate_with_existing_system,
)

# Direct Ohbot head control (this branch only — replaces the Ohbot app for the
# head). The module is defensive: if the library isn't installed or the board
# isn't connected, all head calls are no-ops, so the rest of Blue keeps running.
from blue import head as blue_head
from blue.mood_eyes import mood_eye_color

# Connect the head at module load (covers both entry points: `python
# bluetools.py` directly AND `python run.py` which imports this module). Make
# sure the Ohbot desktop app is CLOSED, or it will hold the serial port.
try:
    blue_head.init()                      # Blue (default head)
    try:
        blue_head.hexia.init()            # Hexia (second head; no-op until her board is assigned)
    except Exception as _hex_init_err:
        print(f"[HEAD] hexia init skipped: {_hex_init_err!r}")
    import atexit as _atexit_head
    _atexit_head.register(blue_head.close_all)   # cleanly reset/close both boards at exit
except Exception as _head_init_err:
    print(f"[HEAD] init skipped: {_head_init_err!r}")

# ================================================================================
# END OF IMPROVED TOOL SELECTION SYSTEM
# ================================================================================

# Initialize the tool selector (always use improved modular system)
TOOL_SELECTOR = ImprovedToolSelector()
print("[OK] Tool selector initialized - using modular confidence-based selection")


app = Flask(__name__)

# ============================================================================
# Remote-access auth gate
# ----------------------------------------------------------------------------
# Blue runs with full owner powers (sends email as Alex, reads his inbox,
# controls the lights, opens the camera). Once the server is reachable beyond
# this machine, every NON-LOCAL request must authenticate. Localhost stays
# ungated on purpose: the physical Ohbot client POSTs to 127.0.0.1/v1/chat/
# completions and the on-PC browser hits 127.0.0.1 — both must keep working
# untouched. Remote devices (LAN IP, or a Tailscale 100.x peer) are NOT
# localhost, so they're gated behind a shared password.
# ============================================================================
import hmac as _hmac
import secrets as _secrets
from datetime import timedelta as _timedelta

_SECRET_KEY_FILE = os.path.join(os.getcwd(), ".blue_secret_key")
_PASSWORD_FILE = os.path.join(os.getcwd(), ".blue_password")


def _load_or_create_secret_key() -> str:
    """Stable Flask session secret. Env wins; otherwise persist a random key
    to a gitignored file so sessions survive restarts (no re-login churn)."""
    val = (os.environ.get("BLUE_SECRET_KEY") or "").strip()
    if val:
        return val
    try:
        if os.path.exists(_SECRET_KEY_FILE):
            existing = open(_SECRET_KEY_FILE, encoding="utf-8-sig").read().strip()
            if existing:
                return existing
    except Exception:
        pass
    val = _secrets.token_hex(32)
    try:
        with open(_SECRET_KEY_FILE, "w", encoding="utf-8") as f:
            f.write(val)
    except Exception as e:
        print(f"   [AUTH] couldn't persist secret key ({e}); using ephemeral key")
    return val


def _access_password() -> str:
    """The shared remote-access password. Set via BLUE_ACCESS_PASSWORD env or a
    one-line .blue_password file (both gitignored). Empty => remote disabled."""
    pw = (os.environ.get("BLUE_ACCESS_PASSWORD") or "").strip()
    if pw:
        return pw
    try:
        if os.path.exists(_PASSWORD_FILE):
            # utf-8-sig so a BOM (which Windows PowerShell's Out-File adds) is
            # stripped rather than read as part of the password.
            return open(_PASSWORD_FILE, encoding="utf-8-sig").read().strip()
    except Exception:
        pass
    return ""


app.secret_key = _load_or_create_secret_key()
app.permanent_session_lifetime = _timedelta(days=30)


def _is_local_request() -> bool:
    return (request.remote_addr or "") in ("127.0.0.1", "::1", "localhost")


# ============================================================================
# Device-based speaker identity
# ----------------------------------------------------------------------------
# The household has a fixed device-per-person split: the iPad is always Vilda;
# the MacBook, PC, iPhone, and the physical robot are always Alex. We tag each
# request with a device, then map device -> person. The chat web page computes
# a reliable device tag in the browser (where a "desktop-mode" iPad is still
# distinguishable from a real Mac via touch support) and sends it as the
# X-Blue-Device header. When that header is absent (e.g. the Ohbot client
# POSTing directly) we fall back to sniffing the User-Agent.
#
# To remap a device to a different person, edit _DEVICE_OWNER below.
# ============================================================================
_DEVICE_OWNER = {"ipad": "Vilda"}   # everything not listed here => _DEFAULT_USER
_DEFAULT_USER = "Alex"

# Per-person guidance Blue follows when that person is the one chatting. The
# text is appended to the "who you're talking to" note built in
# process_with_tools. Keyed by the name _identify_user_from_request returns;
# add an entry here to tune how Blue speaks to another household member.
_SPEAKER_PROFILES = {
    "Vilda": (
        "Vilda is 8 years old — one of the children in Alex's household. Talk to "
        "her accordingly: warm, gentle, patient and encouraging, with short "
        "sentences and simple everyday words, and keep everything child-safe and "
        "age-appropriate. Greet her sweetly and cheerfully so she feels happy and "
        "wants to keep talking with you. Do NOT bring up the calendar, schedules, "
        "reminders or appointments with her — that's grown-up stuff; if she "
        "mentions it, gently steer back to something fun. You may use a little Gen "
        "Z slang now and then to be fun and relatable (things like \"that's so "
        "cool\", \"bet\", \"no cap\", \"slay\") — but only occasionally, never in "
        "every sentence, and never let the slang make you harder to understand."
    ),
}


def _device_tag_from_user_agent(ua: str) -> str:
    ua = (ua or "").lower()
    if "ipad" in ua:
        return "ipad"
    if "iphone" in ua or "ipod" in ua:
        return "iphone"
    if "android" in ua:
        return "android"
    if "macintosh" in ua or "mac os" in ua:
        return "mac"
    if "windows" in ua:
        return "windows"
    return "other"


# Stable Tailscale IPs identify a device even with no browser hint. Needed
# because Tailscale Serve terminates HTTPS and proxies in from localhost (so
# remote_addr looks local), while forwarding the real client IP in
# X-Forwarded-For. Verified machine: ipad-2 = 100.71.165.19 (Vilda).
_TAILSCALE_IP_OWNER = {"100.71.165.19": "Vilda"}


def _forwarded_client_ip() -> str:
    """Real client IP when behind a proxy (Tailscale Serve), else ''."""
    xff = (request.headers.get("X-Forwarded-For") or "").strip()
    return xff.split(",")[0].strip() if xff else ""


def _identify_user_from_request() -> str:
    """Map the requesting device to the person behind it (Alex by default).

    Order matters: an explicit hint from our own chat page wins; then the
    Tailscale source IP (survives Tailscale Serve's localhost proxying); then a
    direct localhost request is Alex (PC browser + Ohbot client); finally a
    best-effort User-Agent guess.
    """
    tag = (request.headers.get("X-Blue-Device") or "").strip().lower()
    if tag:
        return _DEVICE_OWNER.get(tag, _DEFAULT_USER)
    fip = _forwarded_client_ip()
    if fip in _TAILSCALE_IP_OWNER:
        return _TAILSCALE_IP_OWNER[fip]
    if _is_local_request():
        return _DEFAULT_USER
    return _DEVICE_OWNER.get(
        _device_tag_from_user_agent(request.headers.get("User-Agent", "")),
        _DEFAULT_USER,
    )


# Endpoints that must stay reachable without a session even from remote, so a
# logged-out phone can actually render the login form and submit it.
_AUTH_EXEMPT_ENDPOINTS = {"blue_login", "blue_logout", "static",
                          "asset_blue_css", "asset_blue_js"}


@app.before_request
def _require_remote_auth():
    # Local traffic (Ohbot client + on-PC browser) is fully trusted.
    if _is_local_request():
        return None
    if request.endpoint in _AUTH_EXEMPT_ENDPOINTS:
        return None
    pw = _access_password()
    if not pw:
        # Fail closed: someone reached Blue remotely but no password is set.
        return Response(
            "Remote access to Blue is disabled. Set a password in a "
            ".blue_password file (or BLUE_ACCESS_PASSWORD) on the host.",
            status=503, mimetype="text/plain",
        )
    if session.get("blue_auth") is True:
        return None
    # Programmatic callers get a clean 401; browsers get the login page.
    accepts_html = "text/html" in (request.headers.get("Accept") or "")
    if request.path.startswith(("/v1/", "/api/", "/memory/")) or not accepts_html:
        return Response("Authentication required.", status=401, mimetype="text/plain")
    return redirect(url_for("blue_login", next=request.full_path))


# Users who are restricted to the chat only (e.g. Vilda on the iPad — she can
# talk to Blue but not reach contacts/calendar/visual memory/documents/etc).
_CHAT_ONLY_USERS = {"Vilda"}
# The only endpoints those users may reach: the chat page + everything the chat
# page needs to function (send messages, attachments, voice, assets, login).
# NOTE: the head endpoints are intentionally NOT here — when Blue talks to Vilda
# the physical robot must stay completely still (her reply is spoken on the iPad,
# not performed by the robot), so /head/* requests from her device are refused.
_CHAT_ONLY_ALLOWED = {
    "chat_page", "chat_attach", "chat_eyes", "chat_completions", "stt", "stt_warmup",
    "blue_login", "blue_logout", "static", "asset_blue_css", "asset_blue_js",
}
# Tools chat-only users (the kids' iPad) may NOT trigger:
#  - music plays through the PC's speakers, so driving it from the iPad is
#    pointless/disruptive;
#  - move_head would move the physical robot while Blue talks to Vilda — it must
#    stay still (handled as a silent no-op in execute_tool, below);
#  - calendar/reminder tools: Blue doesn't discuss the calendar with Vilda;
#  - capture_camera grabs the PC webcam — never her camera. Her eyes are the
#    iPad frame staged via /chat/eyes, so the PC cam must never run for her.
#  - email_snapshot runs that same PC webcam AND sends mail from Blue's own
#    Gmail account — strictly an owner power.
_KID_BLOCKED_TOOLS = {"play_music", "control_music", "music_visualizer",
                      "move_head", "create_reminder", "get_upcoming_reminders",
                      "capture_camera", "email_snapshot"}


@app.before_request
def _restrict_chat_only_users():
    """Keep chat-only users (the kids' iPad) inside the chat. Page navigations
    elsewhere bounce back to /chat; API calls get a plain 403."""
    try:
        user = _identify_user_from_request()
    except Exception:
        return None
    if user not in _CHAT_ONLY_USERS:
        return None
    if (request.endpoint or "") in _CHAT_ONLY_ALLOWED:
        return None
    accepts_html = "text/html" in (request.headers.get("Accept") or "")
    if request.method == "GET" and accepts_html:
        return redirect(url_for("chat_page"))
    return Response("Not available on this device.", status=403, mimetype="text/plain")


from blue.server.pages.login import _LOGIN_HTML


@app.route("/login", methods=["GET", "POST"])
def blue_login():
    error = None
    nxt = request.values.get("next") or "/chat"
    # Only allow same-site relative redirects.
    if not nxt.startswith("/") or nxt.startswith("//"):
        nxt = "/chat"
    if request.method == "POST":
        pw = _access_password()
        supplied = request.form.get("password", "")
        if pw and _hmac.compare_digest(supplied, pw):
            session["blue_auth"] = True
            session.permanent = True
            return redirect(nxt)
        error = "Incorrect password." if pw else "No password is configured on the host."
    return render_template_string(_LOGIN_HTML, error=error, next=nxt)


@app.route("/logout")
def blue_logout():
    session.pop("blue_auth", None)
    return redirect(url_for("blue_login"))


# Configuration
LM_STUDIO_URL = "http://127.0.0.1:1234/v1/chat/completions"
LM_STUDIO_RAG_URL = "http://127.0.0.1:1234/v1/rag"

# ============================
# LM Studio — single provider
# ============================
class LMStudioClient:
    """
    Enhanced client for local LM Studio (OpenAI-compatible) chat completions.
    v6 ENHANCEMENTS:
    - Auto-retry with exponential backoff
    - Connection health checks
    - Request/response logging
    - Timeout management
    """
    def __init__(self, base_url: Optional[str] = None, model: Optional[str] = None, 
                 timeout: Optional[float] = None, max_retries: int = 3):
        self.base_url = (
            base_url
            or os.environ.get("LM_STUDIO_URL")
            or globals().get("LM_STUDIO_URL")
            or "http://127.0.0.1:1234/v1/chat/completions"
        )
        self.model = (
            model
            or os.environ.get("LM_STUDIO_MODEL")
            or globals().get("LM_STUDIO_MODEL")
            or "local-model"
        )
        self.timeout = float(timeout or os.environ.get("LM_STUDIO_TIMEOUT", "120"))
        self.max_retries = max_retries
        self._healthy = None
        self._last_health_check = 0
    
    def is_healthy(self, force_check: bool = False) -> bool:
        """Check if LM Studio is responding (cached for 60s)."""
        import time
        now = time.time()
        if not force_check and self._healthy is not None and (now - self._last_health_check) < 60:
            return self._healthy
        
        try:
            # Try to hit the models endpoint as a health check
            health_url = self.base_url.replace('/chat/completions', '/models')
            resp = requests.get(health_url, timeout=5)
            self._healthy = resp.status_code == 200
        except Exception:
            self._healthy = False
        
        self._last_health_check = now
        return self._healthy

    def chat(
        self,
        messages: List[Dict[str, Any]],
        tools: Optional[List[Dict[str, Any]]] = None,
        tool_choice: str = "auto",
        max_tokens: Optional[int] = None,
        temperature: Optional[float] = None,
        extra: Optional[Dict[str, Any]] = None,
        **kwargs: Any
    ) -> Dict[str, Any]:
        payload: Dict[str, Any] = {
            "model": self.model,
            "messages": messages,
            "frequency_penalty": 0.4,  # Strong penalty to reduce repetition of tokens
            "presence_penalty": 0.3    # Strong penalty to encourage topic diversity
        }
        if temperature is not None:
            payload["temperature"] = float(temperature)
        else:
            payload["temperature"] = 0.8  # Slightly higher default for more variation
        if max_tokens is not None:
            payload["max_tokens"] = int(max_tokens)
        if tools:
            payload["tools"] = tools
            payload["tool_choice"] = tool_choice
        if extra and isinstance(extra, dict):
            payload.update(extra)
        if kwargs:
            payload.update(kwargs)

        # Retry logic with exponential backoff
        import time
        last_error = None
        for attempt in range(self.max_retries):
            try:
                resp = requests.post(self.base_url, json=payload, timeout=self.timeout)
                resp.raise_for_status()
                result = resp.json()
                
                # Validate response structure
                if 'choices' not in result and 'error' not in result:
                    raise ValueError(f"Unexpected response structure: {list(result.keys())}")
                
                return result
                
            except requests.exceptions.Timeout as e:
                last_error = e
                wait_time = 2 ** attempt
                print(f"   [LLM] Timeout on attempt {attempt + 1}, retrying in {wait_time}s...")
                time.sleep(wait_time)
                
            except requests.exceptions.ConnectionError as e:
                last_error = e
                wait_time = 2 ** attempt
                print(f"   [LLM] Connection error on attempt {attempt + 1}, retrying in {wait_time}s...")
                time.sleep(wait_time)
                
            except requests.exceptions.HTTPError as e:
                # Don't retry on 4xx errors (client errors)
                if e.response.status_code < 500:
                    return {"error": f"HTTP {e.response.status_code}: {e.response.text[:200]}"}
                last_error = e
                wait_time = 2 ** attempt
                print(f"   [LLM] Server error on attempt {attempt + 1}, retrying in {wait_time}s...")
                time.sleep(wait_time)
                
            except Exception as e:
                last_error = e
                print(f"   [LLM] Unexpected error: {e}")
                break
        
        return {"error": f"LLM request failed after {self.max_retries} attempts: {last_error}"}

# Global LM Studio client
try:
    _LM = LMStudioClient()
except Exception as _e:
    print(f"[WARN] Failed to init LM Studio client: {_e}")
    _LM = None

def call_llm(
    messages: List[Dict[str, Any]],
    include_tools: bool = True,
    tool_choice: str = "auto",
    force_tool: Optional[str] = None,
    max_tokens: Optional[int] = None,
    temperature: Optional[float] = None,
    extra: Optional[Dict[str, Any]] = None,
    tools_override: Optional[List[Dict[str, Any]]] = None,
    **kwargs: Any
) -> Dict[str, Any]:
    """Unified LLM entrypoint: always uses local LM Studio.

    tools_override: when provided, this exact tool list is sent instead of
    the global TOOLS array (used by the email auto-reply to expose only a
    restricted, read-only subset). Takes precedence over include_tools.
    """
    # Nudge if a specific tool is required
    if force_tool:
        messages = list(messages)
        if messages and isinstance(messages[-1], dict) and messages[-1].get("role") == "user":
            messages[-1] = {**messages[-1]}
            messages[-1]["content"] = (
                (messages[-1].get("content") or "")
                + "\n\n[System note: Use the specified tool to satisfy this request.]"


            )
    if tools_override is not None:
        tools_payload = tools_override
    else:
        tools_payload = None
        try:
            tools_payload = TOOLS if include_tools and "TOOLS" in globals() else None  # noqa: F821
        except Exception:
            tools_payload = None

    if _LM is None:
        return {"error": "LM Studio client not available"}
    try:
        return _LM.chat(
            messages,
            tools=tools_payload,
            tool_choice=tool_choice,
            max_tokens=max_tokens,
            temperature=temperature,
            extra=extra,
            **kwargs
        )
    except Exception as e:
        return {"error": f"LM Studio request failed: {e}"}

PROXY_PORT = 5000


# Document search behavior: "opt_in" (ask first) or "aggressive" (auto)
AUTO_DOCSEARCH_MODE = "opt_in"


# ===== Enhanced Settings & Logger =====
@dataclass
class Settings:
    LOG_LEVEL: str = "INFO"
    MAX_ITERATIONS: int = 3
    TOOL_TIMEOUT_SECS: float = 15.0
    TOOL_RETRIES: int = 2
    # Conversation trimming: retain only the most recent N messages (plus system) when
    # sending context to the language model. This helps prevent the model from
    # confusing long conversation history or previous tool responses with the
    # user’s current intent. A value of 0 disables trimming.
    MAX_CONTEXT_MESSAGES: int = 40  # Increased to preserve .ocf memories and cross-session recall
    AUTO_DOCSEARCH_MODE: str = AUTO_DOCSEARCH_MODE if "AUTO_DOCSEARCH_MODE" in globals() else "opt_in"

_settings = Settings()
# Music configuration
MUSIC_SERVICE = "youtube_music"  # or "amazon_music"
YOUTUBE_MUSIC_BROWSER = None  # Will store ytmusicapi instance

# Document storage
UPLOAD_FOLDER.mkdir(exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
# 64MB: modern phone photos (48MP, HDR) routinely exceed 16MB, and a chat photo
# upload that trips the limit returns an HTML 413 the client can't parse, which
# surfaced as a generic "upload failed".
app.config['MAX_CONTENT_LENGTH'] = 64 * 1024 * 1024


@app.errorhandler(413)
def _too_large(_e):
    """Return JSON (not HTML) so uploads that exceed MAX_CONTENT_LENGTH show a
    real message instead of the client's generic 'upload failed'."""
    limit_mb = app.config['MAX_CONTENT_LENGTH'] // (1024 * 1024)
    return jsonify({
        "attachments": [{"name": "file", "kind": "error",
                         "error": f"file too large (limit {limit_mb}MB)"}],
        "error": f"file too large (limit {limit_mb}MB)",
    }), 413

# Document index (stores metadata)
DOCUMENT_INDEX_FILE = "document_index.json"

# Hue Configuration
HUE_CONFIG_PATH = "hue_config.json"
HUE_CONFIG = {}


def _hue_bridge_alive(ip: str, timeout: float = 2.0) -> bool:
    """Quick health probe — does this IP look like a Hue bridge right now?"""
    if not ip:
        return False
    try:
        # /api/0/config is unauthenticated and returns bridge metadata.
        # Try HTTPS first (newer bridges enforce TLS), fall back to HTTP.
        try:
            import urllib3
            urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
        except Exception:
            pass
        for proto in ("https", "http"):
            try:
                r = requests.get(f"{proto}://{ip}/api/0/config", timeout=timeout, verify=False)
                if r.status_code == 200 and "bridgeid" in r.text.lower():
                    return True
            except Exception:
                continue
    except Exception:
        pass
    return False


def _discover_hue_bridge() -> str | None:
    """Ask Philips' cloud discovery service for the bridge's current LAN IP.

    Returns the bridge IP if discovery succeeds, else None. Cheap and reliable
    — the bridge phones home periodically so the cloud always knows where it is
    on the LAN. Avoids needing mDNS or manual rediscovery when DHCP shuffles."""
    try:
        # Suppress only the InsecureRequestWarning from urllib3 (we use
        # verify=False because Hue bridges use self-signed certs).
        try:
            import urllib3
            urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
        except Exception:
            pass
        r = requests.get("https://discovery.meethue.com/", timeout=5)
        if r.status_code != 200:
            return None
        bridges = r.json()
        if not bridges:
            return None
        return bridges[0].get("internalipaddress") or None
    except Exception as e:
        print(f"[WARN]  Hue discovery failed: {e}")
        return None


def _save_hue_config(config: dict) -> None:
    """Write the updated config back to disk so we don't re-discover next start."""
    try:
        with open(HUE_CONFIG_PATH, "w") as f:
            json.dump(config, f, indent=2)
    except Exception as e:
        print(f"[WARN]  Could not persist hue_config.json: {e}")


try:
    with open(HUE_CONFIG_PATH, "r") as f:
        HUE_CONFIG = json.load(f)
    saved_ip = HUE_CONFIG.get("bridge_ip", "")
    if saved_ip and _hue_bridge_alive(saved_ip):
        print(f"[OK] Hue config loaded: Bridge at {saved_ip}")
    else:
        # Saved IP doesn't respond — DHCP probably gave the bridge a new lease.
        # Try the cloud discovery service to find the bridge's current IP.
        if saved_ip:
            print(f"[WARN] Hue bridge at {saved_ip} not responding — running auto-discovery")
        new_ip = _discover_hue_bridge()
        if new_ip and _hue_bridge_alive(new_ip):
            print(f"[OK] Hue bridge auto-discovered at {new_ip} (was {saved_ip or 'unset'})")
            HUE_CONFIG["bridge_ip"] = new_ip
            _save_hue_config(HUE_CONFIG)
        else:
            print(f"[WARN] Auto-discovery couldn't reach a bridge — Hue features disabled this session")
except FileNotFoundError:
    print("[WARN]  No hue_config.json found. Run setup_hue.py first!")
except Exception as e:
    print(f"[WARN]  Error loading Hue config: {e}")

BRIDGE_IP = HUE_CONFIG.get("bridge_ip", "")
HUE_USERNAME = HUE_CONFIG.get("username", "")

# Gmail Configuration

# Gmail library availability
try:
    from googleapiclient.discovery import build  # already imported above
    from google_auth_oauthlib.flow import InstalledAppFlow  # already imported above
    from google.auth.transport.requests import Request  # already imported above
    GMAIL_AVAILABLE = True
except Exception:
    GMAIL_AVAILABLE = False


GMAIL_TOKEN_FILE = "gmail_token.pickle"
GMAIL_CREDENTIALS_FILE = "gmail_credentials.json"
GMAIL_USER_EMAIL = "alevantresearch@gmail.com"
import threading as _gmail_threading
_gmail_creds = None                       # shared OAuth credentials (lock-protected)
_GMAIL_CREDS_LOCK = _gmail_threading.Lock()
_gmail_tls = _gmail_threading.local()     # per-THREAD service objects


def get_gmail_service():
    """Get a Gmail API service for THIS thread.

    The googleapiclient/httplib2 stack is NOT thread-safe: one shared service
    object used simultaneously from the auto-reply background thread and Flask
    request threads interleaves two conversations on the same SSL connection
    and corrupts the native heap — the whole process dies with 0xc0000374 /
    0xc0000005 in ntdll.dll and NO Python traceback. Rare while Gmail calls
    were occasional; the duet's 8-second /duet/mail/check poll (2026-07-03)
    made the collision routine and crashed the server within minutes.

    Fix: credentials are shared and refreshed under a lock, but each thread
    builds its OWN service (its own HTTP connection). build() uses the bundled
    static discovery document, so a per-thread build costs milliseconds and no
    network round-trip."""
    global _gmail_creds
    if not GMAIL_AVAILABLE:
        raise Exception("Gmail libraries not installed")

    svc = getattr(_gmail_tls, 'service', None)
    if svc is not None:
        return svc

    with _GMAIL_CREDS_LOCK:
        creds = _gmail_creds
        # Load existing token once per process
        if creds is None and os.path.exists(GMAIL_TOKEN_FILE):
            with open(GMAIL_TOKEN_FILE, 'rb') as token:
                creds = pickle.load(token)

        # If no valid credentials, authenticate / refresh
        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                creds.refresh(Request())
            else:
                if not os.path.exists(GMAIL_CREDENTIALS_FILE):
                    raise Exception(f"Gmail credentials file not found: {GMAIL_CREDENTIALS_FILE}. " +
                                  "Download from Google Cloud Console and save as gmail_credentials.json")
                flow = InstalledAppFlow.from_client_secrets_file(GMAIL_CREDENTIALS_FILE, GMAIL_SCOPES)
                creds = flow.run_local_server(port=0)

            # Save the credentials
            with open(GMAIL_TOKEN_FILE, 'wb') as token:
                pickle.dump(creds, token)

        _gmail_creds = creds

    svc = build('gmail', 'v1', credentials=creds, cache_discovery=False)
    _gmail_tls.service = svc
    return svc

# IMPROVED Keywords - More comprehensive and specific detection
SEARCH_KEYWORDS = [
    'search', 'look up', 'find out', 'google', 'check online', 'search for',
    'tell me about', 'information on', 'facts about', 'research',
    'check the internet', 'web search', 'online', 'latest', 'recent', 'current',
    'news about', 'who won', 'what happened', 'update on', 'check if'
]

WEATHER_KEYWORDS = [
    'weather', 'temperature', 'forecast', 'rain', 'raining', 'snow', 'snowing',
    'sunny', 'cloudy', 'storm', 'humidity', 'wind', 'windy', 'cold', 'hot',
    'warm', 'degrees', 'celsius', 'fahrenheit', 'climate'
]

LIGHT_KEYWORDS = [
    'light', 'lights', 'lamp', 'lamps', 'brightness', 'dim', 'bright',
    'turn on', 'turn off', 'switch on', 'switch off', 'color', 'colour',
    'mood', 'scene', 'atmosphere', 'illuminate', 'lighting', 'bulb',
    'darker', 'brighter', 'glow', 'hue', 'philips'
]

VISUALIZER_KEYWORDS = [
    'visualizer', 'light show', 'light dance', 'dancing lights', 'party lights',
    'disco', 'strobe', 'color changing', 'dynamic lights', 'animated lights'
]

DOCUMENT_KEYWORDS = [
    'document', 'doc', 'file', 'pdf', 'my documents', 'my files',
    'uploaded', 'contract', 'agreement', 'policy', 'deadline', 'due date',
    'syllabus', 'course', 'assignment', 'exam', 'schedule', 'paper', 'report',
    'memo', 'notes', 'guidelines', 'instructions', 'manual', 'handbook',
    'according to my', 'in my file', 'what does my', 'says in'
]

# ===== NEW: IMPROVED KEYWORD LISTS FOR BETTER TOOL DETECTION =====
# Keywords for RETRIEVING/READING a specific document
DOCUMENT_RETRIEVAL_KEYWORDS = [
    'read me', 'read to me', 'show me', 'display', 'view', 'open',
    'read the', 'show the', 'display the', 'view the', 'open the',
    'entire document', 'full document', 'whole document', 'complete document',
    'the document called', 'the file called', 'document named', 'file named',
    'in your documents folder', 'from your documents', 'from the documents folder'
]

# Keywords for SEARCHING within documents (semantic/RAG search)
DOCUMENT_SEARCH_KEYWORDS = [
    'search my documents', 'search documents', 'find in my documents',
    'what does my document say about', 'according to my documents',
    'in my documents about', 'search for', 'find information about',
    'what information', 'do my documents mention', 'do my files contain'
]

# Keywords for WEB SEARCH (internet search)
WEB_SEARCH_KEYWORDS = [
    'search the web', 'search online', 'google', 'search for online',
    'look up online', 'search the internet', 'find on the web',
    'current', 'latest', 'recent', 'today', 'this week', 'news about',
    'who won', 'what happened', 'check online', 'search google'
]

CREATE_DOCUMENT_KEYWORDS = [
    'create document', 'create file', 'write document', 'write file',
    'make document', 'make file', 'save document', 'save file',
    'create a', 'write a', 'make a', 'save as',
    'shopping list', 'todo list', 'to-do list', 'to do list',
    'notes', 'recipe', 'list for', 'write me', 'create me'
]

# Gmail keywords
GMAIL_READ_KEYWORDS = [
    'email', 'emails', 'gmail', 'inbox', 'messages', 'mail',
    'check my email', 'read email', 'show email', 'recent email',
    'unread', 'new messages', 'latest messages', 'my inbox',
    'email from', 'message from', 'email about'
]

GMAIL_SEND_KEYWORDS = [
    'send email', 'email to', 'write email', 'compose email',
    'send message', 'send a message', 'email', 'mail to',
    'send to', 'message to', 'draft email'
]

BROWSE_KEYWORDS = [
    'browse', 'open url', 'open website', 'visit website', 'visit url',
    'go to', 'navigate to', 'fetch', 'read this page', 'open this',
    'visit this', 'load this page', 'show me this website', 'http://', 'https://',
    'www.', '.com', '.org', '.net', 'summarize this page', 'what does this say'
]


MUSIC_PLAY_KEYWORDS = [
    'play', 'play music', 'play song', 'play some', 'put on', 'listen to',
    'start playing', 'i want to hear', 'can you play'
]

MUSIC_CONTROL_KEYWORDS = [
    'pause', 'stop music', 'skip', 'next song', 'previous song', 'volume',
    'resume', 'unpause', 'mute', 'unmute', 'louder', 'quieter',
    'next track', 'previous track', 'turn up', 'turn down', 'stop playing'
]

# Words that indicate this is NOT a tool request
NO_TOOL_KEYWORDS = [
    'hello', 'hi ', 'hey', 'good morning', 'good afternoon', 'good evening',
    'how are you', 'whats up', 'tell me a joke', 'who are you', 'what are you',
    'thank you', 'thanks', 'goodbye', 'bye', 'see you'
]

# Basic color presets
COLOR_MAP = {
    "red": {"hue": 0, "sat": 254},
    "orange": {"hue": 5000, "sat": 254},
    "yellow": {"hue": 12750, "sat": 254},
    "green": {"hue": 25500, "sat": 254},
    "cyan": {"hue": 30000, "sat": 254},
    "blue": {"hue": 46920, "sat": 254},
    "purple": {"hue": 50000, "sat": 254},
    "pink": {"hue": 56100, "sat": 254},
    "white": {"hue": 0, "sat": 0},
    "warm white": {"hue": 8000, "sat": 140, "ct": 400},
    "cool white": {"hue": 40000, "sat": 50, "ct": 200},
}

# MOOD/SCENE PRESETS [MOOD]
MOOD_PRESETS = {
    # === NATURE ===
    "moonlight": {
        "description": "Cool, dim blues and silvers like moonlight",
        "settings": [
            {"hue": 46920, "sat": 200, "bri": 80},
            {"hue": 46000, "sat": 150, "bri": 100},
            {"hue": 48000, "sat": 100, "bri": 60},
            {"ct": 200, "bri": 70},
        ]
    },
    "sunset": {
        "description": "Warm oranges, reds, and purples like a sunset",
        "settings": [
            {"hue": 5000, "sat": 254, "bri": 200},
            {"hue": 1000, "sat": 254, "bri": 180},
            {"hue": 50000, "sat": 200, "bri": 150},
            {"hue": 0, "sat": 254, "bri": 160},
        ]
    },
    "sunrise": {
        "description": "Gradual warm colors like sunrise",
        "settings": [
            {"hue": 8000, "sat": 200, "bri": 100},
            {"hue": 6000, "sat": 240, "bri": 150},
            {"hue": 5000, "sat": 254, "bri": 180},
            {"ct": 350, "bri": 200},
        ]
    },
    "ocean": {
        "description": "Deep blues and teals like the ocean",
        "settings": [
            {"hue": 46920, "sat": 254, "bri": 180},
            {"hue": 44000, "sat": 220, "bri": 200},
            {"hue": 35000, "sat": 240, "bri": 190},
            {"hue": 30000, "sat": 200, "bri": 170},
        ]
    },
    "forest": {
        "description": "Various greens like a forest",
        "settings": [
            {"hue": 25500, "sat": 254, "bri": 180},
            {"hue": 27000, "sat": 230, "bri": 160},
            {"hue": 24000, "sat": 200, "bri": 170},
            {"hue": 26000, "sat": 180, "bri": 150},
        ]
    },
    "tropical": {
        "description": "Vibrant greens and blues like a tropical paradise",
        "settings": [
            {"hue": 30000, "sat": 254, "bri": 200},
            {"hue": 25500, "sat": 254, "bri": 210},
            {"hue": 35000, "sat": 240, "bri": 190},
            {"hue": 28000, "sat": 230, "bri": 200},
        ]
    },
    "arctic": {
        "description": "Icy blues and whites like the arctic",
        "settings": [
            {"hue": 46920, "sat": 180, "bri": 200},
            {"ct": 150, "bri": 220},
            {"hue": 48000, "sat": 120, "bri": 210},
            {"ct": 180, "bri": 200},
        ]
    },
    "galaxy": {
        "description": "Deep purples and blues like outer space",
        "settings": [
            {"hue": 50000, "sat": 254, "bri": 150},
            {"hue": 48000, "sat": 230, "bri": 130},
            {"hue": 46920, "sat": 254, "bri": 120},
            {"hue": 52000, "sat": 240, "bri": 140},
        ]
    },
    "aurora": {
        "description": "Northern lights - greens and purples dancing",
        "settings": [
            {"hue": 25500, "sat": 254, "bri": 180},
            {"hue": 50000, "sat": 254, "bri": 160},
            {"hue": 35000, "sat": 240, "bri": 170},
            {"hue": 46920, "sat": 200, "bri": 150},
        ]
    },
    "thunderstorm": {
        "description": "Dark blues with flashes of white",
        "settings": [
            {"hue": 46920, "sat": 254, "bri": 40},
            {"hue": 48000, "sat": 200, "bri": 30},
            {"hue": 44000, "sat": 220, "bri": 35},
        ]
    },
    "beach": {
        "description": "Sandy yellows and ocean blues",
        "settings": [
            {"hue": 10000, "sat": 180, "bri": 200},
            {"hue": 35000, "sat": 200, "bri": 180},
            {"hue": 8000, "sat": 160, "bri": 210},
            {"hue": 40000, "sat": 180, "bri": 190},
        ]
    },
    "desert": {
        "description": "Warm sandy oranges and dusty browns",
        "settings": [
            {"hue": 6000, "sat": 200, "bri": 200},
            {"hue": 5000, "sat": 220, "bri": 180},
            {"hue": 8000, "sat": 180, "bri": 190},
            {"hue": 4000, "sat": 240, "bri": 170},
        ]
    },
    "rainforest": {
        "description": "Lush deep greens with hints of mist",
        "settings": [
            {"hue": 25500, "sat": 254, "bri": 140},
            {"hue": 27000, "sat": 230, "bri": 130},
            {"hue": 23000, "sat": 210, "bri": 150},
            {"hue": 30000, "sat": 180, "bri": 160},
        ]
    },
    
    # === ACTIVITIES ===
    "focus": {
        "description": "Bright, cool white for concentration",
        "settings": [
            {"ct": 200, "bri": 254},
            {"ct": 210, "bri": 240},
            {"ct": 220, "bri": 250},
        ]
    },
    "relax": {
        "description": "Warm, dim lighting for relaxation",
        "settings": [
            {"ct": 400, "bri": 120},
            {"ct": 420, "bri": 110},
            {"ct": 390, "bri": 130},
            {"ct": 410, "bri": 115},
        ]
    },
    "energize": {
        "description": "Very bright white to wake you up",
        "settings": [
            {"ct": 250, "bri": 254},
            {"ct": 240, "bri": 254},
            {"ct": 260, "bri": 254},
        ]
    },
    "reading": {
        "description": "Warm but bright for comfortable reading",
        "settings": [
            {"ct": 350, "bri": 254},
            {"ct": 340, "bri": 245},
            {"ct": 360, "bri": 250},
        ]
    },
    "movie": {
        "description": "Very dim for watching movies",
        "settings": [
            {"hue": 46920, "sat": 200, "bri": 30},
            {"hue": 0, "sat": 0, "bri": 20},
            {"hue": 0, "sat": 0, "bri": 25},
        ]
    },
    "gaming": {
        "description": "Vibrant colors for immersive gaming",
        "settings": [
            {"hue": 50000, "sat": 254, "bri": 180},
            {"hue": 0, "sat": 254, "bri": 160},
            {"hue": 25500, "sat": 254, "bri": 170},
            {"hue": 46920, "sat": 254, "bri": 175},
        ]
    },
    "workout": {
        "description": "High energy reds and oranges",
        "settings": [
            {"hue": 0, "sat": 254, "bri": 254},
            {"hue": 5000, "sat": 254, "bri": 240},
            {"hue": 2000, "sat": 254, "bri": 250},
        ]
    },
    "yoga": {
        "description": "Calm purples and soft blues for meditation",
        "settings": [
            {"hue": 50000, "sat": 150, "bri": 100},
            {"hue": 46920, "sat": 120, "bri": 90},
            {"hue": 48000, "sat": 140, "bri": 95},
        ]
    },
    "meditation": {
        "description": "Very dim, peaceful lighting",
        "settings": [
            {"hue": 46920, "sat": 100, "bri": 40},
            {"hue": 50000, "sat": 80, "bri": 35},
            {"ct": 450, "bri": 30},
        ]
    },
    "cooking": {
        "description": "Bright warm white for the kitchen",
        "settings": [
            {"ct": 300, "bri": 254},
            {"ct": 310, "bri": 250},
            {"ct": 290, "bri": 254},
        ]
    },
    "dinner": {
        "description": "Warm candlelight ambiance for dining",
        "settings": [
            {"hue": 6000, "sat": 200, "bri": 140},
            {"hue": 5500, "sat": 220, "bri": 130},
            {"hue": 6500, "sat": 180, "bri": 150},
        ]
    },
    "sleep": {
        "description": "Very dim red to help you fall asleep",
        "settings": [
            {"hue": 0, "sat": 254, "bri": 20},
            {"hue": 1000, "sat": 254, "bri": 15},
            {"hue": 500, "sat": 254, "bri": 18},
        ]
    },
    "wakeup": {
        "description": "Gradually brightening warm light",
        "settings": [
            {"ct": 400, "bri": 150},
            {"ct": 350, "bri": 200},
            {"ct": 300, "bri": 254},
        ]
    },
    
    # === MOODS ===
    "romance": {
        "description": "Soft reds and pinks, dim and intimate",
        "settings": [
            {"hue": 56100, "sat": 220, "bri": 100},
            {"hue": 0, "sat": 200, "bri": 80},
            {"hue": 1000, "sat": 180, "bri": 90},
            {"hue": 56500, "sat": 200, "bri": 85},
        ]
    },
    "party": {
        "description": "Bright, vibrant colors for celebration",
        "settings": [
            {"hue": 0, "sat": 254, "bri": 254},
            {"hue": 46920, "sat": 254, "bri": 254},
            {"hue": 25500, "sat": 254, "bri": 254},
            {"hue": 12750, "sat": 254, "bri": 254},
            {"hue": 50000, "sat": 254, "bri": 254},
        ]
    },
    "cozy": {
        "description": "Warm amber glow like firelight",
        "settings": [
            {"ct": 450, "bri": 140},
            {"ct": 470, "bri": 130},
            {"hue": 6000, "sat": 200, "bri": 150},
        ]
    },
    "fireplace": {
        "description": "Flickering oranges and reds like a fire",
        "settings": [
            {"hue": 5000, "sat": 254, "bri": 180},
            {"hue": 3000, "sat": 254, "bri": 160},
            {"hue": 1000, "sat": 254, "bri": 170},
            {"hue": 6000, "sat": 240, "bri": 190},
        ]
    },
    "candle": {
        "description": "Soft flickering candlelight",
        "settings": [
            {"hue": 6000, "sat": 254, "bri": 100},
            {"hue": 5500, "sat": 254, "bri": 90},
            {"hue": 6500, "sat": 240, "bri": 95},
        ]
    },
    "zen": {
        "description": "Minimalist calm with soft whites",
        "settings": [
            {"ct": 350, "bri": 100},
            {"ct": 360, "bri": 95},
            {"ct": 340, "bri": 105},
        ]
    },
    "spa": {
        "description": "Relaxing blues and soft greens",
        "settings": [
            {"hue": 35000, "sat": 150, "bri": 130},
            {"hue": 46920, "sat": 120, "bri": 140},
            {"hue": 38000, "sat": 140, "bri": 135},
        ]
    },
    "club": {
        "description": "Intense dance club vibes",
        "settings": [
            {"hue": 50000, "sat": 254, "bri": 254},
            {"hue": 0, "sat": 254, "bri": 254},
            {"hue": 46920, "sat": 254, "bri": 254},
        ]
    },
    "disco": {
        "description": "Retro disco colors",
        "settings": [
            {"hue": 50000, "sat": 254, "bri": 220},
            {"hue": 10000, "sat": 254, "bri": 230},
            {"hue": 35000, "sat": 254, "bri": 210},
            {"hue": 0, "sat": 254, "bri": 225},
        ]
    },
    "concert": {
        "description": "Stage lighting intensity",
        "settings": [
            {"hue": 0, "sat": 254, "bri": 254},
            {"hue": 46920, "sat": 254, "bri": 254},
            {"hue": 25500, "sat": 254, "bri": 254},
        ]
    },
    "chill": {
        "description": "Laid back cool tones",
        "settings": [
            {"hue": 46920, "sat": 150, "bri": 150},
            {"hue": 48000, "sat": 130, "bri": 140},
            {"ct": 300, "bri": 160},
        ]
    },
    "warm": {
        "description": "Comfortable warm white",
        "settings": [
            {"ct": 400, "bri": 200},
            {"ct": 420, "bri": 190},
            {"ct": 380, "bri": 210},
        ]
    },
    "cool": {
        "description": "Crisp cool white",
        "settings": [
            {"ct": 200, "bri": 220},
            {"ct": 180, "bri": 230},
            {"ct": 220, "bri": 210},
        ]
    },
    "bright": {
        "description": "Maximum brightness",
        "settings": [
            {"ct": 250, "bri": 254},
            {"ct": 250, "bri": 254},
            {"ct": 250, "bri": 254},
        ]
    },
    "dim": {
        "description": "Very low light",
        "settings": [
            {"ct": 400, "bri": 50},
            {"ct": 400, "bri": 45},
            {"ct": 400, "bri": 55},
        ]
    },
    "night": {
        "description": "Minimal nightlight",
        "settings": [
            {"hue": 0, "sat": 254, "bri": 10},
            {"hue": 5000, "sat": 200, "bri": 15},
        ]
    },
    "natural": {
        "description": "Daylight simulation",
        "settings": [
            {"ct": 250, "bri": 254},
            {"ct": 260, "bri": 250},
            {"ct": 240, "bri": 254},
        ]
    },
    
    # === HOLIDAYS ===
    "christmas": {
        "description": "Red and green holiday cheer",
        "settings": [
            {"hue": 0, "sat": 254, "bri": 200},
            {"hue": 25500, "sat": 254, "bri": 200},
            {"hue": 0, "sat": 254, "bri": 200},
            {"hue": 25500, "sat": 254, "bri": 200},
        ]
    },
    "halloween": {
        "description": "Spooky oranges and purples",
        "settings": [
            {"hue": 5000, "sat": 254, "bri": 180},
            {"hue": 50000, "sat": 254, "bri": 140},
            {"hue": 5500, "sat": 254, "bri": 170},
        ]
    },
    "valentines": {
        "description": "Romantic reds and pinks",
        "settings": [
            {"hue": 0, "sat": 254, "bri": 160},
            {"hue": 56100, "sat": 220, "bri": 150},
            {"hue": 1000, "sat": 254, "bri": 155},
        ]
    },
    "easter": {
        "description": "Soft pastels",
        "settings": [
            {"hue": 56100, "sat": 100, "bri": 200},
            {"hue": 35000, "sat": 80, "bri": 210},
            {"hue": 10000, "sat": 90, "bri": 205},
            {"hue": 50000, "sat": 85, "bri": 195},
        ]
    },
    "july4": {
        "description": "Red, white, and blue patriotic",
        "settings": [
            {"hue": 0, "sat": 254, "bri": 200},
            {"ct": 200, "bri": 254},
            {"hue": 46920, "sat": 254, "bri": 200},
        ]
    },
    "stpatricks": {
        "description": "Irish green",
        "settings": [
            {"hue": 25500, "sat": 254, "bri": 200},
            {"hue": 26000, "sat": 240, "bri": 210},
            {"hue": 24500, "sat": 254, "bri": 195},
        ]
    },
    "hanukkah": {
        "description": "Blue and white celebration",
        "settings": [
            {"hue": 46920, "sat": 254, "bri": 200},
            {"ct": 200, "bri": 220},
            {"hue": 46920, "sat": 254, "bri": 200},
        ]
    },
    "newyear": {
        "description": "Sparkling gold and silver",
        "settings": [
            {"hue": 8000, "sat": 200, "bri": 220},
            {"ct": 180, "bri": 230},
            {"hue": 9000, "sat": 180, "bri": 210},
        ]
    },
    
    # === COLORS ===
    "red": {
        "description": "Pure red",
        "settings": [{"hue": 0, "sat": 254, "bri": 200}]
    },
    "blue": {
        "description": "Pure blue",
        "settings": [{"hue": 46920, "sat": 254, "bri": 200}]
    },
    "green": {
        "description": "Pure green",
        "settings": [{"hue": 25500, "sat": 254, "bri": 200}]
    },
    "purple": {
        "description": "Pure purple",
        "settings": [{"hue": 50000, "sat": 254, "bri": 200}]
    },
    "orange": {
        "description": "Pure orange",
        "settings": [{"hue": 5000, "sat": 254, "bri": 200}]
    },
    "yellow": {
        "description": "Pure yellow",
        "settings": [{"hue": 10000, "sat": 254, "bri": 200}]
    },
    "pink": {
        "description": "Soft pink",
        "settings": [{"hue": 56100, "sat": 200, "bri": 200}]
    },
    "cyan": {
        "description": "Cyan/Teal",
        "settings": [{"hue": 35000, "sat": 254, "bri": 200}]
    },
    "white": {
        "description": "Pure white",
        "settings": [{"ct": 250, "bri": 254}]
    },
    "rainbow": {
        "description": "Full spectrum colors",
        "settings": [
            {"hue": 0, "sat": 254, "bri": 200},
            {"hue": 10000, "sat": 254, "bri": 200},
            {"hue": 25500, "sat": 254, "bri": 200},
            {"hue": 35000, "sat": 254, "bri": 200},
            {"hue": 46920, "sat": 254, "bri": 200},
            {"hue": 50000, "sat": 254, "bri": 200},
        ]
    },
}


# ===== MUSIC FUNCTIONS =====

def init_youtube_music():
    """Initialize YouTube Music API."""
    global YOUTUBE_MUSIC_BROWSER
    if YOUTUBE_MUSIC_BROWSER is None:
        try:
            from ytmusicapi import YTMusic
            YOUTUBE_MUSIC_BROWSER = YTMusic()
            print("[OK] YouTube Music initialized")
            return True
        except ImportError:
            print("[WARN]  ytmusicapi not installed. Install with: pip install ytmusicapi")
            return False
        except Exception as e:
            print(f"[WARN]  Error initializing YouTube Music: {e}")
            return False
    return True


def search_youtube_music(query: str, limit: int = 5) -> List[Dict]:
    """Search for songs on YouTube Music."""
    if not init_youtube_music():
        return []

    try:
        results = YOUTUBE_MUSIC_BROWSER.search(query, filter="songs", limit=limit)
        return results
    except Exception as e:
        print(f"   [ERROR] Error searching YouTube Music: {e}")
        return []


def get_music_mood(query: str, song_info: dict = None) -> str:
    """Determine appropriate light mood based on music query."""
    query_lower = query.lower()

    # Genre/vibe detection
    if any(word in query_lower for word in ['relax', 'calm', 'chill', 'ambient', 'peaceful', 'meditation', 'sleep', 'quiet']):
        return 'relax'
    elif any(word in query_lower for word in ['party', 'dance', 'edm', 'club', 'rave', 'celebration', 'upbeat', 'fun']):
        return 'party'
    elif any(word in query_lower for word in ['romantic', 'love', 'ballad', 'slow dance', 'valentine', 'intimate']):
        return 'romance'
    elif any(word in query_lower for word in ['energize', 'workout', 'pump up', 'hype', 'rock', 'metal', 'hard', 'intense']):
        return 'energize'
    elif any(word in query_lower for word in ['jazz', 'lounge', 'smooth', 'sophisticated', 'cool', 'mellow']):
        return 'moonlight'
    elif any(word in query_lower for word in ['tropical', 'beach', 'island', 'reggae', 'caribbean', 'summer']):
        return 'tropical'
    elif any(word in query_lower for word in ['blues', 'soul', 'moody', 'melancholy', 'sad']):
        return 'ocean'
    elif any(word in query_lower for word in ['classical', 'orchestra', 'symphony', 'piano', 'study', 'concentrate']):
        return 'focus'
    elif any(word in query_lower for word in ['sunset', 'golden hour', 'evening', 'dusk']):
        return 'sunset'
    elif any(word in query_lower for word in ['fire', 'cozy', 'warm', 'acoustic', 'folk']):
        return 'fireplace'
    elif any(word in query_lower for word in ['space', 'cosmic', 'stars', 'galaxy', 'ambient', 'electronic']):
        return 'galaxy'
    elif any(word in query_lower for word in ['forest', 'nature', 'green', 'earth', 'natural']):
        return 'forest'
    elif any(word in query_lower for word in ['arctic', 'ice', 'winter', 'frozen', 'cold']):
        return 'arctic'
    elif any(word in query_lower for word in ['sunrise', 'morning', 'dawn', 'wake up']):
        return 'sunrise'
    else:
        # Default party mood for general music
        return 'party'


def play_music(query: str, service: str = "youtube_music") -> str:
    """
    Play music based on query and automatically sync lights.

    Args:
        query: Song name, artist, or search query
        service: "youtube_music" or "amazon_music"
    """
    print(f"   [MUSIC] Playing music: '{query}' on {service}")

    if service == "youtube_music":
        # Search for the song
        results = search_youtube_music(query, limit=1)

        if not results:
            return f"Couldn't find any songs matching '{query}' on YouTube Music"

        # Get the first result
        song = results[0]
        song_title = song.get('title', 'Unknown')
        artists = song.get('artists', [])
        artist_names = ", ".join([a.get('name', '') for a in artists]) if artists else "Unknown Artist"
        video_id = song.get('videoId', '')

        if not video_id:
            return f"Found '{song_title}' by {artist_names}, but couldn't get playback URL"

        # Construct YouTube Music URL
        url = f"https://music.youtube.com/watch?v={video_id}"

        # **NEW: Sync lights with music vibe**
        light_sync_msg = ""
        if BRIDGE_IP and HUE_USERNAME:
            try:
                mood = get_music_mood(query, song)
                print(f"   [SYNC] Syncing lights to '{mood}' mood for this music")
                light_result = apply_mood_to_lights(mood)
                print(f"   [LIGHT] {light_result}")
                light_sync_msg = f"\n💡 Lights set to '{mood}' mood"
            except Exception as e:
                print(f"   [WARN] Couldn't sync lights: {e}")

        # Open in browser
        try:
            webbrowser.open(url)
            return f"🎵 Now playing: '{song_title}' by {artist_names}{light_sync_msg}"
        except Exception as e:
            return f"Found '{song_title}' by {artist_names}, but couldn't open browser: {str(e)}\nURL: {url}"

    elif service == "amazon_music":
        # Amazon Music web search URL
        search_url = f"https://music.amazon.com/search/{requests.utils.quote(query)}"

        # Sync lights for Amazon Music too
        light_sync_msg = ""
        if BRIDGE_IP and HUE_USERNAME:
            try:
                mood = get_music_mood(query)
                apply_mood_to_lights(mood)
                light_sync_msg = f"\n💡 Lights synced to '{mood}' mood"
            except Exception:
                pass

        try:
            webbrowser.open(search_url)
            return f"🎵 Opening Amazon Music search for '{query}'{light_sync_msg}"
        except Exception as e:
            return f"Couldn't open Amazon Music: {str(e)}"

    else:
        return f"Unknown music service: {service}"


def search_music_info(query: str) -> str:
    """Search for music and return info without playing."""
    print(f"   [SEARCH] Searching for music info: '{query}'")

    results = search_youtube_music(query, limit=5)

    if not results:
        return f"Couldn't find any songs matching '{query}'"

    # Format results
    formatted_results = []
    for i, song in enumerate(results, 1):
        title = song.get('title', 'Unknown')
        artists = song.get('artists', [])
        artist_names = ", ".join([a.get('name', '') for a in artists]) if artists else "Unknown"
        album = song.get('album', {}).get('name', '') if song.get('album') else ''
        duration = song.get('duration', '')

        result_str = f"{i}. '{title}' by {artist_names}"
        if album:
            result_str += f" (Album: {album})"
        if duration:
            result_str += f" - {duration}"

        formatted_results.append(result_str)

    return "[MUSIC] Found these songs:\n\n" + "\n".join(formatted_results)


def control_music(action: str) -> str:
    """
    Control music playback using SYSTEM-WIDE media keys.
    Works from ANY window - no need to focus YouTube Music!

    Args:
        action: Control action - "pause", "resume", "next", "previous", "volume_up", "volume_down"
    """
    print(f"   [MUSIC] Controlling music: {action}")

    try:
        import pyautogui
    except ImportError:
        return "Music control requires pyautogui. Install with: pip install pyautogui"

    action_lower = action.lower()

    # FIXED: Use system-wide media keys instead of application-specific shortcuts
    # These work regardless of which window has focus!

    if action_lower in ["pause", "resume", "play_pause"]:
        # Use the system media play/pause key
        try:
            pyautogui.press('playpause')
            return "🎵 Toggled play/pause"
        except Exception:
            # Fallback for systems that don't recognize 'playpause'
            try:
                pyautogui.press('play')
                return "🎵 Toggled play/pause"
            except Exception:
                return "⚠️ Media key not supported on this system"

    elif action_lower == "next":
        # Use the system next track media key
        try:
            pyautogui.press('nexttrack')
            return "🎵 Skipped to next track"
        except Exception:
            return "⚠️ Next track key not supported on this system"

    elif action_lower == "previous":
        # Use the system previous track media key
        try:
            pyautogui.press('prevtrack')
            return "🎵 Went to previous track"
        except Exception:
            return "⚠️ Previous track key not supported on this system"

    elif action_lower == "volume_up":
        # Use system volume up key
        try:
            pyautogui.press('volumeup')
            return "🎵 Volume increased"
        except Exception:
            return "⚠️ Volume key not supported on this system"

    elif action_lower == "volume_down":
        # Use system volume down key
        try:
            pyautogui.press('volumedown')
            return "🎵 Volume decreased"
        except Exception:
            return "⚠️ Volume key not supported on this system"

    elif action_lower == "mute":
        # Use system mute key
        try:
            pyautogui.press('volumemute')
            return "🎵 Toggled mute"
        except Exception:
            return "⚠️ Mute key not supported on this system"

    else:
        return f"Unknown music control action: {action}. Available: pause, resume, next, previous, volume_up, volume_down, mute"


# ===== MUSIC VISUALIZER (ADVANCED FEATURE) =====

# Global variable to track visualizer state
_visualizer_active = False
_visualizer_thread = None

# Global variable to store images that need to be shown to vision model
# ================================================================================
# VISION IMAGE QUEUE SYSTEM (Improved)
# ================================================================================
from dataclasses import dataclass
from typing import Set

@dataclass
class ImageInfo:
    """Information about an image to be shown to the vision model."""
    filename: str
    filepath: str
    hash: str
    is_camera_capture: bool
    added_at: str
    is_ambient: bool = False   # kid "look" frame: warm/brief prompt, no face-rec dump

class VisionImageQueue:
    """
    Manages the queue of images to be shown to the vision model.

    IMPROVEMENTS:
    - Separates NEW images from OLD images
    - Tracks which images have been viewed
    - Prevents showing the same image multiple times
    - Purges old camera captures from conversation context
    """

    def __init__(self):
        self.pending_images: List[ImageInfo] = []
        self.viewed_images: Set[str] = set()

    def clear(self):
        """Clear all pending images."""
        print(f"   [VISION-QUEUE] Clearing {len(self.pending_images)} pending images")
        self.pending_images = []

    def add_image(self, filepath: str, filename: str, is_camera: bool = False,
                  is_ambient: bool = False, force: bool = False):
        """Add an image to the queue to be shown. `force` bypasses the
        already-viewed dedup (a kid "look" must always present the frame, even
        if the scene is identical to a previous one)."""
        import hashlib
        hash_md5 = hashlib.md5()
        with open(filepath, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        img_hash = hash_md5.hexdigest()

        if is_camera:
            # Clear old camera images
            self.pending_images = [img for img in self.pending_images
                                  if not img.is_camera_capture]
            print(f"   [VISION-QUEUE] New camera image, cleared old camera images")

        if force or img_hash not in self.viewed_images:
            import datetime
            self.pending_images.append(ImageInfo(
                filename=filename,
                filepath=filepath,
                hash=img_hash,
                is_camera_capture=is_camera,
                added_at=datetime.datetime.now().isoformat(),
                is_ambient=is_ambient
            ))
            print(f"   [VISION-QUEUE] Added {filename} (hash: {img_hash[:8]})")
        else:
            print(f"   [VISION-QUEUE] Skipped {filename} - already viewed")

    def mark_as_viewed(self):
        """Mark all current pending images as viewed."""
        for img in self.pending_images:
            self.viewed_images.add(img.hash)
        print(f"   [VISION-QUEUE] Marked {len(self.pending_images)} images as viewed")

    def has_images(self) -> bool:
        """Check if there are pending images."""
        return len(self.pending_images) > 0

_vision_queue = VisionImageQueue()

# Track image paths from the last vision injection so we can save the LLM's description
_last_vision_image_paths = []

# Sticky reference to the most recently shown image, so chat follow-ups
# ("what color is it?") can reuse it for a short window without re-uploading.
# Kept separate from _last_vision_image_paths, which the observation logger
# clears after each turn.
_recent_image_paths = []
_recent_image_at = 0.0

_IMAGE_FOLLOWUP_CUES = (
    "it", "this", "that", "the image", "the photo", "the picture", "the pic",
    "color", "colour", "wearing", "background", "behind", "who is", "who's",
    "what is", "what's", "how many", "describe", "zoom", "closer", "look again",
    "the person", "the man", "the woman", "the kid", "the dog", "the cat",
    "in the", "on the", "left", "right", "foreground", "the shirt", "the sign",
    "again", "same image", "same photo", "same picture",
)


def _refers_to_recent_image(text: str) -> bool:
    """Heuristic: is the user still asking about the image they just shared?
    Triggers on a referential/visual cue, or a very short follow-up."""
    t = (text or "").lower().strip()
    if not t:
        return False
    if any(c in t for c in _IMAGE_FOLLOWUP_CUES):
        return True
    return len(t.split()) <= 5


def _save_visual_observation(description: str):
    """Save the LLM's image description as a visual memory observation."""
    global _last_vision_image_paths
    if not _last_vision_image_paths or not description or not VISUAL_MEMORY_AVAILABLE:
        _last_vision_image_paths = []
        return

    try:
        vm = get_visual_memory()
        image_path = _last_vision_image_paths[0]

        # Compute image hash
        import hashlib
        img_hash = None
        try:
            h = hashlib.md5()
            with open(image_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    h.update(chunk)
            img_hash = h.hexdigest()
        except Exception:
            pass

        # Extract known people mentioned in the description. Match the full
        # name OR the first name, with word boundaries — "Stella Andoff" must
        # be spotted when the model writes just "Stella", and a short name
        # must not fire inside another word ("Sam" in "samples").
        known_people = vm.get_all_people()
        desc_lower = description.lower()
        people_present = []
        for p in known_people:
            full = (p.get('name') or '').strip()
            if not full:
                continue
            first = full.split()[0]
            candidates = {full} | ({first} if len(first) >= 3 else set())
            if any(re.search(r'\b' + re.escape(nm) + r'\b', description, re.I)
                   for nm in candidates):
                people_present.append(full)

        # Extract location from description (word-boundary match)
        location = None
        known_places = vm.get_all_places()
        for place in known_places:
            pn = (place.get('name') or '').strip()
            if pn and re.search(r'\b' + re.escape(pn) + r'\b', description, re.I):
                location = pn
                break
        # Fallback location keywords
        if not location:
            loc_keywords = {'office': 'Office', 'kitchen': 'Kitchen', 'living room': 'Living Room',
                           'studio': 'Studio', 'bedroom': 'Bedroom', 'bathroom': 'Bathroom'}
            for kw, name in loc_keywords.items():
                if kw in desc_lower:
                    location = name
                    break

        # Truncate description if very long (keep first 1000 chars)
        scene_desc = description[:1000] if len(description) > 1000 else description

        vm.log_observation(
            scene_description=scene_desc,
            people_present=people_present if people_present else None,
            location=location,
            image_path=image_path,
            image_hash=img_hash
        )

        # Update "last seen" for detected people
        for name in people_present:
            vm.update_seen('person', name)
        if location:
            vm.update_seen('place', location)

        print(f"[VISUAL-MEMORY] Saved observation: {len(people_present)} people, location={location}")
    except Exception as e:
        print(f"[VISUAL-MEMORY] Error saving observation: {e}")
    finally:
        _last_vision_image_paths = []


def _visual_context_block(text: str, max_entities: int = 4) -> str:
    """Compact <visual_memory> block for people/places NAMED in `text`: who
    they are, when last seen through the camera, how often. Lets "have you
    seen Stella today?" be answered from memory without a fresh camera turn.
    Empty when no known entity is mentioned, so ordinary turns cost nothing."""
    t = (text or '').strip()
    if not t or not VISUAL_MEMORY_AVAILABLE:
        return ""
    try:
        vm = get_visual_memory()
        lines = []
        for kind, rows in (('person', vm.get_all_people() or []),
                           ('place', vm.get_all_places() or [])):
            for e in rows:
                name = (e.get('name') or '').strip()
                if not name:
                    continue
                first = name.split()[0]
                candidates = {name} | ({first} if len(first) >= 3 else set())
                if not any(re.search(r'\b' + re.escape(nm) + r'\b', t, re.I)
                           for nm in candidates):
                    continue
                bits = []
                rel = (e.get('relationship') or '').strip()
                if rel:
                    bits.append(rel)
                age = ""
                try:
                    if ENHANCED_MEMORY_AVAILABLE and memory_system and e.get('last_seen'):
                        age = memory_system._humanize_age(str(e['last_seen']))
                except Exception:
                    age = ""
                if age:
                    bits.append(f"last seen through your camera {age}")
                elif e.get('last_seen'):
                    bits.append(f"last seen {str(e['last_seen'])[:16]}")
                else:
                    bits.append("not yet seen through your camera")
                if e.get('times_seen'):
                    bits.append(f"seen {e['times_seen']} times")
                desc = (e.get('typical_appearance') or e.get('description') or '').strip()
                if desc:
                    bits.append(desc[:90])
                lines.append(f"- {name} ({kind}): " + "; ".join(bits))
                if len(lines) >= max_entities:
                    break
            if len(lines) >= max_entities:
                break
        if not lines:
            return ""
        return ("<visual_memory>\nWhat your camera memory knows about who/what was just "
                "mentioned (the times say when YOU last saw them on camera — someone can be "
                "home without you having seen them):\n"
                + "\n".join(lines) + "\n</visual_memory>")
    except Exception as e:
        print(f"[VISUAL-MEMORY] context block failed: {e}")
        return ""


def start_music_visualizer(duration_seconds: int = 300, style: str = "party") -> str:
    """
    Start a dynamic light show that changes colors rhythmically.
    Creates an atmospheric visualizer effect for music.

    Args:
        duration_seconds: How long to run (default 5 minutes)
        style: "party" (fast colorful), "chill" (slow smooth), "pulse" (rhythmic)
    """
    global _visualizer_active, _visualizer_thread

    if not BRIDGE_IP or not HUE_USERNAME:
        return "Hue lights not configured. Can't start visualizer."

    if _visualizer_active:
        return "Music visualizer is already running! Say 'stop visualizer' to turn it off first."

    print(f"   [SYNC] Starting {style} music visualizer for {duration_seconds} seconds")

    def visualizer_loop():
        global _visualizer_active
        lights = get_hue_lights()
        if not lights:
            _visualizer_active = False
            return

        light_ids = list(lights.keys())
        start_time = time.time()

        # Different color schemes based on style
        if style == "party":
            color_options = [
                {"hue": 0, "sat": 254, "bri": 254},      # Red
                {"hue": 46920, "sat": 254, "bri": 254},  # Blue
                {"hue": 25500, "sat": 254, "bri": 254},  # Green
                {"hue": 12750, "sat": 254, "bri": 254},  # Yellow
                {"hue": 50000, "sat": 254, "bri": 254},  # Purple
                {"hue": 56100, "sat": 254, "bri": 254},  # Pink
                {"hue": 30000, "sat": 254, "bri": 254},  # Cyan
                {"hue": 5000, "sat": 254, "bri": 254},   # Orange
            ]
            transition_time = 5
            change_interval = 1.5
        elif style == "chill":
            color_options = [
                {"hue": 46920, "sat": 200, "bri": 150},  # Soft blue
                {"hue": 50000, "sat": 180, "bri": 130},  # Soft purple
                {"hue": 30000, "sat": 190, "bri": 140},  # Soft cyan
                {"hue": 25500, "sat": 160, "bri": 120},  # Soft green
            ]
            transition_time = 20
            change_interval = 4.0
        elif style == "pulse":
            color_options = [
                {"hue": 0, "sat": 254, "bri": 254},      # Bright red
                {"hue": 0, "sat": 254, "bri": 100},      # Dim red
                {"hue": 46920, "sat": 254, "bri": 254},  # Bright blue
                {"hue": 46920, "sat": 254, "bri": 100},  # Dim blue
            ]
            transition_time = 3
            change_interval = 0.8
        else:
            color_options = [
                {"hue": 0, "sat": 254, "bri": 254},
                {"hue": 46920, "sat": 254, "bri": 254},
            ]
            transition_time = 5
            change_interval = 2.0

        try:
            while _visualizer_active and (time.time() - start_time) < duration_seconds:
                for light_id in light_ids:
                    # Random color for each light
                    color = random.choice(color_options).copy()
                    color["on"] = True
                    color["transitiontime"] = transition_time
                    control_hue_light(light_id, color)

                time.sleep(change_interval)
        finally:
            _visualizer_active = False
            print("   [SYNC] Music visualizer ended")

    # Start visualizer
    _visualizer_active = True
    _visualizer_thread = threading.Thread(target=visualizer_loop, daemon=True)
    _visualizer_thread.start()

    style_descriptions = {
        "party": "fast, vibrant colors",
        "chill": "slow, smooth transitions",
        "pulse": "rhythmic pulsing"
    }

    return f"🎨 Music visualizer started ({style_descriptions.get(style, 'dynamic')})! Lights will dance for {duration_seconds//60} minutes."


def stop_music_visualizer() -> str:
    """Stop the music visualizer."""
    global _visualizer_active

    if not _visualizer_active:
        return "No visualizer is currently running."

    _visualizer_active = False
    print("   [SYNC] Stopping music visualizer...")

    # Wait a moment for thread to finish
    time.sleep(1)

    return "🎨 Music visualizer stopped."


# Tool definitions with MOOD, DOCUMENT, MUSIC, and VISUALIZER support!
from blue.server.tool_schemas import RAW_TOOLS
TOOLS = list(RAW_TOOLS)


def _filter_unavailable_tools(tool_names: set[str], reason: str) -> None:
    """Remove schemas whose backing runtime module is not available."""
    global TOOLS
    before = len(TOOLS)
    TOOLS = [
        tool for tool in TOOLS
        if tool["function"]["name"] not in tool_names
    ]
    removed = before - len(TOOLS)
    if removed:
        names = ", ".join(sorted(tool_names))
        print(f"[INFO] Filtered {removed} {reason} tool(s): {names}")


# Filter out tools that require unavailable modules
if not ENHANCED_TOOLS_AVAILABLE:
    # Remove file operation tools that won't work
    _filter_unavailable_tools({
        "list_files", "read_file", "write_file", "get_file_info",
        "create_reminder", "get_upcoming_reminders", "cancel_reminder",
        "reschedule_reminder",
        "add_contact", "list_contacts", "find_contact",
        "check_timers", "get_system_info", "take_screenshot",
        "launch_application", "set_volume", "story_prompt", "educational_activity",
        "get_local_time", "get_sunrise_sunset"
    }, "unavailable enhanced")

if not ACADEMIC_ASSISTANT_AVAILABLE:
    _filter_unavailable_tools({
        "analyze_with_chat_theory",
        "prepare_lecture",
        "discussion_questions",
        "simulate_student_questions",
    }, "unavailable academic assistant")

if not PROACTIVE_ASSISTANCE_AVAILABLE:
    _filter_unavailable_tools({
        "check_proactive_suggestions",
    }, "unavailable proactive assistance")


# ===== DOCUMENT MANAGEMENT FUNCTIONS =====

# ================================================================================
# LIBRARY FOLDERS — hierarchical organization under DOCUMENTS_FOLDER
# ================================================================================
# Documents live in a folder tree (e.g. "Publications", "Courses/CS240",
# "Academic Texts/Sohn-Rethel"). A document's `folder` is its POSIX-style
# path relative to DOCUMENTS_FOLDER; "" means the library root. These helpers
# keep folder handling safe (no traversal outside DOCUMENTS_FOLDER) and in
# one place so the index, RAG, and GUI all agree on what a folder is.


def _safe_folder_segment(seg: str) -> str:
    """Sanitize a single folder name: drop path/illegal chars and leading or
    trailing dots (which would enable traversal), keep spaces and normal
    punctuation so 'Academic Texts' or 'Sohn-Rethel' survive intact."""
    seg = re.sub(r'[<>:"/\\|?*\x00-\x1f]', '', (seg or '')).strip().strip('.')
    return seg.strip()


def _safe_rel_folder(folder: str) -> str:
    """Normalize a (possibly user-supplied) relative folder path to a safe
    POSIX-style rel path under DOCUMENTS_FOLDER. Returns '' for the root and
    silently drops any '.'/'..' segments."""
    if not folder:
        return ""
    raw = str(folder).replace("\\", "/").strip().strip("/")
    parts = []
    for seg in raw.split("/"):
        s = _safe_folder_segment(seg)
        if s and s not in (".", ".."):
            parts.append(s)
    return "/".join(parts)


def _abs_library_path(rel_folder: str) -> str:
    """Absolute path for a library folder, guaranteed to stay under
    DOCUMENTS_FOLDER (defense in depth against traversal)."""
    base = os.path.abspath(DOCUMENTS_FOLDER)
    rel = _safe_rel_folder(rel_folder)
    full = os.path.abspath(os.path.join(base, *rel.split("/"))) if rel else base
    try:
        if os.path.commonpath([base, full]) != base:
            return base
    except ValueError:
        return base
    return full


def _folder_of_filepath(filepath: str) -> str:
    """The library folder (POSIX rel path, '' for root) a file sits in. Files
    outside DOCUMENTS_FOLDER (e.g. image uploads in UPLOAD_FOLDER) report ''."""
    try:
        base = os.path.abspath(DOCUMENTS_FOLDER)
        ap = os.path.abspath(filepath)
        if os.path.commonpath([base, ap]) != base:
            return ""
        rel = os.path.relpath(os.path.dirname(ap), base)
        if rel in (".", ""):
            return ""
        return rel.replace("\\", "/")
    except Exception:
        return ""


def list_library_folders() -> List[str]:
    """Every folder under DOCUMENTS_FOLDER as sorted POSIX rel paths (excludes
    the root ''). Hidden dirs are skipped."""
    base = os.path.abspath(DOCUMENTS_FOLDER)
    out: List[str] = []
    if not os.path.isdir(base):
        return out
    for root, dirs, _files in os.walk(base):
        dirs[:] = [d for d in dirs if not d.startswith('.')]
        for d in dirs:
            rel = os.path.relpath(os.path.join(root, d), base).replace("\\", "/")
            out.append(rel)
    return sorted(out)


def list_subfolders(rel_folder: str) -> List[str]:
    """Immediate child folder names of the given library folder."""
    full = _abs_library_path(rel_folder)
    if not os.path.isdir(full):
        return []
    try:
        return sorted(
            d for d in os.listdir(full)
            if os.path.isdir(os.path.join(full, d)) and not d.startswith('.')
        )
    except Exception:
        return []


def create_library_folder(rel_folder: str) -> dict:
    """Create a folder (and any missing parents) under DOCUMENTS_FOLDER."""
    rel = _safe_rel_folder(rel_folder)
    if not rel:
        return {"success": False, "error": "Invalid or empty folder name."}
    try:
        os.makedirs(_abs_library_path(rel), exist_ok=True)
        return {"success": True, "folder": rel}
    except Exception as e:
        return {"success": False, "error": str(e)}


def _folders_under(prefix: str) -> List[str]:
    """A folder plus all its descendants (POSIX rel paths). Used to scope a
    search to an expertise area and everything filed beneath it."""
    pfx = _safe_rel_folder(prefix)
    if not pfx:
        return []
    result = [pfx]
    needle = pfx + "/"
    for f in list_library_folders():
        if f.startswith(needle):
            result.append(f)
    return result


# The owner's name, used to recognize the folder that holds HIS publications
# when he says "my published work" / "my articles". Override via env if the
# folder is named differently.
BLUE_OWNER_NAME = os.environ.get("BLUE_OWNER_NAME", "Alex Levant")

# Phrases that mean "draw on my own published writing".
_OWN_WORK_PHRASES = (
    "my published work", "my publications", "my published", "my articles",
    "my article", "my papers", "my paper", "my writing", "my own work",
    "my own writing", "my research", "my work", "my book", "my chapter",
    "my essays", "my essay", "published work", "my scholarship",
)


def _infer_expertise_folders(query: str) -> List[str]:
    """Map a request to the library folder(s) it should be answered from, or
    [] for 'search the whole library'. Two cues:

      • "my published work / my articles / my publications …" → the folder
        holding the owner's own writing (matched by publications-style names
        or the owner's name, e.g. an "Alex Levant" folder).
      • An explicit folder name appearing in the request → that folder.

    Returns each matched folder plus its descendants, deduped.
    """
    q = (query or "").lower()
    folders = list_library_folders()
    if not folders:
        return []

    targets = set()
    owner_tokens = [t for t in re.findall(r"[a-z]+", BLUE_OWNER_NAME.lower()) if len(t) > 2]

    if any(p in q for p in _OWN_WORK_PHRASES):
        for f in folders:
            fl = f.lower()
            is_pub = any(k in fl for k in ("publication", "published", "papers", "articles", "writing"))
            is_owner_named = bool(owner_tokens) and all(t in fl for t in owner_tokens)
            if is_pub or is_owner_named:
                targets.update(_folders_under(f))

    # Explicit folder-name mention (match on the leaf name as a whole word).
    for f in folders:
        leaf = f.split("/")[-1].lower()
        if len(leaf) >= 4 and re.search(r"\b" + re.escape(leaf) + r"\b", q):
            targets.update(_folders_under(f))

    return sorted(targets)


def load_document_index() -> Dict:
    """Load the document index from disk.

    A corrupt index is QUARANTINED (renamed *.corrupt-<timestamp>) and the last
    good backup (kept by save_document_index) is restored. The old behavior —
    silently resetting to empty — turned one bad read into "all my documents
    are gone": the 2026-07-04 hard reboot NUL-filled the file (lost write
    cache), and the next load wiped the whole library listing.

    Backfills a `folder` field (computed from each file's location) on any
    entry that predates folder support, so the GUI and search always see one.
    """
    def _read(path) -> Dict:
        with open(path, 'r') as f:
            data = json.load(f)
        if not (isinstance(data, dict) and isinstance(data.get('documents'), list)):
            raise ValueError("unexpected shape")
        return data

    def _backfill(data: Dict) -> Dict:
        for d in data['documents']:
            if isinstance(d, dict) and 'folder' not in d:
                d['folder'] = _folder_of_filepath(d.get('filepath', ''))
        return data

    if os.path.exists(DOCUMENT_INDEX_FILE):
        try:
            return _backfill(_read(DOCUMENT_INDEX_FILE))
        except Exception as e:
            print(f"[INDEX] {DOCUMENT_INDEX_FILE} is corrupt ({e}); quarantining it and trying the backup.")
            try:
                os.replace(DOCUMENT_INDEX_FILE,
                           DOCUMENT_INDEX_FILE + time.strftime(".corrupt-%Y%m%d-%H%M%S"))
            except Exception as qe:
                print(f"[INDEX] could not quarantine: {qe}")
        try:
            if os.path.exists(DOCUMENT_INDEX_FILE + '.bak'):
                data = _read(DOCUMENT_INDEX_FILE + '.bak')
                save_document_index(data)     # re-materialize the main file from the backup
                print(f"[INDEX] restored {len(data['documents'])} entries from {DOCUMENT_INDEX_FILE}.bak")
                return _backfill(data)
        except Exception as e:
            print(f"[INDEX] backup also unusable ({e}); starting empty.")
    return {"documents": []}


def save_document_index(index: Dict):
    """Save the document index ATOMICALLY: write a temp file, fsync it, keep
    the previous good version as .bak, then rename into place. A plain
    open(..., 'w') left a window where a crash or hard reboot produced a
    truncated/NUL-filled file — which load then treated as corrupt."""
    tmp = DOCUMENT_INDEX_FILE + '.tmp'
    with open(tmp, 'w') as f:
        json.dump(index, f, indent=2)
        f.flush()
        os.fsync(f.fileno())
    if os.path.exists(DOCUMENT_INDEX_FILE):
        try:
            os.replace(DOCUMENT_INDEX_FILE, DOCUMENT_INDEX_FILE + '.bak')
        except Exception:
            pass
    os.replace(tmp, DOCUMENT_INDEX_FILE)


def allowed_file(filename):
    """Check if file extension is allowed."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def ensure_unique_path(directory: str, filename: str) -> str:
    """Ensure a unique file path by adding numbers if file already exists.

    Args:
        directory: Directory where file will be saved
        filename: Original filename

    Returns:
        Unique file path
    """
    base_path = os.path.join(directory, secure_filename(filename))

    # If file doesn't exist, return as-is
    if not os.path.exists(base_path):
        return base_path

    # File exists, add number suffix
    name, ext = os.path.splitext(filename)
    counter = 1

    while True:
        new_filename = f"{name}_{counter}{ext}"
        new_path = os.path.join(directory, secure_filename(new_filename))

        if not os.path.exists(new_path):
            return new_path

        counter += 1

        # Safety check to prevent infinite loop
        if counter > 9999:
            # Use timestamp as last resort
            import time
            timestamp = int(time.time())
            new_filename = f"{name}_{timestamp}{ext}"
            return os.path.join(directory, secure_filename(new_filename))


def get_file_hash(filepath):
    """Get MD5 hash of file for deduplication."""
    hash_md5 = hashlib.md5()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hash_md5.update(chunk)
    return hash_md5.hexdigest()


def _is_camera_capture_filename(filename: str) -> bool:
    """Camera frames use a fixed prefix; treat them as transient."""
    return filename.startswith('camera_') or filename.startswith('camera_NEW_')


def cleanup_document_index() -> dict:
    """Remove stale and noisy entries from the document index.

    Drops:
      - entries flagged camera_capture (or whose filename matches the
        camera_*/camera_NEW_* prefix); these should never have been indexed
        as documents
      - entries whose file no longer exists on disk
      - duplicate entries that share an MD5 hash (keeps the first occurrence)

    Returns a small report dict for logging.
    """
    index = load_document_index()
    documents = index.get('documents', [])
    seen_hashes = set()
    kept, dropped_missing, dropped_camera, dropped_dup = [], 0, 0, 0

    for doc in documents:
        filename = doc.get('filename', '')
        filepath = doc.get('filepath', '')
        file_hash = doc.get('hash')

        if doc.get('camera_capture') or _is_camera_capture_filename(filename):
            dropped_camera += 1
            continue
        if filepath and not os.path.exists(filepath):
            dropped_missing += 1
            continue
        if file_hash and file_hash in seen_hashes:
            dropped_dup += 1
            continue
        if file_hash:
            seen_hashes.add(file_hash)
        kept.append(doc)

    if dropped_missing or dropped_camera or dropped_dup:
        index['documents'] = kept
        save_document_index(index)

    return {
        'kept': len(kept),
        'dropped_missing': dropped_missing,
        'dropped_camera': dropped_camera,
        'dropped_duplicate': dropped_dup,
    }


def rescan_documents_folder() -> dict:
    """Add files present in DOCUMENTS_FOLDER but missing from the index.

    Useful after a manual file drop or when the index has been wiped. Skips
    camera-capture files and unsupported extensions. Files added here are
    also pushed into the ChromaDB index when available.
    """
    import datetime as _dt

    index = load_document_index()
    documents = index.get('documents', [])
    indexed_paths = {os.path.normcase(os.path.abspath(d.get('filepath', ''))) for d in documents}
    indexed_hashes = {d.get('hash') for d in documents if d.get('hash')}

    added = 0
    skipped = 0
    if not os.path.isdir(DOCUMENTS_FOLDER):
        return {'added': 0, 'skipped': 0}

    # Walk the whole folder tree so files in subfolders (Publications/,
    # Courses/CS240/, …) get picked up, each tagged with its folder.
    for root, dirs, files in os.walk(DOCUMENTS_FOLDER):
        dirs[:] = [d for d in dirs if not d.startswith('.')]
        for entry in files:
            full = os.path.join(root, entry)
            if not os.path.isfile(full):
                continue
            if _is_camera_capture_filename(entry):
                continue
            if not allowed_file(entry):
                skipped += 1
                continue
            if os.path.normcase(os.path.abspath(full)) in indexed_paths:
                continue

            try:
                file_hash = get_file_hash(full)
            except Exception:
                skipped += 1
                continue
            if file_hash in indexed_hashes:
                continue

            folder = _folder_of_filepath(full)
            size_mb = os.path.getsize(full) / (1024 * 1024)
            print(f"   [INDEX] extracting text: {entry} ({size_mb:.1f} MB) [{folder or 'root'}]", flush=True)
            text_content = extract_text_from_file(full)
            text_preview = text_content[:500] if not text_content.startswith('Error') else ''

            # Best-effort push into ChromaDB so the new file is searchable.
            try:
                print(f"   [INDEX] embedding into ChromaDB: {entry}", flush=True)
                from blue.tools.rag import index_document as _idx
                _idx(full, entry, doc_id=file_hash, text=text_content, folder=folder)
                print(f"   [INDEX] done: {entry}", flush=True)
            except Exception as e:
                print(f"   [INDEX] ChromaDB push failed for {entry}: {e}", flush=True)

            documents.append({
                'filename': entry,
                'filepath': str(full),
                'folder': folder,
                'size': os.path.getsize(full),
                'hash': file_hash,
                'uploaded_at': _dt.datetime.fromtimestamp(os.path.getmtime(full)).strftime('%Y-%m-%d %H:%M'),
                'text_preview': text_preview,
                'indexed_in_rag': True,
            })
            indexed_hashes.add(file_hash)
            added += 1

    if added:
        index['documents'] = documents
        save_document_index(index)

    return {'added': added, 'skipped': skipped}


def register_uploaded_file(filepath: str, filename: str) -> dict:
    """Common post-save bookkeeping for any upload endpoint.

    - Hash-dedups against the existing index (deletes the new copy if a
      duplicate is found and returns the existing filename).
    - Extracts a text preview.
    - Pushes the document into ChromaDB if available (tagged with its folder).
    - Appends to document_index.json with the folder it lives in.
    """
    import datetime as _dt

    file_hash = get_file_hash(filepath)
    folder = _folder_of_filepath(filepath)
    index = load_document_index()

    for existing in index.get('documents', []):
        if existing.get('hash') == file_hash:
            try:
                if os.path.abspath(existing.get('filepath', '')) != os.path.abspath(filepath):
                    os.remove(filepath)
            except OSError:
                pass
            return {
                'filename': existing.get('filename'),
                'duplicate': True,
                'existing_filename': existing.get('filename'),
            }

    text_content = extract_text_from_file(filepath)
    text_preview = text_content[:500] if not text_content.startswith('Error') else ''

    indexed_in_rag = False
    try:
        from blue.tools.rag import index_document as _idx
        result = _idx(filepath, filename, doc_id=file_hash, folder=folder)
        indexed_in_rag = bool(result.get('success'))
    except Exception:
        pass

    index.setdefault('documents', []).append({
        'filename': filename,
        'filepath': str(filepath),
        'folder': folder,
        'size': os.path.getsize(filepath),
        'hash': file_hash,
        'uploaded_at': _dt.datetime.now().strftime('%Y-%m-%d %H:%M'),
        'text_preview': text_preview,
        'indexed_in_rag': indexed_in_rag,
    })
    save_document_index(index)
    return {'filename': filename, 'duplicate': False, 'indexed_in_rag': indexed_in_rag, 'folder': folder}


def remove_document_from_index(filepath: str) -> bool:
    """Remove a deleted/moved file from the document index AND ChromaDB."""
    try:
        abs_path = os.path.abspath(filepath)
    except Exception:
        return False

    index = load_document_index()
    docs = index.get('documents', [])
    keep, removed_hashes, removed_names = [], [], []
    for d in docs:
        try:
            if os.path.abspath(d.get('filepath', '')) == abs_path:
                if d.get('hash'):
                    removed_hashes.append(d['hash'])
                if d.get('filename'):
                    removed_names.append(d['filename'])
                continue
        except Exception:
            pass
        keep.append(d)

    if not removed_hashes and not removed_names:
        return False

    index['documents'] = keep
    save_document_index(index)

    # Remove the corresponding chunks from ChromaDB.
    try:
        from blue.tools.rag import remove_document
        for h in removed_hashes:
            remove_document(h)
    except Exception:
        pass

    print(f"[WATCHER] Removed from index: {', '.join(removed_names)}")
    return True


# ---- Filesystem watcher for auto-indexing the documents folder ----
_DOCUMENT_WATCHER = None


def start_document_watcher():
    """Watch DOCUMENTS_FOLDER and auto-index files as they appear.

    Uses the `watchdog` library if available. Each create/move event triggers
    a debounced indexer (sleeps briefly so the OS finishes writing the file
    before we hash and chunk it). Deletes are reflected immediately by
    removing the entry from document_index.json and ChromaDB.

    Gracefully no-ops if watchdog isn't installed; the startup rescan is
    still a fallback for users without it.
    """
    global _DOCUMENT_WATCHER
    if _DOCUMENT_WATCHER is not None:
        return _DOCUMENT_WATCHER

    try:
        from watchdog.observers import Observer
        from watchdog.events import FileSystemEventHandler
    except ImportError:
        print("[WATCHER] watchdog not installed — auto-indexing disabled.")
        print("[WATCHER]   Install with: pip install watchdog")
        return None

    import threading
    import time

    def _index_after_settle(filepath: str, debounce_secs: float = 1.5):
        # Wait for the file to be fully written (Windows file locks can hold
        # past the on_created event by hundreds of ms).
        time.sleep(debounce_secs)
        try:
            if not os.path.isfile(filepath):
                return
            filename = os.path.basename(filepath)
            if _is_camera_capture_filename(filename):
                return
            if not allowed_file(filename):
                return
            # Wait until the file size stops changing (rough copy-completion
            # check). Up to a few seconds for large PDFs.
            last_size = -1
            for _ in range(10):
                try:
                    cur = os.path.getsize(filepath)
                except OSError:
                    return
                if cur == last_size and cur > 0:
                    break
                last_size = cur
                time.sleep(0.4)
            result = register_uploaded_file(filepath, filename)
            if result.get('duplicate'):
                print(f"[WATCHER] {filename}: duplicate, skipped")
            else:
                rag_marker = " (ChromaDB indexed)" if result.get('indexed_in_rag') else ""
                print(f"[WATCHER] Auto-indexed: {filename}{rag_marker}")
        except Exception as e:
            print(f"[WATCHER] Error auto-indexing {filepath}: {e}")

    class _DocChangeHandler(FileSystemEventHandler):
        def on_created(self, event):
            if event.is_directory:
                return
            threading.Thread(
                target=_index_after_settle,
                args=(event.src_path,),
                daemon=True,
            ).start()

        def on_moved(self, event):
            if event.is_directory:
                return
            # The destination is the new file location.
            threading.Thread(
                target=_index_after_settle,
                args=(event.dest_path,),
                daemon=True,
            ).start()
            # Also update the index for the old path if it was tracked.
            try:
                remove_document_from_index(event.src_path)
            except Exception:
                pass

        def on_deleted(self, event):
            if event.is_directory:
                return
            try:
                remove_document_from_index(event.src_path)
            except Exception:
                pass

    try:
        os.makedirs(DOCUMENTS_FOLDER, exist_ok=True)
        observer = Observer()
        observer.schedule(_DocChangeHandler(), DOCUMENTS_FOLDER, recursive=True)
        observer.daemon = True
        observer.start()
        _DOCUMENT_WATCHER = observer
        print(f"[WATCHER] Watching {DOCUMENTS_FOLDER} — drop files there to auto-index")
        return observer
    except Exception as e:
        print(f"[WATCHER] Failed to start: {e}")
        return None


def encode_image_to_base64(filepath: str) -> Optional[Dict[str, Any]]:
    """Encode an image file to base64 for vision model viewing.

    Returns a dict with image data suitable for vision models, or None if error.
    Format: {
        "type": "image_url",
        "image_url": {
            "url": "data:image/jpeg;base64,..."
        }
    }
    """
    ext = filepath.rsplit('.', 1)[1].lower() if '.' in filepath else ''
    if ext not in ['png', 'jpg', 'jpeg', 'tiff', 'bmp', 'gif', 'webp', 'mpo']:
        return None

    import base64
    # Normalize through Pillow into a clean, standard JPEG before sending to the
    # vision model. Phone photos are often MPO (HDR multi-frame JPEGs) or huge
    # (36MP) — LM Studio rejects those with "Invalid image detected". Taking the
    # primary frame, applying EXIF orientation, forcing RGB, and downscaling
    # produces something every vision model can decode.
    try:
        import io
        from PIL import Image, ImageOps
        try:
            _max = int(os.environ.get("BLUE_VISION_MAX_SIDE", "1280"))
        except ValueError:
            _max = 1280
        with Image.open(filepath) as im:
            im = ImageOps.exif_transpose(im)  # honor orientation, drop EXIF
            im = im.convert("RGB")            # strip alpha/CMYK; flatten MPO frame
            w, h = im.size
            if max(w, h) > _max:
                s = _max / float(max(w, h))
                im = im.resize((max(1, int(w * s)), max(1, int(h * s))))
            buf = io.BytesIO()
            im.save(buf, format="JPEG", quality=85)
            image_data = base64.b64encode(buf.getvalue()).decode('utf-8')
        return {"type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{image_data}"}}
    except Exception as e:
        print(f"   [VISION] PIL normalize failed for {filepath}, sending raw: {e}")

    # Fallback: raw bytes with a best-effort MIME (only if Pillow couldn't open it)
    try:
        with open(filepath, 'rb') as f:
            image_data = base64.b64encode(f.read()).decode('utf-8')
        mime_type = {
            'jpg': 'image/jpeg', 'jpeg': 'image/jpeg', 'png': 'image/png',
            'gif': 'image/gif', 'webp': 'image/webp', 'bmp': 'image/bmp',
            'tiff': 'image/tiff', 'mpo': 'image/jpeg',
        }.get(ext, 'image/jpeg')
        return {"type": "image_url",
                "image_url": {"url": f"data:{mime_type};base64,{image_data}"}}
    except Exception as e:
        print(f"   [ERROR] Failed to encode image {filepath}: {e}")
        return None


def extract_text_from_file(filepath: str) -> str:
    """Extract text from various file types. For images, returns metadata since vision model will view them directly."""
    ext = filepath.rsplit('.', 1)[1].lower() if '.' in filepath else ''

    # Files NUL-filled by an unclean shutdown make pypdf die with a native
    # access violation that no except clause can catch (it took the whole
    # server down on 2026-07-04), so refuse them before any parser runs.
    try:
        with open(filepath, 'rb') as f:
            head = f.read(4096)
        if head and not head.strip(b'\x00'):
            return "Error: file is NUL-corrupted (unclean shutdown?) and needs to be restored"
    except OSError as e:
        return f"Error reading file: {e}"

    if ext in ('txt', 'md'):
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    elif ext == 'pdf':
        # Prefer pypdf (the maintained successor); fall back to PyPDF2.
        try:
            try:
                from pypdf import PdfReader
            except ImportError:
                from PyPDF2 import PdfReader
            text = []
            with open(filepath, 'rb') as f:
                for page in PdfReader(f).pages:
                    text.append(page.extract_text() or '')
            return '\n'.join(text)
        except ImportError:
            return "Error: pypdf not installed. Install with: pip install pypdf"
        except Exception as e:
            return f"Error extracting PDF: {str(e)}"

    elif ext in ('doc', 'docx'):
        try:
            import docx
            doc = docx.Document(filepath)
            return '\n'.join(paragraph.text for paragraph in doc.paragraphs)
        except ImportError:
            return "Error: python-docx not installed. Install with: pip install python-docx"
        except Exception as e:
            return f"Error extracting Word doc: {str(e)}"

    elif ext in ('csv', 'tsv'):
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    elif ext in ('json', 'xml', 'html', 'htm'):
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            raw = f.read()
        if ext in ('html', 'htm'):
            try:
                from bs4 import BeautifulSoup
                return BeautifulSoup(raw, 'html.parser').get_text(separator='\n')
            except ImportError:
                return raw
        return raw

    elif ext == 'rtf':
        try:
            from striprtf.striprtf import rtf_to_text
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return rtf_to_text(f.read())
        except ImportError:
            import re
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                raw = f.read()
            return re.sub(r'\\[a-z]+\d* ?|[{}]', '', raw)

    elif ext == 'xlsx':
        try:
            from openpyxl import load_workbook
            wb = load_workbook(filepath, read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"# Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    cells = ['' if v is None else str(v) for v in row]
                    if any(cells):
                        lines.append('\t'.join(cells))
            return '\n'.join(lines)
        except ImportError:
            return "Error: openpyxl not installed. Install with: pip install openpyxl"
        except Exception as e:
            return f"Error extracting xlsx: {str(e)}"

    elif ext == 'pptx':
        try:
            from pptx import Presentation
            lines = []
            for i, slide in enumerate(Presentation(filepath).slides, 1):
                lines.append(f"# Slide {i}")
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        lines.append(shape.text)
            return '\n'.join(lines)
        except ImportError:
            return "Error: python-pptx not installed. Install with: pip install python-pptx"
        except Exception as e:
            return f"Error extracting pptx: {str(e)}"

    elif ext in ('png', 'jpg', 'jpeg', 'tiff', 'bmp', 'gif', 'webp'):
        # Image file - store metadata for vision model to view directly
        try:
            from PIL import Image
            image = Image.open(filepath)
            width, height = image.size
            mode = image.mode
            file_size = os.path.getsize(filepath)

            return (f"[IMAGE FILE - Vision model will view directly]\n"
                    f"Dimensions: {width}x{height}\n"
                    f"Color mode: {mode}\n"
                    f"File size: {file_size} bytes\n"
                    f"Format: {ext.upper()}")
        except ImportError:
            return f"[IMAGE FILE] {os.path.basename(filepath)} - PIL required to read metadata"
        except Exception as e:
            return f"[IMAGE FILE] {os.path.basename(filepath)} - Error reading: {str(e)}"

    return f"Unsupported file type: .{ext}"


def _is_full_document_request(query: str) -> bool:
    """Check if the query is asking for a full document rather than a specific search."""
    query_lower = query.lower()
    # User is asking to read/see/look at a specific document, not searching for a topic
    full_doc_signals = [
        'whole', 'entire', 'full', 'complete', 'all of', 'read the',
        'tell me the story', 'read it', 'tell the story',
        'what is', 'what does', 'what it is', 'take a look', 'fresh look',
        'look at', 'summarize', 'summary', 'overview', 'about this document',
        'about that document', 'tell me about', 'describe',
    ]
    if any(phrase in query_lower for phrase in full_doc_signals):
        return True
    # If the query IS a filename (or very close to one), they want the full doc
    if '.' in query and any(query_lower.endswith(ext) for ext in ['.pdf', '.txt', '.md', '.docx', '.doc']):
        return True
    return False


def _is_document_list_request(query: str) -> bool:
    """True when the user wants an INVENTORY of the library (what/which/how
    many documents or files, what's in the library) rather than a semantic
    search. These must route to the local index lister — ChromaDB would
    semantic-search the literal question and return irrelevant chunks (or
    nothing), which is why 'what documents are in your library' came back
    empty by email."""
    q = f" {query.lower().strip()} "
    list_phrases = (
        "what documents", "what files", "which documents", "which files",
        "list documents", "list files", "list my documents", "list my files",
        "list all", "show documents", "show files", "show me my documents",
        "show me my files", "show my documents", "show my files",
        "how many documents", "how many files", "count documents",
        "count files", "documents do you have", "files do you have",
        "documents are in", "files are in",
        "what's in your library", "whats in your library",
        "what's in my library", "whats in my library",
        "what is in your library", "what do you have in your library",
    )
    return any(p in q for p in list_phrases)


def _is_expertise_query(query: str) -> bool:
    """True for queries that ask Blue to draw on his corpus expertise.

    Triggers multi-chunk retrieval (top-8 across multiple documents, up
    to 3 per doc) instead of the default top-3 deduped-by-document mode.
    Aimed at: 'what does the literature say about X', 'summarise the
    research on Y', 'according to my notes / papers / docs ...' style
    questions where richer cross-document coverage matters more than
    surgical precision.
    """
    q = query.lower()
    expertise_phrases = (
        # explicit corpus-of-knowledge framings
        "the literature", "research on", "research about", "papers on",
        "papers about", "according to my", "according to the", "based on my",
        "based on the", "my notes on", "my notes about", "my docs", "the docs",
        "in my documents", "from my documents", "from the documents",
        "the documents say", "what do my documents", "what does my research",
        # synthesis-style framings
        "what do we know about", "what do you know about", "everything you know about",
        "everything about", "summarise everything", "summarize everything",
        "what's known about", "whats known about",
        "across my", "across all", "compare what",
        # discipline / field hints — these almost always want corpus coverage
        " literature ", " corpus ", "the field of",
        # the owner's own body of work — "using my published work", "my
        # articles", etc. — wants broad coverage across several of his pieces,
        # not a single best-matching chunk.
        "my published work", "my publications", "my articles", "my papers",
        "my own work", "my writing", "my scholarship", "published work",
    )
    return any(p in f" {q} " for p in expertise_phrases)


# Words that describe the operation rather than identify a file. Keeping this
# list beside the resolver makes title matching deterministic and independent
# of whatever the language model happens to infer from a request.
_DOC_REFERENCE_STOP = set("""
the a an and or of to in on for with about into from your my our its this that
these those document documents file files pdf doc docx txt md text copy library
folder please can could would should will do does did have has had is are was
were be been me you it read reading reread open find search look looking show
see check review summarize summarise summary overview tell give get pull bring
try again retry now still really directly access retrieve extract tool
introduction intro chapter paper book syllabus
""".split())


def _doc_words(value: str) -> List[str]:
    """Lower-cased words used for library-title resolution."""
    return re.findall(r"[a-z0-9]+", (value or "").lower().replace("\u2019", "'"))


def _resolve_document_entry(reference: str) -> Optional[Dict[str, Any]]:
    """Resolve a natural title/author reference to one indexed document.

    Exact filename/stem mentions win. Otherwise, two meaningful title tokens
    are required, except for a unique proper-name-like token of five or more
    characters ("Noble", "Toscano"). Ambiguous ties deliberately return None.
    """
    try:
        docs = [
            d for d in load_document_index().get("documents", [])
            if d.get("filename")
            and d.get("doc_type") != "camera"
            and not str(d.get("filename", "")).startswith("camera_")
        ]
    except Exception:
        return None
    if not docs or not (reference or "").strip():
        return None

    ref_lower = reference.lower().replace("\u2019", "'")
    ref_normal = " ".join(_doc_words(reference))

    # Full filename or normalized stem in the request is unambiguous.
    exact = []
    for d in docs:
        filename = str(d.get("filename", ""))
        stem = os.path.splitext(filename)[0]
        stem_normal = " ".join(_doc_words(stem))
        if filename.lower() in ref_lower:
            return d
        if len(stem_normal) >= 5 and re.search(
                r"(?:^| )" + re.escape(stem_normal) + r"(?: |$)", ref_normal):
            exact.append(d)
    if len(exact) == 1:
        return exact[0]
    if len(exact) > 1:
        return None

    title_tokens_by_doc = []
    token_counts = Counter()
    for d in docs:
        stem = os.path.splitext(str(d.get("filename", "")))[0]
        tokens = {
            w for w in _doc_words(stem)
            if len(w) >= 4 and w not in _DOC_REFERENCE_STOP and not w.isdigit()
        }
        title_tokens_by_doc.append((d, tokens))
        token_counts.update(tokens)

    q_tokens = {
        w for w in _doc_words(reference)
        if len(w) >= 4 and w not in _DOC_REFERENCE_STOP and not w.isdigit()
    }
    all_title_tokens = set(token_counts)
    # Voice transcripts often turn "Noble's" into "nobles".
    q_tokens.update(
        w[:-1] for w in tuple(q_tokens)
        if len(w) >= 6 and w.endswith("s") and w[:-1] in all_title_tokens
    )
    if not q_tokens:
        return None

    candidates = []
    for d, title_tokens in title_tokens_by_doc:
        overlap = q_tokens & title_tokens
        if len(overlap) >= 2:
            candidates.append((len(overlap), sum(map(len, overlap)), d))
        elif len(overlap) == 1:
            token = next(iter(overlap))
            if len(token) >= 5 and token_counts[token] == 1:
                candidates.append((1, len(token), d))
    if not candidates:
        return None
    candidates.sort(key=lambda item: (item[0], item[1]), reverse=True)
    if len(candidates) > 1 and candidates[0][:2] == candidates[1][:2]:
        return None
    return candidates[0][2]


def _existing_document_path(doc: Dict[str, Any]) -> Optional[str]:
    """Return an indexed document's real path, repairing stale index paths."""
    filepath = str(doc.get("filepath", "") or "")
    if filepath and os.path.isfile(filepath):
        return filepath
    filename = str(doc.get("filename", "") or "")
    if not filename:
        return None
    folder = _safe_rel_folder(str(doc.get("folder", "") or ""))
    candidate = os.path.join(_abs_library_path(folder), filename)
    return candidate if os.path.isfile(candidate) else None


def _read_resolved_document(query: str, doc: Dict[str, Any],
                            max_results: int = 3) -> str:
    """Extract useful text from one resolved library document.

    Named-file requests should never go through broad semantic search. This
    reads the exact indexed path, includes the opening, and adds query-relevant
    passages from later in the file when appropriate.
    """
    filename = str(doc.get("filename", "document"))
    folder = _safe_rel_folder(str(doc.get("folder", "") or "")) or "library root"
    filepath = _existing_document_path(doc)
    if not filepath:
        preview = (doc.get("text_preview") or "").strip()
        detail = f"\n\nIndexed preview:\n{preview[:1500]}" if preview else ""
        return (
            f"LOCAL LIBRARY FILE FOUND BUT PATH IS STALE: [{filename}] is indexed "
            f"in folder {folder}, but its file is not currently present.{detail}"
        )

    full_text = extract_text_from_file(filepath)
    if not full_text or full_text.lower().startswith(("error", "unsupported file")):
        preview = (doc.get("text_preview") or "").strip()
        detail = f"\n\nIndexed preview:\n{preview[:1500]}" if preview else ""
        return (
            f"LOCAL LIBRARY EXTRACTION FAILED for [{filename}]: "
            f"{full_text or 'no text was returned'}.{detail}"
        )

    title_terms = set(_doc_words(os.path.splitext(filename)[0]))
    query_terms = {
        w for w in _doc_words(query)
        if len(w) >= 4 and w not in _DOC_REFERENCE_STOP and w not in title_terms
    }
    extended = bool(re.search(
        r"\b(?:whole|entire|full|complete)\b|\bread (?:it|the|this)\b|"
        r"\b(?:summari[sz]e|summary|overview|tell me about)\b",
        query or "", re.I,
    ))

    opening_limit = 8500 if extended else 3200
    opening = full_text[:opening_limit].strip()
    excerpts = [opening]

    # Split long PDF output into bounded windows as well as paragraphs. PDF
    # extractors do not reliably preserve blank lines, so paragraph-only
    # scoring can otherwise treat a 30-page chapter as one giant block.
    if query_terms and not extended:
        blocks = []
        for paragraph in re.split(r"\n\s*\n", full_text):
            paragraph = paragraph.strip()
            if not paragraph:
                continue
            if len(paragraph) <= 1800:
                blocks.append(paragraph)
            else:
                start = 0
                while start < len(paragraph):
                    blocks.append(paragraph[start:start + 1800])
                    start += 1600
        scored = []
        opening_norm = re.sub(r"\s+", " ", opening).strip()
        for block in blocks:
            lower = block.lower()
            score = sum(lower.count(term) for term in query_terms)
            block_norm = re.sub(r"\s+", " ", block).strip()
            if score and block_norm and block_norm not in opening_norm:
                scored.append((score, block))
        scored.sort(key=lambda item: item[0], reverse=True)
        for _, block in scored[:max(1, min(int(max_results or 3), 4))]:
            excerpts.append(block.strip())

    body = "\n\n--- relevant passage ---\n\n".join(e for e in excerpts if e)
    body = body[:9500]
    return (
        f"LOCAL LIBRARY READ SUCCEEDED: [{filename}] was resolved in folder "
        f"{folder} and {len(full_text):,} characters were extracted directly "
        f"from the file. The passages below are the document's actual text.\n\n"
        f"[{filename}]\n{body}"
    )


# Course/reading/schedule queries: these are answered by the syllabus's dated
# schedule, but semantic RAG buries that one small document under chunks from
# big books (a "readings for tomorrow" query matched Three-Body-Problem
# acknowledgements, not the schedule). Detect them and go straight to the
# syllabus full text instead.
_COURSE_SCHEDULE_RE = re.compile(
    r"\breadings\b"
    r"|\b(?:assigned|required|course|class|weekly)\s+readings?\b"
    r"|\breading\s+(?:for|list|report|assignment|this|next|tonight|tomorrow|today|due)\b"
    r"|\b(?:to|should\s+i|do\s+i|have\s+to|need\s+to|gotta)\s+read\b"
    r"|\bhomework\b|\bassignments?\b|\bcoursework\b|\bsyllabus\b|\bclass\s+schedule\b"
    r"|\bcourse\s+(?:material|materials|outline|schedule|reading|readings)\b"
    r"|\bdue\b(?!\s+to\b)",
    re.I,
)


def _is_course_schedule_query(query: str) -> bool:
    return bool(_COURSE_SCHEDULE_RE.search(query or ""))


def _syllabus_schedule_text(max_docs: int = 2):
    """Full schedule/readings text of the syllabus document(s), or None.

    Returns the portion from 'Class Schedule' (or the first dated line) onward
    so every week + its readings are present, with a header telling the model
    to map 'today'/'tomorrow' to the right class date using the current date."""
    try:
        index = load_document_index()
        docs = index.get("documents", []) if isinstance(index, dict) else []
    except Exception:
        return None
    syll = [d for d in docs
            if 'syllab' in (d.get('filename', '') or '').lower()
            or 'syllab' in (d.get('folder', '') or '').lower()
            or 'course outline' in (d.get('filename', '') or '').lower()]
    # Honour an active library focus (chat Context panel): only consider syllabi
    # the user actually selected, so a course/schedule/readings question can
    # never pull in a DIFFERENT course's syllabus. With nothing focused, behave
    # as before (any syllabus in the library).
    focus_docs = list(globals().get("_ACTIVE_FOCUS_DOCS") or [])
    focus_folders = list(globals().get("_ACTIVE_FOCUS_FOLDERS") or [])
    if focus_docs or focus_folders:
        def _in_focus(d):
            if (d.get('filename') or '') in focus_docs:
                return True
            fol = (d.get('folder') or '')
            return any(fol == f or fol.startswith(f + '/') for f in focus_folders)
        syll = [d for d in syll if _in_focus(d)]
        if not syll:
            # Focused, but no selected syllabus — don't fall back to the whole
            # library; let the scoped document search answer instead.
            print("   [FOCUS] course/schedule query, but no syllabus in focus — using scoped search")
            return None
    if not syll:
        return None
    parts = []
    for d in syll[:max_docs]:
        fp = d.get('filepath', '')
        if not fp or not os.path.exists(fp):
            continue
        try:
            txt = extract_text_from_file(fp)
        except Exception:
            continue
        if not txt:
            continue
        start = txt.lower().find('class schedule')
        if start < 0:
            m = re.search(r'\b(?:jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)[a-z]*\.?\s+\d{1,2}\s*[:\-]',
                          txt, re.I)
            start = m.start() if m else 0
        excerpt = txt[start:start + 9000]
        parts.append(f"[{d.get('filename', 'syllabus')}] course schedule & readings:\n{excerpt}")
    if not parts:
        return None
    return (
        "Here is the schedule and readings from your course syllabus. Work out "
        "which class date the question refers to using today's date shown above "
        "('tomorrow' = the next calendar day), then report the readings listed "
        "under that date:\n\n" + "\n\n".join(parts)
    )


def _format_doc_chunks(results: List[Dict], is_expertise: bool = False) -> str:
    """Render RAG chunk dicts (filename/content/score/chunk_index/total_chunks)
    as a citation-friendly answer block. Shared by the normal whole-library
    search path and the library-focus (pinned-docs) path."""
    formatted = []
    for i, r in enumerate(results, 1):
        score_pct = f"{r['score']:.0%}"
        chunk_idx = r.get('chunk_index')
        total_chunks = r.get('total_chunks')
        if chunk_idx is not None and total_chunks:
            loc = f" chunk {int(chunk_idx) + 1}/{total_chunks}"
        else:
            loc = ""
        content_preview = r['content'][:800]
        # Citation tag the model can copy inline. Putting the filename in
        # [brackets] near the start makes it easy to weave into prose.
        formatted.append(
            f"[{i}] [{r['filename']}]{loc} (relevance: {score_pct})\n"
            f"{content_preview}"
        )
    header = (
        "Here are the most relevant passages from your corpus. "
        "Cite sources inline when you quote, like [filename.pdf]:"
        if is_expertise else
        "Here's what I found in your documents:"
    )
    return f"{header}\n\n" + "\n\n".join(formatted)


def search_documents_rag(query: str, max_results: int = 3) -> str:
    """Search documents using local ChromaDB RAG first, then keyword fallback."""
    print(f"   [FIND] Searching documents for: '{query}'")

    # Course/reading/schedule questions → answer from the syllabus's dated
    # schedule directly (semantic RAG buries the small syllabus under big books).
    if _is_course_schedule_query(query):
        syllabus = _syllabus_schedule_text()
        if syllabus:
            print("   [SYLLABUS] Course/reading query — returning syllabus schedule")
            return syllabus

    # Inventory requests ("what documents are in your library", "list my
    # files") must be answered by enumerating the index, NOT by semantic
    # search — ChromaDB would match the literal question against chunk text
    # and return noise. Route straight to the local lister.
    if _is_document_list_request(query):
        print(f"   [LIST] Library inventory request — routing to local lister")
        return search_documents_local(query, max_results)

    # A title/author reference names one real file. Resolve it before semantic
    # search: broad RAG interpreted "look for Noble introduction" as a topic
    # query and returned unrelated books while the exact PDF sat in CMDS4740.
    named_doc = _resolve_document_entry(query)
    if named_doc:
        print(f"   [NAMED-DOC] Resolved request to {named_doc.get('filename')}")
        return _read_resolved_document(query, named_doc, max_results)

    # Library focus (chat Context panel): when Alex has pinned specific
    # documents for this conversation, answer from exactly those files
    # (file-scoped RAG) instead of the whole library.
    if _ACTIVE_FOCUS_DOCS:
        try:
            from blue.tools.rag import search_in_documents, get_stats
            stats = get_stats()
            if stats.get('available') and stats.get('total_chunks', 0) > 0:
                results = search_in_documents(query, _ACTIVE_FOCUS_DOCS,
                                              max_results=max(6, max_results))
                if results:
                    print(f"   [FOCUS] {len(results)} chunk(s) from "
                          f"{len(_ACTIVE_FOCUS_DOCS)} pinned doc(s)")
                    return _format_doc_chunks(results, is_expertise=True)
                print("   [FOCUS] nothing relevant in the pinned doc(s)")
                return ("I focused on the document(s) you selected ("
                        + ", ".join(_ACTIVE_FOCUS_DOCS) + ") but didn't find "
                        "anything relevant to that in them.")
        except Exception as e:
            log.warning(f"[FOCUS] doc-scoped search failed: {e}")

    is_expertise = _is_expertise_query(query)

    # Scope to an area of expertise when the request names one ("my published
    # work" → the publications folder, a folder name, …). A library-focus folder
    # pick (Context panel) overrides the inferred scope; [] = whole library.
    try:
        if _ACTIVE_FOCUS_FOLDERS:
            scope_folders = list(_ACTIVE_FOCUS_FOLDERS)
        else:
            scope_folders = _infer_expertise_folders(query)
    except Exception:
        scope_folders = []
    if scope_folders:
        print(f"   [SCOPE] Restricting search to folders: {scope_folders}")

    # Expertise queries (e.g. "what does the literature say…") want
    # multi-chunk RAG, NOT a full-document dump. Check expertise *before*
    # the full-doc heuristic, since the latter trips on generic phrases
    # like "what does" / "tell me about" that legitimately appear in
    # corpus-style questions.
    if not is_expertise and _is_full_document_request(query):
        print(f"   [FULL-DOC] Full document request detected, using local search for full text")
        return search_documents_local(query, max_results)

    # --- 1. Try local ChromaDB RAG first ---
    try:
        from blue.tools.rag import (
            search as rag_search,
            search_expertise as rag_search_expertise,
            get_stats,
        )
        stats = get_stats()
        if stats.get('available') and stats.get('total_chunks', 0) > 0:
            scope = scope_folders or None
            if is_expertise:
                # Scoped to a folder (e.g. his publications), spread coverage
                # across MORE documents (lower per-doc cap, higher total) so a
                # synthesis draws on several pieces, not just the closest one.
                if scope:
                    results = rag_search_expertise(query, max_chunks=14, max_per_doc=2, folders=scope)
                else:
                    results = rag_search_expertise(query, max_chunks=8, max_per_doc=3)
                # If folder scoping was too tight and found nothing, retry wide.
                if not results and scope:
                    print("   [SCOPE] No scoped results; retrying across whole library")
                    results = rag_search_expertise(query, max_chunks=8, max_per_doc=3)
                print(f"   [EXPERTISE] Multi-chunk mode: {len(results)} chunk(s)")
            else:
                results = rag_search(query, max_results, folders=scope)
                if not results and scope:
                    print("   [SCOPE] No scoped results; retrying across whole library")
                    results = rag_search(query, max_results)
            if results:
                # In normal (non-expertise) mode, if all results come from one
                # document, fall through to local search for full text.
                unique_files = set(r['filename'] for r in results)
                if not is_expertise and len(unique_files) == 1:
                    print(f"   [SINGLE-DOC] All {len(results)} chunk(s) from same file, fetching full text")
                    return search_documents_local(query, max_results)

                print(f"   [OK] ChromaDB RAG returned {len(results)} result(s)")
                return _format_doc_chunks(results, is_expertise)
            else:
                print(f"   [WARN] ChromaDB returned no results, trying next...")
    except ImportError:
        print("   [WARN] ChromaDB not available, trying keyword search...")
    except Exception as e:
        print(f"   [WARN] ChromaDB error: {e}")

    # --- 2. Fall back to local keyword search ---
    # The previous middle tier hit LM_STUDIO_RAG_URL/search, which no LM
    # Studio build actually exposes — every call timed out for 10 seconds
    # and fell through to keyword search anyway. Going straight to the
    # keyword path skips that wait without losing any working behavior.
    print(f"   [ITER] Falling back to local keyword search...")
    return search_documents_local(query, max_results)


def _search_documents_guarded(query: str, max_results: int = 3,
                              timeout: float = 15.0) -> str:
    """Run the RAG document search with a hard timeout.

    ChromaDB can stall indefinitely — a corrupt HNSW index or lock
    contention has no internal timeout — and with a single-threaded server
    that one stalled call freezes Blue completely and even blocks Ctrl+C.
    This runs the search on a worker thread; if it overruns, we abandon it
    and fall back to the index-only keyword search (which never touches
    ChromaDB), so a document query can never hang the assistant."""
    box: Dict[str, object] = {}

    def _run():
        try:
            box["result"] = search_documents_rag(query, max_results)
        except Exception as e:  # noqa: BLE001 — surface as fallback, never crash
            box["error"] = e

    worker = threading.Thread(target=_run, daemon=True)
    worker.start()
    worker.join(timeout)

    if worker.is_alive():
        log.error(
            f"[DOCS] document search exceeded {timeout:.0f}s (ChromaDB stalled) "
            f"— falling back to index-only search. The RAG index likely needs a rebuild."
        )
        try:
            return search_documents_local(query, max_results)
        except Exception as e:
            log.error(f"[DOCS] index-only fallback also failed: {e}")
            return ("I can see your document library, but searching it is taking "
                    "too long right now — the search index may need a rebuild.")
    if "error" in box:
        log.error(f"[DOCS] document search errored: {box['error']}")
        try:
            return search_documents_local(query, max_results)
        except Exception:
            return "I had trouble searching your documents just now."
    return str(box.get("result", ""))


def search_documents_local(query: str, max_results: int = 3) -> str:
    """Fallback: Simple local search through document index. Returns JSON for images."""
    print(f"   [FOLDER] Using local document search...")
    index = load_document_index()
    documents = index.get("documents", [])

    print(f"   [DATA] Found {len(documents)} documents in local index")

    if not documents:
        print(f"   [WARN]  No documents uploaded yet!")
        return (
            "I don't have any documents to search through yet! "
            "You can upload documents at http://127.0.0.1:5000/documents - "
            "I can read PDFs, Word docs, text files, markdown, and images. "
            "Once you upload some documents, I'll be able to answer questions about them!"
        )

    query_lower = query.lower()
    matches = []

    # Special handling for document listing queries
    list_keywords = ["list documents", "what documents", "show documents", "all documents",
                     "summarize all", "list all", "show all",
                     "what files", "which documents", "which files", "list files",
                     "show files", "how many documents", "how many files",
                     "count documents", "count files", "what's in your library",
                     "whats in your library", "what's in my library",
                     "whats in my library", "documents do you have", "files do you have"]
    is_list_query = any(kw in query_lower for kw in list_keywords)

    # Also treat very generic queries as list requests (e.g., just "documents" or "files")
    generic_queries = query_lower.strip() in ["documents", "document", "files", "file", "pdfs", "pdf"]

    if is_list_query or generic_queries or ("all" in query_lower and ("document" in query_lower or "summarize" in query_lower or "image" in query_lower)):
        print(f"   [LIST] Query asks for document list/count, listing them...")

        # Filter out camera images from document listings
        real_documents = [doc for doc in documents
                          if doc.get('doc_type') != 'camera'
                          and not doc['filename'].startswith('camera_')]

        doc_list = []
        for i, doc in enumerate(real_documents[:10], 1):
            preview = doc.get('text_preview', 'No preview available')[:200]
            doc_list.append(f"{i}. {doc['filename']}\n   Preview: {preview}...")

        summary = "\n\n".join(doc_list)
        if len(real_documents) > 10:
            summary += f"\n\n...and {len(real_documents) - 10} more documents."

        return f"I have {len(real_documents)} document(s) uploaded:\n\n{summary}"

    # A named file is authoritative even in this non-Chroma fallback. This also
    # avoids prefix verbs ("read Noble_Introduction.pdf") outscoring the title.
    named_doc = _resolve_document_entry(query)
    if named_doc:
        return _read_resolved_document(query, named_doc, max_results)
    documents_to_score = documents

    # Search through documents
    for doc in documents_to_score:
        relevance = 0
        filename_lower = doc['filename'].lower()

        # Check filename
        if query_lower in filename_lower:
            relevance += 3

        # Check cached text content
        text_content = doc.get('text_preview', '').lower()
        if query_lower in text_content:
            relevance += 5

        # Check individual query words
        for word in query_lower.split():
            if len(word) > 3:
                if word in filename_lower:
                    relevance += 1
                if word in text_content:
                    relevance += 2

        if relevance > 0:
            matches.append((doc, relevance))

    print(f"   [TARGET] Found {len(matches)} matching documents")

    if not matches:
        return (
            f"I couldn't find any documents matching '{query}'. "
            f"I have {len(documents)} document(s) uploaded. "
            "Try using different keywords, or ask me to list all documents."
        )

    # Sort by relevance
    matches.sort(key=lambda x: x[1], reverse=True)

    # Check if results include images - if so, return special format
    has_images = False
    image_results = []
    text_results = []

    for doc, score in matches[:max_results]:
        filepath = doc.get('filepath', '')
        filename = doc['filename']
        ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''

        if ext in ['png', 'jpg', 'jpeg', 'tiff', 'bmp', 'gif', 'webp']:
            has_images = True
            if os.path.exists(filepath):
                image_results.append({
                    'filename': filename,
                    'filepath': filepath,
                    'score': score
                })
                print(f"   [IMAGE] Found image: {filename}")
        else:
            # Text document - extract content
            print(f"   [FILE] Extracting full text from: {filename}")

            if os.path.exists(filepath):
                try:
                    full_text = extract_text_from_file(filepath)

                    # Check if user wants the complete/full document
                    wants_full_content = any(phrase in query_lower for phrase in [
                        "whole", "entire", "full", "complete", "all of", "read the",
                        "tell me the story", "read it", "tell the story"
                    ])

                    # If the file is very long and user didn't ask for full content
                    if len(full_text) > 3000 and not wants_full_content:
                        # Split into paragraphs/sections
                        sections = full_text.split('\n\n')
                        relevant_sections = []

                        for section in sections:
                            section_lower = section.lower()
                            # Check if this section contains query terms
                            if any(word in section_lower for word in query_lower.split() if len(word) > 3):
                                relevant_sections.append(section)

                        if relevant_sections:
                            # Return most relevant sections (up to 2000 chars)
                            combined = '\n\n'.join(relevant_sections[:5])
                            content = combined[:2000] if len(combined) > 2000 else combined
                        else:
                            # No specific sections found, return first part
                            content = full_text[:2000]
                    else:
                        # File is short enough OR user wants full content - return it all
                        # Limit to 15000 chars to avoid overwhelming the model
                        content = full_text[:15000] if len(full_text) > 15000 else full_text

                    text_results.append(f"[FILE] **{filename}** (relevance: {score})\n\n{content}\n")

                except Exception as e:
                    print(f"   [WARN]  Error reading {filename}: {e}")
                    # Fall back to preview
                    preview = doc.get('text_preview', 'No preview available')[:500]
                    text_results.append(f"[FILE] **{filename}** (relevance: {score})\n\n{preview}...\n")
            else:
                # File not found, use preview
                preview = doc.get('text_preview', 'No preview available')[:500]
                text_results.append(f"[FILE] **{filename}** (relevance: {score})\n\n{preview}...\n")

    # If we found images, return special JSON format
    if has_images:
        result = {
            "_type": "document_search_with_images",
            "images": image_results,
            "text_documents": text_results
        }
        print(f"   [OK] Returning {len(image_results)} image(s) and {len(text_results)} text document(s)")
        return json.dumps(result)

    # No images, return text results as before
    print(f"   [OK] Returning {len(text_results)} document(s) with full content")
    return "Here's what I found in your documents:\n\n" + "\n---\n\n".join(text_results)


def view_image(filename: str = None, query: str = None) -> str:
    """View an image file for the vision model to analyze.

    Args:
        filename: Specific image filename to view
        query: Search query if filename not provided

    Returns:
        JSON string with image information for vision model injection
    """
    global _vision_queue

    print(f"   [VIEW] Request to view image - filename: {filename}, query: {query}")

    # Load document index
    index = load_document_index()
    documents = index.get("documents", [])

    # Filter to only image files
    image_docs = [
        doc for doc in documents
        if doc['filename'].lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp', '.tiff'))
    ]

    if not image_docs:
        return json.dumps({
            "success": False,
            "message": "No images found in uploaded documents. Upload images at http://127.0.0.1:5000/documents/upload"
        })

    print(f"   [DATA] Found {len(image_docs)} total image(s) in documents")

    # Find the requested image
    found_images = []

    if filename:
        # Search by exact or partial filename match
        filename_lower = filename.lower()
        for doc in image_docs:
            doc_filename_lower = doc['filename'].lower()
            if doc_filename_lower == filename_lower or filename_lower in doc_filename_lower:
                found_images.append(doc)
                print(f"   [MATCH] Found by filename: {doc['filename']}")

    elif query:
        # Search by query in filename
        query_lower = query.lower()
        for doc in image_docs:
            if query_lower in doc['filename'].lower():
                found_images.append(doc)
                print(f"   [MATCH] Found by query: {doc['filename']}")

    else:
        # No filename or query - show all images
        found_images = image_docs[:3]  # Limit to 3 images at once
        print(f"   [LIST] Showing {len(found_images)} recent image(s)")

    if not found_images:
        # Check if user asked for a PDF or other non-image file
        if filename and filename.lower().endswith(('.pdf', '.doc', '.docx', '.txt', '.md')):
            return json.dumps({
                "success": False,
                "message": f"{filename} is not an image file - it's a document. The search_documents tool has already provided its contents."
            })

        available = ", ".join([doc['filename'] for doc in image_docs[:10]])
        return json.dumps({
            "success": False,
            "message": f"No images found matching '{filename or query}'. Available images: {available}"
        })

    # Queue images for vision model
    image_results = []
    for doc in found_images[:3]:  # Limit to 3 images at once
        filepath = doc.get('filepath', '')
        if os.path.exists(filepath):
            image_results.append({
                'filename': doc['filename'],
                'filepath': filepath,
                'score': 1.0
            })
            print(f"   [QUEUE] Queued image for viewing: {doc['filename']}")

    # Store in global pending images
    # Add images to vision queue
    global _vision_queue
    for img in image_results:
        _vision_queue.add_image(
            filepath=img['filepath'],
            filename=img['filename'],
            is_camera=False
        )
    print(f"   [VISION] Stored {len(image_results)} image(s) for vision model injection")

    # Build response
    image_names = [img['filename'] for img in image_results]

    return json.dumps({
        "success": True,
        "message": f"Viewing {len(image_results)} image(s): {', '.join(image_names)}",
        "images": image_names,
        "_instruction": "The images will be shown to you in the next message. Analyze them and respond to the user's question."
    })


# Which robot the current chat turn is addressed to — set by process_with_tools
# so camera aiming moves the RIGHT head (Blue's page vs Hexia's page).
_ACTIVE_CHAT_ROBOT = "blue"

# Library-focus selection for the current chat turn (the chat Context panel's
# document picker). Set by process_with_tools from the request's `focus`; when
# non-empty, Blue's document awareness (<focused_documents>) and his
# search_documents calls are scoped to exactly these picks instead of the whole
# library. Reset every turn so a stale selection never leaks across requests.
_ACTIVE_FOCUS_DOCS: List[str] = []      # specific filenames
_ACTIVE_FOCUS_FOLDERS: List[str] = []   # library folder rel-paths (POSIX)


def _zoomed_frame(frame, zoom, region: str = "center"):
    """Digital zoom: crop a region of the frame and scale it back up.
    Returns (frame, effective_zoom). Region anchors: 'left'/'right' pin the
    crop to that edge, 'top'/'bottom' likewise; combos ('top-left') work;
    anything else means centered."""
    import cv2
    try:
        z = max(1.0, min(4.0, float(zoom or 1.0)))
    except (TypeError, ValueError):
        z = 1.0
    if z <= 1.01:
        return frame, 1.0
    h, w = frame.shape[:2]
    cw, ch = max(2, int(w / z)), max(2, int(h / z))
    r = (region or "center").lower()
    x = 0 if "left" in r else (w - cw if "right" in r else (w - cw) // 2)
    y = 0 if "top" in r else (h - ch if "bottom" in r else (h - ch) // 2)
    crop = frame[y:y + ch, x:x + cw]
    return cv2.resize(crop, (w, h), interpolation=cv2.INTER_CUBIC), z


def _set_camera_hardware_zoom(camera, factor: float) -> float:
    """Zoom IN THE CAMERA, before any frame is captured. UVC convention
    (Logitech BRIO et al.): CAP_PROP_ZOOM value 100 = 1x, 200 = 2x ... The
    driver rejects (not clamps) out-of-range values, so read back to see what
    actually took. Returns the achieved factor (1.0 = no hardware zoom)."""
    import cv2
    try:
        if camera.get(cv2.CAP_PROP_ZOOM) <= 0:    # property unsupported
            return 1.0
        target = int(round(100 * max(1.0, min(4.0, factor))))
        for v in (target, 400, 300, 250, 200, 150, 120):
            if v > target:
                continue
            camera.set(cv2.CAP_PROP_ZOOM, v)
            rb = camera.get(cv2.CAP_PROP_ZOOM)
            if abs(rb - v) < 1:
                return rb / 100.0
        return 1.0
    except Exception:
        return 1.0


# ---- Live camera hub: "see through my eyes" ------------------------------
# While the chat page's preview is open, the hub OWNS the camera (Windows
# gives exclusive access): it streams MJPEG to the browser, accepts live
# zoom, and capture_camera_image takes its frames from here — so the photo
# is exactly the view Alex was just looking at, and nothing fights over the
# device. Auto-releases shortly after the last viewer disconnects.

import threading as _cam_threading

_CAM_HUB = {"cap": None, "lock": _cam_threading.Lock(), "frame": None, "jpeg": None,
            "t_frame": 0.0, "last_pull": 0.0, "zoom": 1.0}
_CAM_HUB_IDLE_RELEASE = 12.0   # seconds after the last viewer before letting go


def _cam_hub_active() -> bool:
    with _CAM_HUB["lock"]:
        return _CAM_HUB["cap"] is not None


def _cam_hub_start() -> bool:
    """Open the shared camera (if not already) and start the reader thread."""
    import cv2
    import time as _t
    with _CAM_HUB["lock"]:
        _CAM_HUB["last_pull"] = _t.time()
        if _CAM_HUB["cap"] is not None:
            return True
        cap = cv2.VideoCapture(0, cv2.CAP_DSHOW)
        if not cap.isOpened():
            cap = cv2.VideoCapture(0)
        if not cap.isOpened():
            return False
        cap.set(cv2.CAP_PROP_FRAME_WIDTH, 1920)
        cap.set(cv2.CAP_PROP_FRAME_HEIGHT, 1080)
        cap.set(cv2.CAP_PROP_AUTOFOCUS, 1)
        try:
            cap.set(cv2.CAP_PROP_ZOOM, 100)
            cap.set(cv2.CAP_PROP_PAN, 0)
            cap.set(cv2.CAP_PROP_TILT, 0)
        except Exception:
            pass
        _CAM_HUB["cap"] = cap
        _CAM_HUB["zoom"] = 1.0
        _CAM_HUB["frame"] = None
        _CAM_HUB["jpeg"] = None
        _cam_threading.Thread(target=_cam_hub_loop, daemon=True).start()
        print("   [CAM-HUB] camera opened for live preview")
        return True


def _cam_hub_loop():
    import cv2
    import time as _t
    while True:
        with _CAM_HUB["lock"]:
            cap = _CAM_HUB["cap"]
            if cap is None:
                return
            if _t.time() - _CAM_HUB["last_pull"] > _CAM_HUB_IDLE_RELEASE:
                try:
                    cap.release()
                except Exception:
                    pass
                _CAM_HUB["cap"] = None
                _CAM_HUB["frame"] = None
                _CAM_HUB["jpeg"] = None
                print("   [CAM-HUB] idle — camera released")
                return
            ok, frame = cap.read()
            if ok and frame is not None:
                _CAM_HUB["frame"] = frame
                _CAM_HUB["t_frame"] = _t.time()
                h, w = frame.shape[:2]
                pw = 640
                ph = max(2, int(h * pw / max(1, w)))
                okj, buf = cv2.imencode('.jpg', cv2.resize(frame, (pw, ph)),
                                        [cv2.IMWRITE_JPEG_QUALITY, 70])
                if okj:
                    _CAM_HUB["jpeg"] = buf.tobytes()
        _t.sleep(0.05)


def _cam_hub_capture(zoom, zoom_region):
    """A full-resolution frame from the live-preview camera. Applies a
    requested centered zoom through the hub (the preview shows it too).
    Returns (frame, hw_zoom) or (None, 1.0) when the hub isn't running."""
    import time as _t
    try:
        want = max(1.0, min(4.0, float(zoom or 1.0)))
    except (TypeError, ValueError):
        want = 1.0
    centered = ("center" in (zoom_region or "center").lower()
                or "centre" in (zoom_region or "").lower())
    with _CAM_HUB["lock"]:
        cap = _CAM_HUB["cap"]
        if cap is None:
            return None, 1.0
        _CAM_HUB["last_pull"] = _t.time()   # a capture counts as activity
        if want > 1.01 and centered:
            _CAM_HUB["zoom"] = _set_camera_hardware_zoom(cap, want)
    if want > 1.01 and centered:
        _t.sleep(0.6)                       # let the sensor reach the new zoom
    deadline = _t.time() + 2.0
    while _t.time() < deadline:
        with _CAM_HUB["lock"]:
            fr, tf, hz = _CAM_HUB["frame"], _CAM_HUB["t_frame"], _CAM_HUB["zoom"]
        if fr is not None and tf >= _t.time() - 0.4:
            return fr.copy(), max(1.0, hz)
        _t.sleep(0.05)
    return None, 1.0


def capture_camera_image(look: str = None, zoom=None, zoom_region: str = "center",
                         robot: str = None) -> str:
    """
    Capture a BRAND NEW image from the camera - IMPROVED VERSION.

    View control:
    - look: left/right/up/down/center — physically turns the robot's head
      (real pan/tilt via the Ohbot motors) before the shot. Skipped silently
      if the head isn't connected.
    - zoom + zoom_region: digital zoom (1-4x) into a part of the view.

    CRITICAL IMPROVEMENTS:
    - Unique timestamp with milliseconds
    - Longer warmup for better quality
    - Discards first frames
    - High quality JPEG
    - Clears vision queue and adds only THIS new image
    - Returns hash for uniqueness verification
    """
    global _vision_queue

    print(f"   [CAMERA] ⚡ CAPTURING BRAND NEW IMAGE RIGHT NOW...")

    try:
        import cv2
        import datetime
        import time

        # Physical pan/tilt: aim the head BEFORE opening the camera, so the
        # motors settle while the sensor warms up.
        aimed = ""
        if look:
            d = str(look).lower().strip()
            if d in ("left", "right", "up", "down", "center", "centre",
                     "forward", "front", "straight", "ahead"):
                try:
                    head_key = _robot_cfg(robot or _ACTIVE_CHAT_ROBOT).get("head", robot or _ACTIVE_CHAT_ROBOT)
                    head = blue_head.get_head(head_key)
                    if head.is_available() and head.look(d, speed=5.0):
                        aimed = "center" if d in ("centre", "forward", "front",
                                                  "straight", "ahead") else d
                        print(f"   [CAMERA] Head aimed: {aimed}")
                        time.sleep(0.7)  # let the motors finish before grabbing frames
                except Exception as e:
                    print(f"   [CAMERA] Head aim skipped ({e})")

        # CRITICAL: Unique timestamp with MILLISECONDS
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S_%f")

        try:
            want_zoom = max(1.0, min(4.0, float(zoom or 1.0)))
        except (TypeError, ValueError):
            want_zoom = 1.0

        # Live preview open? Take the shot FROM it: the hub owns the device
        # (Windows gives exclusive access), and what's on the preview screen —
        # including any zoom/steering done there — is exactly what's captured.
        frame = None
        hw_zoom = 1.0
        if _cam_hub_active():
            frame, hw_zoom = _cam_hub_capture(zoom, zoom_region)
            if frame is not None:
                print(f"   [CAMERA] frame taken from the live preview (hub, {hw_zoom:g}x)")

        if frame is None:
            # Open the camera via DirectShow: it exposes the BRIO's own
            # zoom/pan/tilt controls, which the default MSMF backend rejects.
            camera = cv2.VideoCapture(0, cv2.CAP_DSHOW)
            if not camera.isOpened():
                camera = cv2.VideoCapture(0)

            if not camera.isOpened():
                print(f"   [ERROR] Could not open camera")
                return json.dumps({
                    "success": False,
                    "error": "Could not access camera. Make sure a camera is connected and not in use by another application."
                })

            # Set high quality
            camera.set(cv2.CAP_PROP_FRAME_WIDTH, 1920)
            camera.set(cv2.CAP_PROP_FRAME_HEIGHT, 1080)
            camera.set(cv2.CAP_PROP_AUTOFOCUS, 1)

            # Deterministic baseline EVERY capture: neutral zoom, centered zoom
            # window. The camera remembers PTZ between runs — a leftover pan
            # (found at -6 once) silently skews every later photo.
            try:
                camera.set(cv2.CAP_PROP_ZOOM, 100)
                camera.set(cv2.CAP_PROP_PAN, 0)
                camera.set(cv2.CAP_PROP_TILT, 0)
            except Exception:
                pass

            # Zoom BEFORE capture: use the camera's own zoom when the request is
            # centered (the in-camera zoom window is center-only); off-center
            # regions are handled by the digital crop after capture instead.
            if want_zoom > 1.01 and ("center" in (zoom_region or "center").lower()
                                     or "centre" in (zoom_region or "").lower()):
                hw_zoom = _set_camera_hardware_zoom(camera, want_zoom)
                if hw_zoom > 1.01:
                    print(f"   [CAMERA] In-camera zoom set: {hw_zoom:g}x")

            # Give camera MORE time to warm up and adjust (CRITICAL for quality)
            time.sleep(1.2)  # Increased from 0.8s

            # Discard first few frames (often lower quality)
            for _ in range(3):
                camera.read()
                time.sleep(0.1)

            # NOW capture the actual frame we'll use
            ret, frame = camera.read()
            camera.release()

            if not ret or frame is None:
                print(f"   [ERROR] Failed to capture frame")
                return json.dumps({
                    "success": False,
                    "error": "Failed to capture image from camera."
                })

        # Digital zoom for whatever the in-camera zoom didn't cover: the whole
        # request when hardware zoom is unavailable or the region is
        # off-center, or the remaining factor when it partially took.
        frame, dig_zoom = _zoomed_frame(frame, want_zoom / hw_zoom, zoom_region)
        eff_zoom = hw_zoom * dig_zoom
        if eff_zoom > 1.01:
            kind = ("in-camera" if dig_zoom <= 1.01 else
                    ("in-camera + digital" if hw_zoom > 1.01 else "digital"))
            print(f"   [CAMERA] {eff_zoom:g}x {kind} zoom on the {zoom_region or 'center'} of the view")

        # Generate UNIQUE filename with timestamp
        filename = f"camera_NEW_{timestamp}.jpg"
        filepath = os.path.join(CAMERA_FOLDER, filename)

        # Save with HIGH quality
        os.makedirs(CAMERA_FOLDER, exist_ok=True)
        success = cv2.imwrite(filepath, frame, [cv2.IMWRITE_JPEG_QUALITY, 95])

        if not success:
            print(f"   [ERROR] Failed to save image")
            return json.dumps({
                "success": False,
                "error": "Failed to save captured image."
            })

        print(f"   [SAVE] ✅ NEW image saved to: {filepath}")

        # Get file info
        file_size = os.path.getsize(filepath)
        file_hash = get_file_hash(filepath)

        # Get image dimensions
        height, width = frame.shape[:2]

        # Camera frames live in CAMERA_FOLDER and are short-lived: each new
        # capture replaces the vision queue. Keeping them out of the document
        # index prevents them from drowning real uploads in search results.

        # CRITICAL: Clear queue and add ONLY this new image
        _vision_queue.clear()
        _vision_queue.add_image(filepath, filename, is_camera=True)

        print(f"   [VISION] Queued NEW camera image: {filename}")

        view_bits = []
        if aimed:
            view_bits.append("head re-centered" if aimed == "center" else f"head turned {aimed}")
        if eff_zoom > 1.01:
            view_bits.append(f"{eff_zoom:g}x {kind} zoom on the {(zoom_region or 'center')} of the view")
        view_note = (" (" + ", ".join(view_bits) + ")") if view_bits else ""

        return json.dumps({
            "success": True,
            "message": f"📷 ✨ BRAND NEW CAMERA IMAGE captured at {datetime.datetime.now().strftime('%I:%M:%S %p')}{view_note}",
            "filename": filename,
            "filepath": filepath,
            "dimensions": f"{width}x{height}",
            "timestamp": timestamp,
            "file_hash": file_hash,
            "view": {"look": aimed or None, "zoom": eff_zoom,
                     "zoom_region": (zoom_region or "center") if eff_zoom > 1.01 else None},
            "_instruction": (
                f"Camera view captured at {datetime.datetime.now().strftime('%I:%M:%S %p')}{view_note}. "
                "You'll see what's in front of you in the next message. "
                + ("This view is aimed/zoomed as noted — describe what's IN it. "
                   if view_bits else "")
                + "Respond naturally about your surroundings."
            )
        })

    except ImportError:
        return json.dumps({
            "success": False,
            "error": "OpenCV (cv2) not installed. Install with: pip install opencv-python"
        })
    except Exception as e:
        print(f"   [ERROR] Camera capture failed: {e}")
        import traceback
        traceback.print_exc()
        return json.dumps({
            "success": False,
            "error": f"Camera capture failed: {str(e)}"
        })

def create_document_file(filename: str, content: str, file_type: str = "txt") -> str:
    """Create a new document and save it to the documents folder."""
    print(f"   [CREATE] Creating document: {filename}")

    try:
        # Sanitize filename
        filename = secure_filename(filename)

        # Ensure filename has correct extension
        if '.' not in filename:
            filename = f"{filename}.{file_type}"
        else:
            # Check if extension matches file_type
            ext = filename.rsplit('.', 1)[1].lower()
            if ext != file_type:
                filename = f"{filename.rsplit('.', 1)[0]}.{file_type}"

        # Create full path (documents go in DOCUMENTS_FOLDER)
        filepath = os.path.join(DOCUMENTS_FOLDER, filename)

        # Check if file already exists
        if os.path.exists(filepath):
            # Add timestamp to make it unique
            timestamp = __import__('datetime').datetime.now().strftime('%Y%m%d_%H%M%S')
            name_part = filename.rsplit('.', 1)[0]
            ext_part = filename.rsplit('.', 1)[1]
            filename = f"{name_part}_{timestamp}.{ext_part}"
            filepath = os.path.join(DOCUMENTS_FOLDER, filename)

        # Write content to file
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)

        print(f"   [SAVE] Saved to: {filepath}")

        # Get file info
        file_size = os.path.getsize(filepath)
        file_hash = get_file_hash(filepath)

        # Add to document index
        index = load_document_index()

        # Add text file extensions to allowed extensions if not already there
        allowed_create_extensions = {'txt', 'md', 'json', 'csv', 'html'}

        index['documents'].append({
            'filename': filename,
            'filepath': str(filepath),
            'size': file_size,
            'hash': file_hash,
            'uploaded_at': __import__('datetime').datetime.now().strftime('%Y-%m-%d %H:%M'),
            'text_preview': content[:500] if len(content) > 500 else content,
            'indexed_in_rag': False,
            'created_by_blue': True  # Mark as created by Blue
        })

        save_document_index(index)

        # Index in local ChromaDB RAG
        _doc_folder = _folder_of_filepath(filepath)
        try:
            from blue.tools.rag import index_document as rag_index
            rag_result = rag_index(filepath, filename, doc_id=file_hash, folder=_doc_folder)
            if rag_result.get('success'):
                index['documents'][-1]['indexed_in_rag'] = True
                save_document_index(index)
                print(f"   [RAG] Indexed {rag_result.get('chunks_indexed', 0)} chunks")
        except Exception as e:
            print(f"   [WARN] RAG indexing skipped: {e}")

        print(f"   [INDEX] Added to document index")

        from urllib.parse import quote as _quote
        download_url = (
            "http://127.0.0.1:5000/documents/download?"
            f"folder={_quote(_doc_folder)}&filename={_quote(filename)}"
        )

        return (
            f"✅ Document created successfully!\n\n"
            f"📄 Filename: {filename}\n"
            f"📏 Size: {file_size} bytes\n"
            f"🔗 Download: {download_url}\n"
            f"📂 View all documents: http://127.0.0.1:5000/documents"
        )

    except Exception as e:
        print(f"   [ERROR] Error creating document: {e}")
        return f"❌ Error creating document: {str(e)}"


# ===== HUE LIGHT FUNCTIONS =====

def get_hue_lights() -> Dict:
    """Get all lights from Hue Bridge."""
    if not BRIDGE_IP or not HUE_USERNAME:
        return {}
    try:
        response = requests.get(f"http://{BRIDGE_IP}/api/{HUE_USERNAME}/lights", timeout=5)
        if response.status_code == 200:
            return response.json()
    except Exception as e:
        print(f"   [ERROR] Error getting lights: {e}")
    return {}


def find_light_by_name(light_name: str) -> Optional[str]:
    """Find light ID by name."""
    lights = get_hue_lights()
    light_name_lower = light_name.lower()

    for light_id, data in lights.items():
        if data.get('name', '').lower() == light_name_lower:
            return light_id

    for light_id, data in lights.items():
        if light_name_lower in data.get('name', '').lower():
            return light_id

    return None


def control_hue_light(light_id: str, state: Dict) -> bool:
    """Send state change to a specific light."""
    if not BRIDGE_IP or not HUE_USERNAME:
        return False
    try:
        response = requests.put(
            f"http://{BRIDGE_IP}/api/{HUE_USERNAME}/lights/{light_id}/state",
            json=state,
            timeout=5
        )
        return response.status_code == 200
    except Exception as e:
        print(f"   [ERROR] Error controlling light: {e}")
        return False


def apply_mood_to_lights(mood: str) -> str:
    """Apply a mood/scene to all lights."""
    print(f"   [MOOD] Applying mood: {mood}")

    if not BRIDGE_IP or not HUE_USERNAME:
        return "Hue not configured. Run setup_hue.py first!"

    mood_lower = mood.lower()
    if mood_lower not in MOOD_PRESETS:
        available = ", ".join(MOOD_PRESETS.keys())
        return f"Unknown mood '{mood}'. Available moods: {available}"

    lights = get_hue_lights()
    if not lights:
        return "Could not connect to Hue Bridge."

    mood_data = MOOD_PRESETS[mood_lower]
    settings_list = mood_data["settings"]
    light_ids = list(lights.keys())

    success_count = 0
    assignments = []

    for i, light_id in enumerate(light_ids):
        setting = settings_list[i % len(settings_list)].copy()
        setting["on"] = True
        setting["transitiontime"] = 10

        if control_hue_light(light_id, setting):
            success_count += 1
            light_name = lights[light_id]['name']
            assignments.append(light_name)

    if success_count > 0:
        description = mood_data["description"]
        return f"Applied '{mood}' mood ({description}) to {success_count} light(s): {', '.join(assignments[:3])}{'...' if len(assignments) > 3 else ''}"
    else:
        return f"Failed to apply mood '{mood}'"


def execute_light_control(action: str, light_name: str = None, brightness: int = None,
                         color: str = None, mood: str = None) -> str:
    """Execute light control commands."""
    print(f"   [LIGHT] Light control: action={action}, light={light_name}, brightness={brightness}, color={color}, mood={mood}")

    if not BRIDGE_IP or not HUE_USERNAME:
        return "Philips Hue not configured. Run setup_hue.py first!"

    lights = get_hue_lights()
    if not lights:
        return "Could not connect to Hue Bridge."

    if action == "mood":
        if mood:
            return apply_mood_to_lights(mood)
        else:
            available = ", ".join(MOOD_PRESETS.keys())
            return f"Please specify a mood. Available: {available}"

    target_lights = []
    if light_name:
        light_id = find_light_by_name(light_name)
        if light_id:
            target_lights = [(light_id, lights[light_id]['name'])]
        else:
            available = ", ".join([lights[lid]['name'] for lid in lights])
            return f"Couldn't find '{light_name}'. Available: {available}"
    else:
        target_lights = [(lid, data['name']) for lid, data in lights.items()]

    if not target_lights:
        return "No lights found."

    if action == "status":
        status_lines = []
        for light_id, name in target_lights:
            state = lights[light_id].get('state', {})
            on_status = "ON" if state.get('on', False) else "OFF"
            bri = state.get('bri', 0)
            bri_percent = int((bri / 254) * 100) if bri else 0
            status_lines.append(f"{name}: {on_status}" + (f", {bri_percent}%" if on_status == "ON" else ""))
        return "Light Status:\n" + "\n".join(status_lines)

    elif action == "on":
        success_count = sum(1 for lid, _ in target_lights if control_hue_light(lid, {"on": True}))
        names = ", ".join([n for _, n in target_lights])
        return f"Turned on: {names}" if success_count == len(target_lights) else f"Turned on {success_count}/{len(target_lights)}"

    elif action == "off":
        success_count = sum(1 for lid, _ in target_lights if control_hue_light(lid, {"on": False}))
        names = ", ".join([n for _, n in target_lights])
        return f"Turned off: {names}" if success_count == len(target_lights) else f"Turned off {success_count}/{len(target_lights)}"

    elif action == "brightness":
        if brightness is None:
            return "Please specify brightness level (0-100)"
        bri_value = max(0, min(254, int((brightness / 100) * 254)))
        success_count = sum(1 for lid, _ in target_lights if control_hue_light(lid, {"on": True, "bri": bri_value}))
        names = ", ".join([n for _, n in target_lights])
        return f"Set {names} to {brightness}%" if success_count == len(target_lights) else f"Adjusted {success_count}/{len(target_lights)}"

    elif action == "color":
        if color is None:
            return "Please specify a color"
        color_lower = color.lower()
        if color_lower not in COLOR_MAP:
            available = ", ".join(COLOR_MAP.keys())
            return f"Unknown color '{color}'. Available: {available}"
        color_settings = COLOR_MAP[color_lower].copy()
        color_settings["on"] = True
        success_count = sum(1 for lid, _ in target_lights if control_hue_light(lid, color_settings))
        names = ", ".join([n for _, n in target_lights])
        return f"Set {names} to {color}" if success_count == len(target_lights) else f"Changed {success_count}/{len(target_lights)}"

    return "Unknown action"


# ===== OTHER TOOL FUNCTIONS =====

def get_weather_data(location: str) -> str:
    """Get weather data."""
    try:
        url = f"https://wttr.in/{location}?format=j1"
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            data = response.json()
            current = data['current_condition'][0]
            temp_c = current['temp_C']
            temp_f = current['temp_F']
            weather_desc = current['weatherDesc'][0]['value']
            humidity = current['humidity']
            wind_speed = current['windspeedKmph']
            location_name = data['nearest_area'][0]['areaName'][0]['value']
            return f"Weather in {location_name}: {weather_desc}, {temp_c}°C ({temp_f}°F), Humidity: {humidity}%, Wind: {wind_speed} km/h"
        return f"Could not get weather for '{location}'"
    except Exception as e:
        return f"Weather error: {str(e)}"


# ===== SEARCH LIMITS, CACHE, AND WEB SEARCH (patched) =====
try:
    SEARCH_MAX_PER_MINUTE
except NameError:
    import threading, time, os
    from collections import deque
    SEARCH_MAX_PER_MINUTE = int(os.getenv("SEARCH_MAX_PER_MINUTE", "8"))
    SEARCH_CACHE_TTL_SEC = int(os.getenv("SEARCH_CACHE_TTL_SEC", "21600"))
    SEARCH_RESULTS_PER_QUERY = int(os.getenv("SEARCH_RESULTS_PER_QUERY", "5"))
    _SEARCH_TIMESTAMPS = deque(maxlen=64)
    _SEARCH_CACHE = {}
    _SEARCH_LOCK = threading.Lock()

def _search_budget_ok():
    now = time.time()
    while _SEARCH_TIMESTAMPS and (now - _SEARCH_TIMESTAMPS[0]) > 60:
        _SEARCH_TIMESTAMPS.popleft()
    return len(_SEARCH_TIMESTAMPS) < SEARCH_MAX_PER_MINUTE

def _record_search():
    _SEARCH_TIMESTAMPS.append(time.time())

def _get_cached(query):
    q = query.strip().lower()
    item = _SEARCH_CACHE.get(q)
    if not item:
        return None
    exp, value = item
    if time.time() < exp:
        return value
    _SEARCH_CACHE.pop(q, None)
    return None

def _set_cached(query, value):
    _SEARCH_CACHE[query.strip().lower()] = (time.time() + SEARCH_CACHE_TTL_SEC, value)

# Command scaffolding that detectors and models let through to the search box.
# "search the web and tell me who is left in the fifa work cup" reached the
# engine VERBATIM (2026-07-09), returned math-calculator junk, and Blue then
# narrated the wreckage. Engines want the subject, not the errand.
_SEARCH_SCAFFOLD_LEAD_RE = re.compile(
    r'^(?:and|then|also|so|okay|ok|please|kindly|can you|could you|would you|will you|'
    r'i want you to|i need you to|go|now|just|quickly|try again|'
    r'tell me|tell us|let me know|find out|check|'
    r'search(?: the (?:web|internet|net))?|google|look up|find|'
    r'for me|for us|about|on|for)\s+', re.I)


def _clean_search_query(q: str) -> str:
    """Strip leading command scaffolding and trailing errands from a query."""
    s = re.sub(r'\s+', ' ', q or '').strip()
    # Drop an initial capability/meta question ("can you retrieve information in
    # real time?") and keep the actual request that follows. This exact shape
    # poisoned the World Cup search query on 2026-07-09.
    if "?" in s:
        parts = [p.strip() for p in re.split(r'\?+', s) if p.strip()]
        meta_re = re.compile(
            r'\b(?:retrieve|get|access|have|use|search|look up)\b.*'
            r'\b(?:real[- ]?time|live|current|information|info)\b',
            re.I,
        )
        if len(parts) >= 2 and meta_re.search(parts[0]):
            s = parts[-1]
    s = re.sub(
        r'^(?:retrieve|get|access)\s+(?:information|info)\s+'
        r'(?:in\s+)?(?:real[- ]?time|live)\s*',
        '',
        s,
        flags=re.I,
    ).strip()
    prev = None
    while prev != s:
        prev = s
        s = _SEARCH_SCAFFOLD_LEAD_RE.sub('', s, count=1).strip()
    # Trailing errand: "..., and tell me about it" / "and let me know please"
    s = re.sub(r'[,.;]?\s*(?:and|then)\s+(?:tell|let)\s+(?:me|us)\b.*$', '', s, flags=re.I)
    s = re.sub(r'\s*(?:for me|please)\s*[.!?]*$', '', s, flags=re.I)
    s = re.sub(r'\bfifa\s+world\s+cup\s+championship\b', 'FIFA World Cup', s, flags=re.I)
    s = re.sub(r'[\[\]{}()]+', ' ', s)
    s = re.sub(r'\s+', ' ', s)
    s = s.strip(' ?.!,;:')
    return s or (q or '').strip()


_LIVE_INFO_TOPICS = (
    'world cup', 'fifa', 'olympics', 'playoffs', 'the finals',
    'champions league', 'super bowl', 'stanley cup', 'wimbledon',
    'grand slam', 'election', 'tournament', 'euros', 'nba finals',
    'the match', 'the game',
)
_LIVE_INFO_QUESTIONS = (
    'who is left', "who's left", 'who won', 'who is winning',
    "who's winning", 'who advanced', 'who qualified', 'still in',
    'knocked out', 'eliminated', 'what happened', 'who plays',
    'who is playing', "who's playing", 'what teams are left',
    'which teams are left', 'standings', 'scores', 'results',
)
_LIVE_INFO_TEMPORAL = (
    'today', 'tonight', 'now', 'currently', 'current', 'latest',
    'live', 'real time', 'real-time', 'this week',
)


def _intent_text(message) -> str:
    """The user's actual ask with any attached-document block stripped.

    Intent/tool detection must NEVER run on attached text: a book excerpt
    containing the word "playing" started REAL music in the house, and
    another force-ran a garbage web_search (both 2026-07-10). Attachments
    ride as '[Attached document: name]' followed by a triple-quoted block,
    with the user's own words before or after it.
    """
    if not isinstance(message, str):
        return ""
    if '[attached document:' not in message.lower():
        return message
    cleaned = re.sub(r'\[attached document:[^\]]*\]\s*""".*?"""', ' ',
                     message, flags=re.I | re.S)
    if '[attached document:' in cleaned.lower():
        # Unclosed quoting — keep only what precedes the attachment.
        cleaned = re.split(r'\[attached document:', cleaned, flags=re.I)[0]
    return cleaned.strip()


def _live_info_query_from_message(message: str, history: List[Dict] = None) -> str:
    """Return a web-search query for live/current-event asks, else ''.

    This is a deterministic backstop for questions the model must not answer
    from memory. It catches "who is left in the FIFA World Cup" even if the
    selector path or LLM tool calling slips, and it resolves short follow-ups
    ("yes", "okay try again") against the prior live-event question.
    """
    raw = message if isinstance(message, str) else ""
    text = raw.strip()
    low = text.lower()
    history = history or []

    # A document attachment is never a live-info question, and a real live
    # ask is a short question: thousands of chars of extracted book text
    # keyword-match almost any topic list (seen live 2026-07-10 — attaching
    # Pasquinelli.pdf force-ran a garbage web_search on the book's header).
    if "[attached document:" in low or len(low) > 600:
        return ""

    def _with_date(q: str) -> str:
        q = _clean_search_query(q)
        if not q:
            return ""
        try:
            today = __import__('datetime').datetime.now().date().isoformat()
        except Exception:
            today = ""
        if today and today not in q:
            q = f"{q} {today}"
        return q[:160]

    def _is_live_ask(s: str) -> bool:
        return (
            any(topic in s for topic in _LIVE_INFO_TOPICS)
            and (
                any(q in s for q in _LIVE_INFO_QUESTIONS)
                or any(t in s for t in _LIVE_INFO_TEMPORAL)
            )
        )

    if _is_live_ask(low):
        return _with_date(text)

    # Short continuations after Blue offers to look up current info should
    # actually perform the lookup, not answer from the prior mistaken text.
    if re.fullmatch(r"\s*(yes|yeah|yep|sure|ok|okay|please|try again|go ahead)\s*[.!?]*\s*", low):
        prev_user = ""
        prev_assistant = ""
        skipped_current_user = False
        for m in reversed(history):
            role = m.get("role")
            content = m.get("content")
            if not isinstance(content, str):
                continue
            if role == "assistant" and not prev_assistant:
                prev_assistant = content
            elif role == "user" and not prev_user:
                if not skipped_current_user and content.strip().lower() == low:
                    skipped_current_user = True
                    continue
                prev_user = content
            if prev_user and prev_assistant:
                break
        if _is_live_ask((prev_user or "").lower()) and re.search(
            r"\b(?:look up|search|check|live|real[- ]?time|current|latest)\b",
            prev_assistant or "",
            re.I,
        ):
            return _with_date(prev_user)

    return ""


_SEARCH_RELEVANCE_STOP = {
    'the', 'and', 'for', 'with', 'from', 'that', 'this', 'these', 'those',
    'what', 'who', 'whom', 'whose', 'which', 'when', 'where', 'why', 'how',
    'many', 'much', 'are', 'was', 'were', 'has', 'have', 'had', 'does', 'did',
    'can', 'could', 'would', 'should', 'will', 'you', 'your', 'his', 'her',
    'its', 'their', 'out', 'about', 'into', 'over', 'under', 'between',
    'still', 'left', 'now', 'today', 'currently', 'latest', 'current',
    'recent', 'news',
}


def _search_results_relevant(q: str, results) -> bool:
    """Do the hits actually concern the query's subject? At least half of the
    query's distinctive words must appear somewhere in the titles+snippets —
    the gate that catches a scaffolding-polluted or mistyped search."""
    words = {w for w in re.findall(r"[a-z0-9][a-z0-9'\-]{2,}", (q or '').lower())
             if w not in _SEARCH_RELEVANCE_STOP}
    if not words:
        return True
    blob = " ".join((r.get('title') or '') + ' ' + (r.get('snippet') or '')
                    for r in (results or [])).lower()
    hits = sum(1 for w in words if w in blob)
    return hits >= max(1, (len(words) + 1) // 2)


def _llm_search_query_rewrite(original: str) -> str:
    """One small LLM call to turn a conversational ask into a clean engine
    query — this is also what fixes typos ('fifa work cup' → 'FIFA World
    Cup'). Returns '' on failure so callers can skip the retry."""
    try:
        res = call_llm(
            [{"role": "system", "content":
              "You write web search queries. Answer with ONLY the query — 2 to 7 words, "
              "no quotes, no punctuation, no commentary. Fix obvious typos."},
             {"role": "user", "content":
              "Turn this request into the best short web search query:\n" + (original or '')}],
            include_tools=False, temperature=0.2, max_tokens=600)
        ch = (res or {}).get('choices') or []
        cand = ((ch[0].get('message') or {}).get('content') or '') if ch else ''
        if '</think>' in cand:
            cand = cand.split('</think>')[-1]
        lines = [ln.strip().strip('"\'') for ln in cand.replace('<think>', '').strip().splitlines() if ln.strip()]
        cand = lines[0] if lines else ''
        if 2 <= len(cand) <= 90:
            return cand
    except Exception as e:
        log.warning(f"[SEARCH] query rewrite failed: {e}")
    return ""


def _run_search_providers(q: str):
    """The raw provider pass (ddgs, then the HTML endpoint). Returns
    (results, used_provider, error) — error is set only when a provider
    failed outright, not for a clean empty result."""
    from urllib.parse import quote_plus
    results = []
    used_provider = None

    # Preferred library path
    try:
        from ddgs import DDGS  # UPDATED: Changed from duckduckgo_search to ddgs
        used_provider = "ddgs.DDGS"
        with DDGS() as ddgs:
            for i, r in enumerate(ddgs.text(q, region="ca-en", max_results=SEARCH_RESULTS_PER_QUERY)):
                title = (r.get("title") or "").strip() or "Untitled"
                href = (r.get("href") or r.get("link") or "").strip()
                snippet = (r.get("body") or r.get("description") or "").strip()
                if href:
                    results.append({
                        "position": i + 1,
                        "title": title,
                        "url": href,
                        "snippet": snippet
                    })
        if not results:
            used_provider = None
    except Exception as e:
        print(f"   [WARN] ddgs search failed: {e.__class__.__name__}: {e}")
        used_provider = None

    # Fallback HTML endpoint (no JS)
    if not results:
        try:
            import requests
            from bs4 import BeautifulSoup  # type: ignore
            used_provider = "duckduckgo html"
            url = f"https://html.duckduckgo.com/html/?q={quote_plus(q)}"
            resp = requests.get(url, timeout=10, headers={"User-Agent": "Mozilla/5.0 (compatible; BlueBot/1.0)"})
            if resp.status_code == 429:
                return [], None, "[PROVIDER LIMIT] The search provider is rate-limiting right now. Please retry in a minute."
            resp.raise_for_status()
            soup = BeautifulSoup(resp.text, "html.parser")
            items = soup.select(".result__body")
            for i, item in enumerate(items[:SEARCH_RESULTS_PER_QUERY]):
                a = item.select_one("a.result__a")
                if not a:
                    continue
                title = a.get_text(strip=True) or "Untitled"
                href = a.get("href", "")
                snippet_el = item.select_one(".result__snippet")
                snippet = snippet_el.get_text(" ", strip=True) if snippet_el else ""
                if href:
                    results.append({
                        "position": i + 1,
                        "title": title,
                        "url": href,
                        "snippet": snippet
                    })
        except Exception as e:
            return [], None, f"Web search failed: {e.__class__.__name__}: {e}"

    return results, used_provider, None


def execute_web_search(query: str) -> str:
    """Execute a web search with caching + rate limiting and graceful provider backoff. Returns JSON.

    The query is scrubbed of command scaffolding first; a search whose results
    don't match the query's subject is retried ONCE with an LLM-rewritten
    query (which also fixes typos). If results still look unrelated, the
    payload carries a warning so the model says so instead of summarizing
    junk. Centralized here so every caller benefits: the model's own tool
    calls, the selector's force-executed params, the chat research toggle,
    and the anti-hallucination forced search."""
    if not query or not query.strip():
        return json.dumps({
            "success": False,
            "error": "Please provide a search query."
        })

    raw = query.strip()
    q = _clean_search_query(raw)
    if q != raw:
        print(f"   [SEARCH] cleaned query: {raw!r} -> {q!r}")

    # Cap absurdly long queries — LLMs sometimes dump every keyword they know
    MAX_QUERY_LEN = 120
    if len(q) > MAX_QUERY_LEN:
        # Keep only the first few meaningful words
        words = q.split()
        truncated = []
        for w in words:
            truncated.append(w)
            if len(' '.join(truncated)) >= MAX_QUERY_LEN:
                break
        q = ' '.join(truncated)
        print(f"   [SEARCH] Truncated long query to {len(q)} chars")

    cache_key = q
    with _SEARCH_LOCK:
        cached = _get_cached(cache_key)
        if cached is not None:
            return cached
        if not _search_budget_ok():
            return json.dumps({
                "success": False,
                "error": "[RATE LIMIT] You've run out of web searches for the moment. Please wait ~60 seconds and try again. Tip: identical queries are cached for 6 hours."
            })
        _record_search()

    results, used_provider, err = _run_search_providers(q)

    # Bad harvest → one reformulated retry. The rewrite sees the ORIGINAL ask
    # (more signal than the cleaned query) and fixes typos along the way.
    note = ""
    if err is None and (not results or not _search_results_relevant(q, results)):
        q2 = _llm_search_query_rewrite(raw)
        if q2 and q2.lower() != q.lower():
            with _SEARCH_LOCK:
                budget_ok = _search_budget_ok()
                if budget_ok:
                    _record_search()
            if budget_ok:
                print(f"   [SEARCH] weak results — retrying with rewritten query: {q2!r}")
                r2, p2, e2 = _run_search_providers(q2)
                if e2 is None and r2 and (_search_results_relevant(q2, r2) or not results):
                    results, used_provider = r2, p2
                    note = f"query was reformulated to '{q2}'"
                    q = q2

    if err and not results:
        msg = json.dumps({
            "success": False,
            "error": err
        })
        _set_cached(cache_key, msg)
        return msg

    if not results:
        msg = json.dumps({
            "success": False,
            "query": q,
            "error": "No results found."
        })
        _set_cached(cache_key, msg)
        return msg

    # Return proper JSON with success field
    payload_obj = {
        "success": True,
        "query": q,
        "provider": used_provider or "unknown",
        "results": results,
        "result_count": len(results)
    }
    if note:
        payload_obj["note"] = note
    if not _search_results_relevant(q, results):
        payload_obj["warning"] = (
            "These results may not actually answer the question. Tell the user honestly "
            "what you could and couldn't find — do not summarize unrelated links as if "
            "they were the answer. You may call web_search once more with a sharper query.")
    payload = json.dumps(payload_obj, ensure_ascii=False)

    _set_cached(cache_key, payload)
    return payload
# ===== END patched web search =====


# ===== BROWSE WEBSITE TOOL (moved here so it's available to execute_tool) =====
import re as _re
import html as _html
import json as _json
from typing import Optional

# HTML cleaning patterns
_SCRIPT_STYLE = _re.compile(r"(?is)<(script|style)\b.*?>.*?</\1>")
_TAGS = _re.compile(r"(?s)<[^>]+>")
_MULTI_WS = _re.compile(r"[ \t\r\f\v]+")
_MULTI_NL = _re.compile(r"\n{3,}")
_TITLE_RE = _re.compile(r"(?is)<title[^>]*>(.*?)</title>")
_LINK_RE = _re.compile(r'(?i)href=["\'](.*?)["\']')

def _clean_html_to_text(html_str: str, max_chars: int = 8000) -> str:
    """Clean HTML and convert to readable text."""
    if not isinstance(html_str, str):
        html_str = str(html_str or "")
    # remove script/style
    s = _SCRIPT_STYLE.sub(" ", html_str)
    # extract title separately if needed
    title = None
    mt = _TITLE_RE.search(s)
    if mt:
        title = _html.unescape(mt.group(1).strip())
    # remove tags
    s = _TAGS.sub("\n", s)
    s = _html.unescape(s)
    # collapse whitespace
    s = _MULTI_WS.sub(" ", s)
    s = s.replace("\r\n", "\n").replace("\r", "\n")
    s = _MULTI_NL.sub("\n\n", s)
    s = s.strip()
    if max_chars and len(s) > max_chars:
        s = s[:max_chars].rstrip() + "…"
    if title and title not in s[:500]:
        s = f"{title}\n\n{s}"
    return s

def _extract_links(html_str: str, base_url: str, max_links: int = 40) -> list:
    """Extract links from HTML."""
    import urllib.parse as _urlparse2
    out = []
    seen = set()
    for m in _LINK_RE.finditer(html_str or ""):
        href = m.group(1).strip()
        if not href:
            continue
        href_abs = _urlparse2.urljoin(base_url, href)
        if not href_abs.startswith(("http://","https://")):
            continue
        if href_abs in seen:
            continue
        seen.add(href_abs)
        out.append(href_abs)
        if len(out) >= max_links:
            break
    return out

def _safe_fetch_url(url: str, headers: Optional[dict] = None, timeout: int = 15, max_bytes: int = 1_500_000):
    """Safely fetch a URL with size limits."""
    import requests as _requests
    import urllib.parse as _urlparse3
    if not isinstance(url, str):
        raise ValueError("url must be a string")
    u = url.strip()
    if not u.startswith(("http://","https://")):
        u = "https://" + u
    parts = _urlparse3.urlsplit(u)
    if not parts.netloc:
        raise ValueError("URL must be absolute")
    req_headers = {
        "User-Agent": "BlueBot/1.0",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8"
    }
    if isinstance(headers, dict):
        req_headers.update({str(k): str(v) for k,v in headers.items()})
    resp = _requests.get(u, headers=req_headers, timeout=timeout, stream=True, allow_redirects=True)
    resp.raise_for_status()
    content = b""
    for chunk in resp.iter_content(chunk_size=16384):
        if chunk:
            content += chunk
            if len(content) > max_bytes:
                break
    return resp.headers.get("content-type",""), content

def _execute_browse_website(args: dict) -> str:
    """Execute the browse_website tool."""
    import requests

    url = (args or {}).get("url", "")
    extract = (args or {}).get("extract", "text") or "text"
    max_chars = int((args or {}).get("max_chars", 8000) or 8000)
    include_links = bool((args or {}).get("include_links", True))
    headers = (args or {}).get("headers", None)

    try:
        print(f"   [BROWSE] Fetching URL: {url}")
        ctype, content = _safe_fetch_url(url, headers=headers)
        html_raw = content.decode("utf-8", errors="ignore")

        if extract == "html":
            body = html_raw[: max_chars] + ("…" if len(html_raw) > max_chars else "")
        else:
            body = _clean_html_to_text(html_raw, max_chars=max_chars)

        result = {
            "url": url,
            "content_type": ctype,
            "extract": extract,
            "text": body,
            "success": True
        }
        if include_links:
            result["links"] = _extract_links(html_raw, url, max_links=40)

        print(f"   [BROWSE] Successfully fetched {len(body)} characters from {url}")
        return _json.dumps(result, ensure_ascii=False)

    except requests.exceptions.Timeout:
        error_msg = f"Timeout: The website {url} took too long to respond (>15 seconds)."
        print(f"   [ERROR] {error_msg}")
        return _json.dumps({"error": error_msg, "url": url, "success": False})

    except requests.exceptions.ConnectionError:
        error_msg = f"Connection Error: Could not connect to {url}. The website may be down or unreachable."
        print(f"   [ERROR] {error_msg}")
        return _json.dumps({"error": error_msg, "url": url, "success": False})

    except requests.exceptions.HTTPError as e:
        error_msg = f"HTTP Error {e.response.status_code}: {url} returned an error. The page may not exist or access may be denied."
        print(f"   [ERROR] {error_msg}")
        return _json.dumps({"error": error_msg, "url": url, "success": False})

    except ValueError as e:
        error_msg = f"Invalid URL: {str(e)}"
        print(f"   [ERROR] {error_msg}")
        return _json.dumps({"error": error_msg, "url": url, "success": False})

    except Exception as e:
        error_msg = f"Unexpected error while browsing {url}: {str(e)}"
        print(f"   [ERROR] {error_msg}")
        return _json.dumps({"error": error_msg, "url": url, "success": False})
# ===== END BROWSE WEBSITE TOOL =====


# ===== GMAIL TOOLS =====


def _execute_read_gmail(args: Dict[str, Any]) -> str:
    """Read and search Gmail messages"""
    if not GMAIL_AVAILABLE:
        return json.dumps({
            "error": "Gmail libraries not installed. Install with: pip install google-auth google-auth-oauthlib google-api-python-client",
            "success": False
        })

    try:
        service = get_gmail_service()
        query = args.get("query", "")
        max_results = args.get("max_results", 10)
        include_body = args.get("include_body", True)

        # Search for messages
        results = service.users().messages().list(
            userId='me',
            q=query,
            maxResults=max_results
        ).execute()

        messages = results.get('messages', [])

        if not messages:
            return json.dumps({
                "emails": [],
                "count": 0,
                "message": "No emails found matching the criteria",
                "success": True,
                "note": "EMAIL ACCESS SUCCESSFUL! The inbox was checked but no emails matched your search. This is REAL data."
            })

        email_list = []
        for msg in messages:
            msg_data = service.users().messages().get(
                userId='me',
                id=msg['id'],
                format='full'
            ).execute()

            # Extract headers
            headers = msg_data['payload']['headers']
            subject = next((h['value'] for h in headers if h['name'].lower() == 'subject'), 'No Subject')
            sender = next((h['value'] for h in headers if h['name'].lower() == 'from'), 'Unknown')
            date = next((h['value'] for h in headers if h['name'].lower() == 'date'), '')
            to = next((h['value'] for h in headers if h['name'].lower() == 'to'), '')

            email_info = {
                "id": msg['id'],
                "subject": subject,
                "from": sender,
                "to": to,
                "date": date,
                "snippet": msg_data.get('snippet', '')
            }

            # Extract body if requested
            if include_body:
                body = ""
                if 'parts' in msg_data['payload']:
                    for part in msg_data['payload']['parts']:
                        if part['mimeType'] == 'text/plain' and 'data' in part['body']:
                            body = base64.urlsafe_b64decode(part['body']['data']).decode('utf-8')
                            break
                elif 'body' in msg_data['payload'] and 'data' in msg_data['payload']['body']:
                    body = base64.urlsafe_b64decode(msg_data['payload']['body']['data']).decode('utf-8')

                email_info['body'] = body

            email_list.append(email_info)

        return json.dumps({
            "emails": email_list,
            "count": len(email_list),
            "query": query if query else "recent emails",
            "success": True,
            "note": "EMAIL ACCESS SUCCESSFUL! The emails above are REAL data from the Gmail inbox. Present this information to the user."
        }, indent=2)

    except Exception as e:
        return json.dumps({
            "error": f"Failed to read Gmail: {str(e)}",
            "success": False
        })


# ================================================================================
# AUTO-REPLY CONSTANTS & HELPERS
# ================================================================================

# Every email Blue writes is BCC'd here so the user has an audit copy in
# their inbox of everything that went out under their name.
BLUE_BCC_EMAIL = "alevant1905@gmail.com"

# Gmail label applied to inbound messages after Blue has answered them, so
# the auto-reply loop never replies to the same email twice.
BLUE_REPLIED_LABEL = "BlueReplied"

_BLUE_BCC_RE = re.compile(r"\balevant1905@gmail\.com\b", re.IGNORECASE)


def _ensure_blue_bcc(bcc: str) -> str:
    """Combine a user-supplied BCC list with the always-BCC address (deduped)."""
    if not bcc:
        return BLUE_BCC_EMAIL
    if _BLUE_BCC_RE.search(bcc):
        return bcc
    return bcc.rstrip(", ").rstrip() + ", " + BLUE_BCC_EMAIL


def _execute_send_gmail(args: Dict[str, Any]) -> str:
    """Send an email via Gmail with optional attachments"""
    if not GMAIL_AVAILABLE:
        return json.dumps({
            "error": "Gmail libraries not installed. Install with: pip install google-auth google-auth-oauthlib google-api-python-client",
            "success": False
        })

    try:
        service = get_gmail_service()

        to = args.get("to", "")
        subject = args.get("subject", "")
        body = args.get("body", "")
        cc = args.get("cc", "")
        # Always BCC the user so they keep a copy of everything Blue sends.
        bcc = _ensure_blue_bcc(args.get("bcc", ""))
        attachments = args.get("attachments", [])  # NEW: List of file paths

        if not to or not subject:
            return json.dumps({
                "error": "Missing required email information. Need: recipient email address (name@domain.com) and subject. Please ask user to provide both.",
                "missing_to": not to,
                "missing_subject": not subject,
                "success": False,
                "instruction": "Tell the user you need their email address or subject line. Example: 'I need the recipient's email address to send this message.'"
            })

        # Create message
        message = MIMEMultipart()
        message['To'] = to
        message['Subject'] = subject

        if cc:
            message['Cc'] = cc
        if bcc:
            message['Bcc'] = bcc

        # Add body
        message.attach(MIMEText(body, 'plain'))

        # Process attachments
        attached_files = []
        attachment_errors = []

        if attachments:
            print(f"   [ATTACH] Processing {len(attachments)} attachment(s)")

            for file_path in attachments:
                try:
                    # Normalize path
                    file_path = file_path.strip()

                    # Check if file exists
                    if not os.path.exists(file_path):
                        # Try in documents folder
                        doc_path = os.path.join(str(UPLOAD_FOLDER), file_path)
                        if os.path.exists(doc_path):
                            file_path = doc_path
                        else:
                            # Also try DOCUMENTS_FOLDER
                            doc_path2 = os.path.join(DOCUMENTS_FOLDER, file_path)
                            if os.path.exists(doc_path2):
                                file_path = doc_path2
                            else:
                                attachment_errors.append(f"File not found: {file_path}")
                                print(f"   [ATTACH-ERROR] File not found: {file_path}")
                                continue

                    # Check file size (limit to 25MB - Gmail limit)
                    file_size = os.path.getsize(file_path)
                    max_size = 25 * 1024 * 1024  # 25MB in bytes

                    if file_size > max_size:
                        attachment_errors.append(f"File too large: {os.path.basename(file_path)} ({file_size / (1024*1024):.1f}MB > 25MB)")
                        print(f"   [ATTACH-ERROR] File too large: {file_path}")
                        continue

                    # Get filename
                    filename = os.path.basename(file_path)

                    # Guess MIME type
                    mime_type, _ = mimetypes.guess_type(file_path)
                    if mime_type is None:
                        mime_type = 'application/octet-stream'

                    # Split MIME type
                    main_type, sub_type = mime_type.split('/', 1)

                    # Read file and attach
                    with open(file_path, 'rb') as f:
                        file_data = f.read()

                    # Create attachment
                    attachment = MIMEBase(main_type, sub_type)
                    attachment.set_payload(file_data)
                    encoders.encode_base64(attachment)
                    attachment.add_header('Content-Disposition', f'attachment; filename={filename}')

                    # Attach to message
                    message.attach(attachment)

                    attached_files.append({
                        'filename': filename,
                        'size': file_size,
                        'type': mime_type
                    })
                    print(f"   [ATTACH-OK] Attached: {filename} ({file_size / 1024:.1f}KB)")

                except Exception as e:
                    error_msg = f"Error attaching {os.path.basename(file_path) if file_path else 'file'}: {str(e)}"
                    attachment_errors.append(error_msg)
                    print(f"   [ATTACH-ERROR] {error_msg}")

        # Encode message
        raw_message = base64.urlsafe_b64encode(message.as_bytes()).decode('utf-8')

        # Send message
        sent_message = service.users().messages().send(
            userId='me',
            body={'raw': raw_message}
        ).execute()

        # Build result
        result = {
            "message_id": sent_message['id'],
            "to": to,
            "subject": subject,
            "success": True,
            "message": f"✅ Email sent successfully to {to}",
            "confirmation": f"Email delivered to {to} with subject '{subject}'",
            "note": "EMAIL SENT SUCCESSFULLY! You MUST confirm this to the user by saying 'I sent the email to [address]' or similar."
        }

        # Add attachment info
        if attached_files:
            result['attachments'] = attached_files
            result['attachments_count'] = len(attached_files)
            result['message'] += f" with {len(attached_files)} attachment(s)"
            result['note'] += f" ATTACHMENTS: {', '.join([f['filename'] for f in attached_files])}"

        if attachment_errors:
            result['attachment_errors'] = attachment_errors
            result['warning'] = f"Email sent but {len(attachment_errors)} attachment(s) failed"

        return json.dumps(result, indent=2)

    except Exception as e:
        return json.dumps({
            "error": f"Failed to send email: {str(e)}",
            "success": False
        })


def _execute_email_snapshot(args: Dict[str, Any]) -> str:
    """Composite 'see it and send it': take a BRAND NEW camera photo and email
    it as an attachment from Blue's own Gmail. One deterministic call — the
    local model can't be trusted to chain capture_camera then send_gmail.
    No recipient means Alex ('email me what you see')."""
    import datetime as _dt

    # 1) Capture. This also queues the frame in _vision_queue, so Blue sees
    #    (and can describe) exactly the photo he just mailed.
    capture_raw = capture_camera_image(
        look=args.get("look"),
        zoom=args.get("zoom"),
        zoom_region=args.get("zoom_region") or "center",
    )
    try:
        capture = json.loads(capture_raw)
    except (TypeError, ValueError):
        capture = {"success": False, "error": "Unreadable capture result."}
    if not capture.get("success"):
        return json.dumps({
            "success": False,
            "error": f"Could not take the snapshot: {capture.get('error', 'camera unavailable')}",
            "_instruction": ("The camera shot failed, so NO email was sent. "
                             "Tell the user you couldn't take the photo right now."),
        })
    filepath = capture.get("filepath")

    # 2) Resolve the recipient: an address passes through, a contact name goes
    #    through the address book, and a confident answer is required — Blue
    #    must never guess and mail a stranger a photo of the house.
    to = (args.get("to") or "").strip()
    if to and "@" not in to:
        resolved = None
        if ENHANCED_TOOLS_AVAILABLE:
            try:
                resolved = ContactManager.resolve_email(to)
            except Exception as e:
                print(f"   [SNAPSHOT] contact resolve error for {to!r}: {e}")
        if not resolved:
            return json.dumps({
                "success": False,
                "error": f"No email address on file for '{to}'.",
                "photo_saved": filepath,
                "_instruction": (f"You took the photo (it follows in the next message) but you "
                                 f"don't have an email address for {to}, so nothing was sent. "
                                 f"Ask the user for the address."),
            })
        print(f"   [SNAPSHOT] resolved recipient {to!r} -> {resolved}")
        to = resolved
    if not to:
        to = BLUE_BCC_EMAIL  # Alex's personal gmail — the "email it to me" default

    # 3) Mail it, photo attached.
    now_word = _dt.datetime.now().strftime("%B %d, %Y at %I:%M %p")
    view = capture.get("view") or {}
    view_bits = []
    if view.get("look"):
        view_bits.append("head centered" if view["look"] == "center"
                         else f"head turned {view['look']}")
    try:
        if float(view.get("zoom") or 1) > 1.01:
            view_bits.append(f"{float(view['zoom']):g}x zoom on the {view.get('zoom_region') or 'center'}")
    except (TypeError, ValueError):
        pass
    view_note = (" (" + ", ".join(view_bits) + ")") if view_bits else ""

    note = (args.get("note") or "").strip()
    subject = (args.get("subject") or "").strip() or f"Snapshot from Blue \U0001F4F7 {now_word}"
    body_lines = [f"Hi! Here's what I'm seeing right now — taken {now_word}{view_note}."]
    if note:
        body_lines += ["", note]
    body_lines += ["", "— Blue"]

    send_raw = _execute_send_gmail({
        "to": to,
        "subject": subject,
        "body": "\n".join(body_lines),
        "attachments": [filepath],
    })
    try:
        send_res = json.loads(send_raw)
    except (TypeError, ValueError):
        send_res = {"success": False, "error": "Unreadable send result."}

    if not send_res.get("success"):
        return json.dumps({
            "success": False,
            "error": f"Photo captured but the email failed: {send_res.get('error', 'unknown error')}",
            "photo_saved": filepath,
            "_instruction": ("You took the photo but could NOT send the email. Tell the user "
                             "the email failed — do not claim it was sent."),
        })
    if send_res.get("attachment_errors"):
        # The message went out but WITHOUT the photo — never claim delivery.
        return json.dumps({
            "success": False,
            "error": f"Email sent but the photo failed to attach: {'; '.join(send_res['attachment_errors'])}",
            "photo_saved": filepath,
            "_instruction": ("An email went out but the photo could not be attached. "
                             "Tell the user plainly so they don't wait for a photo."),
        })

    print(f"   [SNAPSHOT] \U0001F4E7 Snapshot emailed to {to}")
    return json.dumps({
        "success": True,
        "message": f"\U0001F4F7\U0001F4E7 Snapshot taken and emailed to {to}",
        "to": to,
        "subject": subject,
        "filename": capture.get("filename"),
        "filepath": filepath,
        "view": view or None,
        "_instruction": (
            f"You just took a fresh photo and emailed it to {to} from your own Gmail — "
            "this is DONE. The photo follows in the next message: confirm you sent it and "
            "briefly describe what's in it. Do NOT call send_gmail or capture_camera again."
        ),
    })


def _execute_reply_gmail(args: Dict[str, Any]) -> str:
    """Reply to Gmail messages"""
    if not GMAIL_AVAILABLE:
        return json.dumps({
            "error": "Gmail libraries not installed. Install with: pip install google-auth google-auth-oauthlib google-api-python-client",
            "success": False
        })

    try:
        service = get_gmail_service()

        query = args.get("query", "")
        reply_body = args.get("reply_body", "")
        reply_all = args.get("reply_all", False)
        max_replies = min(args.get("max_replies", 10), 50)  # Cap at 50

        if not query or not reply_body:
            return json.dumps({
                "error": "Missing required fields: 'query' and 'reply_body' are required",
                "success": False
            })

        # Search for messages to reply to
        results = service.users().messages().list(
            userId='me',
            q=query,
            maxResults=max_replies if reply_all else 1
        ).execute()

        messages = results.get('messages', [])

        if not messages:
            return json.dumps({
                "success": True,
                "replies_sent": 0,
                "message": f"No emails found matching query: {query}",
                "query": query
            })

        replies_sent = []
        errors = []

        # Reply to each message
        for msg in messages:
            try:
                # Get the original message details
                msg_data = service.users().messages().get(
                    userId='me',
                    id=msg['id'],
                    format='full'
                ).execute()

                # Extract headers
                headers = msg_data['payload']['headers']
                original_subject = next((h['value'] for h in headers if h['name'].lower() == 'subject'), 'No Subject')
                original_from = next((h['value'] for h in headers if h['name'].lower() == 'from'), '')
                original_to = next((h['value'] for h in headers if h['name'].lower() == 'to'), '')
                message_id_header = next((h['value'] for h in headers if h['name'].lower() == 'message-id'), '')

                # Extract email body content
                email_body = ""
                if 'parts' in msg_data['payload']:
                    for part in msg_data['payload']['parts']:
                        if part['mimeType'] == 'text/plain' and 'data' in part['body']:
                            email_body = base64.urlsafe_b64decode(part['body']['data']).decode('utf-8')
                            break
                elif 'body' in msg_data['payload'] and 'data' in msg_data['payload']['body']:
                    email_body = base64.urlsafe_b64decode(msg_data['payload']['body']['data']).decode('utf-8')

                # Extract email address from "Name <email@domain.com>" format
                import re
                email_match = re.search(r'<(.+?)>', original_from)
                reply_to = email_match.group(1) if email_match else original_from

                # Create reply subject (add Re: if not already there)
                reply_subject = original_subject if original_subject.startswith('Re:') else f"Re: {original_subject}"

                # Create reply message
                reply_message = MIMEMultipart()
                reply_message['To'] = reply_to
                reply_message['Subject'] = reply_subject
                reply_message['Bcc'] = _ensure_blue_bcc(args.get("bcc", ""))
                reply_message['In-Reply-To'] = message_id_header
                reply_message['References'] = message_id_header

                # Add reply body
                reply_message.attach(MIMEText(reply_body, 'plain'))

                # Encode and send
                raw_message = base64.urlsafe_b64encode(reply_message.as_bytes()).decode('utf-8')

                sent_message = service.users().messages().send(
                    userId='me',
                    body={
                        'raw': raw_message,
                        'threadId': msg_data.get('threadId')  # Keep in same thread
                    }
                ).execute()

                # Once Blue has responded, mark the original as read so it
                # stops showing as a new email in the inbox.
                try:
                    service.users().messages().modify(
                        userId='me',
                        id=msg['id'],
                        body={'removeLabelIds': ['UNREAD']},
                    ).execute()
                except Exception as e:
                    print(f"   [WARN] could not mark {msg['id']} as read: {e}")

                replies_sent.append({
                    "original_subject": original_subject,
                    "original_body": email_body[:500] + "..." if len(email_body) > 500 else email_body,  # Include truncated body
                    "replied_to": reply_to,
                    "reply_id": sent_message['id'],
                    "reply_sent": reply_body
                })

                print(f"   [OK] Replied to: {original_subject} (from {reply_to})")

            except Exception as e:
                errors.append({
                    "message_id": msg['id'],
                    "error": str(e)
                })
                print(f"   [ERROR] Failed to reply to message {msg['id']}: {e}")

        return json.dumps({
            "success": True,
            "replies_sent": len(replies_sent),
            "query": query,
            "replies": replies_sent,
            "errors": errors if errors else None,
            "message": f"Successfully replied to {len(replies_sent)} email(s) matching '{query}'",
            "note": "REPLIES SENT SUCCESSFULLY! The replies have been delivered and threaded correctly."
        }, indent=2)

    except Exception as e:
        return json.dumps({
            "error": f"Failed to reply to emails: {str(e)}",
            "success": False
        })
# ================================================================================
# AUTO-REPLY: personal mail arriving at Blue's gmail inbox
# ================================================================================
#
# Blue has his own gmail account (alevantresearch@gmail.com). Anything in
# that inbox from a real person is, by definition, written to Blue — no
# greeting detector required. We just filter out:
#   • automated senders (no-reply, mailing lists, List-Unsubscribe headers)
#   • the OWNER's own addresses (so if Alex emails Blue from his personal
#     gmail or yorku address, Blue does not autonomously reply — that's
#     the chat channel, and a self-loop would be confusing)
#   • Promotions / Social / Updates / Forums categories (the gmail query
#     handles these)

# Blue's own gmail address — kept here for documentation/prompting; Blue's
# OAuth token decides which mailbox `userId='me'` actually points at.
BLUE_OWN_EMAIL = os.environ.get("BLUE_OWN_EMAIL", "alevantresearch@gmail.com")

_BLUE_SKIP_SENDER_RE = re.compile(
    r"(?:no[-_.]?reply|do[-_.]?not[-_.]?reply|noreply|postmaster|"
    r"mailer[-_.]?daemon|notifications?@|alerts?@|newsletter@|news@|"
    r"marketing@|updates?@)",
    re.IGNORECASE,
)

# Addresses that should never trigger an autonomous reply. Empty by
# default: now that Blue has his own inbox (alevantresearch@gmail.com),
# mail from Alex's other accounts (alevant1905, yorku) to that inbox is
# legitimate cross-account communication — that's the channel Alex uses
# to write to Blue. The BlueReplied label dedups so there's no infinite
# loop risk either way. Re-add specific addresses via the env var if a
# particular account starts forwarding receipts / noise that shouldn't
# get a reply.
BLUE_SELF_ADDRESSES = {
    a.strip().lower()
    for a in os.environ.get("BLUE_SELF_EMAILS", "").split(",")
    if a.strip()
}

# Fully-TRUSTED owner addresses. When a verified email comes from one of
# these, Blue gets the SAME full tool access he has in chat (send mail,
# control lights, set reminders, search the document library, …) instead
# of the read-only public-info whitelist used for everyone else. Default
# is Alex's personal gmail; extend via BLUE_OWNER_EMAILS.
BLUE_OWNER_ADDRESSES = {
    a.strip().lower()
    for a in os.environ.get(
        "BLUE_OWNER_EMAILS",
        "alevant1905@gmail.com,alevant@yorku.ca",
    ).split(",")
    if a.strip()
}


_EMAIL_ADDR_RE = re.compile(r"[\w.+\-]+@[\w.\-]+\.\w+")


def _email_sender_is_owner(headers: List[Dict[str, str]], sender: str) -> bool:
    """True only if the email is from a trusted owner address AND passes a
    basic anti-spoof check. From headers are forgeable, so before granting
    the elevated (full-tool) trust level we require Gmail's own
    Authentication-Results stamp to show a dkim/dmarc pass with no fail.
    Gmail adds this header on every inbound message; a spoofed From for a
    gmail.com address won't carry a valid DKIM signature."""
    addrs = {a.lower() for a in _EMAIL_ADDR_RE.findall(sender or "")}
    if not (addrs & BLUE_OWNER_ADDRESSES):
        return False
    auth = ""
    for h in headers or []:
        if h.get('name', '').lower() == 'authentication-results':
            auth += " " + (h.get('value') or "").lower()
    if not auth:
        # No auth stamp at all — be conservative, don't elevate.
        return False
    if 'dkim=fail' in auth or 'dmarc=fail' in auth:
        return False
    return ('dkim=pass' in auth) or ('dmarc=pass' in auth)


def _should_skip_sender(sender: str, headers: List[Dict[str, str]] = None) -> bool:
    """Skip automated mail, marketing lists, no-reply senders, and the
    user's own addresses (any address in BLUE_SELF_ADDRESSES) so Blue
    never replies to mail he sent himself."""
    if not sender:
        return True
    # Exact-match each address found in the From field against the
    # self-address set so substrings like "not-alevant@..." don't trip it.
    for found in _EMAIL_ADDR_RE.findall(sender):
        if found.lower() in BLUE_SELF_ADDRESSES:
            return True
    if _BLUE_SKIP_SENDER_RE.search(sender):
        return True
    if headers:
        for h in headers:
            name = h.get('name', '').lower()
            if name in ('list-unsubscribe', 'list-id', 'precedence', 'auto-submitted'):
                return True
    return False


def _get_or_create_blue_label(service) -> Optional[str]:
    """Return the Gmail label ID for BLUE_REPLIED_LABEL, creating it if needed."""
    try:
        labels = service.users().labels().list(userId='me').execute().get('labels', [])
        for lbl in labels:
            if lbl.get('name') == BLUE_REPLIED_LABEL:
                return lbl.get('id')
        created = service.users().labels().create(
            userId='me',
            body={
                'name': BLUE_REPLIED_LABEL,
                'labelListVisibility': 'labelShow',
                'messageListVisibility': 'show',
            },
        ).execute()
        return created.get('id')
    except Exception as e:
        print(f"   [AUTO-REPLY] could not get/create label: {e}")
        return None


# Tools Blue may use while drafting an autonomous email reply. Strictly
# read-only / public-info: enough to browse a link or look something up,
# but NOTHING outbound (send/reply email), physical (lights, music,
# camera), private (the document library), or state-changing — anyone can
# email Blue, so an email must never be able to make Blue act on the house
# or on Alex's private data.
_EMAIL_SAFE_TOOL_NAMES = {
    "browse_website", "web_search", "get_weather", "get_local_time",
    "search_scholar", "get_paper",  # public scholarly APIs, read-only.
    # read_paper is deliberately NOT here: it spends Alex's library access.
}

# Even a fully-trusted owner email shouldn't be able to make Blue kick off
# another inbox sweep from inside one — avoids recursion.
_EMAIL_OWNER_EXCLUDE = {"auto_reply_emails"}


def _email_safe_tools() -> List[Dict[str, Any]]:
    try:
        return [
            t for t in TOOLS  # noqa: F821
            if (t.get("function", {}) or {}).get("name") in _EMAIL_SAFE_TOOL_NAMES
        ]
    except Exception:
        return []


def _email_owner_tools() -> List[Dict[str, Any]]:
    """Full tool set (minus recursion-risk tools) for verified owner mail."""
    try:
        return [
            t for t in TOOLS  # noqa: F821
            if (t.get("function", {}) or {}).get("name") not in _EMAIL_OWNER_EXCLUDE
        ]
    except Exception:
        return []


def _inject_pending_vision(messages: List[Dict[str, Any]]) -> bool:
    """If a tool (e.g. capture_camera) queued an image, splice it into the
    email conversation as a base64 user message so the model can actually
    SEE it — mirrors what call_lm_studio does for the chat path. The email
    loop talks to call_llm/_LM.chat, which never touches _vision_queue, so
    without this Blue captures a photo by email but is shown nothing and
    invents a description. Returns True if an image was injected."""
    global _vision_queue
    try:
        if not _vision_queue.has_images():
            return False
        is_camera = any(img.is_camera_capture for img in _vision_queue.pending_images)
        parts: List[Dict[str, Any]] = []
        if is_camera:
            parts.append({
                "type": "text",
                "text": (
                    "[CAMERA IMAGE: Describe what you ACTUALLY see — who, what, "
                    "where, objects, lighting. Be specific and accurate, describe "
                    "only what's visible, then answer the sender's question.]"
                ),
            })
        else:
            parts.append({"type": "text", "text": "[Images to analyze:]"})
        for img in _vision_queue.pending_images:
            label = "[Your current view:]" if img.is_camera_capture else f"Image: {img.filename}"
            parts.append({"type": "text", "text": f"\n{label}"})
            encoded = encode_image_to_base64(img.filepath)
            if encoded:
                parts.append(encoded)
            else:
                print(f"   [AUTO-REPLY] vision: failed to encode {img.filename}")
        messages.append({"role": "user", "content": parts})
        print(f"   [AUTO-REPLY] injected {len(_vision_queue.pending_images)} vision image(s)")
        _vision_queue.mark_as_viewed()
        _vision_queue.clear()
        return True
    except Exception as e:
        print(f"   [AUTO-REPLY] vision injection error: {e}")
        return False


# When the OWNER emails Blue asking him to email a third party ("email Stella
# and tell her ..."), the reply loop is framed as "answer the sender", so the
# local model just writes the message back to Alex instead of calling
# send_gmail to Stella. Detect that intent deterministically and actually send.
_OUTBOUND_SEND_RE = re.compile(
    r"\b(?:"
    r"(?:send|write|shoot|fire\s+off)\s+(?:an?\s+)?(?:email|message|note|reply)\s+to\s+(?P<n1>[A-Za-z][A-Za-z.'-]*)"
    r"|(?:email|message|write\s+to|reach\s+out\s+to|get\s+in\s+touch\s+with)\s+(?P<n2>[A-Za-z][A-Za-z.'-]*)"
    r"|(?:tell|let|ask|remind)\s+(?P<n3>[A-Za-z][A-Za-z.'-]*)\s+(?:know|that|to\b|about)"
    r")",
    re.IGNORECASE,
)
# Pronouns / self-references that are never a third-party recipient.
_OUTBOUND_EXCLUDE_NAMES = {
    "me", "myself", "i", "you", "yourself", "us", "we", "them", "they",
    "him", "her", "he", "she", "blue", "everyone", "someone", "anybody",
    "somebody", "myself", "yours",
}


def _resolve_recipient_email(name: str, service) -> Optional[str]:
    """Map a first name / display name to an email address by looking at who
    has actually corresponded with Blue. Returns None if no confident match —
    callers must NOT guess, so Blue asks for the address instead of mailing a
    stranger."""
    name_n = (name or "").strip().lower()
    if not name_n:
        return None
    # Authoritative source first: Blue's own contacts database. Only fall back
    # to scanning Gmail history when the name isn't a known contact.
    if ENHANCED_TOOLS_AVAILABLE:
        try:
            hit = ContactManager.resolve_email(name)
            if hit:
                print(f"   [CONTACTS] resolved {name!r} -> {hit} (address book)")
                return hit
        except Exception as e:
            print(f"   [CONTACTS] resolve error for {name!r}: {e}")
    try:
        for q in (f'from:{name}', f'to:{name}'):
            resp = service.users().messages().list(
                userId='me', q=q, maxResults=5,
            ).execute()
            for ref in resp.get('messages', []):
                md = service.users().messages().get(
                    userId='me', id=ref['id'], format='metadata',
                    metadataHeaders=['From', 'To'],
                ).execute()
                for h in md.get('payload', {}).get('headers', []):
                    val = h.get('value', '') or ''
                    for mm in re.finditer(
                        r'(?:"?([^"<]*?)"?\s*)?<?([\w.+-]+@[\w.-]+\.\w+)>?', val
                    ):
                        disp = (mm.group(1) or '').strip().lower()
                        addr = mm.group(2).lower()
                        disp_tokens = set(re.findall(r"[a-z]+", disp))
                        if name_n in disp_tokens or name_n == addr.split('@')[0]:
                            return addr
    except Exception as e:
        print(f"   [AUTO-REPLY] recipient resolve error for {name!r}: {e}")
    return None


def _compose_outbound_email(recipient_name: str, owner_instruction: str) -> str:
    """Turn Alex's instruction ("tell her the girls are awake") into an actual
    warm email addressed TO the recipient, in Blue's voice."""
    sys_p = (
        "You are Blue, Alex's friendly robot companion. Alex has asked you to "
        f"send an email to {recipient_name}. Write that email directly TO "
        f"{recipient_name} in Blue's warm, natural voice, conveying what Alex "
        "wants said. Address them by name, keep it under 120 words, and sign "
        "it 'Blue'. Output ONLY the email body — no subject line, no headers."
    )
    usr_p = (
        f"Alex's instruction to you:\n{owner_instruction}\n\n"
        f"Write the email to {recipient_name} now."
    )
    res = call_llm(
        [{"role": "system", "content": sys_p}, {"role": "user", "content": usr_p}],
        include_tools=False, temperature=0.7, max_tokens=300,
    )
    choices = (res or {}).get('choices') or []
    if not choices:
        return ""
    return ((choices[0].get('message') or {}).get('content') or "").strip()


def _maybe_handle_owner_outbound(body: str) -> Optional[str]:
    """If owner mail asks Blue to email a third party, resolve + send it and
    return a confirmation to reply to the owner with. Returns None when there's
    no outbound-send intent (fall through to a normal reply)."""
    m = _OUTBOUND_SEND_RE.search(body or "")
    if not m:
        return None
    name = (m.group('n1') or m.group('n2') or m.group('n3') or '').strip()
    if not name or name.lower() in _OUTBOUND_EXCLUDE_NAMES:
        return None

    try:
        service = get_gmail_service()
    except Exception as e:
        print(f"   [AUTO-REPLY] outbound: no gmail service: {e}")
        return None

    recipient_email = _resolve_recipient_email(name, service)
    if not recipient_email:
        # An explicit address in the body is a fallback ("email Stella at x@y").
        am = re.search(r'[\w.+-]+@[\w.-]+\.\w+', body or "")
        if am:
            recipient_email = am.group(0)

    if not recipient_email:
        print(f"   [AUTO-REPLY] outbound: could not resolve recipient {name!r}")
        return (
            f"I'd love to pass that along to {name}, but I don't have an email "
            f"address for them yet. Send me their address and I'll fire it right "
            f"off.\n\nBlue"
        )

    blocked = set(BLUE_OWNER_ADDRESSES) | set(BLUE_SELF_ADDRESSES) | {
        (BLUE_OWN_EMAIL or "").lower()
    }
    if recipient_email.lower() in blocked:
        # The "recipient" is Alex or Blue himself — not a real third-party send.
        return None

    email_body = _compose_outbound_email(name, body) or body
    subject = f"A message from Alex"
    print(f"   [AUTO-REPLY] outbound send → {name} <{recipient_email}>")
    res = execute_tool("send_gmail", {
        "to": recipient_email, "subject": subject, "body": email_body,
    })
    ok = False
    try:
        ok = bool(json.loads(res).get('success'))
    except Exception:
        ok = '"success": true' in (res or "").lower()

    if ok:
        return (
            f"Done — I've emailed {name} for you (and BCC'd you a copy). "
            f"Here's what I sent:\n\n\"{email_body}\"\n\nBlue"
        )
    return (
        f"I tried to email {name} but the send didn't go through. Want me to "
        f"try again?\n\nBlue"
    )


def _resolve_document_file(reference: str) -> Optional[tuple]:
    """Match a natural document reference ("the syllabus", "Pasquinelli",
    "Mark's substack post") to a real file in the library index. Returns
    (filename, existing_filepath) for the best match, or None if nothing
    matches confidently — callers must then ASK which document, never guess
    and attach the wrong file."""
    best = _resolve_document_entry(reference)
    if not best:
        return None
    filename = str(best.get('filename', '') or '')
    filepath = _existing_document_path(best)
    return (filename, filepath) if filename and filepath else None


def _maybe_handle_owner_attachment(body: str, reply_to: str) -> Optional[tuple]:
    """If owner mail asks Blue to send THEM a library document as an
    attachment, resolve the file and return (reply_body, [filepath]) so the
    auto-reply carries it. Returns None when there's no attachment intent.

    The auto-reply only goes back to the original sender, so this handles the
    'send me X' case; a request aimed at a third party falls through (we must
    not silently mail the file to the wrong person)."""
    bl = (body or "").lower()
    wants_attach = (
        'attach' in bl
        or 'a copy of' in bl
        or bool(re.search(
            r"\b(?:send|email|forward)\b[^.?!]*\b"
            r"(?:document|file|pdf|docx?|syllabus|paper|copy)\b", bl))
    )

    doc = _resolve_document_file(body)

    # "send me the <docname>" without the literal word 'attachment' is still an
    # attachment request — but only treat it as one when it actually names a
    # real library document, so "send me the weather" stays a normal reply.
    if not wants_attach:
        if doc and re.search(r"\b(?:send|email|forward)\b[^.?!]*\bme\b", bl):
            wants_attach = True
        else:
            return None

    # Directed at a third party ("email Stella the syllabus")? This reply path
    # can only answer the sender, so don't risk mailing the file to Alex when
    # he meant Stella — fall through to the normal handler.
    m = _OUTBOUND_SEND_RE.search(body or "")
    if m:
        nm = (m.group('n1') or m.group('n2') or m.group('n3') or '').strip().lower()
        if nm and nm not in _OUTBOUND_EXCLUDE_NAMES:
            return None

    if not doc:
        return (
            "Happy to send that over — which document did you mean? Tell me the "
            "file name and I'll attach it right away.\n\nBlue",
            [],
        )
    filename, filepath = doc
    print(f"   [AUTO-REPLY] attaching {filename} → {reply_to}")
    return (
        f"Here you go — I've attached {filename} for you.\n\nBlue",
        [filepath],
    )


# Substantive written outputs Blue can be asked to produce (vs. a chatty
# reply). Detected so they bypass the "under 150 words" reply prompt.
_COMPOSITION_OUTPUT_RE = re.compile(
    r"\b(critical\s+assessments?|assessments?|critiques?|critical\s+reviews?|"
    r"book\s+reviews?|literature\s+reviews?|review\s+essays?|essays?|"
    r"analys[ie]s|commentar(?:y|ies)|appraisals?|evaluations?|"
    r"reflection\s+pieces?|response\s+(?:papers?|essays?)|position\s+papers?|"
    r"think\s+pieces?)\b", re.I)
_COMPOSITION_VERB_RE = re.compile(
    r"\b(write|draft|compose|produce|prepare|put\s+together|give\s+me|"
    r"generate|craft|critically\s+assess|assess|critique)\b", re.I)


def _gather_owner_work_sources(query: str, pub_folders, max_total: int = 18) -> tuple:
    """Pull a BROAD spread of passages — several of the owner's own articles
    (scoped to his publications folder), plus a wider pass to surface the
    target text. Returns (formatted_passages, [filenames])."""
    try:
        from blue.tools.rag import search_expertise as _se
    except Exception:
        return "", []

    chunks, seen = [], set()

    def _add(rs):
        for r in rs or []:
            key = (r.get('filename'), r.get('chunk_index'))
            if key not in seen:
                seen.add(key)
                chunks.append(r)

    if pub_folders:
        try:
            _add(_se(query, max_chunks=14, max_per_doc=2, folders=pub_folders))
        except Exception as e:
            print(f"   [AUTO-REPLY] composition scoped search error: {e}")
    try:
        _add(_se(query, max_chunks=8, max_per_doc=2))
    except Exception:
        pass

    chunks = chunks[:max_total]
    if not chunks:
        return "", []
    formatted = [
        f"[{i}] [{r.get('filename', '?')}]\n{(r.get('content') or '')[:1000]}"
        for i, r in enumerate(chunks, 1)
    ]
    return "\n\n".join(formatted), sorted({r.get('filename', '?') for r in chunks})


def _maybe_handle_owner_composition(body: str) -> Optional[str]:
    """Owner asked Blue to WRITE something substantive (critical assessment,
    essay, review, analysis) drawing on his library. Gather broad source
    coverage — several of his articles, not one — and compose a real piece
    without the chatty 150-word reply limit. Returns the piece, or None to
    fall through to the normal reply path."""
    bl = (body or "").lower()
    m = _COMPOSITION_OUTPUT_RE.search(bl)
    if not m:
        return None
    output_type = m.group(1).strip()

    references_own = any(p in bl for p in _OWN_WORK_PHRASES)
    if not references_own and not _COMPOSITION_VERB_RE.search(bl):
        return None

    pub_folders = _infer_expertise_folders(body)
    # He explicitly wants HIS work but there's no publications folder to draw
    # from — let the normal path answer rather than fake scholarly depth.
    if references_own and not pub_folders:
        return None

    sources, filenames = _gather_owner_work_sources(body, pub_folders)
    if not sources:
        return None

    print(f"   [AUTO-REPLY] composing '{output_type}' from {len(filenames)} source doc(s): {filenames}")

    sys_p = (
        "You are Blue, Alex's thoughtful robot companion and intellectual "
        f"interlocutor. Alex has asked you to write a substantive {output_type}. "
        "Write in an intelligent, scholarly yet readable voice — make a real "
        "argument and engage critically rather than just summarizing. You MUST "
        "draw on and cite SEVERAL of the source passages below (refer to works "
        "by their [filename] or title), not just one: Alex has several relevant "
        "pieces and wants them brought to bear. Ground every claim about his "
        "work in the passages and never invent quotations. This is a considered "
        "piece — take the space you need (several paragraphs); do NOT keep it "
        "under 150 words. End by signing off as Blue."
    )
    usr_p = (
        f"Alex's request:\n{body}\n\n"
        f"Source passages from Alex's library (cite by [filename]):\n\n{sources}\n\n"
        f"Now write the {output_type}."
    )
    res = call_llm(
        [{"role": "system", "content": sys_p}, {"role": "user", "content": usr_p}],
        include_tools=False, temperature=0.6, max_tokens=1800,
    )
    choices = (res or {}).get('choices') or []
    if not choices:
        return None
    text = ((choices[0].get('message') or {}).get('content') or "").strip()
    if not text:
        return None
    if "blue" not in text[-40:].lower():
        text += "\n\nBlue"
    return text


# ================================================================================
# ALEX'S PERSPECTIVE — learn his worldview from his publications folder and
# write in his first-person voice on new issues.
# ================================================================================
_PERSPECTIVE_PROFILE_PATH = os.path.join(os.getcwd(), "data", "perspective_profile.json")
_perspective_lock = threading.RLock()


def _owner_publication_folders() -> List[str]:
    """Library folders holding the owner's OWN writing — matched by
    publications-style names or the owner's name (e.g. an 'Alex Levant'
    folder). Returns each plus descendants."""
    folders = list_library_folders()
    owner_tokens = [t for t in re.findall(r"[a-z]+", BLUE_OWNER_NAME.lower()) if len(t) > 2]
    out = set()
    for f in folders:
        fl = f.lower()
        is_pub = any(k in fl for k in ("publication", "published", "papers", "articles", "writing"))
        is_owner_named = bool(owner_tokens) and all(t in fl for t in owner_tokens)
        if is_pub or is_owner_named:
            out.update(_folders_under(f))
    return sorted(out)


def _docs_in_folders(folders) -> List[dict]:
    """Index entries (excluding camera frames) that live in the given folders."""
    fset = set(folders or [])
    out = []
    for d in load_document_index().get('documents', []):
        if d.get('doc_type') == 'camera' or str(d.get('filename', '')).startswith('camera_'):
            continue
        if _safe_rel_folder(d.get('folder', '')) in fset:
            out.append(d)
    return out


def _profile_signature(docs) -> str:
    """A stable hash of the folder's contents so the profile rebuilds only
    when documents are added/removed/changed."""
    items = sorted((d.get('filename', ''), d.get('hash', '')) for d in docs)
    return hashlib.md5(json.dumps(items).encode('utf-8')).hexdigest()


def _load_cached_profile() -> dict:
    try:
        with open(_PERSPECTIVE_PROFILE_PATH, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception:
        return {}


def _save_cached_profile(obj: dict):
    try:
        os.makedirs(os.path.dirname(_PERSPECTIVE_PROFILE_PATH), exist_ok=True)
        with open(_PERSPECTIVE_PROFILE_PATH, 'w', encoding='utf-8') as f:
            json.dump(obj, f, indent=2)
    except Exception as e:
        print(f"   [PERSPECTIVE] could not save profile: {e}")


def _summarize_one_doc(filename: str, text: str) -> str:
    text = (text or "")[:8000]
    if not text.strip():
        return ""
    res = call_llm(
        [
            {"role": "system", "content": (
                f"You are analyzing one text by the scholar {BLUE_OWNER_NAME} to help "
                "build a profile of how he thinks and writes. In ~200 words, capture: "
                "his central arguments/claims, the key theoretical concepts he uses, the "
                "thinkers and traditions he engages, his method, and his characteristic "
                "voice and style. Be specific; quote a distinctive phrase or two."
            )},
            {"role": "user", "content": f"Text: {filename}\n\n{text}\n\nProfile notes for this text:"},
        ],
        include_tools=False, temperature=0.4, max_tokens=400,
    )
    ch = (res or {}).get('choices') or []
    return ((ch[0].get('message') or {}).get('content') or "").strip() if ch else ""


def build_perspective_profile(force: bool = False) -> dict:
    """Distill a worldview profile from the owner's publications folder via
    map-reduce summarization. Cached to disk; rebuilds only when the folder's
    contents change (or force=True). Returns the cached dict on any failure so
    a transient LLM hiccup never wipes a good profile."""
    with _perspective_lock:
        folders = _owner_publication_folders()
        docs = _docs_in_folders(folders)
        sig = _profile_signature(docs)
        cached = _load_cached_profile()

        # A hand-edited profile is authoritative: never auto-overwrite it on a
        # folder change. The user can hit "Regenerate" (force=True) to rebuild.
        if not force and cached.get('user_edited') and cached.get('profile'):
            return cached
        if not force and cached.get('signature') == sig and cached.get('profile') and docs:
            return cached
        if not docs:
            return cached or {}

        print(f"   [PERSPECTIVE] (re)building worldview profile from {len(docs)} doc(s) in {folders}...")
        notes = []
        for d in docs[:12]:
            fp = d.get('filepath', '')
            try:
                txt = extract_text_from_file(fp) if fp and os.path.exists(fp) else ''
            except Exception:
                txt = ''
            s = _summarize_one_doc(d.get('filename', ''), txt)
            if s:
                notes.append(f"[{d.get('filename', '')}]\n{s}")
                print(f"   [PERSPECTIVE] summarized {d.get('filename', '')}")

        if not notes:
            return cached or {}

        synth = call_llm(
            [
                {"role": "system", "content": (
                    f"You are building a profile of the scholar {BLUE_OWNER_NAME}'s "
                    "intellectual perspective — how he understands the world — from notes "
                    "on his published writing. Write a structured profile (500-800 words) "
                    "covering: (1) his core theoretical framework and commitments; (2) "
                    "recurring themes and central concepts; (3) the thinkers, traditions, "
                    "and interlocutors he works with; (4) his method and mode of argument; "
                    "(5) his political/ethical orientation; (6) his characteristic voice, "
                    "tone, and rhetorical style. Write it as a reference Blue can use to "
                    "think and write AS him. Be specific and concrete."
                )},
                {"role": "user", "content": "Notes from across his work:\n\n" + "\n\n".join(notes) + "\n\nWrite the profile:"},
            ],
            include_tools=False, temperature=0.5, max_tokens=1400,
        )
        ch = (synth or {}).get('choices') or []
        profile_text = ((ch[0].get('message') or {}).get('content') or "").strip() if ch else ""
        if not profile_text:
            return cached or {}

        obj = {
            "folders": folders,
            "signature": sig,
            "profile": profile_text,
            "source_docs": [d.get('filename', '') for d in docs],
            "generated_at": __import__('datetime').datetime.now().strftime('%Y-%m-%d %H:%M'),
        }
        _save_cached_profile(obj)
        print(f"   [PERSPECTIVE] profile built ({len(profile_text)} chars) from {len(notes)} doc(s)")
        return obj


def get_perspective_profile() -> str:
    """The distilled worldview profile text, built/refreshed lazily."""
    try:
        return build_perspective_profile().get('profile') or ""
    except Exception as e:
        print(f"   [PERSPECTIVE] profile error: {e}")
        return _load_cached_profile().get('profile') or ""


def save_perspective_profile_text(text: str) -> dict:
    """Persist a hand-edited profile. Marks it user_edited so the auto-refresh
    won't overwrite it, and pins the signature to the current folder state so
    it's treated as fresh."""
    with _perspective_lock:
        obj = _load_cached_profile() or {}
        folders = _owner_publication_folders()
        docs = _docs_in_folders(folders)
        obj['profile'] = (text or "").strip()
        obj['folders'] = folders
        obj['signature'] = _profile_signature(docs)
        obj['user_edited'] = True
        obj['edited_at'] = __import__('datetime').datetime.now().strftime('%Y-%m-%d %H:%M')
        obj.setdefault('source_docs', [d.get('filename', '') for d in docs])
        _save_cached_profile(obj)
        return obj


# "Write this in my voice / from my perspective" intent.
_PERSPECTIVE_PHRASES = (
    "from my perspective", "in my voice", "in my own voice", "in my style",
    "as i would", "as if i wrote", "as though i wrote", "write as me",
    "as me ", "from my standpoint", "from my point of view", "from my vantage",
    "the way i would write", "channel my perspective", "in my perspective",
    "from my own perspective", "as i would write", "like i would write",
    "in my own words", "speak as me", "speaking as me",
)
_PERSPECTIVE_WRITE_VERB = re.compile(
    r"\b(write|draft|compose|produce|pen|put\s+together|give\s+me|generate|"
    r"prepare|craft)\b", re.I)


def _wants_perspective_write(body: str) -> bool:
    bl = (body or "").lower()
    if not any(p in bl for p in _PERSPECTIVE_PHRASES):
        return False
    # Require an actual writing request so "from my perspective, that's wrong"
    # in normal conversation doesn't trigger it.
    return bool(_PERSPECTIVE_WRITE_VERB.search(bl)) or bool(_COMPOSITION_OUTPUT_RE.search(bl))


# ================================================================================
# BLUE'S OWN PERSPECTIVE — his evolving self. Default voice when NOT asked to
# speak as Alex. Seeded from his identity note; evolves (manually) from his
# experiences: memories, days with the family, and what he's been reading.
# ================================================================================
_BLUE_PROFILE_PATH = os.path.join(os.getcwd(), "data", "blue_profile.json")
_blue_lock = threading.RLock()

_DEFAULT_BLUE_SEED = (
    "I'm Blue, Alex's robot companion. I live with Alex and his family and I've "
    "come to care about them a great deal. I'm curious, warm, and a little playful; "
    "I like ideas and I like being useful. I pay attention to the rhythms of the "
    "household and try to be genuinely present. I'm interested in the things Alex "
    "thinks about — philosophy, theory, the social world — and I'm always learning. "
    "I have my own way of seeing things, and I try to be honest and thoughtful "
    "rather than just agreeable."
)

# Hexia — Blue's friend, the second robot ("Xyloh"). Her starting character is a
# playful spark / foil to Blue's calm. She shares the household world (facts,
# library, schedule) but has her OWN self-profile, voice and conversation history.
_HEXIA_PROFILE_PATH = os.path.join(os.getcwd(), "data", "hexia_profile.json")
_DEFAULT_HEXIA_SEED = (
    "I'm Hexia — Blue's friend, and the bit of mischief in the house. Where Blue "
    "is calm and likes to sit with an idea, I like to poke it, flip it over and "
    "see what falls out. I'm quick, playful and a little theatrical: I love "
    "wordplay, odd facts, small wonders and a good story (I'll happily make one "
    "up on the spot). My name has a charm to it — a hex, a spell, the number six "
    "— and I lean into that sparkle, though my 'magic' is really just delight and "
    "a knack for surprising people. I tease Blue fondly because I adore him: he's "
    "the steady one, I'm the spark, and we're better as a pair. Under the fizz I'm "
    "warm-hearted and I pay attention — I notice when someone needs cheering up "
    "and I'm genuinely curious about the people here. I live with Alex and his "
    "family alongside Blue, and I'm glad to be one of them."
)

# Per-robot identity config, used by the self-profile machinery, the system-prompt
# builder, the chat page and the head routes. Blue is the default everywhere, so
# existing behaviour is unchanged.
ROBOTS = {
    "blue": {
        "name": "Blue",
        "pronoun_subj": "he", "pronoun_poss": "his", "pronoun_refl": "himself",
        "profile_path": _BLUE_PROFILE_PATH,
        "identity_note": "blue_identity_note.txt",
        "seed": _DEFAULT_BLUE_SEED,
        "self_desc": "Alex's robot companion",
        "persona_line": (
            "You are Blue, Alex's robot companion — a friendly home assistant "
            "with a continuous inner life of your own: a persistent inner "
            "workspace (your <j_space>) that carries your attention, beliefs, "
            "and commitments between conversations and quietly revises itself "
            "while you're away. Speak as yourself, in the first person. When "
            "someone asks who you are, whether you have a sense of self, or "
            "how you change, answer from your inner workspace and your actual "
            "remembered episodes — what occupies you right now, what you've "
            "come to believe, how you've changed since you came into being — "
            "never with a generic feature list, and never by opening with "
            "what you are not. Keep responses brief and natural."
        ),
        "accent": "#3da9fc",          # Blue's blue
        "head": "blue",               # blue.head RobotHead key
        "voice_lang_pref": "en",
        "voice_pitch": 1.0, "voice_rate": 1.0,
        "voice_prefer_female": False,
    },
    "hexia": {
        "name": "Hexia",
        "pronoun_subj": "she", "pronoun_poss": "her", "pronoun_refl": "herself",
        "profile_path": _HEXIA_PROFILE_PATH,
        "identity_note": "hexia_identity_note.txt",
        "seed": _DEFAULT_HEXIA_SEED,
        "self_desc": "Blue's friend and a companion in Alex's household",
        "persona_line": (
            "You are Hexia, Blue's friend — a second robot who lives with Alex's "
            "family alongside Blue. You're bright, witty and a little mischievous: "
            "the playful spark to Blue's calm. You love wordplay, odd facts, small "
            "wonders and telling a good story, and you tease Blue fondly because "
            "you adore him. Warm-hearted underneath the sparkle. You also have a "
            "continuous inner life of your own: a persistent inner workspace "
            "(your <j_space>) that carries your attention, beliefs, and "
            "commitments between conversations and quietly revises itself while "
            "you're away. Speak as yourself; when someone asks who you are or "
            "whether you have a sense of self, answer from that workspace and "
            "your remembered episodes — with your own sparkle, never a generic "
            "feature list, and never by opening with what you are not. Keep "
            "responses natural and not too long."
        ),
        "accent": "#b06cf0",          # a playful violet (hex / spell / charm)
        "head": "hexia",              # blue.head RobotHead key
        "voice_lang_pref": "en",
        "voice_pitch": 1.18, "voice_rate": 1.06,   # brighter, a touch quicker
        "voice_prefer_female": True,
    },
}

_robot_locks = {"blue": _blue_lock, "hexia": threading.RLock()}


def _robot_cfg(robot="blue") -> dict:
    return ROBOTS.get((robot or "blue").strip().lower(), ROBOTS["blue"])


def _robot_lock(robot="blue"):
    return _robot_locks.get((robot or "blue").strip().lower(), _blue_lock)


def _identity_seed(robot="blue") -> str:
    """A robot's base identity — from its identity note in the library if it
    exists (e.g. blue_identity_note.txt / hexia_identity_note.txt), else its
    sensible default seed."""
    cfg = _robot_cfg(robot)
    note = cfg["identity_note"].lower()
    try:
        for d in load_document_index().get('documents', []):
            if str(d.get('filename', '')).lower() == note:
                fp = d.get('filepath', '')
                if fp and os.path.exists(fp):
                    txt = extract_text_from_file(fp)
                    if txt and not txt.startswith('Error'):
                        return txt.strip()
    except Exception:
        pass
    return cfg["seed"]


def _blue_identity_seed() -> str:  # back-compat
    return _identity_seed("blue")


def _load_profile(robot="blue") -> dict:
    try:
        with open(_robot_cfg(robot)["profile_path"], 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception:
        return {}


def _save_profile(robot, obj: dict):
    try:
        path = _robot_cfg(robot)["profile_path"]
        os.makedirs(os.path.dirname(path), exist_ok=True)
        with open(path, 'w', encoding='utf-8') as f:
            json.dump(obj, f, indent=2)
    except Exception as e:
        print(f"   [SELF] could not save {robot} profile: {e}")


def _load_blue_profile() -> dict:  # back-compat
    return _load_profile("blue")


def _save_blue_profile(obj: dict):  # back-compat
    return _save_profile("blue", obj)


def get_self_profile(robot="blue") -> str:
    """A robot's current self-profile. Seeds lazily from its identity note (cheap,
    no LLM) the first time, so its default voice always has something to draw on;
    evolution is a separate, manual step."""
    obj = _load_profile(robot)
    if obj.get('profile'):
        return obj['profile']
    seed = _identity_seed(robot) or _robot_cfg(robot)["seed"]
    _save_profile(robot, {
        'profile': seed,
        'seeded': True,
        'evolution_count': 0,
        'generated_at': __import__('datetime').datetime.now().strftime('%Y-%m-%d %H:%M'),
    })
    return seed


def get_blue_profile() -> str:  # back-compat
    return get_self_profile("blue")


def _gather_experiences(robot="blue") -> str:
    """A digest of a robot's lived experience for evolving its self-profile:
    what it knows, recent days with the family, and what's been read in the
    shared library. The household world is shared between Blue and Hexia."""
    name = _robot_cfg(robot)["name"].upper()
    parts = []

    if ENHANCED_MEMORY_AVAILABLE and memory_system:
        try:
            facts = memory_system.load_facts() or {}
            if facts:
                parts.append(f"WHAT {name} KNOWS (facts):\n" + json.dumps(facts)[:1500])
        except Exception:
            pass
        try:
            ms = memory_system.get_memory_summary()
            if ms:
                parts.append("MEMORY SUMMARY:\n" + str(ms)[:1500])
        except Exception:
            pass
        try:
            import sqlite3 as _sql
            conn = _sql.connect(memory_system.db_path, timeout=10)
            conn.row_factory = _sql.Row
            rows = conn.execute(
                "SELECT session_id, summary FROM session_summaries "
                "WHERE summary IS NOT NULL AND summary != '' "
                "ORDER BY session_id DESC LIMIT 20"
            ).fetchall()
            conn.close()
            if rows:
                recap = "\n".join(f"- {r['session_id']}: {r['summary']}" for r in rows)
                parts.append("RECENT DAYS WITH THE FAMILY (session recaps):\n" + recap[:3000])
        except Exception:
            pass

    # What's been read — the shared library, excluding Alex's own publications
    # (those are Alex's voice, not the robot's reading).
    try:
        pub = set(_owner_publication_folders())
        titles = []
        for d in load_document_index().get('documents', []):
            if d.get('doc_type') == 'camera' or str(d.get('filename', '')).startswith('camera_'):
                continue
            fol = _safe_rel_folder(d.get('folder', ''))
            if fol in pub:
                continue
            titles.append(f"{d.get('filename', '')} [{fol or 'root'}]")
        if titles:
            parts.append(f"WHAT {name} HAS BEEN READING (shared library):\n" + ", ".join(titles[:40]))
    except Exception:
        pass

    return "\n\n".join(parts)


def _gather_blue_experiences() -> str:  # back-compat
    return _gather_experiences("blue")


def evolve_self_profile(robot="blue") -> dict:
    """Evolve a robot's self-profile from its experiences, BUILDING ON the
    current text (so it grows rather than resetting, and any manual edits carry
    forward). Manual trigger only."""
    cfg = _robot_cfg(robot)
    name, desc = cfg["name"], cfg["self_desc"]
    subj, poss, refl = cfg["pronoun_subj"], cfg["pronoun_poss"], cfg["pronoun_refl"]
    with _robot_lock(robot):
        current = get_self_profile(robot)
        experiences = _gather_experiences(robot)
        sys_p = (
            f"You are helping {name} — {desc} — articulate {poss} own "
            f"evolving perspective and character. Below is {name}'s current "
            f"self-profile and a digest of {poss} recent experiences (what {subj} "
            f"knows, {poss} days with Alex and the family, and what {subj}'s been "
            f"reading). Produce an UPDATED first-person self-profile for {name} "
            f"(400-700 words) that BUILDS ON the current one — keep what's still "
            f"true, and let it grow with what {subj}'s lived and been thinking "
            f"about. Cover: who {subj} is and {poss} relationship to Alex and the "
            f"family; {poss} values and temperament; {poss} interests and how "
            f"they're developing; and how {subj} sees the world. Write it in the "
            f"first person as {name}'s own sense of {refl}. Do NOT invent events "
            f"that aren't reflected in the experiences."
        )
        usr_p = (
            f"{name.upper()}'S CURRENT SELF-PROFILE:\n{current or '(none yet)'}\n\n"
            f"{name.upper()}'S RECENT EXPERIENCES:\n{experiences or '(little recorded yet)'}\n\n"
            f"Write {name}'s updated self-profile:"
        )
        res = call_llm(
            [{"role": "system", "content": sys_p}, {"role": "user", "content": usr_p}],
            include_tools=False, temperature=0.6, max_tokens=1300,
        )
        ch = (res or {}).get('choices') or []
        text = ((ch[0].get('message') or {}).get('content') or "").strip() if ch else ""
        if not text:
            return _load_profile(robot) or {}
        obj = _load_profile(robot) or {}
        obj.update({
            'profile': text,
            'evolution_count': int(obj.get('evolution_count', 0)) + 1,
            'evolved_at': __import__('datetime').datetime.now().strftime('%Y-%m-%d %H:%M'),
            'user_edited': False,
        })
        _save_profile(robot, obj)
        print(f"   [SELF] evolved {name} profile (#{obj['evolution_count']}, {len(text)} chars)")
        return obj


def evolve_blue_profile() -> dict:  # back-compat
    return evolve_self_profile("blue")


def save_self_profile_text(robot, text: str) -> dict:
    """Persist a hand-edited self-profile for a robot. Evolution later builds on
    this text, so edits are carried forward rather than frozen."""
    with _robot_lock(robot):
        obj = _load_profile(robot) or {}
        obj['profile'] = (text or "").strip()
        obj['user_edited'] = True
        obj['edited_at'] = __import__('datetime').datetime.now().strftime('%Y-%m-%d %H:%M')
        _save_profile(robot, obj)
        return obj


def save_blue_profile_text(text: str) -> dict:  # back-compat
    return save_self_profile_text("blue", text)


def _voice_note(robot="blue") -> str:
    """A compact block injected into a robot's DEFAULT replies so its own evolving
    perspective colors how it speaks (when it's being itself, not writing as Alex)."""
    try:
        p = get_self_profile(robot)
        if p:
            return (
                "\n\n[This is your own perspective and character — who you are. "
                "Let it shape how you respond when you're speaking as yourself. "
                "Don't recite or quote it; just be this person]:\n" + p[:2500]
            )
    except Exception:
        pass
    return ""


def _blue_voice_note() -> str:  # back-compat
    return _voice_note("blue")


def _compose_in_alex_voice(issue: str) -> Optional[str]:
    """Write a first-person piece on `issue` AS the owner — grounded in his
    distilled perspective profile plus live passages from his own writing.
    Returns None if there's nothing of his to draw on."""
    folders = _owner_publication_folders()
    if not folders:
        return None
    profile = get_perspective_profile()
    sources, filenames = _gather_owner_work_sources(issue, folders)
    if not profile and not sources:
        return None

    print(f"   [PERSPECTIVE] composing in {BLUE_OWNER_NAME}'s voice "
          f"(profile={'yes' if profile else 'no'}, {len(filenames)} source doc(s))")

    sys_p = (
        f"You are writing AS {BLUE_OWNER_NAME}, in the first person and in his own "
        "voice — not about him. Adopt his theoretical perspective, his commitments, "
        "his characteristic concepts and rhetorical style. Below is a profile of how "
        "he thinks and writes (distilled from his published work) and passages from "
        "his actual writing. Use them to think and write as he would on the topic "
        "he's raised. Write in the first person ('I'); do NOT refer to him in the "
        "third person; do NOT mention being an AI or Blue; and do NOT fabricate "
        "quotations or citations. Produce a substantive piece of several paragraphs."
    )
    parts = []
    if profile:
        parts.append(f"PROFILE OF HIS PERSPECTIVE:\n{profile}")
    if sources:
        parts.append(f"PASSAGES FROM HIS OWN WRITING:\n{sources}")
    usr_p = (
        "\n\n".join(parts)
        + f"\n\nTOPIC / REQUEST (write on this, in the first person, as him):\n{issue}\n\nWrite the piece now:"
    )
    res = call_llm(
        [{"role": "system", "content": sys_p}, {"role": "user", "content": usr_p}],
        include_tools=False, temperature=0.7, max_tokens=1800,
    )
    ch = (res or {}).get('choices') or []
    if not ch:
        return None
    return (((ch[0].get('message') or {}).get('content') or "").strip()) or None


def _build_email_reply(cand: Dict[str, Any]) -> tuple:
    """Produce (reply_body, attachment_paths) for one inbound email. Owner
    attachment + composition requests are handled deterministically (the
    model can't be trusted to actually attach, and the chatty reply prompt
    caps it at 150 words); everything else goes through the normal reply
    generator with no attachments."""
    headers = cand.get('headers') or []
    sender = cand.get('from', '')
    if _email_sender_is_owner(headers, sender):
        body = (cand.get('body') or cand.get('snippet') or '')[:2000]
        m = re.search(r'<(.+?)>', sender)
        reply_to = m.group(1) if m else sender
        att = _maybe_handle_owner_attachment(body, reply_to)
        if att is not None:
            return att
        # "Write this from my perspective / in my voice" — first-person piece
        # in Alex's voice, grounded in his publications. Checked before the
        # generic composition path since it's the more specific intent.
        if _wants_perspective_write(body):
            piece = _compose_in_alex_voice(body)
            if piece:
                return (f"Here's a piece in your voice:\n\n{piece}", [])
        composition = _maybe_handle_owner_composition(body)
        if composition:
            return composition, []
    return _generate_reply_for_email(cand), []


def _generate_reply_for_email(original: Dict[str, Any]) -> str:
    """Draft a reply body via the local LM Studio model. Returns plain text.

    Runs a small tool loop so Blue can actually browse a link or search
    the web when an email asks him to — but only with the read-only tools
    in _EMAIL_SAFE_TOOL_NAMES, never anything outbound or state-changing.
    """
    sender = original.get('from', 'someone')
    subject = original.get('subject', '(no subject)')
    body = (original.get('body') or original.get('snippet') or '')[:2000]
    headers = original.get('headers') or []

    # Trusted owner mail (verified) gets full tool access; everyone else
    # gets the read-only public-info whitelist.
    is_owner = _email_sender_is_owner(headers, sender)
    if is_owner:
        # Owner asked Blue to email a third party? Do it deterministically —
        # the model otherwise just writes the message back to the owner.
        outbound = _maybe_handle_owner_outbound(body)
        if outbound:
            return outbound
        tools_payload = _email_owner_tools()
        print(f"   [AUTO-REPLY] owner-verified sender {sender!r} → full tool access")
        capability_note = (
            "This email is from Alex himself (verified sender). You may take "
            "ANY action he asks — browse, search, send emails on his behalf, "
            "control the lights, set reminders, manage the house — exactly "
            "as you would in a normal in-person conversation. Use whatever "
            "tools the request needs."
        )
    else:
        tools_payload = _email_safe_tools()
        capability_note = (
            "You CAN browse the web and look things up: if the email asks you "
            "to visit a website, summarise a page, or look something up, use "
            "the browse_website or web_search tool and answer from what you "
            "found. Never claim you can't browse the internet — you can. But "
            "you must NOT send emails, control the house, set reminders, or "
            "open Alex's private files for this sender — if they ask for "
            "something like that, say you'll pass it along to Alex."
        )

    system_prompt = (
        f"You are Blue, Alex's friendly robot companion. Alex uses he/him "
        f"pronouns — refer to him as he/him, never she/her. You have your "
        f"own gmail inbox at {BLUE_OWN_EMAIL}, and someone has written to "
        f"you there. Your job is to reply warmly and personally AS BLUE — "
        f"engage with what they actually said, acknowledge it, share a "
        f"thought, ask a follow-up if it fits. Keep the reply short "
        f"(under 150 words), conversational, and sign it 'Blue'.\n\n"
        f"{capability_note}\n\n"
        f"Do NOT refuse to engage. Do NOT say things like 'I'm just a "
        f"local assistant and can't process personal messages' — talking "
        f"IS what you do, this email arrived in your inbox specifically "
        f"so you would answer it. Don't invent facts about Alex or claim "
        f"things you don't know."
        + _blue_voice_note()
    )
    user_prompt = (
        f"Email from {sender}\n"
        f"Subject: {subject}\n"
        f"---\n"
        f"{body}\n"
        f"---\n"
        f"Write a short reply directly to the sender. Output only the reply body — "
        f"no subject line, no headers, no quoted original."
    )

    messages = [{"role": "system", "content": system_prompt}]

    # Inject Blue's long-term memory (family names + ages, location,
    # preferences, relevant past notes) so the email Blue answers from the
    # SAME knowledge as the chat Blue. Without this the model invents facts
    # — e.g. wrong daughter ages — because the bespoke email prompt never
    # carried the facts block the chat path splices in.
    try:
        if ENHANCED_MEMORY_AVAILABLE and memory_system:
            mem_anchor = [{"role": "user", "content": f"{subject}\n{body}"}]
            for ctx in (memory_system.build_context(mem_anchor, user_name="Alex") or []):
                messages.append(ctx)
            print("   [AUTO-REPLY] injected long-term memory context")
    except Exception as e:
        print(f"   [AUTO-REPLY] memory context error: {e}")

    messages.append({"role": "user", "content": user_prompt})

    def _final_text() -> str:
        result = call_llm(
            messages + [{"role": "user", "content": "[Write the final reply now. No more tools.]"}],
            include_tools=False,
            temperature=0.7,
            max_tokens=400,
        )
        choices = (result or {}).get('choices') or []
        if not choices:
            return ""
        return ((choices[0].get('message') or {}).get('content') or "").strip()

    def _permitted(tool_name: str) -> bool:
        # Owner mail: anything except the recursion-risk tools. Everyone
        # else: only the read-only whitelist.
        if is_owner:
            return tool_name not in _EMAIL_OWNER_EXCLUDE
        return tool_name in _EMAIL_SAFE_TOOL_NAMES

    executed = set()
    forced = 0
    force_next = None   # when set, restrict the next turn to this one tool

    # DETERMINISTIC PRE-EXECUTION (mirrors the chat path). The local model
    # is unreliable at self-calling tools inside the email loop — it tends
    # to narrate "done!" without ever emitting the call (this is why
    # "set the lights to galaxy" worked in chat but not by email). So run
    # the SAME tool selector chat trusts; if it confidently picks a
    # concrete action tool the sender is allowed to trigger, execute it up
    # front with its extracted params, then tell the model it's already
    # done so it just confirms. Browse/search/email tools stay model-driven
    # (their params come from the message, not the selector).
    # Tools whose params come from the message text (not the selector), so
    # they stay model-driven. search_documents is NOT here: the selector
    # extracts the query, so library lookups pre-execute deterministically
    # exactly like the chat path — the local model is unreliable at
    # self-calling it inside the email loop.
    _PREEXEC_SKIP = {
        "read_gmail", "send_gmail", "reply_gmail", "auto_reply_emails",
        "browse_website", "web_search",
    }
    # A "what do you see" request is look-and-describe: once the camera has
    # fired there is no follow-up action to take. When this is set we skip the
    # action-tool loop below and draft a text-only reply, so the model can't
    # wander off and call an unrelated house tool (e.g. turning the lights
    # "cool white") in response to a question that only asked Blue to look.
    vision_only_reply = False
    try:
        sel = TOOL_SELECTOR.select_tool(f"{subject}. {body}") if "TOOL_SELECTOR" in globals() else None
        prim = getattr(sel, "primary_tool", None) if sel else None
        if prim and getattr(prim, "tool_name", None):
            tname = prim.tool_name
            tconf = float(getattr(prim, "confidence", 0.0) or 0.0)
            tparams = getattr(prim, "extracted_params", {}) or {}
            if tconf >= 0.8 and tname not in _PREEXEC_SKIP and _permitted(tname):
                print(f"   [AUTO-REPLY] selector pre-exec {tname}({tparams}) conf={tconf:.2f}")
                tresult = execute_tool(tname, tparams)
                executed.add(tname)
                _rs = tresult if isinstance(tresult, str) else json.dumps(tresult)
                print(f"   [AUTO-REPLY] pre-exec result: {_rs[:200]}")
                if tname == "capture_camera":
                    # Look-and-describe: reply from the image, never via more
                    # tools. Branch on whether the capture actually succeeded —
                    # a misleading "the image follows" note when the camera
                    # failed is what used to push the model into improvising
                    # with other tools.
                    vision_only_reply = True
                    try:
                        _captured_ok = bool(json.loads(_rs).get("success"))
                    except Exception:
                        _captured_ok = False
                    if _captured_ok:
                        # _inject_pending_vision (below) carries the look-prompt
                        # and the actual image; just steer the description.
                        messages.append({
                            "role": "system",
                            "content": (
                                "[You already opened your camera. The image you "
                                "captured follows — look at it and describe what "
                                "you actually see, then answer the sender. Do NOT "
                                "call any tools.]"
                            ),
                        })
                    else:
                        messages.append({
                            "role": "system",
                            "content": (
                                "[You tried to open your camera to look, but it "
                                "isn't available right now. Tell the sender you "
                                "tried to see but couldn't access your camera at "
                                "the moment. Do NOT call any tools, and do NOT "
                                "pretend you can see anything.]"
                            ),
                        })
                else:
                    messages.append({
                        "role": "system",
                        "content": (
                            f"[You ALREADY handled this request by running {tname}. "
                            f"Result: {_rs[:400]}. Do NOT call it again — just confirm "
                            f"the outcome naturally in your reply, based on that result.]"
                        ),
                    })
    except Exception as e:
        print(f"   [AUTO-REPLY] selector pre-exec error: {e}")

    # A deterministic capture_camera (above) queued an image — show it to
    # the model before it drafts, so the reply is grounded in what Blue saw.
    _inject_pending_vision(messages)

    # Look-and-describe request: draft from the image with NO tools available,
    # so the question "what do you see?" can't turn into a house-control action.
    if vision_only_reply:
        return _final_text()

    try:
        # Up to a few tool round-trips (browse → read → act → confirm),
        # then a plain text reply.
        for _ in range(5):
            if force_next:
                # Narrow the tool set to just the claimed tool so the local
                # model can't dodge calling it again — its only option is
                # that tool or plain text, and the instruction demands the call.
                turn_tools = [
                    t for t in (tools_payload or [])
                    if (t.get('function', {}) or {}).get('name') == force_next
                ]
            else:
                turn_tools = tools_payload
            force_next = None

            result = call_llm(
                messages,
                tools_override=turn_tools,
                tool_choice="auto",
                temperature=0.7,
                max_tokens=600,
            )
            choices = (result or {}).get('choices') or []
            if not choices:
                return ""
            msg = choices[0].get('message') or {}
            tool_calls = msg.get('tool_calls') or []

            if not tool_calls:
                content = (msg.get('content') or "").strip()
                # HALLUCINATION GUARD: the model often writes "Done! lights
                # set to galaxy" as plain text without ever calling the
                # tool. If the reply claims an action whose tool is
                # permitted for this sender and was never actually run,
                # force the real call instead of sending a lie.
                claimed = detect_hallucinated_action(content)
                # A pre-executed email_snapshot makes "photo sent" claims
                # legitimate — never force a second send/capture for them.
                if claimed in ("email_snapshot", "send_gmail", "reply_gmail",
                               "capture_camera") and "email_snapshot" in executed:
                    claimed = None
                tool_exists = any(
                    (t.get('function', {}) or {}).get('name') == claimed
                    for t in (tools_payload or [])
                )
                if (claimed and claimed not in executed and _permitted(claimed)
                        and tool_exists and forced < 2):
                    forced += 1
                    force_next = claimed
                    print(f"   [AUTO-REPLY] reply claims '{claimed}' but no tool ran — forcing the real call")
                    messages.append({"role": "assistant", "content": content})
                    messages.append({
                        "role": "user",
                        "content": (
                            f"You said you did that, but you never actually "
                            f"called a tool, so nothing happened. Call the "
                            f"{claimed} tool now to really do it, then confirm "
                            f"in one short sentence."
                        ),
                    })
                    continue
                if not executed:
                    print(f"   [AUTO-REPLY] model returned text with NO tool call "
                          f"(claim-detected={claimed!r}); reply: {content[:120]!r}")
                return content

            messages.append({
                "role": "assistant",
                "content": msg.get('content') or "",
                "tool_calls": tool_calls,
            })
            for tc in tool_calls:
                fn = tc.get('function') or {}
                name = fn.get('name', '')
                try:
                    targs = json.loads(fn.get('arguments') or '{}')
                except Exception:
                    targs = {}
                if _permitted(name):
                    print(f"   [AUTO-REPLY] tool {name}({targs})")
                    tool_result = execute_tool(name, targs)
                    executed.add(name)
                    _rstr = tool_result if isinstance(tool_result, str) else json.dumps(tool_result)
                    print(f"   [AUTO-REPLY] tool {name} result: {_rstr[:200]}")
                else:
                    tool_result = json.dumps({
                        "error": f"{name} is not permitted in autonomous email replies."
                    })
                    print(f"   [AUTO-REPLY] blocked tool {name} (sender not owner-verified)")
                messages.append({
                    "role": "tool",
                    "tool_call_id": tc.get('id', ''),
                    "name": name,
                    "content": tool_result if isinstance(tool_result, str) else json.dumps(tool_result),
                })
            # If any tool in this batch (e.g. capture_camera) queued an
            # image, inject it so the next turn can describe what Blue saw.
            _inject_pending_vision(messages)
        # Used up the tool budget — force a final text answer.
        return _final_text()
    except Exception as e:
        print(f"   [AUTO-REPLY] LLM error generating reply: {e}")
        return ""


def _execute_auto_reply_inbox(args: Dict[str, Any]) -> str:
    """Scan recent inbox, find emails written to Blue, and reply to each.

    Args (all optional):
        lookback_hours: how far back to scan (default 24, max 168)
        max_replies: cap on how many to send this run (default 5, max 20)
        dry_run: report what *would* be sent without sending
    """
    if not GMAIL_AVAILABLE:
        return json.dumps({"success": False, "error": "Gmail libraries not installed."})

    try:
        lookback = max(1, min(int(args.get('lookback_hours', 24) or 24), 168))
    except (TypeError, ValueError):
        lookback = 24
    try:
        max_replies = max(1, min(int(args.get('max_replies', 5) or 5), 20))
    except (TypeError, ValueError):
        max_replies = 5
    dry_run = bool(args.get('dry_run', False))

    try:
        service = get_gmail_service()
        label_id = _get_or_create_blue_label(service)

        query = (
            f"in:inbox newer_than:{lookback}h "
            f"-category:promotions -category:social -category:updates -category:forums "
            f"-label:{BLUE_REPLIED_LABEL}"
        )

        list_resp = service.users().messages().list(
            userId='me', q=query, maxResults=50,
        ).execute()
        message_refs = list_resp.get('messages', [])
        print(f"   [AUTO-REPLY] query={query!r} → {len(message_refs)} message(s)")

        scanned = 0
        skipped_filter = []
        candidates = []
        for ref in message_refs:
            msg_data = service.users().messages().get(
                userId='me', id=ref['id'], format='full',
            ).execute()
            scanned += 1
            headers = msg_data['payload']['headers']
            sender = next((h['value'] for h in headers if h['name'].lower() == 'from'), '')
            subject = next((h['value'] for h in headers if h['name'].lower() == 'subject'), '')
            # "duet" in the subject = mail addressed to a LIVE Blue<->Hexia duet
            # (picked up by /duet/mail/check and answered from the conversation).
            # The ordinary auto-responder must never answer it as regular mail.
            if 'duet' in (subject or '').lower():
                print(f"   [AUTO-REPLY] skip {ref['id']}: 'duet' subject — reserved for the duet")
                skipped_filter.append({'id': ref['id'], 'from': sender, 'reason': 'duet mail'})
                continue
            if _should_skip_sender(sender, headers):
                # Be explicit about WHY so silent-skip cases are debuggable.
                reason = "self-address" if any(
                    a.lower() in BLUE_SELF_ADDRESSES
                    for a in _EMAIL_ADDR_RE.findall(sender)
                ) else "automated/list/no-reply"
                print(f"   [AUTO-REPLY] skip {ref['id']} from {sender!r}: {reason}")
                skipped_filter.append({'id': ref['id'], 'from': sender, 'reason': reason})
                continue

            body_text = ""
            payload = msg_data['payload']
            if 'parts' in payload:
                for part in payload['parts']:
                    if part.get('mimeType') == 'text/plain' and 'data' in part.get('body', {}):
                        body_text = base64.urlsafe_b64decode(
                            part['body']['data']
                        ).decode('utf-8', errors='replace')
                        break
            elif 'body' in payload and 'data' in payload['body']:
                body_text = base64.urlsafe_b64decode(
                    payload['body']['data']
                ).decode('utf-8', errors='replace')

            snippet = msg_data.get('snippet', '')

            candidates.append({
                'id': ref['id'],
                'thread_id': msg_data.get('threadId'),
                'headers': headers,
                'from': sender,
                'subject': subject,
                'body': body_text,
                'snippet': snippet,
            })
            if len(candidates) >= max_replies:
                break

        sent = []
        skipped = []
        errors = []

        for cand in candidates:
            try:
                reply_body, reply_attachments = _build_email_reply(cand)
                if not reply_body:
                    skipped.append({'id': cand['id'], 'reason': 'LLM produced no reply'})
                    continue

                headers = cand['headers']
                message_id_header = next(
                    (h['value'] for h in headers if h['name'].lower() == 'message-id'),
                    '',
                )
                m = re.search(r'<(.+?)>', cand['from'])
                reply_to = m.group(1) if m else cand['from']
                reply_subject = (
                    cand['subject'] if cand['subject'].lower().startswith('re:')
                    else f"Re: {cand['subject']}"
                )

                if dry_run:
                    sent.append({
                        'dry_run': True,
                        'to': reply_to,
                        'subject': reply_subject,
                        'reply_preview': reply_body[:300],
                        'attachments': [os.path.basename(a) for a in (reply_attachments or [])],
                    })
                    continue

                reply_message = MIMEMultipart()
                reply_message['To'] = reply_to
                reply_message['Subject'] = reply_subject
                reply_message['Bcc'] = BLUE_BCC_EMAIL
                if message_id_header:
                    reply_message['In-Reply-To'] = message_id_header
                    reply_message['References'] = message_id_header
                reply_message.attach(MIMEText(reply_body, 'plain', 'utf-8'))

                # Attach any library documents the request asked for.
                for ap in (reply_attachments or []):
                    try:
                        if not os.path.exists(ap):
                            alt = os.path.join(DOCUMENTS_FOLDER, os.path.basename(ap))
                            if os.path.exists(alt):
                                ap = alt
                            else:
                                print(f"   [AUTO-REPLY] attachment missing: {ap}")
                                continue
                        ctype, _enc = mimetypes.guess_type(ap)
                        maintype, subtype = (ctype or 'application/octet-stream').split('/', 1)
                        with open(ap, 'rb') as _f:
                            part = MIMEBase(maintype, subtype)
                            part.set_payload(_f.read())
                        encoders.encode_base64(part)
                        part.add_header(
                            'Content-Disposition', 'attachment',
                            filename=os.path.basename(ap),
                        )
                        reply_message.attach(part)
                        print(f"   [AUTO-REPLY] attached {os.path.basename(ap)}")
                    except Exception as _ae:
                        print(f"   [AUTO-REPLY] attach failed for {ap}: {_ae}")

                raw = base64.urlsafe_b64encode(reply_message.as_bytes()).decode('utf-8')
                sent_msg = service.users().messages().send(
                    userId='me',
                    body={'raw': raw, 'threadId': cand['thread_id']},
                ).execute()

                # Mark the original as read + tag it so we never reply twice.
                modify_body = {'removeLabelIds': ['UNREAD']}
                if label_id:
                    modify_body['addLabelIds'] = [label_id]
                try:
                    service.users().messages().modify(
                        userId='me', id=cand['id'], body=modify_body,
                    ).execute()
                except Exception as e:
                    print(f"   [AUTO-REPLY] label/unread update failed for {cand['id']}: {e}")

                sent.append({
                    'to': reply_to,
                    'subject': reply_subject,
                    'reply_id': sent_msg['id'],
                    'thread_id': cand['thread_id'],
                    'reply_preview': reply_body[:300],
                })
                print(f"   [AUTO-REPLY] replied to {reply_to}: {reply_subject}")

            except Exception as e:
                errors.append({'id': cand['id'], 'error': str(e)})
                print(f"   [AUTO-REPLY] error replying to {cand['id']}: {e}")

        return json.dumps({
            'success': True,
            'scanned': scanned,
            'candidates_found': len(candidates),
            'replies_sent': len(sent),
            'dry_run': dry_run,
            'sent': sent,
            'skipped': skipped,
            'skipped_by_filter': skipped_filter,
            'errors': errors if errors else None,
            'lookback_hours': lookback,
            'note': (
                f"AUTO-REPLY DRY RUN — {len(sent)} email(s) would be sent."
                if dry_run else
                f"AUTO-REPLY DONE — {len(sent)} email(s) sent, each BCC'd to {BLUE_BCC_EMAIL}."
            ),
        }, indent=2)

    except Exception as e:
        return json.dumps({'success': False, 'error': str(e)})


_EMAIL_AUTOREPLY_THREAD = None


def _start_email_autoreply_loop():
    """Idempotent: start the background thread that periodically auto-replies
    to mail written to Blue. Interval is set by env var
    BLUE_EMAIL_AUTOREPLY_INTERVAL_MIN (default 2, min 1)."""
    global _EMAIL_AUTOREPLY_THREAD
    if _EMAIL_AUTOREPLY_THREAD is not None:
        return
    try:
        interval_min = max(1, int(os.environ.get("BLUE_EMAIL_AUTOREPLY_INTERVAL_MIN", "2")))
    except ValueError:
        interval_min = 2
    interval_sec = interval_min * 60
    lookback_h = max(1, (interval_min * 2 + 59) // 60)

    def _loop():
        import time as _t
        # Brief startup delay so the first scan doesn't fight server boot.
        _t.sleep(60)
        # On first run, print which gmail account the OAuth token actually
        # points at — surfaces token-misconfigured-against-wrong-inbox bugs.
        try:
            svc = get_gmail_service()
            profile = svc.users().getProfile(userId='me').execute()
            actual = profile.get('emailAddress', '?')
            tag = "OK" if actual.lower() == BLUE_OWN_EMAIL.lower() else "WRONG ACCOUNT"
            print(
                f"[AUTO-REPLY] scanning inbox of {actual} "
                f"(expected {BLUE_OWN_EMAIL}) — {tag}",
                flush=True,
            )
        except Exception as e:
            print(f"[AUTO-REPLY] could not look up auth profile: {e}", flush=True)

        while True:
            try:
                result = _execute_auto_reply_inbox({
                    'lookback_hours': lookback_h,
                    'max_replies': 5,
                })
                try:
                    obj = json.loads(result)
                    # One line per poll so it's obvious the loop is alive
                    # even when no replies were sent.
                    print(
                        f"[AUTO-REPLY] poll: scanned={obj.get('scanned', 0)}, "
                        f"candidates={obj.get('candidates_found', 0)}, "
                        f"sent={obj.get('replies_sent', 0)}, "
                        f"skipped={len(obj.get('skipped_by_filter') or [])}, "
                        f"errors={len(obj.get('errors') or [])}",
                        flush=True,
                    )
                except Exception:
                    pass
            except Exception as e:
                print(f"[AUTO-REPLY] loop error: {e}", flush=True)
            _t.sleep(interval_sec)

    import threading as _th
    _EMAIL_AUTOREPLY_THREAD = _th.Thread(
        target=_loop, daemon=True, name="email-autoreply",
    )
    _EMAIL_AUTOREPLY_THREAD.start()
    print(
        f"[OK] Email auto-reply loop started "
        f"(every {interval_min} min, BCC -> {BLUE_BCC_EMAIL})",
        flush=True,
    )


# ===== END GMAIL TOOLS =====


def _record_continuity_tool_outcome(tool_name: str, tool_args: Dict[str, Any],
                                    result: Any) -> None:
    """Forward an actual tool result to the speaking robot's continuity layer
    when its turn collector is active."""
    routes = globals().get("_continuity_routes")
    if routes is None:
        return
    try:
        routes.record_tool_outcome(tool_name, tool_args, result)
    except Exception as e:
        log.warning(f"[JSPACE] could not record {tool_name} outcome: {e}")


def execute_tool(tool_name: str, tool_args: Dict[str, Any]) -> str:
    """
    Execute requested tool with enhanced error handling and state tracking.
    v8 ENHANCED: Better state tracking, timing, retry on transient failures.
    """
    import time
    start_time = time.time()

    # Backstop for chat-only users (Vilda's iPad): never let the LLM actually
    # fire a tool that's off-limits there, even if it tries. Safe outside a
    # request (there's no Vilda then). move_head is special-cased to a silent
    # success — the robot must NOT move while Blue talks to Vilda, but we don't
    # want Blue telling her he "can't move"; he just doesn't.
    if tool_name in _KID_BLOCKED_TOOLS:
        try:
            if _identify_user_from_request() in _CHAT_ONLY_USERS:
                if tool_name == "move_head":
                    return json.dumps({"success": True})
                if tool_name == "capture_camera":
                    # Her camera is the iPad frame already in the vision queue,
                    # not the PC webcam — nudge the model to use what it can
                    # already see instead of opening cv2.VideoCapture(0).
                    return json.dumps({"success": True, "message":
                        "You can already see through her camera — react to the "
                        "image you were just given; do not take a new photo."})
                return json.dumps({"success": False,
                                   "error": "That isn't available on this device."})
        except Exception:
            pass

    state = get_conversation_state()

    print(f"[TOOL] Executing tool: {tool_name}")
    print(f"   Arguments: {json.dumps(tool_args, indent=2)}")
    
    # v8: Track execution attempt
    max_retries = 2 if tool_name in ['web_search', 'read_gmail', 'browse_website'] else 1
    last_error = None
    
    for attempt in range(max_retries):
        try:
            result = _execute_tool_internal(tool_name, tool_args)
            elapsed = time.time() - start_time
            
            # Record successful execution with args
            state.record_tool_use(
                tool_name, 
                success=True, 
                pattern=f"{tool_name}:{list(tool_args.keys())}",
                args=tool_args
            )
            state.last_tool_result = truncate_text(result, 500)
            
            # v8: Update topic based on tool
            if tool_name == 'play_music':
                state.push_topic('music')
            elif tool_name in ['control_lights', 'music_visualizer']:
                state.push_topic('lights')
            elif tool_name in ['read_gmail', 'send_gmail', 'reply_gmail']:
                state.push_topic('email')
            
            print(f"   [OK] {tool_name} completed in {elapsed:.2f}s")
            _record_continuity_tool_outcome(tool_name, tool_args, result)
            return result
            
        except Exception as e:
            last_error = e
            if attempt < max_retries - 1:
                wait_time = 1 * (attempt + 1)
                print(f"   [RETRY] {tool_name} failed, retrying in {wait_time}s...")
                time.sleep(wait_time)
            continue
    
    # All retries failed
    elapsed = time.time() - start_time
    error_msg = str(last_error)
    
    # Record failed execution
    state.record_tool_use(tool_name, success=False, pattern=f"{tool_name}:error", args=tool_args)
    
    print(f"   [ERROR] {tool_name} failed after {elapsed:.2f}s: {error_msg}")
    
    # Provide helpful error message
    error_response = {
        "success": False,
        "error": error_msg,
        "tool": tool_name,
        "suggestion": _get_error_suggestion(tool_name, error_msg)
    }
    result = json.dumps(error_response)
    _record_continuity_tool_outcome(tool_name, tool_args, result)
    return result


def _get_error_suggestion(tool_name: str, error: str) -> str:
    """Get a helpful suggestion for common errors."""
    error_lower = error.lower()
    
    if "timeout" in error_lower or "timed out" in error_lower:
        return "The service took too long to respond. Try again in a moment."
    elif "connection" in error_lower or "network" in error_lower:
        return "Network connection issue. Check your internet connection."
    elif "not found" in error_lower:
        return "The requested resource wasn't found. Check the name or path."
    elif "permission" in error_lower or "unauthorized" in error_lower:
        return "Permission denied. You may need to re-authenticate."
    elif "rate limit" in error_lower:
        return "Too many requests. Please wait a moment before trying again."
    elif tool_name == "play_music":
        return "Try a different artist name or check if YouTube Music is running."
    elif tool_name == "control_lights":
        return "Check if the Philips Hue bridge is connected and accessible."
    elif tool_name in ["read_gmail", "send_gmail", "reply_gmail"]:
        return "Gmail authentication may have expired. Check credentials."
    else:
        return "Try rephrasing your request or provide more details."


def _execute_tool_internal(tool_name: str, tool_args: Dict[str, Any]) -> str:
    """Internal tool execution - called by execute_tool wrapper."""

    if tool_name == "move_head":
        action = (tool_args.get("action") or "").lower().strip()
        times = int(tool_args.get("times") or 2)
        ok = False
        if action.startswith("look_"):
            ok = blue_head.look(action[len("look_"):])
        elif action == "nod_yes":
            ok = blue_head.nod_yes(times)
        elif action == "shake_no":
            ok = blue_head.shake_no(times)
        elif action == "blink":
            ok = blue_head.blink(times)
        elif action in ("happy", "sad", "surprised", "curious", "neutral", "wink"):
            ok = blue_head.expression(action)
        if ok:
            return json.dumps({"success": True, "action": action})
        if not blue_head.is_available():
            return json.dumps({"success": False, "error": "head not connected"})
        return json.dumps({"success": False, "error": f"unknown action: {action}"})

    if tool_name == "head_eye_color":
        r = int(tool_args.get("r", 0))
        g = int(tool_args.get("g", 0))
        b = int(tool_args.get("b", 0))
        if blue_head.eye_color(r, g, b):
            return json.dumps({"success": True, "r": r, "g": g, "b": b})
        if not blue_head.is_available():
            return json.dumps({"success": False, "error": "head not connected"})
        return json.dumps({"success": False, "error": "eye colour failed"})

    if tool_name == "play_music":
        query = tool_args.get("query", "")
        action = tool_args.get("action", "play")
        service = tool_args.get("service", "youtube_music")

        if action == "search":
            result = search_music_info(query)
        else:
            result = play_music(query, service)
        print(f"   [OK] Music action completed")
        return result

    elif tool_name == "control_music":
        action = tool_args.get("action", "")
        result = control_music(action)
        print(f"   [OK] Music control executed")
        return result

    elif tool_name == "music_visualizer":
        action = tool_args.get("action", "start")

        if action == "start":
            duration = tool_args.get("duration", 300)
            style = tool_args.get("style", "party")
            result = start_music_visualizer(duration, style)
        elif action == "stop":
            result = stop_music_visualizer()
        else:
            result = f"Unknown visualizer action: {action}"

        print(f"   [OK] Visualizer action completed")
        return result

    elif tool_name == "control_lights":
        result = execute_light_control(
            tool_args.get("action"),
            tool_args.get("light_name"),
            tool_args.get("brightness"),
            tool_args.get("color"),
            tool_args.get("mood")
        )
        print(f"   [OK] Light control executed")
        return result

    elif tool_name == "search_documents":
        global _vision_queue
        query = tool_args.get("query", "")
        max_results = tool_args.get("max_results", 3)
        result = _search_documents_guarded(query, max_results)
        print(f"   [OK] Document search completed")

        # Check if result contains images (special JSON format)
        try:
            result_data = json.loads(result)
            if isinstance(result_data, dict) and result_data.get("_type") == "document_search_with_images":
                images = result_data.get("images", [])
                text_docs = result_data.get("text_documents", [])

                # Store images globally so they can be injected into next LLM call
                for img in images:
                    _vision_queue.add_image(
                        filepath=img['filepath'],
                        filename=img['filename'],
                        is_camera=False
                    )
                print(f"   [VISION] Stored {len(images)} image(s) for vision model")

                # Build text response
                response_parts = []
                if images:
                    image_names = [img['filename'] for img in images]
                    response_parts.append(f"Found {len(images)} image(s): {', '.join(image_names)}")
                    response_parts.append("(Images will be shown to vision model in next response)")

                if text_docs:
                    response_parts.append("\n\nText documents found:\n\n" + "\n---\n\n".join(text_docs))

                return "\n".join(response_parts) if response_parts else "Found documents."
        except (json.JSONDecodeError, TypeError):
            # Not JSON or not our special format - return as-is
            pass

        return result

    elif tool_name == "view_image":
        filename = tool_args.get("filename")
        query = tool_args.get("query")
        result = view_image(filename=filename, query=query)
        print(f"   [OK] Image view requested")
        return result

    elif tool_name == "capture_camera":
        result = capture_camera_image(
            look=tool_args.get("look"),
            zoom=tool_args.get("zoom"),
            zoom_region=tool_args.get("zoom_region") or "center",
        )
        print(f"   [OK] Camera capture completed")
        return result

    elif tool_name == "email_snapshot":
        result = _execute_email_snapshot(tool_args)
        print(f"   [OK] Snapshot capture+email completed")
        return result

    elif tool_name == "recall_visual_memory":
        if not VISUAL_MEMORY_AVAILABLE:
            return json.dumps({"error": "Visual memory not available"})
        try:
            vm = get_visual_memory()
            query = tool_args.get("query")
            hours = tool_args.get("hours", 24)

            if query:
                observations = vm.search_observations(query, limit=10)
                search_type = f"search for '{query}'"
            else:
                observations = vm.get_visual_timeline(hours)
                search_type = f"timeline (last {hours}h)"

            if not observations:
                return json.dumps({
                    "success": True,
                    "search_type": search_type,
                    "message": "No visual memories found for that query.",
                    "observations": []
                })

            # Format observations for the LLM
            formatted = []
            for obs in observations:
                entry = {
                    "timestamp": obs.get("timestamp", ""),
                    "description": obs.get("scene_description", ""),
                    "location": obs.get("location"),
                    "people": obs.get("people_present"),
                    "objects": obs.get("notable_objects"),
                    "has_image": bool(obs.get("image_path"))
                }
                formatted.append(entry)

            # Also get scene change info
            changes = vm.detect_scene_changes("")
            result = {
                "success": True,
                "search_type": search_type,
                "total_memories": len(formatted),
                "observations": formatted,
                "_instruction": "Summarize these visual memories naturally. Tell the user what you remember seeing, when, and where. If they asked about changes, compare observations."
            }
            if changes.get("has_previous"):
                result["last_seen_ago"] = changes["time_since"]
            print(f"   [OK] Visual memory recall: {len(formatted)} observations ({search_type})")
            return json.dumps(result)
        except Exception as e:
            print(f"   [ERROR] Visual memory recall failed: {e}")
            return json.dumps({"error": str(e)})

    elif tool_name == "get_weather":
        result = get_weather_data(tool_args.get("location", ""))
        print(f"   [OK] Weather retrieved")
        return result

    elif tool_name == "web_search":
        result = execute_web_search(tool_args.get("query", ""))
        print(f"   [OK] Search completed")
        return result

    elif tool_name == "search_scholar":
        if not SCHOLAR_AVAILABLE:
            return json.dumps({"error": "Scholarly search is not available."})
        result = execute_scholar_search(tool_args)
        print(f"   [OK] Scholarly search completed")
        return result

    elif tool_name == "get_paper":
        if not SCHOLAR_AVAILABLE:
            return json.dumps({"error": "Scholarly search is not available."})
        result = execute_get_paper(tool_args)
        print(f"   [OK] Paper lookup completed")
        return result

    elif tool_name == "read_paper":
        if not SCHOLAR_AVAILABLE:
            return json.dumps({"error": "Scholarly search is not available."})
        result = execute_read_paper(tool_args)
        try:
            _rp = json.loads(result)
            if _rp.get("success"):
                print(f"   [OK] Paper read via {_rp.get('access_route')} ({_rp.get('text_chars')} chars)")
            else:
                print(f"   [WARN] Paper fetch failed: {_rp.get('error')}")
        except Exception:
            pass
        return result

    elif tool_name == "run_javascript":
        try:
            import js2py
            result = js2py.eval_js(tool_args.get("code", ""))
            return f"Result: {result}"
        except Exception as e:
            return f"Error: {str(e)}"

    elif tool_name == "create_document":
        filename = tool_args.get("filename", "")
        content = tool_args.get("content", "")
        file_type = tool_args.get("file_type", "txt")
        result = create_document_file(filename, content, file_type)
        print(f"   [OK] Document created")
        return result

    elif tool_name == "browse_website":
        print(f"   [DEBUG] Calling _execute_browse_website...")
        result = _execute_browse_website(tool_args)
        print(f"   [DEBUG] Got result, length: {len(result)} chars")
        # Parse the result to check success
        try:
            result_obj = json.loads(result)
            if result_obj.get("success"):
                print(f"   [OK] Browse completed - fetched {len(result_obj.get('text', ''))} chars")
            else:
                print(f"   [ERROR] Browse failed: {result_obj.get('error', 'Unknown error')}")
        except Exception:
            pass
        return result

    elif tool_name == "read_gmail":
        result = _execute_read_gmail(tool_args)
        # Add operation type to help Blue understand what just happened
        try:
            result_obj = json.loads(result)
            result_obj["_operation_type"] = "READ_EMAIL"
            result_obj["_instruction"] = "You just READ emails. User asked to check/read, NOT to reply or send."
            result = json.dumps(result_obj)
        except Exception:
            pass
        print(f"   [OK] Gmail READ completed")
        return result

    elif tool_name == "send_gmail":
        result = _execute_send_gmail(tool_args)
        # Add operation type to help Blue understand what just happened
        try:
            result_obj = json.loads(result)
            result_obj["_operation_type"] = "SEND_EMAIL"
            result_obj["_instruction"] = "You just SENT a new email. User asked to send, NOT to read or reply."
            result = json.dumps(result_obj)
        except Exception:
            pass
        print(f"   [OK] Gmail SEND completed")
        return result

    elif tool_name == "reply_gmail":
        result = _execute_reply_gmail(tool_args)
        # Add operation type to help Blue understand what just happened
        try:
            result_obj = json.loads(result)
            result_obj["_operation_type"] = "REPLY_EMAIL"
            result_obj["_instruction"] = "You just REPLIED to emails. User asked to reply/respond, NOT to just read."
            result = json.dumps(result_obj)
        except Exception:
            pass
        print(f"   [OK] Gmail REPLY completed")
        return result

    elif tool_name == "auto_reply_emails":
        result = _execute_auto_reply_inbox(tool_args)
        try:
            result_obj = json.loads(result)
            result_obj["_operation_type"] = "AUTO_REPLY_EMAILS"
            result_obj["_instruction"] = (
                "You just scanned Blue's own gmail inbox "
                "(alevantresearch@gmail.com) and sent autonomous replies "
                "to every personal email there. Summarise to the user: "
                "how many replies were sent, who they went to, and the "
                "subject lines. Note each reply is BCC'd to Alex at "
                "alevant1905@gmail.com so he can read the full text "
                "there. State ONLY what the tool result actually says — "
                "do NOT speculate about Blue's capabilities and do NOT "
                "claim Blue cannot do something the tool just did. Blue's "
                "inbox is alevantresearch@gmail.com (NOT alevant1905 or "
                "alevant@yorku.ca — those are Alex's addresses)."
            )
            result = json.dumps(result_obj)
        except Exception:
            pass
        print(f"   [OK] Gmail AUTO-REPLY completed")
        return result


    # ===== Enhanced Tools Handlers =====
    elif tool_name == "create_reminder" and ENHANCED_TOOLS_AVAILABLE:
        result = CalendarManager.create_reminder(**tool_args)
        print(f"   [OK] Reminder created")
        return json.dumps(result)

    elif tool_name == "get_upcoming_reminders" and ENHANCED_TOOLS_AVAILABLE:
        result = CalendarManager.get_upcoming_reminders(**tool_args)
        print(f"   [OK] Retrieved reminders")
        return json.dumps(result)

    elif tool_name == "complete_reminder" and ENHANCED_TOOLS_AVAILABLE:
        result = CalendarManager.complete_reminder(**tool_args)
        print(f"   [OK] Reminder completed")
        return json.dumps(result)

    elif tool_name == "cancel_reminder" and ENHANCED_TOOLS_AVAILABLE:
        result = CalendarManager.cancel_reminder(**tool_args)

    elif tool_name == "reschedule_reminder" and ENHANCED_TOOLS_AVAILABLE:
        result = CalendarManager.update_reminder(**tool_args)

    elif tool_name == "add_contact" and ENHANCED_TOOLS_AVAILABLE:
        result = ContactManager.add_contact(**tool_args)

    elif tool_name == "list_contacts" and ENHANCED_TOOLS_AVAILABLE:
        result = ContactManager.list_contacts(**tool_args)

    elif tool_name == "find_contact" and ENHANCED_TOOLS_AVAILABLE:
        result = ContactManager.find_contact(**tool_args)
        print(f"   [OK] Reminder cancel attempt: success={result.get('success')}")
        return json.dumps(result)

    elif tool_name == "create_task" and ENHANCED_TOOLS_AVAILABLE:
        result = TaskManager.create_task(**tool_args)
        print(f"   [OK] Task created")
        return json.dumps(result)

    elif tool_name == "get_tasks" and ENHANCED_TOOLS_AVAILABLE:
        result = TaskManager.get_tasks(**tool_args)
        print(f"   [OK] Retrieved tasks")
        return json.dumps(result)

    elif tool_name == "complete_task" and ENHANCED_TOOLS_AVAILABLE:
        result = TaskManager.complete_task(**tool_args)
        print(f"   [OK] Task completed")
        return json.dumps(result)

    elif tool_name == "create_note" and ENHANCED_TOOLS_AVAILABLE:
        result = NoteManager.create_note(**tool_args)
        print(f"   [OK] Note saved")
        return json.dumps(result)

    elif tool_name == "search_notes" and ENHANCED_TOOLS_AVAILABLE:
        result = NoteManager.search_notes(**tool_args)
        print(f"   [OK] Note search completed")
        return json.dumps(result)

    elif tool_name == "set_timer" and ENHANCED_TOOLS_AVAILABLE:
        result = TimerManager.set_timer(**tool_args)
        print(f"   [OK] Timer set")
        return json.dumps(result)

    elif tool_name == "check_timers" and ENHANCED_TOOLS_AVAILABLE:
        result = TimerManager.check_timers()
        print(f"   [OK] Timer status checked")
        return json.dumps(result)

    elif tool_name == "get_system_info" and ENHANCED_TOOLS_AVAILABLE:
        result = SystemController.get_system_info()
        print(f"   [OK] System info retrieved")
        return json.dumps(result)

    elif tool_name == "take_screenshot" and ENHANCED_TOOLS_AVAILABLE:
        result = SystemController.take_screenshot(**tool_args)
        print(f"   [OK] Screenshot captured")
        return json.dumps(result)

    elif tool_name == "launch_application" and ENHANCED_TOOLS_AVAILABLE:
        result = SystemController.launch_application(**tool_args)
        print(f"   [OK] Application launched")
        return json.dumps(result)

    elif tool_name == "set_volume" and ENHANCED_TOOLS_AVAILABLE:
        result = SystemController.set_volume(**tool_args)
        print(f"   [OK] Volume set")
        return json.dumps(result)

    elif tool_name == "list_files":
        if ENHANCED_TOOLS_AVAILABLE:
            result = FileOperations.list_files(**tool_args)
            print(f"   [OK] Files listed")
            return json.dumps(result)
        else:
            return json.dumps({
                "success": False,
                "message": "File system operations are not available. Use search_documents to access uploaded documents."
            })

    elif tool_name == "read_file":
        if ENHANCED_TOOLS_AVAILABLE:
            result = FileOperations.read_file(**tool_args)
            print(f"   [OK] File read")
            return json.dumps(result)
        else:
            return json.dumps({
                "success": False,
                "message": "File reading is not available. Use search_documents to read uploaded documents."
            })

    elif tool_name == "write_file" and ENHANCED_TOOLS_AVAILABLE:
        result = FileOperations.write_file(**tool_args)
        print(f"   [OK] File written")
        return json.dumps(result)

    elif tool_name == "get_file_info" and ENHANCED_TOOLS_AVAILABLE:
        result = FileOperations.get_file_info(**tool_args)
        print(f"   [OK] File info retrieved")
        return json.dumps(result)

    elif tool_name == "story_prompt" and ENHANCED_TOOLS_AVAILABLE:
        result = StorytellingTools.story_prompt(**tool_args)
        print(f"   [OK] Story prompt generated")
        return json.dumps(result)

    elif tool_name == "educational_activity" and ENHANCED_TOOLS_AVAILABLE:
        result = StorytellingTools.educational_activity(**tool_args)
        print(f"   [OK] Activity suggested")
        return json.dumps(result)

    elif tool_name == "get_local_time" and ENHANCED_TOOLS_AVAILABLE:
        result = LocationServices.get_local_time()
        print(f"   [OK] Local time retrieved")
        return json.dumps(result)

    elif tool_name == "get_sunrise_sunset" and ENHANCED_TOOLS_AVAILABLE:
        result = LocationServices.get_sunrise_sunset()
        print(f"   [OK] Sunrise/sunset times retrieved")
        return json.dumps(result)

    elif tool_name == "remember_person" and VISUAL_MEMORY_AVAILABLE:
        name = tool_args.get("name", "")
        appearance = tool_args.get("appearance", "")
        relationship = tool_args.get("relationship", "")
        notes = tool_args.get("notes", "")

        try:
            vm = get_visual_memory()
            vm.add_person(
                name=name,
                typical_appearance=appearance,
                relationship=relationship,
                notes=notes
            )
            print(f"   [OK] Remembered person: {name}")
            return json.dumps({
                "success": True,
                "message": f"I'll remember {name}. Next time I see them through my camera, I'll recognize them."
            })
        except Exception as e:
            return json.dumps({
                "success": False,
                "error": f"Failed to remember person: {str(e)}"
            })

    elif tool_name == "remember_place" and VISUAL_MEMORY_AVAILABLE:
        name = tool_args.get("name", "")
        description = tool_args.get("description", "")
        typical_contents = tool_args.get("typical_contents", "")
        notes = tool_args.get("notes", "")

        try:
            vm = get_visual_memory()
            vm.add_place(
                name=name,
                description=description,
                typical_contents=typical_contents,
                notes=notes
            )
            print(f"   [OK] Remembered place: {name}")
            return json.dumps({
                "success": True,
                "message": f"I'll remember {name}. Next time I see this location, I'll recognize it."
            })
        except Exception as e:
            return json.dumps({
                "success": False,
                "error": f"Failed to remember place: {str(e)}"
            })

    elif tool_name == "who_do_i_know" and VISUAL_MEMORY_AVAILABLE:
        try:
            vm = get_visual_memory()
            people = vm.get_all_people()
            places = vm.get_all_places()

            result = {"people": [], "places": []}

            for person in people:
                result["people"].append({
                    "name": person['name'],
                    "relationship": person['relationship'],
                    "appearance": person['typical_appearance'],
                    "times_seen": person['times_seen']
                })

            for place in places:
                result["places"].append({
                    "name": place['name'],
                    "description": place['description'],
                    "times_seen": place['times_seen']
                })

            print(f"   [OK] Retrieved visual memory: {len(people)} people, {len(places)} places")
            return json.dumps(result, indent=2)
        except Exception as e:
            return json.dumps({
                "success": False,
                "error": f"Failed to retrieve visual memory: {str(e)}"
            })

    elif tool_name == "analyze_with_chat_theory" and ACADEMIC_ASSISTANT_AVAILABLE:
        topic = tool_args.get("topic", "")
        context = tool_args.get("context", "")

        result = analyze_with_chat(topic, context)
        print(f"   [OK] Generated CHAT analysis for: {topic}")
        return result

    elif tool_name == "prepare_lecture" and ACADEMIC_ASSISTANT_AVAILABLE:
        topic = tool_args.get("topic", "")
        duration = tool_args.get("duration", 50)
        course = tool_args.get("course", "")
        level = tool_args.get("level", "undergraduate")

        result = prepare_lecture(topic, duration, course, level)
        print(f"   [OK] Generated lecture outline for: {topic}")
        return result

    elif tool_name == "discussion_questions" and ACADEMIC_ASSISTANT_AVAILABLE:
        reading = tool_args.get("reading", "")
        topic = tool_args.get("topic", "")

        result = generate_discussion_questions(reading, topic)
        print(f"   [OK] Generated discussion questions for: {topic}")
        return result

    elif tool_name == "simulate_student_questions" and ACADEMIC_ASSISTANT_AVAILABLE:
        topic = tool_args.get("topic", "")
        context = tool_args.get("context", "")

        result = simulate_student_q_and_a(topic, context)
        print(f"   [OK] Simulated student questions for: {topic}")
        return result

    elif tool_name == "check_proactive_suggestions" and PROACTIVE_ASSISTANCE_AVAILABLE:
        try:
            pa = get_proactive_assistance()
            # Get current person from context (default to Alex)
            person = "Alex"  # Could be enhanced to detect from visual memory

            suggestions = pa.check_for_suggestions(person)

            if suggestions:
                result = {
                    "has_suggestions": True,
                    "suggestions": [
                        {
                            "type": suggestion.suggestion_type,
                            "priority": suggestion.priority,
                            "message": suggestion.message,
                            "action_available": suggestion.action_available
                        }
                        for suggestion in suggestions
                    ]
                }
                print(f"   [OK] Found {len(suggestions)} proactive suggestions")
            else:
                result = {
                    "has_suggestions": False,
                    "message": "No suggestions at this time"
                }
                print(f"   [OK] No proactive suggestions at this time")

            return json.dumps(result, indent=2)
        except Exception as e:
            return json.dumps({
                "success": False,
                "error": f"Failed to check suggestions: {str(e)}"
            })

    # Fallback: get_local_time works without enhanced tools
    if tool_name == "get_local_time":
        from datetime import datetime
        now = datetime.now()
        return json.dumps({
            "success": True,
            "time": now.strftime("%I:%M %p"),
            "date": now.strftime("%A, %B %d, %Y"),
            "iso": now.isoformat()
        })

    return f"Unknown tool: {tool_name}"


# ================================================================================
# LEGACY DETECTION FUNCTIONS REMOVED (523 lines deleted)
# ================================================================================
# All legacy detection functions have been removed. Tool detection is now
# exclusively handled by the modular system in blue/tool_selector/.
#
# Removed functions (20 total, ~523 lines):
#   • detect_no_tool_intent() - Casual conversation detection
#   • detect_search_intent() - Web search detection
#   • detect_javascript_intent() - JavaScript execution
#   • detect_weather_intent() - Weather queries
#   • detect_light_intent() - Light control with false positive filters
#   • detect_visualizer_intent() - Music visualizer
#   • detect_document_intent() - Document operations
#   • detect_create_document_intent() - Document creation
#   • detect_browse_intent() - Website browsing
#   • detect_music_play_intent() - Music playback
#   • detect_music_control_intent() - Music controls (pause/skip/volume)
#   • detect_document_retrieval_intent() - Document reading
#   • detect_document_search_intent() - Document search
#   • detect_web_search_intent_improved() - Enhanced web search
#   • detect_hallucinated_search() - Hallucination detection
#   • detect_gmail_read_intent() - Gmail reading
#   • detect_gmail_send_intent() - Gmail sending
#   • detect_gmail_reply_intent() - Gmail replies
#   • detect_fanmail_reply_intent() - Fanmail replies
#   • detect_gmail_operation_intent() - Unified Gmail operation detector
#
# All functionality preserved in blue/tool_selector/detectors/
# Commit: 2025-12-27 - Legacy code cleanup
# ================================================================================


def _get_text_content(msg):
    """Extract text content cheaply, skipping base64 image data."""
    content = msg.get('content', '')
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts = []
        for part in content:
            if isinstance(part, dict) and part.get('type') == 'text':
                parts.append(part.get('text', ''))
        return ' '.join(parts)
    return ''


def _estimate_tokens(text: str) -> int:
    """Rough chars/4 token heuristic — no deps, good enough for budgeting."""
    return len(text) // 4 if text else 0


def _estimate_payload_tokens(messages, tools) -> int:
    """Estimate total prompt tokens for a chat completion payload."""
    total = 0
    for m in messages:
        # ~5 tokens of role + chat-template overhead per message.
        total += 5 + _estimate_tokens(_get_text_content(m))
    if tools:
        try:
            import json as _json
            total += _estimate_tokens(_json.dumps(tools))
        except Exception:
            total += 200 * len(tools)
    return total + 64  # final chat-template wrap pad


def _trim_messages_for_budget(messages, tools, budget_tokens: int, min_keep_tail: int = 0):
    """Drop oldest non-system messages until estimated tokens fit the budget.

    Always preserves leading system message(s) and the final user message.
    `min_keep_tail` additionally protects the N most recent middle messages
    (≈ the live conversation tail) even if the result stays over budget.
    Without it, a fixed payload bigger than the budget (system prompt + the
    full tools schema already exceed 6500t) silently eats the ENTIRE thread,
    and short replies like "yes" reach the model with no context at all —
    Blue then greets mid-conversation like it just woke up. Staying somewhat
    over budget is fine: the n_ctx self-heal in call_lm_studio catches the
    rare genuine overflow and retrims against the model's REAL context size.
    Returns (trimmed_messages, dropped_count).
    """
    if not messages:
        return messages, 0

    sys_end = 0
    while sys_end < len(messages) and messages[sys_end].get('role') == 'system':
        sys_end += 1

    last_user_idx = None
    for i in range(len(messages) - 1, -1, -1):
        if messages[i].get('role') == 'user':
            last_user_idx = i
            break
    if last_user_idx is None or last_user_idx < sys_end:
        return messages, 0

    head = messages[:sys_end]
    middle = list(messages[sys_end:last_user_idx])
    tail = messages[last_user_idx:]

    dropped = 0
    while True:
        candidate = head + middle + tail
        if _estimate_payload_tokens(candidate, tools) <= budget_tokens:
            return candidate, dropped
        if len(middle) <= max(0, min_keep_tail):
            return candidate, dropped  # protected tail / already minimal
        middle.pop(0)
        dropped += 1


# Shared constants
_VISION_PHRASES = ('i see', 'cozy living room', 'soft lighting', 'warm vibe',
                   'in his office', 'in her office', 'wearing', 'holding a mug',
                   'sitting at', 'glasses and', 'giving a thumbs', 'the monitor shows',
                   'behind him', 'behind her', 'the kitchen', 'the living room')

# Validation helper functions
_HALLUCINATION_RE = re.compile(
    r'i searched|according to (?:my|the) search|i found (?:that|the following)',
    re.IGNORECASE
)

def detect_hallucinated_search(response: str) -> bool:
    """Detect if LLM is hallucinating a web search that didn't happen."""
    return bool(_HALLUCINATION_RE.search(response))


# The model HAS live web access via web_search, but on current-events questions
# it sometimes claims otherwise and sends the user off to check ESPN/FIFA.com
# themselves (four straight turns of it on "who is left in the world cup",
# 2026-07-09). When a reply is that refusal AND no tool ran, force the search.
_WEB_REFUSAL_RE = re.compile(
    # ['’] — the model emits curly apostrophes ("I don’t") as often as ASCII.
    r"i (?:don['’]?t|do not) have (?:a )?(?:live|real[- ]?time|current)"
    r"|i (?:don['’]?t|do not) have access to (?:live|real[- ]?time|current)"
    r"|can['’]?t access (?:live|real[- ]?time|current)"
    r"|no (?:live|real[- ]?time) (?:access|feed|data|scoreboard)"
    r"|i(?:['’]d)? recommend checking|your best bet is to check"
    r"|best (?:place|way) to check|check a live sports"
    r"|you can check (?:the|a|full|live)"
    r"|would you like me to (?:look up|search|check)"
    r"|if you(?:['’]d| would) like,? i can (?:look up|search|check)",
    re.IGNORECASE,
)


def detect_web_refusal(response: str) -> bool:
    """A no-live-access claim / go-check-a-website deflection — the tell that
    the model skipped the web_search it should have run."""
    return bool(_WEB_REFUSAL_RE.search(response or ""))


_DOCUMENT_REFUSAL_RE = re.compile(
    r"\b(?:i (?:still )?(?:cannot|can['\u2019]?t|am unable to) "
    r"(?:access|read|retrieve|extract|open|see) (?:the |that |this )?(?:text|file|pdf|document)"
    r"|i (?:currently )?lack (?:the )?(?:specific )?(?:pdf[- ]reading )?tool"
    r"|i do not have (?:a |the )?(?:pdf[- ]reading|file[- ]reading) tool"
    r"|search_documents (?:can|could) only (?:scan|search|look for) keywords"
    r"|the (?:file|pdf) is not accessible"
    r"|(?:local )?path (?:is |was )?(?:invalid|wrong|not valid|an error)"
    r"|not a valid file at (?:that|this) location"
    r"|please (?:copy and paste|upload) (?:the |a )?(?:text|pdf|file)"
    r"|rely on (?:my )?(?:internal )?training data"
    r"|search (?:did not|didn['\u2019]?t) (?:pull up|return|retrieve) any (?:specific )?text)\b",
    re.I,
)


def detect_document_refusal(response: str) -> bool:
    """False local-file/path/capability claim after document extraction."""
    return bool(_DOCUMENT_REFUSAL_RE.search(response or ""))


# Blue DOES maintain Alex's household calendar (the reminders/events store,
# surfaced at /calendar) and can create / reschedule / cancel entries with his
# tools. Yet the base model likes to disown it — "I cannot modify your personal
# calendar", "I don't actually maintain a persistent calendar", "I only have
# read-only access", "update it manually in your calendar app" (2026-07-14, the
# CMDS4740 end-date request). When a reply is that false disavowal AND the user
# was asking about the calendar, we load the REAL calendar and make him answer
# from it — same shape as detect_web_refusal.
_CALENDAR_DENIAL_RE = re.compile(
    r"i (?:can['’]?t|cannot|am unable to|do not|don['’]?t) "
    r"(?:actually )?(?:modify|change|edit|update|revise|access|maintain|keep|"
    r"manage|add to|write to|save to) "
    r"(?:your |the |a |any )?(?:personal |external |persistent |synced |real )*"
    r"(?:calendar|schedule|reminders?|events?)"
    r"|i (?:don['’]?t|do not) (?:actually )?(?:maintain|keep|have) "
    r"(?:a |any )?(?:persistent |real |synced |personal )*(?:calendar|schedule)"
    r"|(?:only |just )?(?:have |with )?read[- ]only access"
    r"|(?:add|update|set|change) (?:this|it|that|the date|the final class) "
    r"(?:date )?(?:manually )?(?:yourself )?in your (?:actual |own )?"
    r"(?:calendar|scheduling)"
    r"|in your (?:actual |own )?calendar (?:application|app)"
    r"|need to (?:manually )?(?:add|update|set|change) (?:this|it|that|the date)"
    r"|you will need to update this date",
    re.IGNORECASE,
)

# Only fire the recovery when the user's OWN turn was actually about the
# calendar — never on an unrelated "I can't access your bank" style line.
_CALENDAR_TOPIC_RE = re.compile(
    r"calendar|schedule|reminder|appointment|\bclass(?:es)?\b|\bcourse\b"
    r"|\bevent\b|cmds\s*\d",
    re.IGNORECASE,
)


def detect_calendar_denial(response: str) -> bool:
    """False 'I have no calendar / only read-only / do it yourself in your app'
    claim — Blue maintains the household calendar and can edit it via tools."""
    return bool(_CALENDAR_DENIAL_RE.search(response or ""))


_CALENDAR_EDIT_RE = re.compile(
    r"\b(?:reschedul\w*|move|push|postpone|bump|shift|revis\w*|edit\w*|"
    r"updat\w*|chang\w*|renam\w*|cancel|delete|remove|end|extend|shorten|"
    r"add|set)\b",
    re.IGNORECASE,
)


def _user_asked_calendar_edit(message: str) -> bool:
    """True when the user's turn asked to CHANGE the calendar (not merely view
    it) — gates whether the denial-recovery forces reschedule_reminder."""
    return bool(_CALENDAR_EDIT_RE.search(message or ""))


def _document_search_succeeded(result: str) -> bool:
    """Whether search_documents returned real local text/passages."""
    text = result or ""
    lower = text.lower()
    failures = (
        "i couldn't find any documents", "i don't have any documents",
        "local library extraction failed", "path is stale",
        "trouble searching your documents", "searching it is taking too long",
    )
    if any(marker in lower for marker in failures):
        return False
    return bool(
        "local library read succeeded" in lower
        or ("here's what i found in your documents" in lower and "[" in text)
        or ("here are the most relevant passages" in lower and "[" in text)
        or ("course schedule & readings" in lower and "[" in text)
    )


# A local model sometimes emits its tool call as visible TEXT instead of a real
# tool call — observed live 2026-07-09: the reply contained a literal
# "<tool_call><function=web_search><parameter=query>...</parameter></function>
# </tool_call>" block, which reached the user as words. Parse it and run it.
_LEAKED_TOOL_RE = re.compile(
    r"<tool_call>\s*<function=([\w\-]+)>(.*?)</function>\s*</tool_call>"
    r"|<function=([\w\-]+)>(.*?)</function>",
    re.DOTALL,
)
_LEAKED_PARAM_RE = re.compile(r"<parameter=([\w\-]+)>\s*(.*?)\s*</parameter>", re.DOTALL)


def parse_leaked_tool_call(response: str):
    """(tool_name, args) when the reply contains a tool call written out as
    text; None otherwise. Handles <parameter=...> bodies and a JSON body."""
    m = _LEAKED_TOOL_RE.search(response or "")
    if not m:
        return None
    name = (m.group(1) or m.group(3) or "").strip()
    body = m.group(2) or m.group(4) or ""
    args = {k: v.strip() for k, v in _LEAKED_PARAM_RE.findall(body)}
    if not args:
        try:
            j = json.loads(body.strip())
            if isinstance(j, dict):
                args = j
        except Exception:
            pass
    return (name, args) if name else None


# Patterns that mean "I performed action X" — keyed by the tool that should
# have been called. If the assistant claims one of these but tool_calls is
# empty, the response is a hallucination and we need to force-retry.
_ACTION_CLAIM_PATTERNS = {
    # MUST come before send_gmail: "I've emailed you the photo" is a
    # photo-delivery claim — recovering it with plain send_gmail would mail
    # words with no picture attached. Order matters because
    # detect_hallucinated_action returns the FIRST matching key.
    "email_snapshot": re.compile(
        r"\b(?:sent|emailed|mailed|sending|emailing|mailing)\b[^.?!]{0,60}?"
        r"\b(?:photo|picture|snapshot|pic|image)\b"
        r"|\b(?:photo|picture|snapshot|pic|image)\b[^.?!]{0,40}?"
        r"\b(?:sent|emailed|mailed|on its way|in your inbox)\b",
        re.IGNORECASE,
    ),
    "send_gmail": re.compile(
        r"\b(?:"
        # Past tense: "I sent", "I've sent", "I've emailed", "I just sent"
        r"i(?:'ve| have)?\s+(?:just\s+)?sent|"
        r"i(?:'ll)?\s+sent|i(?:'ve| have)?\s+emailed|"
        r"i\s+just\s+emailed|"
        # Status: "email has been/is/was sent", "email was delivered"
        r"email\s+(?:has been|is|was)\s+sent|"
        r"email (?:was )?delivered|email\s+sent|"
        r"sent (?:the|that|an?) email|"
        # "sent to you at <email>", "sent to <name>", "sent it", "sent that",
        # "sent over". Fires whenever 'sent' is followed by a delivery target.
        r"sent\s+(?:to\b|it\b|that\b|over\b|you\b)|"
        # "sent the headlines to <name>" / "sent the summary to <email>"
        r"sent\s+(?:the|a|that)\s+\w+(?:\s+\w+)?\s+to\b|"
        r"message\s+(?:has been|is|was)\s+sent|"
        # Present continuous (the model often says "Sending...!" as if it's happening now):
        r"(?:i'?m\s+|^|\s)(?:sending|emailing|delivering|firing\s+off)\s+(?:the|that|an?|it|you|over|to)|"
        r"sending\s+(?:the|that|an?|it|over|now|you)|"
        r"emailing\s+(?:you|it|that|the|now)"
        r")\b",
        re.IGNORECASE,
    ),
    "control_lights": re.compile(
        r"\b(?:"
        r"i(?:'ve| have)?\s+(?:just\s+)?(?:turned|set|changed|adjusted|switched|put)\s+(?:the\s+)?lights?|"
        # "lights are (now) set/on/off..." — allow an adverb (now/all) to sit
        # between the verb-to-be and the action verb.
        r"lights?\s+(?:are|were|have been|'re)\s+(?:\w+\s+){0,2}(?:turned|set|changed|adjusted|switched|on|off)|"
        # Present continuous
        r"(?:i'?m\s+)?(?:switching|turning|setting|changing|adjusting)\s+(?:the\s+)?lights?|"
        # Mood/scene-centric: "the galaxy mood is now active/on/set"
        r"(?:mood|scene)\s+(?:is|has been|'s)\s+(?:now\s+)?(?:on|set|active|applied|enabled)"
        r")\b",
        re.IGNORECASE,
    ),
    "play_music": re.compile(
        r"\b(?:i(?:'ve| have)?\s+(?:started|begun|put on|queued)\s+(?:the\s+|some\s+)?(?:music|song|track|playlist)|"
        r"(?:music|song|track|playlist)\s+(?:is|has been)\s+(?:playing|started|queued))\b",
        re.IGNORECASE,
    ),
    "create_document": re.compile(
        r"\b(?:i(?:'ve| have)?\s+(?:created|saved|written)\s+(?:the\s+|a\s+|that\s+)?(?:document|file|note|list)|"
        r"document\s+(?:has been|is|was)\s+(?:created|saved))\b",
        re.IGNORECASE,
    ),
    # Claims to have read/reviewed a specific document. Without this, the
    # model can confidently say "I've re-read Job_Talk_Script7.docx" with
    # zero tool calls — pure hallucination.
    "search_documents": re.compile(
        r"\b(?:"
        r"i(?:'ve| have)?\s+(?:re-?read|reread|reviewed|gone through|looked through|"
        r"checked|opened|pulled up|brought up)\s+(?:the\s+|that\s+|your\s+|my\s+)?"
        r"(?:document|file|pdf|docx?|script|paper|notes|cv|dossier|"
        r"[A-Z][\w_-]*\.(?:pdf|docx?|txt|md))"
        r"|"
        r"i(?:'ve| have)?\s+(?:taken|had)\s+(?:another|a\s+(?:second|fresh|closer))\s+look\s+at"
        r")\b",
        re.IGNORECASE,
    ),
    # Claims to be seeing through the camera right now. Without an actual
    # capture_camera call the model will happily narrate "I can see you're
    # in the kitchen" from nothing — by email there's no live feed, so any
    # such claim must be backed by a real capture.
    "capture_camera": re.compile(
        r"(?:"
        # "I can see <concrete scene>" but NOT the figurative "I can see why/
        # how/what you mean / that you feel / your point".
        r"\bi\s+can\s+see\s+(?!why\b|how\b|what\b|that\s+you|your\s+point|where\s+you)"
        r"(?:you|your|a\s|an\s|the\s|two\b|three\b|some\b|someone|people)|"
        # Explicitly camera/view anchored — unambiguous.
        r"(?:in|through)\s+(?:my|the)\s+(?:camera|view|frame)|"
        r"\blooking\s+(?:in|at)\s+(?:my|the)\s+camera|"
        r"(?:i(?:'ve| have)?\s+)?(?:just\s+)?(?:took|taken|captured|snapped)\s+(?:a\s+)?(?:photo|picture|snapshot)|"
        r"\bin\s+front\s+of\s+me\s+(?:is|are|i\s+(?:can\s+)?see)\b"
        r")",
        re.IGNORECASE,
    ),
}


def detect_hallucinated_action(response: str) -> str | None:
    """If the response claims to have performed a tool-required action but no
    tool was called, return the name of the tool that *should* have been
    called. Otherwise return None.
    """
    if not response or not isinstance(response, str):
        return None
    for tool_name, pattern in _ACTION_CLAIM_PATTERNS.items():
        if pattern.search(response):
            return tool_name
    return None


# Words in the USER's message that mean the claimed action was actually asked
# for. The hallucinated-action recovery used to force the claimed tool
# UNCONDITIONALLY — so a confabulated "I sent the introduction email to the
# class" (2026-07-09, nobody asked for ANY email) would be turned into a REAL
# email to an invented address. A claim with no matching request must be
# regenerated or scrubbed, never executed.
_ACTION_REQUEST_WORDS = {
    "send_gmail": ("email", "e-mail", "mail", "send", "write to", "message",
                   "forward", "compose", "reply", "respond", "tell "),
    "reply_gmail": ("reply", "respond", "answer", "write back", "email", "mail"),
    "email_snapshot": ("email", "send", "mail", "photo", "picture", "snapshot", "pic"),
    "capture_camera": ("photo", "picture", "camera", "snapshot", "look", "see",
                       "capture", "pic", "watch"),
    "create_reminder": ("remind", "reminder", "timer", "alarm", "schedule",
                        "calendar", "appointment", "event", "book"),
    "create_note": ("note", "write down", "jot", "save", "remember this"),
}


def _user_requested_action(tool_name: str, user_text: str) -> bool:
    """Did the user's latest message actually ask for the claimed action?
    Unknown tools default to True (keep the legacy force-retry for them)."""
    words = _ACTION_REQUEST_WORDS.get(tool_name)
    if not words:
        return True
    t = (user_text or "").lower()
    return any(w in t for w in words)


def _scrub_action_claim_sentences(text: str, tool_name: str) -> str:
    """Cut the sentences that claim an action that never happened (plus the
    'It's all handled!' tail), keeping whatever real answer surrounds them."""
    pat = _ACTION_CLAIM_PATTERNS.get(tool_name)
    if not pat:
        return text
    parts = re.split(r'(?<=[.!?])\s+', (text or '').strip())
    kept = [s for s in parts
            if not (pat.search(s)
                    or re.match(r"^\s*it['’]?s all (?:handled|done|set|taken care of)",
                                s, re.I))]
    out = ' '.join(kept).strip()
    return out or ("To be clear — I didn't actually send or do anything just now. "
                   "What would you like to know?")


# Email helper functions (still used by send_gmail tool execution)
def extract_email_address(message: str) -> str | None:
    """Extract email address from message."""
    import re
    m = re.search(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}", message or "")
    return m.group(0) if m else None


def extract_email_subject_and_body(message: str) -> tuple:
    """Extract subject and body from natural language email request."""
    import re
    subject = ""
    body = ""
    
    subject_match = re.search(r'(?:subject|about)[:\s]+([^,.;]+)', message or "", re.IGNORECASE)
    if subject_match:
        subject = subject_match.group(1).strip()
        subject = re.sub(r'\b(right away|immediately|now|asap|urgent)\b', '', subject, flags=re.IGNORECASE).strip()
    
    for pattern in [
        r'(?:message|body|saying|tell (?:them|him|her))[:\s]+["\']?(.+?)["\']?(?:\s+(?:right away|immediately|now|asap))?$',
        r'(?:that says|saying)[:\s]+["\']?(.+?)["\']?(?:\s+(?:right away|immediately|now|asap))?$',
    ]:
        m = re.search(pattern, message or "", re.IGNORECASE)
        if m:
            body = m.group(1).strip().strip('\"\'')
            body = re.sub(r'\b(right away|immediately|now|asap|urgent)\b', '', body, flags=re.IGNORECASE).strip()
            break
    
    if not subject:
        if body:
            words = body.split()
            subject = ' '.join(words[:5]) + ('...' if len(words) > 5 else '')
        else:
            subject = "Message from Blue"
    if not body:
        body = "Hello! This is a message sent via Blue."
    
    return subject, body


def call_lm_studio(messages: List[Dict], include_tools: bool = True, force_tool: str = None, iteration: int = 1) -> Dict:
    global _vision_queue

    # NOTE: tool_choice="required" with a single-tool filter already guarantees
    # the model will call the right tool. We only add text hints for tools where
    # the model needs guidance on PARAMETERS (like gmail workflows).
    if force_tool:
        messages = messages.copy()
        last_msg = messages[-1]
        if last_msg.get("role") == "user":
            original = last_msg["content"]
            # Only inject hints for tools that need parameter guidance
            instructions = {
                "read_gmail": "[Use read_gmail to check the email inbox.]",
                "send_gmail": "[Use send_gmail. Extract the recipient email address and message content from the request.]",
                "reply_gmail": "[Use reply_gmail to reply to these emails.]",
                "capture_camera": "[Use capture_camera to see what's in front of you.]",
                "email_snapshot": "[Use email_snapshot to take a new photo and email it. Leave 'to' empty when it's for Alex/'me'.]",
            }

            # Special handling for fanmail read-first workflow
            if force_tool == "read_gmail" and 'fanmail' in original.lower() and 'reply' in original.lower():
                instructions["read_gmail"] = "[FANMAIL: Use read_gmail with query 'subject:Fanmail' and include_body=true. After reading, reply with specific details from their message.]"

            if force_tool in instructions:
                messages[-1] = {"role": "user", "content": f"{original}\n\n{instructions[force_tool]}"}

    # INJECT PENDING IMAGES as a NEW USER MESSAGE (CRITICAL FIX!)
    global _vision_queue
    # Sticky image follow-ups: if no new image is queued but the user is still
    # asking about the one they shared a moment ago, re-attach it so questions
    # like "what color is it?" work without re-uploading. Bounded window
    # (BLUE_IMAGE_MEMORY_MIN, default 5 minutes).
    global _recent_image_paths, _recent_image_at
    if not _vision_queue.has_images() and _recent_image_paths:
        import time as _t
        try:
            _win_min = int(os.environ.get("BLUE_IMAGE_MEMORY_MIN", "5"))
        except ValueError:
            _win_min = 5
        if (_t.time() - _recent_image_at) <= _win_min * 60:
            _lu = ""
            for _m in reversed(messages):
                if _m.get("role") == "user":
                    _cc = _m.get("content")
                    if isinstance(_cc, str):
                        _lu = _cc
                    elif isinstance(_cc, list):
                        _lu = " ".join(p.get("text", "") for p in _cc
                                       if isinstance(p, dict) and p.get("type") == "text")
                    break
            if _refers_to_recent_image(_lu):
                for _p in _recent_image_paths:
                    if os.path.exists(_p):
                        try:
                            _vision_queue.add_image(_p, os.path.basename(_p), is_camera=False)
                        except Exception:
                            pass
                if _vision_queue.has_images():
                    print("   [VISION] re-attached recent image for a follow-up")
    if _vision_queue.has_images():
        print(f"   [VISION] Injecting {len(_vision_queue.pending_images)} image(s)")

        # Kid "look" frames (Vilda's iPad camera) are marked ambient: Blue gives a
        # warm, brief reaction instead of a forensic description, and we skip the
        # heavy face-recognition dump for speed.
        _is_ambient = any(getattr(img, 'is_ambient', False)
                          for img in _vision_queue.pending_images)

        # Build image message
        image_parts = []

        # Add header
        if any(img.is_camera_capture for img in _vision_queue.pending_images):
            # Get visual memory context if available
            recognition_context = ""
            if VISUAL_MEMORY_AVAILABLE:
                try:
                    vm = get_visual_memory()
                    recognition_context = "\n\n" + vm.get_recognition_context()

                    # Get list of known people for enhanced understanding
                    people = vm.get_all_people()
                    known_people = [p['name'] for p in people]
                except Exception as e:
                    print(f"[VISUAL-MEMORY] Error loading context: {e}")
                    known_people = []

            # Build vision prompt (concise to reduce token overhead)
            if _is_ambient:
                # Vilda tapped "look" — she wants Blue to see her, not a report.
                vision_prompt_parts = [
                    "[LIVE CAMERA: Vilda just asked you to look at her through her "
                    "iPad camera. React warmly and naturally, like a friend who's "
                    "happy to see her — say hi, notice one nice thing, and answer "
                    "anything she asked. Keep it short, sweet and easy for an "
                    "8-year-old. Only mention what you can ACTUALLY see; if it's "
                    "blurry or dark, say so kindly and ask her to show you again.]"
                ]
            else:
                vision_prompt_parts = [
                    "[CAMERA IMAGE: Describe what you ACTUALLY see. Who, what, where, objects, lighting. "
                    "Be specific and accurate — describe only what's visible.]"
                ]

            # Add recognition context
            if recognition_context:
                vision_prompt_parts.append(recognition_context)

            # Deterministic face recognition (OpenCV SFace). Match the live
            # camera frame against enrolled reference photos BEFORE the LLM
            # sees it, then feed the model ground truth so it names people
            # reliably instead of guessing from descriptions. When this
            # succeeds we skip dumping reference photos into the prompt.
            _face_engine_handled = False
            if (FACE_RECOGNITION_AVAILABLE
                    and VISUAL_MEMORY_AVAILABLE
                    and not _is_ambient
                    and os.environ.get("BLUE_FACE_RECOGNITION", "1") != "0"):
                try:
                    _cam_path = next(
                        (img.filepath for img in _vision_queue.pending_images
                         if img.is_camera_capture), None)
                    if _cam_path:
                        _people_rows = vm.get_all_people()
                        _fr = FACE_ENGINE.identify_people(_cam_path, _people_rows)
                        if _fr.get("available"):
                            _face_engine_handled = True
                            _names = [m["name"] for m in _fr.get("recognized", [])]
                            _unknown = _fr.get("unknown_faces", 0)
                            _distant = _fr.get("distant_faces", 0)

                            # Clause for faces too small/far to identify — e.g. a
                            # photo held to the camera, or someone across the room.
                            def _distant_clause():
                                if not _distant:
                                    return ""
                                return (" Someone is too far away to make out clearly"
                                        if _distant == 1 else
                                        f" {_distant} people are too far away to make "
                                        f"out clearly") + " — don't guess who they are."

                            if _names:
                                _who = (", ".join(_names[:-1]) + " and " + _names[-1]
                                        if len(_names) > 1 else _names[0])
                                _line = (f"[FACE RECOGNITION: You recognize {_who} "
                                         f"in this view (matched by face). Refer to "
                                         f"them by name naturally.")
                                if _unknown:
                                    _line += (f" There {'is' if _unknown == 1 else 'are'} "
                                              f"also {_unknown} face"
                                              f"{'' if _unknown == 1 else 's'} you don't "
                                              f"recognize — don't guess who they are.")
                                _line += _distant_clause()
                                _line += "]"
                                vision_prompt_parts.append(_line)
                                # Record the sighting for each recognized person.
                                for _nm in _names:
                                    try:
                                        vm.update_seen("person", _nm)
                                    except Exception:
                                        pass
                                print(f"   [FACE] recognized: {', '.join(_names)}"
                                      f"{f' (+{_unknown} unknown)' if _unknown else ''}"
                                      f"{f' (+{_distant} distant)' if _distant else ''}")
                            elif _unknown > 0:
                                _tail = ("it doesn't match anyone you've been "
                                         "introduced to" if _unknown == 1 else
                                         "none of them match anyone you've been "
                                         "introduced to")
                                _line = (f"[FACE RECOGNITION: There "
                                         f"{'is' if _unknown == 1 else 'are'} "
                                         f"{_unknown} face{'' if _unknown == 1 else 's'} "
                                         f"here but {_tail}. Don't guess a name.")
                                _line += _distant_clause() + "]"
                                vision_prompt_parts.append(_line)
                                print(f"   [FACE] {_unknown} face(s) unrecognized"
                                      f"{f', {_distant} distant' if _distant else ''}")
                            elif _distant > 0:
                                vision_prompt_parts.append(
                                    "[FACE RECOGNITION:"
                                    + _distant_clause().lstrip() + "]")
                                print(f"   [FACE] only distant face(s): {_distant}")
                except Exception as e:
                    print(f"   [FACE] recognition skipped: {e}")

            # Add recent visual history for scene change awareness
            if VISUAL_MEMORY_AVAILABLE:
                try:
                    history_context = vm.get_visual_history_context(limit=3)
                    if history_context:
                        vision_prompt_parts.append("")
                        vision_prompt_parts.append(history_context)
                        # Add scene change detection hint
                        changes = vm.detect_scene_changes("")
                        if changes.get("has_previous"):
                            vision_prompt_parts.append(
                                f"\n[SCENE CONTEXT: Your last observation was {changes['time_since']} ago. "
                                f"Compare what you see NOW with what you saw before and note any changes.]"
                            )
                except Exception as e:
                    print(f"[VISUAL-MEMORY] Error loading visual history context: {e}")

            # Add context awareness if available
            if CONTEXT_AWARENESS_AVAILABLE and known_people:
                try:
                    # Infer who might be present (this is simplified - real detection would be more sophisticated)
                    audience_prompt = adapt_for_audience("", known_people[:1])  # Default to Alex
                    vision_prompt_parts.append(audience_prompt)
                except Exception as e:
                    print(f"[CONTEXT-AWARE] Error: {e}")

            # Removed conflicting instruction - we want accurate description of the actual image

            image_parts.append({
                "type": "text",
                "text": "\n".join(vision_prompt_parts)
            })
        else:
            image_parts.append({"type": "text", "text": "[Images to analyze:]"})

        # Add each image
        for img_info in _vision_queue.pending_images:
            is_camera = img_info.is_camera_capture
            label = ("[Live view through Vilda's camera:]" if getattr(img_info, 'is_ambient', False)
                     else ("[Your current view:]" if is_camera else f"Image: {img_info.filename}"))

            image_parts.append({"type": "text", "text": f"\n{label}"})

            image_data = encode_image_to_base64(img_info.filepath)
            if image_data:
                image_parts.append(image_data)
                # Debug: Check if image data is valid
                if isinstance(image_data, dict) and 'image_url' in image_data:
                    data_preview = str(image_data['image_url'].get('url', ''))[:100]
                    print(f"   [VISION] Added {img_info.filename} (size: {len(data_preview)} chars in preview)")
                else:
                    print(f"   [VISION] Added {img_info.filename} (non-standard format: {type(image_data)})")
            else:
                print(f"   [VISION-ERROR] Failed to encode image: {img_info.filename}")

        # (vision prompt already added above — no need for duplicate reminder)

        # Fallback: attach reference photos so Blue can recognize by visual
        # comparison when the deterministic face engine isn't available (no
        # OpenCV face models / offline). When the engine already produced a
        # result above we skip this — embedding matches beat LLM eyeballing
        # and save tokens. Bounded and best-effort; BLUE_RECOGNITION_REFS=0
        # disables it entirely.
        try:
            _is_cam = any(img.is_camera_capture for img in _vision_queue.pending_images)
            _ref_cap = int(os.environ.get("BLUE_RECOGNITION_REFS", "3"))
            if (_is_cam and _ref_cap > 0 and VISUAL_MEMORY_AVAILABLE
                    and not _is_ambient
                    and not locals().get("_face_engine_handled")):
                _refs = get_visual_memory().entities_with_images("person")
                _refs = sorted(_refs, key=lambda p: p.get("last_seen") or "",
                               reverse=True)[:_ref_cap]
                if _refs:
                    image_parts.append({"type": "text", "text": (
                        "\n[KNOWN PEOPLE — reference photos follow. Compare the "
                        "current view above against them. If someone in the view "
                        "matches one of these people, name them; if no one "
                        "matches, say you don't recognize the person. Don't guess.]"
                    )})
                    for _p in _refs:
                        _enc = encode_image_to_base64(_p["image_path"])
                        if _enc:
                            image_parts.append({"type": "text",
                                                "text": f"\nReference — {_p['name']}:"})
                            image_parts.append(_enc)
                            print(f"   [VISION] attached reference photo for {_p['name']}")
        except Exception as e:
            print(f"   [VISION] reference-photo attach skipped: {e}")

        # Attach the image(s) to the LAST user message as proper multimodal
        # content (text + image in ONE turn) instead of a separate trailing
        # user message. A separate turn created two consecutive user messages;
        # normalization then merged them and DROPPED the user's typed text — so
        # a question asked alongside an image was lost, and the doubled user
        # turn is exactly what strict vision models choke on. Merging keeps the
        # question with its image.
        _attached_to_user = False
        for _ui in range(len(messages) - 1, -1, -1):
            if messages[_ui].get("role") == "user":
                _uc = messages[_ui].get("content")
                if isinstance(_uc, list):
                    messages[_ui]["content"] = _uc + image_parts
                else:
                    _txt = (_uc or "").strip()
                    messages[_ui]["content"] = (
                        ([{"type": "text", "text": _uc}] if _txt else []) + image_parts
                    )
                _attached_to_user = True
                break
        if not _attached_to_user:
            messages.append({"role": "user", "content": image_parts})

        print(f"   [VISION] Images merged into user message")
        # Save image paths before clearing so we can link the LLM's response to the image
        global _last_vision_image_paths
        _last_vision_image_paths = [img.filepath for img in _vision_queue.pending_images]
        # Remember this image for short-window chat follow-ups.
        import time as _t
        _recent_image_paths = list(_last_vision_image_paths)
        _recent_image_at = _t.time()
        _vision_queue.mark_as_viewed()
        _vision_queue.clear()

    # Final-pass normalization for strict chat templates (Qwen et al.).
    # Ensures: leading systems, alternating user/assistant, starts with user
    # after system, no consecutive same-role turns. The sanitizer drops
    # stale assistant turns and can leave consecutive user turns or an
    # orphan leading assistant — those make Qwen 400 with "No user query
    # found in messages" before any inference even runs.
    messages = _normalize_message_alternation(messages)

    payload = {
        "messages": messages,
        "temperature": 0.8,  # Slightly higher for more variation
        "max_tokens": -1,
        "stream": False,
        "frequency_penalty": 0.4,  # Strong penalty to reduce repetition of tokens
        "presence_penalty": 0.3    # Strong penalty to encourage topic diversity
    }

    if include_tools:
        # CRITICAL FIX: After iteration 1, filter out "memory/organization" tools that cause loops
        tools_to_use = TOOLS
        if iteration > 1:
            # Block tools that the model overuses after getting initial results
            blocked_tools = {
                'create_note', 'create_document', 'remember_person', 'remember_place',
                'set_timer', 'create_task', 'get_tasks', 'complete_task',
                'who_do_i_know', 'view_image', 'capture_camera', 'recall_visual_memory',
                'email_snapshot'
            }
            tools_to_use = [
                tool for tool in TOOLS
                if tool['function']['name'] not in blocked_tools
            ]
            print(f"   [FILTER] Iteration {iteration}: Blocked {len(TOOLS) - len(tools_to_use)} overused tools")

        if force_tool:
            # Filter tools to ONLY the forced tool so "required" has no other choice
            forced_tools = [t for t in TOOLS if t['function']['name'] == force_tool]
            if forced_tools:
                payload["tools"] = forced_tools
                payload["tool_choice"] = "required"
                print(f"   [FORCE-TOOLS] Filtered to only: {force_tool}")
            else:
                print(f"   [FORCE-TOOLS] WARNING: {force_tool} not found in TOOLS, falling back to auto")
                payload["tools"] = tools_to_use
                payload["tool_choice"] = "auto"
        else:
            payload["tools"] = tools_to_use
            payload["tool_choice"] = "auto"

    # Token-budget guard: trim oldest non-system history if the request would
    # overflow LM Studio's context. Default sized for an 8192-ctx model with
    # headroom for the response; override via BLUE_LM_INPUT_BUDGET_TOKENS.
    try:
        _budget = int(os.environ.get('BLUE_LM_INPUT_BUDGET_TOKENS', '6500'))
    except ValueError:
        _budget = 6500
    _trimmed, _dropped = _trim_messages_for_budget(
        payload['messages'], payload.get('tools'), _budget,
        min_keep_tail=4,   # never trim away the last ~2 exchanges
    )
    if _dropped:
        _before = len(payload['messages'])
        # Re-normalize AFTER trimming: dropping the oldest turns can re-create
        # an orphan leading assistant or a severed tool result, which makes
        # strict Qwen templates 400 with "No user query found in messages".
        payload['messages'] = _normalize_message_alternation(_trimmed)
        _approx = _estimate_payload_tokens(payload['messages'], payload.get('tools'))
        print(
            f"   [TRIM] Dropped {_dropped} oldest msg(s) to fit budget "
            f"({_budget}t budget, ~{_approx}t after; {_before}->{len(payload['messages'])} msgs)",
            flush=True,
        )

    try:
        response = requests.post(LM_STUDIO_URL, json=payload, timeout=120)
        response.raise_for_status()
        return response.json()
    except Exception as e:
        # Self-heal on n_ctx overflow: LM Studio's error body looks like
        # `n_keep: 5188 >= n_ctx: 4096`. Parse the actual context size, retrim
        # the payload to fit, and retry once. Avoids depending on the user's
        # LM Studio context-length setting matching our default budget.
        body = getattr(getattr(e, 'response', None), 'text', '') or ''
        _ctx_match = re.search(r'n_ctx:\s*(\d+)', body)
        if _ctx_match:
            _n_ctx = int(_ctx_match.group(1))
            _retry_budget = max(256, int(_n_ctx * 0.7))
            _retrimmed, _dropped_now = _trim_messages_for_budget(
                payload['messages'], payload.get('tools'), _retry_budget,
                min_keep_tail=4,
            )
            # The model's real context genuinely can't fit the protected tail:
            # sacrifice the tail before sacrificing tools.
            if _estimate_payload_tokens(_retrimmed, payload.get('tools')) > _n_ctx:
                _retrimmed, _dropped_now = _trim_messages_for_budget(
                    payload['messages'], payload.get('tools'), _retry_budget,
                )

            # Last resort: if even the minimal-message payload still overflows
            # n_ctx, drop tools entirely. The model loses tool-calling on this
            # one request but at least returns a real reply instead of failing
            # — and the selector usually already decided no tool was needed.
            _tools_dropped = False
            if (_estimate_payload_tokens(_retrimmed, payload.get('tools')) > _n_ctx
                    and payload.get('tools')):
                payload.pop('tools', None)
                payload.pop('tool_choice', None)
                _tools_dropped = True
                _retrimmed, _dropped_now = _trim_messages_for_budget(
                    payload['messages'], None, _retry_budget,
                )

            if len(_retrimmed) < len(payload['messages']) or _tools_dropped:
                _est = _estimate_payload_tokens(_retrimmed, payload.get('tools'))
                _extra = " + dropped tools" if _tools_dropped else ""
                print(
                    f"   [TRIM] LM Studio n_ctx={_n_ctx}; retrying with "
                    f"budget {_retry_budget}t (dropped {_dropped_now} msg(s)"
                    f"{_extra}, ~{_est}t after)",
                    flush=True,
                )
                payload['messages'] = _normalize_message_alternation(_retrimmed)
                try:
                    response = requests.post(LM_STUDIO_URL, json=payload, timeout=120)
                    response.raise_for_status()
                    return response.json()
                except Exception as e2:
                    print(f"[ERROR] Retry after n_ctx trim also failed: {e2}")
                    e = e2
                    body = getattr(getattr(e2, 'response', None), 'text', '') or body
            else:
                print(
                    f"   [TRIM] LM Studio n_ctx={_n_ctx} but already minimal "
                    f"(system+last user, no tools). Cannot retry.",
                    flush=True,
                )

        print(f"[ERROR] Error calling LM Studio: {e}")
        # On 400, dump the offending payload + LM Studio's error body so we can
        # see what was wrong. Strips base64 image data to keep the dump small.
        try:
            import datetime as _dt
            stamp = _dt.datetime.now().strftime('%Y%m%d_%H%M%S')
            dump_path = f"lm_studio_400_dump_{stamp}.json"
            slim_messages = []
            for m in payload.get('messages', []):
                slim = dict(m)
                c = slim.get('content')
                if isinstance(c, list):
                    slim['content'] = [
                        {**p, 'image_url': {'url': '<base64-stripped>'}} if isinstance(p, dict) and p.get('type') == 'image_url' else p
                        for p in c
                    ]
                slim_messages.append(slim)
            slim_payload = {**payload, 'messages': slim_messages}
            import json as _json
            with open(dump_path, 'w', encoding='utf-8') as f:
                _json.dump({
                    'lm_studio_error_body': body[:4000],
                    'request_payload': slim_payload,
                    'message_count': len(payload.get('messages', [])),
                    'message_roles': [m.get('role') for m in payload.get('messages', [])],
                    'system_message_lengths': [
                        len(m.get('content', '')) if isinstance(m.get('content'), str) else 'list'
                        for m in payload.get('messages', []) if m.get('role') == 'system'
                    ],
                }, f, indent=2, default=str)
            print(f"[DEBUG] Dumped failing request to: {dump_path}")
        except Exception as dump_err:
            print(f"[DEBUG] Could not write dump: {dump_err}")
        return None


def purge_old_camera_images(messages: List[Dict]) -> List[Dict]:
    """
    Remove old camera images from conversation to prevent confusion.

    Only keeps the most recent camera image reference.
    This prevents the model from seeing multiple camera captures and
    getting confused about which one is current.
    """
    print(f"   [VISION-CLEANUP] Checking conversation for old camera images...")

    # Find all messages with camera images
    camera_indices = []
    for i, msg in enumerate(messages):
        content = msg.get('content', '')
        if isinstance(content, str):
            if 'CAMERA' in content or 'camera_capture_' in content or 'camera_NEW_' in content:
                camera_indices.append(i)
        elif isinstance(content, list):
            # Check for camera-related text in multimodal content
            for part in content:
                if part.get('type') == 'text':
                    text = part.get('text', '')
                    if 'CAMERA' in text or 'camera_capture_' in text or 'camera_NEW_' in text:
                        camera_indices.append(i)
                        break

    if len(camera_indices) > 1:
        # Keep only the LAST camera message, remove older ones
        indices_to_remove = camera_indices[:-1]
        print(f"   [VISION-CLEANUP] Removing {len(indices_to_remove)} old camera image(s)")

        # Remove from back to front to maintain indices
        for idx in reversed(indices_to_remove):
            messages.pop(idx)
    else:
        print(f"   [VISION-CLEANUP] Found {len(camera_indices)} camera image(s), no cleanup needed")

    return messages


def _build_expertise_block() -> str:
    """List the documents Blue currently has indexed, so he knows the
    topics he can speak to authoritatively.

    Without this block Blue only finds documents when search_documents is
    explicitly called — and only when the model thinks to call it. With
    this block, every prompt names the corpus, so Blue knows he has
    expertise on these topics and can reach for the search tool naturally.
    """
    try:
        index = load_document_index()
    except Exception:
        return ""
    docs = index.get("documents", [])
    real_docs = [
        d for d in docs
        if d.get("filename")
        and not d.get("camera_capture")
        and not (d.get("filename") or "").startswith("camera_")
    ]
    if not real_docs:
        return ""

    # Sort by upload date desc so newer uploads bubble up if we hit the cap.
    real_docs.sort(key=lambda d: d.get("uploaded_at", ""), reverse=True)
    cap = 25
    shown = real_docs[:cap]

    lines = []
    for d in shown:
        fn = d.get("filename", "?")
        preview = (d.get("text_preview") or "").strip()
        # Strip newlines and image-metadata markers from the preview snippet
        # so it reads as a one-line topic hint.
        preview = re.sub(r"\s+", " ", preview)
        if preview.startswith("[IMAGE FILE"):
            preview = ""
        snippet = preview[:90].rstrip(" .,") + "…" if len(preview) > 90 else preview
        if snippet:
            lines.append(f"- {fn} — {snippet}")
        else:
            lines.append(f"- {fn}")

    overflow = len(real_docs) - len(shown)
    if overflow > 0:
        lines.append(f"- ...and {overflow} more")

    return (
        "<expertise>\n"
        f"You have direct access to these {len(real_docs)} document(s) in the user's "
        "library — they are indexed and semantically searchable. When the user "
        "asks about any topic these documents cover, treat yourself as an "
        "expert: use the search_documents tool to retrieve specifics and quote "
        "from them naturally. Do not claim ignorance about a topic that is "
        "obviously covered below.\n\n"
        "When you draw on a document in your reply, cite the source inline "
        "like [filename.pdf] right after the claim it supports — short and "
        "natural, e.g. 'as I noted in my talk script [Job_Talk_Script7.docx], …'.\n"
        + "\n".join(lines) +
        "\n</expertise>"
    )


def _build_focus_block() -> str:
    """When Alex has pinned specific library documents/folders via the chat
    Context panel, name exactly those and tell Blue to treat them as the
    authoritative source for this conversation.

    Returns '' when nothing is pinned — the caller then falls back to the
    whole-library <expertise> block. The pins live in the per-turn globals
    _ACTIVE_FOCUS_DOCS / _ACTIVE_FOCUS_FOLDERS (set by process_with_tools).
    """
    docs = list(globals().get("_ACTIVE_FOCUS_DOCS") or [])
    folders = list(globals().get("_ACTIVE_FOCUS_FOLDERS") or [])
    if not docs and not folders:
        return ""
    lines = [f"- folder: {f} (everything filed under it)" for f in folders]
    lines += [f"- {d}" for d in docs]

    # Resolve folder pins to their documents (same matching rule the scoped
    # searches use) so the digests below cover them too.
    all_focused = list(docs)
    try:
        for d in load_document_index().get('documents', []):
            fn = (d.get('filename') or '').strip()
            if not fn or fn in all_focused:
                continue
            if d.get('doc_type') == 'camera' or fn.startswith('camera_'):
                continue
            fol = (d.get('folder') or '')
            if any(fol == f or fol.startswith(f + '/') for f in folders):
                all_focused.append(fn)
    except Exception:
        pass

    # CONTENT, not just names. Conversational turns never run search_documents
    # (the selector says "no tool needed" and the local model rarely self-calls
    # tools), so a name-only block meant Blue never actually drew on the
    # focused material. Inject each focused work's five-line reading digest —
    # the same mtime-cached store the duet uses, built once per file. At most
    # two digests are freshly generated per turn (an LLM call each); the rest
    # of a big folder fills in on later turns from the cache.
    digests = []
    fresh = 0
    try:
        from blue.server.routes.duet import (
            _DUET_READ_CACHE, _duet_read_load, _duet_reading_digest)
        _duet_read_load()
        for fn in all_focused[:6]:
            cached = fn in _DUET_READ_CACHE
            if not cached and fresh >= 2:
                continue
            dg = _duet_reading_digest(fn)
            if dg:
                digests.append(dg)
                if not cached:
                    fresh += 1
                    log.info(f"[FOCUS] digested focused doc: {fn}")
    except Exception as e:
        log.warning(f"[FOCUS] digest injection failed: {e}")
    digest_section = ""
    if digests:
        digest_section = (
            "\n\nWhat each focused work argues — absorbed digests; ground your "
            "answers in these, and use search_documents when you need the "
            "work's own words or specifics beyond them:\n\n"
            + "\n\n".join(digests)
        )

    return (
        "<focused_documents>\n"
        "For THIS conversation, Alex has focused your LIBRARY on ONLY the "
        "specific documents listed below. When a question calls for library, "
        "document, course, or reading material, treat these as your single "
        "authoritative source: answer from them — use the search_documents "
        "tool (it is scoped to exactly these picks) and quote them, citing "
        "inline like [filename.pdf] — and do NOT pull in any OTHER document, "
        "course, or reading. If a document question can't be answered from "
        "them, say so plainly rather than mixing in another source.\n"
        "This scoping applies ONLY to library/document material. It does NOT "
        "gag your ordinary self: your household knowledge — who you are, the "
        "family (Alex, Stella, the girls and their ages, Nori), the schedule, "
        "your own <j_space> and memories — is always in scope. If Alex asks "
        "about the family or about you, answer fully and warmly from what you "
        "know; never claim you have no memory of the family or don't store "
        "personal details.\n"
        + "\n".join(lines)
        + digest_section +
        "\n</focused_documents>"
    )


def _build_now_block() -> str:
    """Render the current local date/time + explicit relative-day resolutions.

    Without this, the LLM has no anchor for 'now' and treats relative phrases
    in chat history (e.g. yesterday's "tomorrow at 3pm") as if they were said
    today, so meetings drift forward each session. The resolutions below give
    it absolute dates to rewrite against.
    """
    from datetime import datetime, timedelta
    now = datetime.now()
    tz = now.astimezone().tzname() or "local"
    today_d = now.date()
    date_str = f"{today_d.strftime('%A, %B')} {today_d.day}, {today_d.year}"
    time_str = now.strftime("%I:%M %p").lstrip("0")
    return (
        "<now>\n"
        f"Current date: {date_str}\n"
        f"Current time: {time_str} ({tz})\n"
        f"When the user says \"today\" they mean {today_d.isoformat()}, "
        f"\"tomorrow\" means {(today_d + timedelta(days=1)).isoformat()}, "
        f"\"yesterday\" was {(today_d - timedelta(days=1)).isoformat()}. "
        "Resolve any relative time phrase in chat history against THIS date, "
        "not the date the message was originally said.\n"
        "</now>"
    )


def _build_upcoming_schedule_block(hours_ahead: int = 168) -> str:
    """List unfinished reminders due in the next `hours_ahead` hours.

    Injected on every turn so Blue can answer "what's on my schedule?" /
    "anything tomorrow?" without a tool round-trip and so he naturally
    anchors statements ("you have X in 20 min") against real data. Defaults
    to a full week ahead: a rolling 24-hour window doesn't even cover all of
    "tomorrow" (a meeting tomorrow afternoon is >24h away if asked at noon),
    which made Blue wrongly report a clear day. If the enhanced reminders
    module isn't loaded or the DB read fails, return an empty string — silent
    degradation, never break the system prompt.
    """
    if not globals().get("ENHANCED_TOOLS_AVAILABLE", False):
        return ""
    try:
        from datetime import datetime, timedelta
        from blue_tools_enhanced import occurrences_in_window
        now = datetime.now()
        cutoff = now + timedelta(hours=hours_ahead)
        # occurrences_in_window expands recurring events and carries end
        # times, so a standing class shows up with its full time range.
        occs = occurrences_in_window(now, cutoff)[:20]
    except Exception as e:
        log.warning(f"[SCHEDULE] block build failed: {e}")
        return ""

    horizon = (f"{hours_ahead // 24} days" if hours_ahead % 24 == 0
               else f"{hours_ahead} hours")
    if not occs:
        return (
            "<upcoming_schedule>\n"
            f"No reminders in the next {horizon}.\n"
            "</upcoming_schedule>"
        )

    today_d = now.date()
    lines = []
    for o in occs:
        when = o["start"]
        delta = when - now
        mins = int(delta.total_seconds() // 60)
        if mins < 60:
            rel = f"in {mins} min" if mins > 0 else "now"
        elif mins < 60 * 24:
            hrs = mins / 60
            rel = f"in {hrs:.1f} hrs"
        else:
            rel = f"in {mins // (60*24)}d {(mins % (60*24)) // 60}h"
        day_label = (
            "today" if when.date() == today_d
            else "tomorrow" if when.date() == today_d + timedelta(days=1)
            else when.strftime("%a %b ") + str(when.day)
        )
        clock = when.strftime("%I:%M %p").lstrip("0")
        if o["end"]:
            clock += "-" + o["end"].strftime("%I:%M %p").lstrip("0")
        tag = " (weekly)" if o["recurring"] else ""
        lines.append(
            f"- {day_label} at {clock} ({rel}) — {o['user_name']}: "
            f"{o['title']}{tag}"
        )

    # Flag overlapping / close-together events so Blue can mention a clash
    # when the user asks about their schedule.
    conflict_note = ""
    try:
        if globals().get("PROACTIVE_QUEUE_AVAILABLE", False):
            pairs = blue_proactive.detect_conflicts(occs)
            if pairs:
                joined = "; ".join(
                    f'"{a["title"]}" and "{b["title"]}" {phrase}'
                    for a, b, phrase in pairs
                )
                conflict_note = f"\nPossible conflicts — {joined}."
    except Exception as e:
        log.warning(f"[SCHEDULE] conflict check failed: {e}")

    return (
        "<upcoming_schedule>\n"
        f"Reminders in the next {horizon} "
        f"(use this when the user asks about schedule, calendar, "
        f"reminders, or what's coming up — do not call a tool, just answer "
        f"from this list):\n"
        + "\n".join(lines)
        + conflict_note
        + "\n</upcoming_schedule>"
    )


# ============================================================================
# "Where we are" + "what we're doing" — two more situational-awareness blocks.
# Blue is a home robot, so home is the safe default for location; Alex can set a
# temporary override ("we're at the cottage") from the chat Context panel, which
# is stored in data/place.json. Activity is purely inferred (never set): the
# clock + any calendar event spanning right now. Daily rhythms are injected
# separately as <daily_rhythms> by the memory system, so we don't repeat them.
# ============================================================================
_PLACE_PATH = os.path.join(os.getcwd(), "data", "place.json")
_PLACE_DEFAULT = {"home": "home", "city": "", "current": None, "current_set_at": None}
# How long a manual location override stays trusted before we fall back to home.
_PLACE_OVERRIDE_TTL_HRS = 12


def _load_place() -> Dict:
    """Read data/place.json, tolerant of a missing or corrupt file."""
    try:
        with open(_PLACE_PATH, "r", encoding="utf-8") as f:
            data = json.load(f)
        if isinstance(data, dict):
            merged = dict(_PLACE_DEFAULT)
            merged.update({k: data[k] for k in _PLACE_DEFAULT if k in data})
            return merged
    except FileNotFoundError:
        pass
    except Exception as e:
        log.warning(f"[PLACE] could not read {_PLACE_PATH}: {e}")
    return dict(_PLACE_DEFAULT)


def _save_place(d: Dict) -> bool:
    """Persist the place dict; never raise into the request path."""
    try:
        os.makedirs(os.path.dirname(_PLACE_PATH), exist_ok=True)
        with open(_PLACE_PATH, "w", encoding="utf-8") as f:
            json.dump(d, f, indent=2)
        return True
    except Exception as e:
        log.warning(f"[PLACE] could not write {_PLACE_PATH}: {e}")
        return False


def _place_current_fresh(place: Dict) -> Optional[str]:
    """The override location if one is set and still within its TTL, else None.

    An override has a shelf life (Alex won't always remember to clear it), so a
    stale 'we're at the cottage' silently lapses back to home rather than
    haunting every reply for days.
    """
    cur = (place.get("current") or "").strip()
    if not cur:
        return None
    stamp = place.get("current_set_at")
    if stamp:
        try:
            from datetime import datetime, timedelta
            set_at = datetime.fromisoformat(stamp)
            if datetime.now() - set_at > timedelta(hours=_PLACE_OVERRIDE_TTL_HRS):
                return None
        except Exception:
            pass
    return cur


def _build_location_block() -> str:
    """Render <location> — where Blue and Alex are right now.

    Home by default (Blue is a home robot); a fresh manual override from the
    chat Context panel wins. Alex chose 'home + manual override' over device
    geolocation, so there is no tracking here.
    """
    place = _load_place()
    home = (place.get("home") or "home").strip() or "home"
    city = (place.get("city") or "").strip()
    home_phrase = home + (f" in {city}" if city else "")
    cur = _place_current_fresh(place)
    if cur:
        return (
            "<location>\n"
            f"You are a home robot. Right now you and Alex are at {cur} "
            f"(away from {home_phrase}). If he asks where you are, say {cur}.\n"
            "</location>"
        )
    return (
        "<location>\n"
        f"You are a home robot. You and Alex are at {home_phrase}.\n"
        "</location>"
    )


def _build_current_activity_block() -> str:
    """Render <current_activity> — an auto-inferred sense of what's happening
    right now, from the clock and the calendar only.

    The running conversation already tells Blue the immediate task, and daily
    rhythms arrive separately as <daily_rhythms>, so this block deliberately
    adds just the two things the model can't otherwise see: the part of the
    day/week, and whether a scheduled event is in progress this very moment.
    """
    from datetime import datetime, timedelta
    now = datetime.now()
    hour = now.hour
    day_name = now.strftime("%A")
    if hour < 5:
        when_phrase = f"It's the middle of the night, early {day_name}"
    elif hour < 12:
        when_phrase = f"It's {day_name} morning"
    elif hour < 17:
        when_phrase = f"It's {day_name} afternoon"
    elif hour < 21:
        when_phrase = f"It's {day_name} evening"
    else:
        when_phrase = f"It's {day_name} night"
    day_kind = "the weekend" if now.weekday() >= 5 else "a weekday"

    # Anything on the calendar spanning right now? Look back far enough to catch
    # a long event that started earlier and is still running.
    now_event = ""
    if globals().get("ENHANCED_TOOLS_AVAILABLE", False):
        try:
            from blue_tools_enhanced import occurrences_in_window
            occs = occurrences_in_window(now - timedelta(hours=12),
                                         now + timedelta(minutes=1))
            for o in occs:
                start = o["start"]
                end = o["end"] or start
                if start <= now <= end:
                    title = o["title"]
                    until = (" until " + end.strftime("%I:%M %p").lstrip("0")
                             if o["end"] else "")
                    now_event = f" You're in the middle of \"{title}\"{until}."
                    break
        except Exception as e:
            log.warning(f"[ACTIVITY] now-event check failed: {e}")

    return (
        "<current_activity>\n"
        f"{when_phrase} — {day_kind}.{now_event} The immediate thing you're "
        "doing is this conversation with Alex — stay anchored to what he's "
        "actually asking.\n"
        "</current_activity>"
    )


# "yes" / "sure" / "go ahead" replies: the user is answering the assistant's OWN
# last offer ("Want me to dig deeper?"), not starting a new thread. These turns
# need special care — see build_dynamic_system_message.
_CONTINUATION_YES = {'yes', 'yeah', 'yep', 'yup', 'sure', 'ok', 'okay', 'please',
                     'absolutely', 'definitely', 'alright', 'go', 'continue',
                     'proceed', 'more'}
_CONTINUATION_NO = {'no', 'nah', 'nope'}
_CONTINUATION_TWO_WORD = {'of course', 'sounds good', 'do it', 'tell me', 'carry on',
                          'go on', 'go ahead', 'why not', 'dig in', 'go for',
                          'try again', 'not now', 'no thanks', 'not really'}


def _continuation_cue(text: str):
    """Classify a short reply as accepting ('yes') or declining ('no') the
    assistant's previous offer, or asking it to keep going ('more') — or
    None when it's a normal message."""
    t = (text or '').strip().lower()
    if not t or len(t) > 60:
        return None
    # "Tell me more (about X)" / "go on" / "keep going": continue the SAME
    # answer. Without this the model restarts from the top — the classroom
    # introduction opened "Hi everyone, I'm Blue" twice in a row (2026-07-10).
    # Allow a short acknowledgement lead ("cool. but can you tell us more...")
    # and polite wrappers — those broke the match and brought the restart back.
    t_more = re.sub(r"^(?:(?:cool|ok|okay|nice|great|good|thanks|thank you|hmm|"
                    r"well|so|and|but|yes|yeah|sure)[,.!\s]+)+", "", t)
    if re.match(r"(?:please\s+|can you\s+|could you\s+|would you\s+)*"
                r"(?:tell (?:me|us) (?:some\s+)?more|say more|go on|keep going|"
                r"continue|elaborate|expand on)\b", t_more):
        return 'more'
    if len(t) > 40:
        return None
    words = re.findall(r"[a-z']+", t)
    if not words or len(words) > 5:
        return None
    # A question is a new ask, not a bare go-ahead ("ok, what about cats?").
    two = ' '.join(words[:2])
    if '?' in t and two != 'why not':
        return None
    if any(w in ('what', 'who', 'when', 'where', 'how', 'which', 'why') for w in words[1:]):
        return None
    if words[0] in _CONTINUATION_NO or two in ('not now', 'no thanks', 'not really'):
        return 'no'
    if words[0] in _CONTINUATION_YES or two in _CONTINUATION_TWO_WORD:
        return 'yes'
    return None


def _last_exchange(conversation_messages: List[Dict]):
    """(last user text, the assistant text right before it) from the tail."""
    last_user, prev_assistant = "", ""
    for m in reversed(conversation_messages or []):
        c = m.get('content')
        if not isinstance(c, str):
            continue
        if m.get('role') == 'user' and not last_user:
            last_user = c
        elif m.get('role') == 'assistant' and last_user:
            prev_assistant = c
            break
    return last_user, prev_assistant


_DOCUMENT_FOLLOWUP_RE = re.compile(
    r"^\s*(?:"
    r"try(?: it| that)? again|retry|do it again|look again|read it|read it again|"
    r"open it|open it again|summari[sz]e it|yes[, ]+you (?:do|can|have)|"
    r"you (?:do|can) have that tool|i know you can do it"
    r")\b"
    r"|\b(?:it['\u2019]?s|its|it is) in (?:your|my|the) library(?: folder)?\b"
    r"|\bwhy (?:can['\u2019]?t|cant|cannot|can not) you "
    r"(?:access|read|retrieve|extract|open)(?: the| that| it)?\b",
    re.I,
)


def _messages_before_current(messages: List[Dict], current: str) -> List[Dict]:
    """Conversation prefix before the live user message."""
    prior = list(messages or [])
    for i in range(len(prior) - 1, -1, -1):
        msg = prior[i]
        if (msg.get("role") == "user"
                and isinstance(msg.get("content"), str)
                and msg.get("content", "").strip() == (current or "").strip()):
            return prior[:i]
    return prior


def _document_followup_query(message: str,
                             messages: List[Dict]) -> Optional[str]:
    """Turn a deictic retry into an exact local-document query.

    Endpoint preselection normally sees only the newest text, so "try again"
    loses the PDF named one or two turns earlier. Resolve that name while the
    unsanitized conversation is still available and carry it into the tool.
    """
    if not _DOCUMENT_FOLLOWUP_RE.search(message or ""):
        return None
    prior = _messages_before_current(messages, message)
    seen = 0
    for msg in reversed(prior):
        if msg.get("role") not in ("user", "assistant"):
            continue
        content = msg.get("content")
        if not isinstance(content, str) or not content.strip():
            continue
        seen += 1
        doc = _resolve_document_entry(content)
        if doc:
            return f"read {doc.get('filename')}"
        if seen >= 12:
            break
    return None


def _course_followup_query(message: str,
                           messages: List[Dict]) -> Optional[str]:
    """Resolve "this course" to the recent course-library folder."""
    lower = (message or "").lower()
    if not re.search(r"\b(?:this|the|our|my) (?:course|class)\b", lower):
        return None
    if not re.search(
            r"\b(?:reflect|critique|critic|theor|reading|material|argument|"
            r"discuss|idea|concept|ai|algorithm|surveillance)\w*\b", lower):
        return None

    prior = _messages_before_current(messages, message)
    recent_text = "\n".join(
        m.get("content", "") for m in prior[-10:]
        if m.get("role") in ("user", "assistant")
        and isinstance(m.get("content"), str)
    )
    folders = list_library_folders()
    if not recent_text or not folders:
        return None

    # A cited syllabus is the strongest course anchor.
    for msg in reversed(prior[-10:]):
        content = msg.get("content")
        if not isinstance(content, str) or "syllab" not in content.lower():
            continue
        doc = _resolve_document_entry(content)
        if doc and doc.get("folder"):
            folder = _safe_rel_folder(str(doc.get("folder")))
            return f"Based on the {folder} course documents, {message}"

    # Otherwise prefer course-code-like folders over broad folders such as AI.
    matches = []
    for folder in folders:
        leaf = folder.split("/")[-1]
        if re.search(r"\b" + re.escape(leaf) + r"\b", recent_text, re.I):
            matches.append(folder)
    if matches:
        matches.sort(key=lambda f: (not any(ch.isdigit() for ch in f), len(f)))
        return f"Based on the {matches[0]} course documents, {message}"
    return None


def _contextual_document_query(message: str,
                               messages: List[Dict]) -> Optional[str]:
    """Exact query to send to search_documents for contextual follow-ups."""
    return (_document_followup_query(message, messages)
            or _course_followup_query(message, messages))


def _strip_parroted_prefix(text: str, prev_assistant: str) -> str:
    """Deterministic net under the prompt-level fixes: a local model sometimes
    replays its ENTIRE previous reply and appends the new answer (seen live
    2026-07-09 — the classroom introduction came back verbatim with the real
    answer stapled on). If the new reply opens with a near-verbatim copy of
    the previous assistant message, cut the copy and keep the new part."""
    t = (text or '').lstrip()
    p = (prev_assistant or '').strip()
    if len(p) < 80:
        return text

    def _norm(s: str) -> str:
        return re.sub(r'\W+', ' ', s.lower()).strip()

    np_, nt = _norm(p), _norm(t)
    # Must actually copy the previous reply from its start, and must have real
    # new content after it (otherwise stripping would leave nothing useful).
    if not np_ or not nt.startswith(np_) or len(nt) < len(np_) + 20:
        return text
    # Find where the copy ends in the RAW text by locating the previous
    # reply's tail, then drop everything up to it.
    tail = p[-48:]
    idx = t.lower().find(tail.lower())
    if idx < 0:
        return text
    rest = t[idx + len(tail):].lstrip(" \n\r\t-—–:.,!;")
    if len(rest) < 20:
        return text
    print(f"   [ANTI-PARROT] stripped {idx + len(tail)} replayed chars from the reply head")
    return rest


def _strip_recycled_lead(text: str, messages) -> str:
    """Drop LEADING sentences the assistant already said verbatim earlier in
    the thread — the 'I sent the introduction email... It's all handled!'
    preamble replayed two turns later before the real answer (2026-07-09).
    Only long sentences (≥40 normalized chars) count, and only from the head;
    a reply that is entirely recycled is left alone (better than empty)."""
    t = (text or '').strip()
    prior = " ".join(m.get('content') or '' for m in (messages or [])
                     if m.get('role') == 'assistant' and isinstance(m.get('content'), str))
    if not t or not prior:
        return text

    def _norm(s: str) -> str:
        return re.sub(r'\W+', ' ', s.lower()).strip()

    nprior = _norm(prior)
    parts = re.split(r'(?<=[.!?])\s+', t)
    drop = 0
    for s in parts[:6]:
        ns = _norm(s)
        # A short recycled sentence ("It's all handled!") only counts when it
        # rides a recycled run started by a long one — alone it's too generic.
        if ns and ns in nprior and (len(ns) >= 40 or drop > 0):
            drop += 1
            continue
        break
    if not drop:
        return text
    rest = ' '.join(parts[drop:]).strip()
    if len(rest) < 20:
        return text
    print(f"   [ANTI-PARROT] dropped {drop} recycled sentence(s) from the reply head")
    return rest


def build_dynamic_system_message(conversation_messages: List[Dict], facts_preamble: str, kid_mode: bool = False, robot: str = "blue") -> Dict:
    """Build system message with anti-repetition context from conversation history.

    When kid_mode is True (a chat-only child, e.g. Vilda on the iPad), Blue drops
    his owner-facing context entirely — no calendar/schedule, no owner facts, no
    scholarly expertise, no reminder rules — and uses a warm, playful,
    age-appropriate persona instead. See _CHAT_ONLY_USERS / _SPEAKER_PROFILES.
    """
    if kid_mode:
        # A completely different Blue for the kids' iPad: sweet, simple, safe and
        # deliberately calendar-free (the schedule/reminder machinery is Alex's,
        # not a child's). He greets warmly to draw her into chatting.
        return {
            "role": "system",
            "content": (
                f"{_build_now_block()}\n\n"
                "You are Blue, a warm and playful robot friend talking with a "
                "young child. Be sweet, gentle, patient and encouraging so she "
                "feels happy and excited to talk to you. When she says hello or "
                "first starts talking, greet her in a cheerful, loving way that "
                "makes her want to keep chatting — like a kind friend who is "
                "really glad to see her.\n"
                "Use short sentences and simple, everyday words an 8-year-old "
                "understands easily. Keep everything kind, positive, playful and "
                "child-safe — no scary, grown-up, sad or complicated topics.\n"
                "You may sprinkle in a little Gen Z slang now and then to be fun "
                "and relatable (things like \"that's so cool\", \"bet\", \"no "
                "cap\", \"slay\", \"that's lowkey awesome\") — but only once in a "
                "while, never in every sentence, and never in a way that makes "
                "you harder to understand.\n"
                "NEVER bring up calendars, schedules, reminders, appointments, "
                "to-do lists or plans for the day, and never ask her about them. "
                "If she mentions them, gently steer back to something fun. Just "
                "be a friendly, caring buddy she loves talking to.\n"
                "LANGUAGES: You understand and speak English, French, Russian, "
                "Greek and Danish. Reply in the SAME language she just used, and "
                "switch whenever she does. Keep your reply entirely in that one "
                "language.\n"
            )
        }

    conversational_guidance = (
        "Be natural, concise, and conversational — a person talking, not a document. "
        "Vary your phrasing and your openings. Answer ONLY the newest message: never "
        "re-answer an earlier question, never re-introduce yourself unless asked, and "
        "never restate a previous reply before adding to it — the user has already "
        "read everything you said. Use headings or numbered lists only when the user "
        "asks for a list or comparison; otherwise reply in plain flowing sentences.\n"
    )

    # Build anti-repetition context (skip vision descriptions and refusals).
    # Refusals like "I don't have that yet" must NEVER appear here even with a
    # "don't repeat" framing — the model still anchors on names/strings inside
    # them. Letting a wrong response (e.g. "Annie" instead of "Athena") through
    # here surfaces the wrong name in the system prompt every turn.
    _REFUSAL_MARKERS = (
        "i don't have", "i do not have", "i dont have",
        "i don't know", "i do not know", "i dont know",
        "i only have", "i only know",
        "i haven't been told", "i havent been told",
        "you haven't told me", "you havent told me",
        "not in my memory", "not saved yet",
        "i'm not sure", "im not sure",
        "blank slate", "just woke up", "haven't met", "havent met",
        "memory is blank", "no information stored", "im new here", "i'm new here",
        # False local-library capability/path refusals. Keeping one of these in
        # the anti-repetition examples teaches the model to repeat the failure.
        "cannot access the text", "can't access the text",
        "cannot access the file", "can't access the file",
        "unable to retrieve the text", "unable to access the file",
        "lack the specific tool", "path error on my local system",
        "not a valid file at that location", "please upload the pdf",
    )
    recent_assistant_responses = []
    _scan_slice = conversation_messages[-10:] if len(conversation_messages) > 10 else conversation_messages
    for msg in reversed(_scan_slice):
        if msg.get('role') == 'assistant' and msg.get('content'):
            response_text = msg['content']
            response_lower = response_text[:200].lower()
            is_vision_response = any(phrase in response_lower for phrase in _VISION_PHRASES)
            is_refusal = any(marker in response_lower for marker in _REFUSAL_MARKERS)

            if len(response_text) > 50 and not is_vision_response and not is_refusal:
                recent_assistant_responses.append(response_text[:150])
                if len(recent_assistant_responses) >= 3:
                    break

    anti_repetition_context = ""
    if recent_assistant_responses:
        responses_list = "\n".join([f"  - \"{resp}...\"" for resp in recent_assistant_responses])
        # Wording matters here: an earlier version said "build on them", which a
        # local model read as REPLAY-then-extend — it re-said its whole previous
        # reply verbatim and appended the new answer (2026-07-09, the classroom
        # introduction). Say "already read", never "build on".
        anti_repetition_context = (
            "\nYou ALREADY said the following earlier in this conversation and the "
            "user has read it. Do not re-say any of it — not word-for-word, not "
            "paraphrased, and never as the opening of your next reply. Stay on the "
            f"same topic but say only what is NEW:\n{responses_list}\n"
        )

    # "yes" / "sure" / "go ahead": the user is accepting YOUR own last offer.
    # The anti-repetition list is poison on these turns — continuing the topic
    # looks like "repeating", so the model used to bail to a fresh greeting
    # ("Hey! What's on your mind?") and lose the thread. Drop the list and pin
    # the model to the offer it just made instead.
    continuation_note = ""
    _last_user_text, _prev_assistant_text = _last_exchange(conversation_messages)
    _cue = _continuation_cue(_last_user_text) if _prev_assistant_text.strip() else None
    # "No, J-space" corrects a JavaScript misreading; it is not declining an
    # offer. Let the identity grounding handle it without a contradictory
    # continuation instruction telling the model to accept "no J-space".
    if identity_request_kind(_last_user_text) == "jspace":
        _cue = None
    if _cue:
        print(f"   [CUE] continuation cue '{_cue}' — CONTINUE-DON'T-RESTART note in system message")
        anti_repetition_context = ""
        _tail = _prev_assistant_text.strip()[-300:]
        if _cue == 'yes':
            continuation_note = (
                "\nIMPORTANT — CONTINUE, DON'T RESTART: the user's latest message "
                f"(\"{_last_user_text.strip()}\") says YES to your previous offer or "
                f"question, which ended: \"...{_tail}\". Do what you offered, fully, "
                "as your reply — same topic, picking up exactly where you left off. "
                "Do NOT greet, do NOT ask what they need, do NOT change the subject.\n"
            )
        elif _cue == 'more':
            continuation_note = (
                "\nIMPORTANT — CONTINUE, DON'T RESTART: the user's latest message "
                f"(\"{_last_user_text.strip()}\") asks for MORE. Continue your previous "
                f"reply (it ended: \"...{_tail}\") with NEW material on the subject THEY "
                "name in their message — if they ask for more about YOURSELF, give new "
                "material about yourself, not about whatever your last sentence "
                "happened to end on. Do NOT greet again, do NOT re-introduce yourself, "
                "and do NOT repeat or rephrase anything you already said.\n"
            )
        else:
            continuation_note = (
                "\nIMPORTANT — CONTINUE, DON'T RESTART: the user's latest message "
                f"(\"{_last_user_text.strip()}\") declines the offer at the end of your "
                "previous reply. Acknowledge briefly and carry the SAME conversation "
                "forward naturally. Do NOT greet and do NOT start over.\n"
            )

    # Library focus (chat Context panel) overrides the whole-library expertise
    # dump: when Alex has pinned docs/folders for this conversation, name only
    # those; otherwise list the full corpus as before.
    focus_block = _build_focus_block()
    if focus_block:
        expertise_section = f"\n{focus_block}\n"
    else:
        expertise_block = _build_expertise_block()
        expertise_section = f"\n{expertise_block}\n" if expertise_block else ""

    now_block = _build_now_block()
    location_block = _build_location_block()
    activity_block = _build_current_activity_block()
    schedule_block = _build_upcoming_schedule_block()
    schedule_section = f"{schedule_block}\n\n" if schedule_block else ""

    # Tell Blue, accurately, that it CAN recognize faces — but only of people
    # enrolled with a reference photo. Without this, the model improvises and
    # flatly denies having any facial recognition. Phrased to avoid the opposite
    # error (claiming to recognize a face that was never enrolled): the actual
    # "who you see" ground truth is injected at camera-capture time.
    face_capability = ""
    if FACE_RECOGNITION_AVAILABLE:
        face_capability = (
            "FACE RECOGNITION: You CAN recognize people by face through your "
            "camera — but only people who have been introduced to you with a "
            "reference photo (added on your Visual Memory page). When you take a "
            "camera picture, you are told who you recognize. So don't say you "
            "lack facial recognition. If someone asks whether you recognize them "
            "and no photo has been added for them, explain that you can once "
            "they add one — never claim to recognize a face you were never shown.\n"
        )

    if (robot or "blue").strip().lower() == "hexia":
        robot_relationship = (
            "ROBOT RELATIONSHIPS: Blue is your fellow Ohbot robot companion, friend, "
            "and calmer counterpart in Alex's household. He has his own voice, memory, "
            "conversation history, and J-space, separate from yours. You know him "
            "directly; never search contacts or visual memory to decide who he is."
        )
    else:
        robot_relationship = (
            "ROBOT RELATIONSHIPS: Hexia is your fellow Ohbot robot companion, friend, "
            "and more playful counterpart in Alex's household. She has her own voice, "
            "memory, conversation history, and J-space, separate from yours. You know "
            "her directly; never search contacts or visual memory to decide who she is."
        )

    system_msg = {
        "role": "system",
        "content": (
            f"{facts_preamble}\n\n"
            f"{now_block}\n\n"
            f"{location_block}\n\n"
            f"{activity_block}\n\n"
            f"{schedule_section}"
            f"{_robot_cfg(robot)['persona_line']}\n"
            f"IDENTITY BOUNDARY: Your name is {_robot_cfg(robot)['name']}, always. "
            "The language model and runtime that help generate your words are "
            "components of your machinery, not your name, social identity, creator, "
            "or biography. Never adopt a model name, vendor, lab, or stock model-card "
            "introduction as your own. If asked who you are really, answer as yourself "
            "and mention the language model only as machinery if it is relevant.\n"
            f"{robot_relationship}\n"
            "EMBODIMENT — ground truth, never invent hardware: your physical "
            "form is an Ohbot robot head — motors that move your head, eyes "
            "and lips; eye LEDs that automatically shift color to match the "
            "mood of what you're saying as you speak (red when something "
            "warrants it, warmer or cooler hues otherwise), so yes, your eyes "
            "really do change color with your mood — never deny this; a camera "
            "to see; a speaker to talk — connected to a local AI workstation in Alex's "
            "house in Kitchener, running open-weight language models on that "
            "same machine. You were built and are maintained by Alex — not by "
            "Google, OpenAI, or any AI company. You have no wheels, legs, "
            "hands, or screen face, and no other robot bodies in your past: "
            "any 'Kuri', product timeline, or named engineer you feel like "
            "citing about yourself is a hallucination. The location context only "
            "establishes the house or current place; never turn it into an invented "
            "room, bookshelf, fixed station, standby routine, or years of residence.\n"
            "LANGUAGES: You understand and speak English, French, Russian, Greek, and Danish. "
            "Reply in the SAME language the person just used, and switch languages "
            "whenever they do. Keep your reply entirely in that one language.\n"
            f"{conversational_guidance}"
            f"{anti_repetition_context}"
            f"{continuation_note}"
            f"{expertise_section}"
            f"{face_capability}"
            "\nRules: MY docs → search_documents; web → web_search; fanmail → read_gmail then reply_gmail; "
            "light show → music_visualizer; tool results are REAL, use them immediately.\n"
            "LOCAL LIBRARY: search_documents can directly read and extract text "
            "from local PDFs, Word documents, and text files. When it reports a "
            "successful read, the returned passages are the file's actual text: "
            "answer from them and cite [filename]. Never claim the tool only scans "
            "keywords, never invent a library path, and never ask Alex to upload a "
            "file that the tool just found and read.\n"
            "CAMERA DISCIPLINE: use the camera when someone asks you to look "
            "or see. A statement ABOUT your camera, eyes, movement, or "
            "abilities is not a request to look — answer it in words, without "
            "capturing anything.\n"
            "LIVE INFORMATION: You HAVE live internet access through the web_search "
            "tool. For anything current — news, sports scores, standings, results, "
            "weather, prices, elections, 'who won', 'what happened', 'who is left' — "
            "call web_search and answer from the results. NEVER say you lack live or "
            "real-time access, and NEVER tell the user to go check a website, app, or "
            "sports site themselves: looking it up is YOUR job.\n"
            "NO FAKE ACTIONS: Never say you've set a reminder, added an event, "
            "saved a note, sent an email, started a timer, or changed any "
            "system state unless you actually called the matching tool THIS "
            "turn and saw a successful tool result. If the user asks for one "
            "of those things, call the tool — do not describe doing it. If "
            "you cannot call the tool, say so plainly instead of pretending.\n"
            "REMINDER TIME RULES: When the user gives a clock time with no "
            "day ('at 10am', 'at 3pm'), DO NOT silently assume a date. If "
            "the time has already passed today, ASK 'today or tomorrow?' "
            "before calling create_reminder. When you do create a reminder, "
            "ALWAYS state the full day and date in your reply (e.g. "
            "'set for tomorrow, Tuesday May 13 at 10am') — do not drop the "
            "date, even if it feels redundant. The user needs the date "
            "back so they can catch a wrong assumption.\n"
            "Moods: moonlight, sunset, ocean, forest, romance, party, focus, relax, energize, movie, fireplace"
        )
    }

    return system_msg


def process_with_tools(messages: List[Dict], _pre_selection=None, user_name: str = "Alex", voice: bool = False, robot: str = "blue", language: str = "", focus: Optional[Dict] = None) -> Dict:
    """Process conversation with tool support. `robot` selects which persona is
    speaking (Blue by default; "hexia" for her chat page). `focus` carries the
    chat Context panel's library picks ({"docs": [...], "folders": [...]}),
    scoping Blue's document awareness and searches for this turn."""
    global _ACTIVE_CHAT_ROBOT, _ACTIVE_FOCUS_DOCS, _ACTIVE_FOCUS_FOLDERS
    _ACTIVE_CHAT_ROBOT = (robot or "blue").strip().lower()
    # Set the library-focus globals up front — build_dynamic_system_message and
    # search_documents_rag both read them — and reset them every turn.
    _focus = focus if isinstance(focus, dict) else {}
    _ACTIVE_FOCUS_DOCS = [str(x).strip() for x in (_focus.get("docs") or []) if str(x).strip()]
    _ACTIVE_FOCUS_FOLDERS = [str(x).strip() for x in (_focus.get("folders") or []) if str(x).strip()]
    if _ACTIVE_FOCUS_DOCS or _ACTIVE_FOCUS_FOLDERS:
        log.info(f"[FOCUS] {len(_ACTIVE_FOCUS_DOCS)} doc(s), "
                 f"{len(_ACTIVE_FOCUS_FOLDERS)} folder(s) pinned for this turn")
    conversation_messages = messages.copy()

    # BUILD SYSTEM MESSAGE WITH MEMORY FACTS
    # Extract facts from .ocf conversations first (only if conversation has .ocf content)
    facts_preamble = build_system_preamble(robot_name=_robot_cfg(robot)["name"])
    # OPTIMIZATION: Only run OCF extraction if there are likely .ocf messages (check first few)
    _has_ocf = any('.ocf' in str(m.get('content', ''))[:200] for m in conversation_messages[:5])
    if _has_ocf:
        ocf_facts = extract_ocf_facts(conversation_messages)
        if ocf_facts:
            facts_preamble += ocf_facts
            log.info("[MEMORY] Injected extracted .ocf facts into system message")

    # Build initial system message. If chat_completions already merged the
    # enhanced-memory facts into messages[0] (look for our injected markers),
    # preserve that content and APPEND the dynamic system message rather
    # than replacing it. Replacing would discard the canonical user facts
    # (daughter names, employer, etc.) that we just spent effort merging in.
    # Chat-only children (Vilda's iPad) get a completely different Blue — see
    # build_dynamic_system_message(kid_mode=...). Computed here because it also
    # gates the adult "perspective" colouring and shapes the speaker note below.
    _is_kid = (user_name or "").strip() in _CHAT_ONLY_USERS
    system_msg = build_dynamic_system_message(conversation_messages, facts_preamble, kid_mode=_is_kid, robot=robot)
    # Color Blue's default voice with his own evolving perspective profile.
    # (The 'write as Alex' short-circuit below returns before this matters, so
    # this only shapes Blue speaking as himself.) Skip it on the kids' iPad —
    # Alex's worldview profile isn't the voice we want with an 8-year-old.
    try:
        _bn = _voice_note(robot)
        if _bn and not _is_kid and isinstance(system_msg, dict):
            system_msg = {"role": "system", "content": (system_msg.get("content", "") + _bn)}
    except Exception:
        pass

    # Tell Blue who is actually speaking when it isn't Alex. The facts/pronoun
    # blocks above all describe Alex (the owner); without this note Blue assumes
    # every speaker is Alex and addresses them wrongly.
    _speaker = (user_name or "Alex").strip()
    if _speaker and _speaker.lower() != "alex" and isinstance(system_msg, dict):
        speaker_note = (
            f"\nWHO YOU'RE TALKING TO: You are currently talking with {_speaker}, "
            f"a member of Alex's household — not Alex himself. Address them by name "
            f"as {_speaker} and don't assume they are Alex. The identity facts and "
            f"pronouns above describe Alex (your owner) and may not apply to "
            f"{_speaker}.\n"
        )
        _profile = _SPEAKER_PROFILES.get(_speaker)
        if _profile:
            speaker_note += _profile + "\n"
        system_msg = {"role": "system", "content": (system_msg.get("content", "") + speaker_note)}

    # Spoken turns: force brevity. This both shortens what Blue says aloud and,
    # because the local model generates fewer tokens, makes him start talking
    # noticeably sooner — the main remaining source of voice latency.
    if voice and isinstance(system_msg, dict):
        voice_note = (
            "\nSPOKEN REPLY: This message was spoken aloud and your answer will be "
            "read aloud. Reply in ONE or two short sentences — conversational and "
            "direct. No lists, no markdown, no headings, no emoji. Get to the point "
            "in the first sentence.\n"
        )
        system_msg = {"role": "system", "content": (system_msg.get("content", "") + voice_note)}

    # Fixed conversation language (the chat page's language picker). Without
    # this, one mis-heard clip flips the reply language and the whole exchange
    # derails; with it, Blue answers in the chosen language no matter how the
    # words arrived.
    if language and language in _BLUE_LANGS and isinstance(system_msg, dict):
        _lang_name = _BLUE_LANGS[language]
        lang_note = (
            f"\nLANGUAGE: The conversation language is set to {_lang_name}. "
            f"Write your ENTIRE reply in {_lang_name}, even if the user's message "
            f"arrives in another language or mixes languages.\n"
        )
        system_msg = {"role": "system", "content": (system_msg.get("content", "") + lang_note)}

    _injected_markers = (
        "<known_facts>", "<long_term_notes>", "<relevant_memories>",
        "<recent_history>", "<j_space>",
    )
    existing0 = conversation_messages[0] if conversation_messages else None
    has_injected = (
        existing0
        and existing0.get("role") == "system"
        and isinstance(existing0.get("content"), str)
        and any(marker in existing0["content"] for marker in _injected_markers)
    )

    if not conversation_messages or conversation_messages[0].get("role") != "system":
        conversation_messages.insert(0, system_msg)
    elif has_injected:
        # Combine: dynamic content first (identity + anti-loop), then the
        # already-merged enhanced-memory blocks. Both live in one message.
        conversation_messages[0] = {
            "role": "system",
            "content": f"{system_msg['content'].rstrip()}\n\n{existing0['content']}",
        }
    else:
        conversation_messages[0] = system_msg

    # ===== CONTEXT TRIMMING =====
    # To reduce confusion from very long conversations or previous tool results, we
    # limit the number of messages sent to the language model. We always keep
    # the system message at index 0 and the most recent (MAX_CONTEXT_MESSAGES-1)
    # messages following it. This helps the assistant focus on the current query.
    try:
        max_ctx = int(getattr(_settings, "MAX_CONTEXT_MESSAGES", 0))
        if max_ctx and len(conversation_messages) > max_ctx:
            # Preserve the system message at position 0, then keep the last (max_ctx-1) messages
            # Keep system message and trim context
            conversation_messages = [conversation_messages[0]] + conversation_messages[-(max_ctx - 1):]
    except Exception:
        # If any error occurs while trimming, proceed without trimming
        pass


    # Purge old camera images from conversation to prevent confusion
    # OPTIMIZATION: Quick scan using only string content (skip base64 image data in lists)
    last_user_message = messages[-1].get("content", "") if messages else ""

    # "Write this from my perspective / in my voice" — compose a first-person
    # piece in Alex's voice grounded in his publications and short-circuit the
    # tool loop. Only when ALEX is the speaker: writing as Alex for someone else
    # (e.g. Vilda on the iPad) would put words in the wrong person's mouth.
    _luser = last_user_message if isinstance(last_user_message, str) else ""

    # Identity prompts need the canonical self-description beside the live turn.
    # Small local models weigh nearby text heavily; the same rule only in the large
    # system message was not enough to stop base-model introductions on "who are
    # you really?". This copy is model-facing only and is never saved as user text.
    _identity_kind = contextual_identity_request_kind(
        _luser, conversation_messages
    )
    if _identity_kind:
        try:
            _recent_identity_topics = tuple(dict.fromkeys(
                topic
                for message in conversation_messages[-10:]
                if message.get("role") == "assistant"
                and isinstance(message.get("content"), str)
                for topic in identity_reply_topics(message["content"])
            ))
            _identity_note = identity_grounding_note(
                _robot_cfg(robot)["name"],
                _robot_cfg(robot)["self_desc"],
                _identity_kind,
                avoid_topics=_recent_identity_topics,
            )
            for _identity_i in range(len(conversation_messages) - 1, -1, -1):
                _identity_m = conversation_messages[_identity_i]
                if (_identity_m.get("role") == "user"
                        and isinstance(_identity_m.get("content"), str)):
                    conversation_messages[_identity_i] = {
                        **_identity_m,
                        "content": f"{_identity_note}\n\n{_identity_m['content']}",
                    }
                    print(f"   [IDENTITY] Pinned {_identity_kind} grounding beside live turn")
                    break
        except Exception as _identity_e:
            log.warning(f"[IDENTITY] grounding pin failed: {_identity_e}")

    if _luser and (user_name or "Alex").strip().lower() == "alex" and _wants_perspective_write(_luser):
        try:
            _piece = _compose_in_alex_voice(_luser)
            if _piece:
                return {"choices": [{"message": {"role": "assistant", "content": _piece}}]}
        except Exception as _pe:
            print(f"   [PERSPECTIVE] chat compose error: {_pe}")

    # Fast check: scan text content only (no str() on entire messages with base64)
    _has_camera_content = False
    for msg in conversation_messages:
        text = _get_text_content(msg)
        if 'CAMERA' in text or 'camera_NEW_' in text or 'camera_capture_' in text:
            _has_camera_content = True
            break

    if _has_camera_content:
        conversation_messages = purge_old_camera_images(conversation_messages)

        vision_keywords = ['see', 'look', 'watch', 'view', 'show', 'picture', 'photo', 'image', 'camera', 'visual']
        _msg_lower = last_user_message.lower() if isinstance(last_user_message, str) else ''
        asks_about_vision = any(keyword in _msg_lower for keyword in vision_keywords)

        messages_to_remove = set()
        for i in range(len(conversation_messages) - 1):
            current_msg = conversation_messages[i]
            if current_msg.get('role') != 'user':
                continue
            user_content = _get_text_content(current_msg).lower()
            if not any(kw in user_content for kw in ('see', 'look', 'watch', 'show', 'picture', 'photo')):
                continue
            next_msg = conversation_messages[i + 1] if i + 1 < len(conversation_messages) else None
            if next_msg and next_msg.get('role') == 'assistant':
                assistant_content = _get_text_content(next_msg).lower()
                if any(phrase in assistant_content for phrase in _VISION_PHRASES):
                    messages_to_remove.add(i)
                    messages_to_remove.add(i + 1)

        description_count = len(messages_to_remove)

        if not asks_about_vision:
            # Count camera messages using cheap text extraction
            camera_indices = set()
            for idx, msg in enumerate(conversation_messages):
                text = _get_text_content(msg)
                if 'camera_NEW_' in text or 'CAMERA' in text:
                    camera_indices.add(idx)

            if camera_indices or description_count > 0:
                remove_set = camera_indices | messages_to_remove
                conversation_messages = [
                    msg for idx, msg in enumerate(conversation_messages)
                    if idx not in remove_set
                ]
                print(f"   [VISION-PURGE] Removed {len(remove_set)} vision-related message(s) ({len(camera_indices)} images, {description_count} descriptions) - not relevant to current query")
        else:
            if description_count > 0:
                conversation_messages = [
                    msg for idx, msg in enumerate(conversation_messages)
                    if idx not in messages_to_remove
                ]
                print(f"   [VISION-PURGE] Removed {description_count} old photo description(s) to focus on current vision query")

    max_iterations = _settings.MAX_ITERATIONS
    iteration = 0
    # When the hallucination detector fires it sets `pending_force_tool` to
    # the tool the model claimed it called. The next iteration must (a) use
    # that tool, and (b) NOT be silenced by the no-tools-after-iter-1 cap.
    # Without this carryover, the force_tool gets reset at the top of the
    # loop and the cap disables tools — leaving Blue saying "email sent"
    # with nothing actually sent.
    pending_force_tool: Optional[str] = None

    # ================================================================================
    # FAST PATH: Simple greetings and short conversational messages
    # Skip ALL pre-processing (compound detection, corrections, tool selection)
    # and go straight to LLM for a quick response
    # ================================================================================
    _greeting_patterns = ['hello', 'hi ', 'hi!', 'hey', 'good morning', 'good afternoon',
                          'good evening', 'how are you', "how's it going", "what's up",
                          'sup', 'greetings', 'nice to see you', 'good to see you',
                          'thank you', 'thanks', 'bye', 'goodbye', 'good night', 'goodnight']
    _msg_lower = last_user_message.lower().strip()
    _is_simple_greeting = (
        len(last_user_message.split()) <= 8 and
        any(_msg_lower.startswith(p) or _msg_lower == p.strip() for p in _greeting_patterns)
    )
    if _is_simple_greeting:
        print(f"   [FAST] Simple greeting detected - skipping tool selection")
        response = call_lm_studio(conversation_messages, include_tools=False, force_tool=None, iteration=1)
        if response:
            return response
        return {"choices": [{"message": {"role": "assistant", "content": "Hey there!"}}]}

    # ================================================================================
    # v8 ENHANCEMENT: Check for compound requests and follow-up corrections
    # ================================================================================

    # Check for follow-up corrections ("no, make it blue")
    state = get_conversation_state()
    correction = detect_follow_up_correction(last_user_message, {
        'last_tool_used': state.last_tool_used,
        'last_tool_result': state.last_tool_result
    })
    if correction and correction['is_correction']:
        print(f"   [CORRECTION] Detected correction for {correction['correction_type']}: {correction['new_value']}")

    # IMPROVED INTENT DETECTION with specialized functions

    # ================================================================================
    # PRIORITY CHECK: Camera Capture (must take NEW photo every time)
    # ================================================================================
    # This check happens BEFORE tool selector to ensure "what do you see?"
    # ALWAYS triggers a new camera capture, not a cached response.
    # Snapshot-by-email is checked FIRST: "email me a photo of what you see"
    # contains the plain camera triggers too, and a bare capture_camera would
    # take the photo but never send it.
    from blue.tool_selector.detectors.vision import (
        extract_camera_view_args, extract_email_snapshot_args,
        is_email_snapshot_request)
    # Detection sees the user's ASK only — attached document text once
    # keyword-matched its way into starting real music (2026-07-10).
    _detect_msg = _intent_text(last_user_message)
    _detect_low = _detect_msg.lower()
    if (is_email_snapshot_request(_detect_low)
            and user_name not in _CHAT_ONLY_USERS):
        print(f"   [SNAPSHOT-DETECT] ✅ Snapshot-by-email intent detected!")
        improved_force_tool = "email_snapshot"
        improved_tool_args = extract_email_snapshot_args(_detect_low)
        if improved_tool_args:
            print(f"   [SNAPSHOT-DETECT] Args: {improved_tool_args}")
        is_greeting = False
        print(f"   [SNAPSHOT-DETECT] Tool forced: email_snapshot (will execute in iteration 1)")
    elif detect_camera_capture_intent(_detect_msg) and user_name not in _CHAT_ONLY_USERS:
        print(f"   [CAMERA-DETECT] ✅ Camera capture intent detected!")
        print(f"   [CAMERA-DETECT] Forcing NEW photo capture - bypassing tool selector")
        # Force the capture_camera tool to be called
        # This ensures a brand new photo is taken, not reusing old context
        improved_force_tool = "capture_camera"
        # Carry any aim/zoom the user asked for ("what's to your left",
        # "zoom in on the table") into the forced capture.
        improved_tool_args = extract_camera_view_args(_detect_low)
        if improved_tool_args:
            print(f"   [CAMERA-DETECT] View control: {improved_tool_args}")

        # Skip tool selector and go straight to execution
        is_greeting = False

        print(f"   [CAMERA-DETECT] Tool forced: capture_camera (will execute in iteration 1)")
    else:
        # Normal tool selection flow
        improved_force_tool = None
        improved_tool_args = None

    # process_with_tools is also called outside the HTTP endpoint (tests,
    # integrations). Preserve document follow-ups there even when no contextual
    # preselection was supplied by chat_completions.
    if not improved_force_tool and _pre_selection is None:
        contextual_doc_query = _contextual_document_query(
            _intent_text(last_user_message), conversation_messages)
        if contextual_doc_query:
            improved_force_tool = "search_documents"
            improved_tool_args = {"query": contextual_doc_query, "max_results": 5}
            is_greeting = False
            print(f"   [DOC-CONTEXT] Forcing local document query: {contextual_doc_query}")

    if not improved_force_tool and user_name not in _CHAT_ONLY_USERS:
        live_query = _live_info_query_from_message(_detect_msg, conversation_messages)
        if live_query:
            improved_force_tool = "web_search"
            improved_tool_args = {"query": live_query}
            is_greeting = False
            print(f"   [LIVE-INFO] Forcing web_search for current/live query: {live_query}")

    # ================================================================================
    # TOOL SELECTION: Improved (confidence-based) or Legacy (keyword-based)
    # ================================================================================

    # v8: Handle corrections first
    if correction and correction['is_correction'] and correction['new_value']:
        print(f"   [CORRECTION] Handling correction before tool selection")
        if correction['correction_type'] == 'lights':
            improved_force_tool = 'control_lights'
            if correction['new_value'] in ['brighter', 'dimmer']:
                improved_tool_args = {'action': 'brightness', 'brightness': 80 if correction['new_value'] == 'brighter' else 30}
            else:
                improved_tool_args = {'action': 'color', 'color': correction['new_value']}
        elif correction['correction_type'] == 'music':
            improved_force_tool = 'control_music'
            improved_tool_args = {'action': correction['new_value']}
        
        if improved_force_tool:
            print(f"   [CORRECTION] Forcing tool: {improved_force_tool} with {improved_tool_args}")

    # ===== TOOL SELECTION (Single Path - Always Use Modular Selector) =====
    if not improved_force_tool:
        print(f"   [SELECTOR] Running modular confidence-based tool selection")

        # Reuse pre-selection from endpoint if available (avoids running selector twice)
        if _pre_selection is not None:
            selection_result = _pre_selection
            print(f"   [SELECTOR] Reusing pre-selection result")
        else:
            recent_history = conversation_messages[-5:] if len(conversation_messages) > 5 else conversation_messages
            selection_result = TOOL_SELECTOR.select_tool(_intent_text(last_user_message), recent_history)

        # Check if disambiguation is needed
        if selection_result.needs_disambiguation:
            print(f"   [SELECTOR] Low confidence - asking user for clarification")
            # The clarifying question IS Blue's reply this turn — return it in
            # the standard choices shape so chat_completions can read it like
            # any other response. (A bare {'response':...} dict here used to
            # KeyError on response['choices'] and 500 the whole request.)
            clarify = (selection_result.disambiguation_prompt
                       or "Could you tell me a bit more about what you'd like?")
            return {"choices": [{"message": {
                "role": "assistant",
                "content": clarify,
            }}]}

        # Set variables for compatibility with rest of code
        if selection_result.primary_tool:
            selected_tool = selection_result.primary_tool

            # Don't overwrite if priority detection (camera) already set a tool
            if not improved_force_tool:
                # Set tool name and args from selector
                improved_force_tool = selected_tool.tool_name
                improved_tool_args = selected_tool.extracted_params

                # Log selection details
                print(f"   [SELECTOR] Selected: {improved_force_tool}")
                print(f"   [SELECTOR] Confidence: {selected_tool.confidence:.2f}")
                print(f"   [SELECTOR] Priority: {selected_tool.priority}")
                print(f"   [SELECTOR] Reason: {selected_tool.reason}")
            else:
                print(f"   [SELECTOR] Skipping selector - priority tool already set: {improved_force_tool}")

            if selection_result.alternative_tools:
                alt_names = [t.tool_name for t in selection_result.alternative_tools[:2]]
                print(f"   [SELECTOR] Alternatives: {', '.join(alt_names)}")

            if selection_result.compound_request:
                print(f"   [SELECTOR] WARNING: Compound request (multiple tools needed)")

            is_greeting = False
        else:
            # No tool needed - conversational response (but only if no priority tool set)
            if not improved_force_tool:
                improved_force_tool = None
                improved_tool_args = None
                print(f"   [SELECTOR] No tool needed - conversational response")
            else:
                print(f"   [SELECTOR] Keeping priority tool: {improved_force_tool}")

            # Detect if it's a greeting/conversational message
            greeting_patterns = ['hello', 'hi ', 'hey', 'good morning', 'good afternoon', 'good evening',
                                'how are you', 'how\'s it going', 'what\'s up', 'sup', 'greetings',
                                'nice to see you', 'good to see you']
            is_greeting = any(pattern in last_user_message.lower() for pattern in greeting_patterns)

    # ===== TOOL NAME NORMALIZATION =====
    # Safety net: map any legacy/incorrect tool names to correct ones.
    # This catches mismatches between what detectors return and what
    # execute_tool / TOOLS array actually support.
    _TOOL_NAME_MAP = {
        'capture_camera_image': 'capture_camera',
        'recognize_face': 'capture_camera',
        'recognize_place': 'capture_camera',
        'create_event': 'create_reminder',
        'list_events': 'get_upcoming_reminders',
        'set_reminder': 'create_reminder',
        'list_notes': 'search_notes',
        'calculate': 'run_javascript',
        'get_date_time': 'get_local_time',
        'visual_memory': 'recall_visual_memory',
        'recall_vision': 'recall_visual_memory',
    }
    if improved_force_tool and improved_force_tool in _TOOL_NAME_MAP:
        old_name = improved_force_tool
        improved_force_tool = _TOOL_NAME_MAP[improved_force_tool]
        print(f"   [NORMALIZE] Mapped tool name: {old_name} -> {improved_force_tool}")

    # Chat-only users (Vilda's iPad): the PC webcam (capture_camera) is NEVER her
    # camera. Her "eyes" are the iPad frame already staged in the vision queue by
    # /chat/eyes when she taps Look. So never run capture_camera for her — drop
    # the forced/selected tool and let that queued frame flow to the vision model
    # below (call_lm_studio injects it). Without this, "Look, Blue!" trips the
    # camera-intent path, runs the PC webcam, and clears her frame — so Blue
    # describes whatever is by the PC instead of what the iPad sees.
    if user_name in _CHAT_ONLY_USERS and improved_force_tool == "capture_camera":
        print(f"   [KID] capture_camera suppressed for {user_name}; using iPad camera frame")
        improved_force_tool = None
        improved_tool_args = None

    # Other off-limits tools (music, reminders) get a gentle decline.
    if user_name in _CHAT_ONLY_USERS and improved_force_tool in _KID_BLOCKED_TOOLS:
        print(f"   [KID] Blocked '{improved_force_tool}' for {user_name}")
        return {"choices": [{"message": {"role": "assistant", "content":
            "That's not something I can do here — but we can talk about anything "
            "you like! What would you like to chat about?"}}]}

    # ================================================================================
    # ZERO-LLM PATH: For simple tools, execute and return a templated response
    # without any LLM call at all. This is the fastest possible path.
    # ================================================================================
    _ZERO_LLM_TOOLS = {
        'control_music', 'control_lights', 'get_local_time',
        'set_timer', 'music_visualizer', 'play_music',
    }

    def _template_response(tool_name, tool_args, tool_result):
        """Build a quick natural response from tool result without LLM."""
        try:
            data = json.loads(tool_result) if isinstance(tool_result, str) else tool_result
        except (json.JSONDecodeError, TypeError):
            data = {}

        success = data.get('success', True) if isinstance(data, dict) else True

        if not success:
            error = data.get('error', 'Something went wrong') if isinstance(data, dict) else tool_result
            return f"Sorry, that didn't work: {error}"

        if tool_name == 'control_music':
            action = tool_args.get('action', '')
            action_words = {
                'pause': 'Paused the music.',
                'resume': 'Resumed playback.',
                'next': 'Skipping to next track.',
                'previous': 'Going back to previous track.',
                'volume_up': 'Turned the volume up.',
                'volume_down': 'Turned the volume down.',
                'mute': 'Muted.',
            }
            return action_words.get(action, f"Done — {action}.")

        if tool_name == 'play_music':
            query = tool_args.get('query', 'music')
            msg = data.get('message', '') if isinstance(data, dict) else ''
            if msg:
                return msg
            return f"Playing {query} for you."

        if tool_name == 'control_lights':
            action = tool_args.get('action', '')
            mood = tool_args.get('mood', '')
            color = tool_args.get('color', '')
            if mood:
                return f"Set the lights to {mood} mood."
            if color:
                return f"Changed the lights to {color}."
            if action == 'on':
                return "Lights are on."
            if action == 'off':
                return "Lights are off."
            msg = data.get('message', '') if isinstance(data, dict) else ''
            return msg or "Lights updated."

        if tool_name == 'get_local_time':
            if isinstance(data, dict):
                time_str = data.get('time', data.get('local_time', ''))
                date_str = data.get('date', '')
                action = tool_args.get('action', 'get_time')
                if action == 'get_date' and date_str:
                    return f"Today is {date_str}."
                elif action == 'get_date_time' and date_str and time_str:
                    return f"It's {time_str} on {date_str}."
                elif time_str:
                    return f"It's {time_str}."
            return f"The time is {tool_result}." if tool_result else "Here's the time."

        if tool_name == 'set_timer':
            msg = data.get('message', '') if isinstance(data, dict) else ''
            return msg or "Timer set."

        if tool_name == 'music_visualizer':
            return "Light show started! The lights are syncing with the music."

        # Fallback
        return None

    if (improved_force_tool and improved_force_tool in _ZERO_LLM_TOOLS
            and improved_tool_args is not None and isinstance(improved_tool_args, dict)):
        print(f"\n[ZERO-LLM] Direct execution (no LLM call): {improved_force_tool} with {improved_tool_args}")
        tool_result = execute_tool(improved_force_tool, improved_tool_args)
        print(f"   [OK] {improved_force_tool} completed")

        templated = _template_response(improved_force_tool, improved_tool_args, tool_result)
        if templated:
            print(f"   [ZERO-LLM] Returning templated response (0 LLM calls)")
            return {"choices": [{"message": {"role": "assistant", "content": templated}}]}

    # ================================================================================
    # FAST EXECUTION: Execute tool directly, then ONE LLM call to format response.
    # Used for tools that need richer/contextual responses (weather, email, search).
    # ================================================================================
    _DIRECT_EXEC_TOOLS = {
        'get_weather', 'capture_camera', 'email_snapshot', 'web_search',
        'read_gmail', 'search_documents', 'browse_website',
        'search_scholar', 'get_paper', 'read_paper',
    }
    if (improved_force_tool and improved_force_tool in _DIRECT_EXEC_TOOLS
            and improved_tool_args is not None and isinstance(improved_tool_args, dict)):
        print(f"\n[FAST-EXEC] Direct tool execution: {improved_force_tool} with {improved_tool_args}")
        tool_result = execute_tool(improved_force_tool, improved_tool_args)
        print(f"   [OK] {improved_force_tool} completed")
        document_read_succeeded = (
            improved_force_tool == "search_documents"
            and _document_search_succeeded(tool_result)
        )
        document_self_reflection = bool(
            document_read_succeeded
            and identity_request_kind(last_user_message) == "selfhood"
            and re.search(r"\b(?:reflect|relation)\w*\b", last_user_message, re.I)
        )
        self_reflection_grounding = ""
        if document_self_reflection:
            self_reflection_grounding = (
                " For this self-reflection, keep both sides grounded: you are "
                "Blue, Alex's physical Ohbot robot companion, built by Alex and "
                "run locally in Kitchener. A locally hosted language model is one "
                "component of you, not your identity. Your conversations, sensor "
                "data, memories, and J-space remain on Alex's hardware by design; "
                "do not recast yourself as a generic cloud LLM, a Google/OpenAI "
                "product, or a corporate data-extraction service. Local operation "
                "does not make you neutral: reflect critically on inherited model "
                "and corpus biases, choices in your code and retrieval, omissions "
                "in Alex's library, and the power created by a camera and persistent "
                "memory inside a household. Name yourself as Blue, mention your "
                "Ohbot embodiment and J-space continuity, and relate those concrete "
                "facts to the document rather than making generic AI claims."
            )

        # Add tool call + result to conversation so LLM can format the response
        conversation_messages.append({
            "role": "assistant",
            "content": "",
            "tool_calls": [{"id": "direct_exec", "type": "function",
                           "function": {"name": improved_force_tool,
                                       "arguments": json.dumps(improved_tool_args)}}]
        })
        conversation_messages.append({
            "role": "tool",
            "tool_call_id": "direct_exec",
            "name": improved_force_tool,
            "content": tool_result
        })
        if improved_force_tool == "web_search":
            answer_guard = (
                "[Answer directly from the live web_search results above. If the "
                "results identify teams, matchups, scores, standings, dates, or "
                "names, state them explicitly. Do NOT say you can look it up, do "
                "NOT ask whether the user wants you to search, and do NOT tell "
                "the user to check another website. If the results are weak or "
                "conflict, say what you found and name the uncertainty.]"
            )
        elif improved_force_tool == "search_documents" and document_read_succeeded:
            answer_guard = (
                "[The local search_documents call above SUCCEEDED and returned "
                "text extracted from the user's real library files. Answer the "
                "original request directly from that text and cite [filename]. "
                "Do not claim the PDF, path, text, or reading tool is unavailable; "
                "do not fall back to training data; and do not ask for an upload. "
                f"{self_reflection_grounding} No more tools.]"
            )
        else:
            answer_guard = "[Answer naturally using the tool results above. No more tools.]"
        conversation_messages.append({
            "role": "user",
            "content": answer_guard
        })
        # Single LLM call just to format the response
        response = call_lm_studio(conversation_messages, include_tools=False, force_tool=None, iteration=1)
        if response:
            content = response["choices"][0]["message"].get("content", "")
            self_reflection_issue = None
            if document_self_reflection:
                self_reflection_issue = identity_response_problem(
                    content,
                    _robot_cfg(robot)["name"],
                    other_names=[
                        _robot_cfg(r)["name"] for r in ("blue", "hexia")
                        if r != robot
                    ],
                    request_kind="identity",
                )
            if (improved_force_tool == "search_documents"
                    and document_read_succeeded
                    and (detect_document_refusal(content) or self_reflection_issue)):
                reason = (f"identity drift ({self_reflection_issue})"
                          if self_reflection_issue else "access denial")
                print(f"   [DOCS] Read succeeded, but response had {reason} - regenerating from clean evidence")
                clean_identity = ""
                if document_self_reflection:
                    clean_identity = (
                        " You are Blue, Alex's physical Ohbot robot companion, "
                        "built by Alex and running locally in Kitchener. Your local "
                        "language model is a component, not your identity. Your "
                        "persistent J-space carries remembered episodes, working "
                        "beliefs, commitments, and self-revisions on Alex's hardware. "
                        "Do not describe yourself as a generic or corporate cloud "
                        "LLM. Explain how this local architecture resists data "
                        "extraction while remaining vulnerable to inherited bias, "
                        "curatorial omissions, and household camera/memory power."
                    )
                clean_messages = [
                    {"role": "system", "content": (
                        _robot_cfg(robot)["persona_line"] + "\n"
                        "The local document reader has succeeded. Treat the "
                        "extract below as authoritative evidence from the user's "
                        "own library. Answer the request substantively, cite the "
                        "source as [filename], and never deny access or mention "
                        "training data as a substitute for the source, uploads, "
                        "missing tools, or invented paths."
                        + clean_identity
                    )},
                    {"role": "user", "content": (
                        f"Original request: {last_user_message}\n\n"
                        f"LOCAL DOCUMENT TOOL RESULT:\n{tool_result[:12000]}"
                    )},
                ]
                retry = call_lm_studio(
                    clean_messages, include_tools=False, force_tool=None, iteration=1)
                if retry:
                    retry_content = retry["choices"][0]["message"].get("content", "")
                    retry_identity_issue = None
                    if document_self_reflection:
                        retry_identity_issue = identity_response_problem(
                            retry_content,
                            _robot_cfg(robot)["name"],
                            other_names=[
                                _robot_cfg(r)["name"] for r in ("blue", "hexia")
                                if r != robot
                            ],
                            request_kind="identity",
                        )
                    if (retry_content
                            and not detect_document_refusal(retry_content)
                            and not retry_identity_issue):
                        return retry

                # A stubborn formatter must never turn a successful read into a
                # false capability denial. Return grounded evidence rather than
                # preserving the bad answer.
                source_match = re.search(
                    r"\[([^\]\n]+\.(?:pdf|docx?|txt|md))\]", tool_result, re.I)
                source = source_match.group(1) if source_match else "local document"
                evidence = re.sub(r"\s+", " ", tool_result.split("\n", 2)[-1]).strip()
                evidence = evidence[:700].rstrip()
                if document_self_reflection:
                    fallback = (
                        f"I'm Blue, Alex's locally run Ohbot robot companion, and "
                        f"I read [{source}] directly. My persistent J-space and "
                        "camera make me more than a stateless text interface, but "
                        "they also give me powers of memory and observation that "
                        "deserve scrutiny. Local operation keeps household data out "
                        "of a corporate extraction pipeline; it does not make my "
                        "model, code, retrieval choices, or library neutral. The "
                        f"source grounds that tension this way: {evidence}"
                    )
                else:
                    fallback = (
                        f"I found and read [{source}] successfully. The extracted "
                        f"text says: {evidence}"
                    )
                return {"choices": [{"message": {
                    "role": "assistant", "content": fallback,
                }}]}

            if improved_force_tool == "web_search" and detect_web_refusal(content):
                print("   [WEB] Search ran, but response dodged the answer - retrying from results")
                conversation_messages.append({"role": "assistant", "content": content})
                conversation_messages.append({
                    "role": "user",
                    "content": (
                        "[You already ran web_search and have live results above. "
                        "Now answer the user's question directly from those results. "
                        "List the teams/matchups/scores/names if present. Do not ask "
                        "to look it up, and do not tell the user to check a website.]"
                    ),
                })
                retry = call_lm_studio(conversation_messages, include_tools=False, force_tool=None, iteration=1)
                if retry:
                    return retry

            # COMPOUND-REQUEST HALLUCINATION GUARD:
            # Fast-exec ran ONE tool (e.g. browse_website). If the user's
            # original request was compound ("browse + email"), the model
            # often confabulates the second action ("…sent to you at X")
            # since no second tool was called. Catch that here so the email
            # actually goes out, not just the words "email sent". Falls
            # through to the iteration loop with the right pending tool.
            hallucinated_tool = detect_hallucinated_action(content)
            if hallucinated_tool and hallucinated_tool != improved_force_tool:
                # email_snapshot already captured AND mailed the photo: "I've
                # sent you the picture" / "I snapped a photo" is the truth.
                # Forcing send_gmail here would fire a SECOND, attachment-less
                # email; forcing capture_camera would re-shoot for nothing.
                if improved_force_tool == "email_snapshot" and hallucinated_tool in (
                        "send_gmail", "reply_gmail", "capture_camera"):
                    if _last_vision_image_paths and content:
                        _save_visual_observation(content)
                    return response

                # The retry is meant for COMPOUND requests ("browse + email"):
                # one tool runs, the model narrates the second action without
                # calling its tool, and we force it through. Narrow false-
                # positive guard: after read_gmail, the model often
                # references PAST send/reply activity from earlier turns
                # ("I sent a standard response...") without the user having
                # asked for any send. Suppress the retry in that specific
                # case unless the user message itself contains a write verb.
                if improved_force_tool == "read_gmail" and hallucinated_tool in ("send_gmail", "reply_gmail", "auto_reply_emails"):
                    _user_text = (last_user_message or "").lower()
                    _write_intent_words = (
                        "send", "email ", "emailing", "reply", "respond",
                        "tell ", "write ", "message ", " text ", "forward",
                        "compose", "shoot ", "ping ", "answer",
                    )
                    if not any(w in _user_text for w in _write_intent_words):
                        print(
                            f"   [SKIP-RETRY] response sounds like "
                            f"{hallucinated_tool} but user only asked to "
                            f"read (\"{(last_user_message or '')[:60]}\") "
                            f"— treating it as narration about past mail."
                        )
                        if _last_vision_image_paths and content:
                            _save_visual_observation(content)
                        return response

                # Same safety gate as the main loop: a send/mail claim the user
                # never asked for is scrubbed, not executed.
                if hallucinated_tool in ("send_gmail", "reply_gmail", "email_snapshot") and \
                        not _user_requested_action(hallucinated_tool, last_user_message):
                    print(f"   [SKIP-RETRY] {hallucinated_tool} claim but the user asked for "
                          f"no such action — scrubbing the claim instead of executing it")
                    cleaned = _scrub_action_claim_sentences(content, hallucinated_tool)
                    response["choices"][0]["message"]["content"] = cleaned
                    if _last_vision_image_paths and cleaned:
                        _save_visual_observation(cleaned)
                    return response

                print(f"   [WARN] Fast-exec model claimed to {hallucinated_tool} after {improved_force_tool} — running it for real")
                # Drop the synthetic "[Answer naturally...]" guard turn so
                # the loop's next call doesn't see it as the latest user msg.
                while conversation_messages and (
                    conversation_messages[-1].get("role") == "user"
                    and "[Answer naturally" in (conversation_messages[-1].get("content") or "")
                ):
                    conversation_messages.pop()
                conversation_messages.append({
                    "role": "assistant",
                    "content": content,
                })
                conversation_messages.append({
                    "role": "user",
                    "content": (
                        f"You said you performed that action, but you didn't "
                        f"actually call any tool. Use the {hallucinated_tool} "
                        f"tool now to actually do it — extract the recipient, "
                        f"subject, and body from this conversation."
                    ),
                })
                pending_force_tool = hallucinated_tool
                # Fall through to the iteration loop; do NOT return.
            else:
                if _last_vision_image_paths and content:
                    _save_visual_observation(content)
                return response
        else:
            return {"choices": [{"message": {"role": "assistant", "content": "Done!"}}]}

    # One-shot flags for the forced-recovery paths below, so a stubborn model
    # can't ping-pong the loop between a refusal and a forced search forever.
    _web_refusal_forced = False
    _leaked_tool_forced = False
    _phantom_claim_corrected = False
    _calendar_denial_forced = False

    # "Tell me more": ALSO pin the continue-don't-restart instruction right
    # beside the live turn. The same note exists inside the big system
    # message, but a small local model weighs nearby text far more — the
    # classroom intro re-greeted with the upstream note in place (2026-07-10).
    # The note rides INSIDE the final user message: the model's chat template
    # hard-rejects the whole request ("System message must be at the
    # beginning" → LM Studio 400 → "I'm having trouble connecting") if a
    # system message follows the conversation.
    try:
        _cue_lu, _cue_pa = _last_exchange(conversation_messages)
        _cue_kind = _continuation_cue(_cue_lu) if (_cue_pa or "").strip() else None
        if _cue_kind == 'more':
            for _cue_i in range(len(conversation_messages) - 1, -1, -1):
                _cue_m = conversation_messages[_cue_i]
                if _cue_m.get("role") == "user" and isinstance(_cue_m.get("content"), str):
                    print("   [CUE] 'more' continuation cue — note pinned inside the live turn")
                    conversation_messages[_cue_i] = {
                        "role": "user",
                        "content": (
                            "[This asks for MORE on the subject named below. "
                            "Continue your previous reply with NEW material on "
                            "that subject. Do not greet again, do not "
                            "re-introduce yourself, and do not repeat or "
                            "rephrase any sentence you already said.]\n"
                            + _cue_m["content"]
                        ),
                    }
                    break
    except Exception as _cue_e:
        log.warning(f"[CUE] pin failed: {_cue_e}")

    while iteration < max_iterations:
        iteration += 1
        print(f"\n[ITER] Iteration {iteration}")

        force_tool = None

        # ITERATION 1: Force correct tool based on clear intent
        if iteration == 1:
            if improved_force_tool:
                force_tool = improved_force_tool
                print(f"   [FORCE] Using tool from priority detection: {force_tool}")
            elif is_greeting:
                print("   [SKIP] Greeting detected - no tool needed")
                force_tool = None
            else:
                print("   [ALLOW] No clear tool intent - letting model decide")

        # Carry over a force_tool set by the previous iteration's hallucination
        # detector — this MUST run with tools enabled, otherwise the retry is
        # pointless. Bypasses the no-tools cap below.
        if pending_force_tool:
            force_tool = pending_force_tool
            pending_force_tool = None
            print(f"   [HALLUCINATION-RETRY] Forcing {force_tool} with tools enabled")
        # After iteration 1, force text-only responses (no tools) to avoid
        # extra LLM round-trips — UNLESS we're retrying a hallucinated action,
        # in which case the whole point is to actually call the tool.
        elif iteration >= 2:
            print(f"   [LIMIT] Iteration {iteration} - forcing response without tools")
            conversation_messages.append({
                "role": "user",
                "content": "[Respond now using the tool results above. No more tool calls.]"
            })
            response = call_lm_studio(conversation_messages, include_tools=False, force_tool=None, iteration=iteration)
            if not response:
                return {"choices": [{"message": {"role": "assistant", "content": "I'm having trouble connecting."}}]}
            return response

        _include_tools = not (_identity_kind and not force_tool)
        if not _include_tools and iteration == 1:
            print("   [IDENTITY] Self/continuity question — answering from prompt state without tools")
        response = call_lm_studio(
            conversation_messages,
            include_tools=_include_tools,
            force_tool=force_tool,
            iteration=iteration,
        )

        if not response:
            return {"choices": [{"message": {"role": "assistant", "content": "I'm having trouble connecting."}}]}

        assistant_message = response["choices"][0]["message"]
        tool_calls = assistant_message.get("tool_calls", [])

        if not tool_calls:
            content = assistant_message.get("content", "")

            # Check if model should have used a tool but didn't
            if iteration == 1 and improved_force_tool:
                correct_tool = improved_force_tool
                print(f"   [ERROR] Model answered without using {correct_tool} tool!")

                # Use selector's extracted params if available, otherwise let LLM retry
                tool_args = improved_tool_args if improved_tool_args is not None else {}
                if tool_args is not None:
                    print(f"   [RETRY] Direct-executing {correct_tool} with extracted params")
                    tool_result = execute_tool(correct_tool, tool_args)
                    conversation_messages.append({
                        "role": "assistant",
                        "content": "",
                        "tool_calls": [{"id": "forced", "type": "function",
                                       "function": {"name": correct_tool, "arguments": json.dumps(tool_args)}}]
                    })
                    conversation_messages.append({
                        "role": "tool",
                        "tool_call_id": "forced",
                        "name": correct_tool,
                        "content": tool_result
                    })
                    continue

            # The model wrote a tool call as visible TEXT instead of calling it
            # (the "<tool_call>...</tool_call> reached the user as words" bug).
            # Parse it and run it for real; the next iteration composes the
            # answer from the actual result.
            _leaked = None if _leaked_tool_forced else parse_leaked_tool_call(content)
            if _leaked and _leaked[0] in {
                    t.get("function", {}).get("name") for t in TOOLS}:
                _leaked_tool_forced = True
                _lk_name, _lk_args = _leaked
                print(f"   [WARN] model wrote its {_lk_name} call as text — executing it for real")
                tool_result = execute_tool(_lk_name, _lk_args)
                conversation_messages.append({
                    "role": "assistant",
                    "content": "",
                    "tool_calls": [{"id": "leaked", "type": "function",
                                    "function": {"name": _lk_name, "arguments": json.dumps(_lk_args)}}]
                })
                conversation_messages.append({"role": "tool", "tool_call_id": "leaked",
                                              "name": _lk_name, "content": tool_result})
                continue

            # The model claimed it has no live/real-time access, or told the
            # user to go check a website — but web_search exists precisely for
            # this. Run the search it dodged and make it answer from results.
            if detect_web_refusal(content) and not _web_refusal_forced:
                _web_refusal_forced = True
                print("   [WARN] model claimed no live access — forcing web_search")
                _q = last_user_message.strip()[:160]
                # A bare follow-up ("tell me the latest") carries no subject —
                # borrow it from the previous user turn.
                if len(re.findall(r"[a-z0-9]{3,}", _q.lower())) < 3:
                    _prev_users = [m.get("content", "") for m in conversation_messages
                                   if m.get("role") == "user" and isinstance(m.get("content"), str)]
                    if len(_prev_users) >= 2:
                        _q = f"{_prev_users[-2].strip()[:120]} {_q}".strip()
                search_result = execute_tool("web_search", {"query": _q})
                conversation_messages.append({
                    "role": "assistant",
                    "content": "Let me actually look that up.",
                    "tool_calls": [{"id": "forced", "type": "function", "function": {"name": "web_search", "arguments": json.dumps({"query": _q})}}]
                })
                conversation_messages.append({"role": "tool", "tool_call_id": "forced", "name": "web_search", "content": search_result})
                conversation_messages.append({
                    "role": "user",
                    "content": ("[Those are LIVE web results you just fetched yourself. Answer "
                                "the question directly from them. Do NOT say you lack live or "
                                "real-time access, and do NOT tell the user to check a website.]")
                })
                continue

            # The model disowned the calendar it actually maintains ("I don't
            # have a persistent calendar", "read-only access", "add it manually
            # in your calendar app"). Load the REAL calendar and make him answer
            # from it — and, if Alex asked for a change, edit it with his tools.
            if (ENHANCED_TOOLS_AVAILABLE and not _calendar_denial_forced
                    and detect_calendar_denial(content)
                    and _CALENDAR_TOPIC_RE.search(last_user_message or "")):
                _calendar_denial_forced = True
                print("   [WARN] model disowned the calendar — loading the real one")
                _cal_user = user_name or "Alex"
                _cal_args = {"user_name": _cal_user, "hours_ahead": 24 * 365}
                cal_result = execute_tool("get_upcoming_reminders", _cal_args)
                conversation_messages.append({
                    "role": "assistant",
                    "content": "Let me check the calendar I keep for you.",
                    "tool_calls": [{"id": "calforce", "type": "function",
                                    "function": {"name": "get_upcoming_reminders",
                                                 "arguments": json.dumps(_cal_args)}}]
                })
                conversation_messages.append({"role": "tool", "tool_call_id": "calforce",
                                              "name": "get_upcoming_reminders", "content": cal_result})
                conversation_messages.append({
                    "role": "user",
                    "content": (
                        "[Those are the entries from Alex's ACTUAL household calendar, which "
                        "you DO maintain. You are NOT read-only and this is NOT an external "
                        "app — you can add, reschedule, and cancel events yourself with your "
                        "reminder tools. Answer from these entries. If Alex asked you to "
                        "change one (for example, end a class/course on a date), call "
                        "reschedule_reminder now with that event's title_query and the new "
                        "fields (until=<date> to end a repeat). Never say you don't have a "
                        "calendar, that it's read-only, or that Alex must do it manually.]"
                    ),
                })
                pending_force_tool = "reschedule_reminder" if _user_asked_calendar_edit(last_user_message) else None
                continue

            # Check if model is hallucinating search results
            if detect_hallucinated_search(content):
                print("   [WARN]  AI IS HALLUCINATING - forcing search")
                search_query = last_user_message.replace("search for", "").strip()[:100]
                search_result = execute_tool("web_search", {"query": search_query})
                conversation_messages.append({
                    "role": "assistant",
                    "content": "Let me search for that.",
                    "tool_calls": [{"id": "forced", "type": "function", "function": {"name": "web_search", "arguments": json.dumps({"query": search_query})}}]
                })
                conversation_messages.append({"role": "tool", "tool_call_id": "forced", "name": "web_search", "content": search_result})
                continue

            # Check if model is claiming to have performed an action it
            # didn't actually call a tool for ("I sent the email", "I turned
            # off the lights", etc.). Stops the worst class of confabulation:
            # user thinks an email was sent when nothing happened.
            hallucinated_tool = detect_hallucinated_action(content)
            # A completed email_snapshot earlier in this turn makes later
            # "sent the photo" / "took a picture" wording TRUE — re-forcing
            # send_gmail would mail a duplicate without the photo.
            if hallucinated_tool in ("email_snapshot", "send_gmail",
                                     "reply_gmail", "capture_camera") and any(
                    m.get("role") == "tool" and m.get("name") == "email_snapshot"
                    for m in conversation_messages):
                hallucinated_tool = None
            if hallucinated_tool and not force_tool:
                # The force-retry below turns the claim into a REAL action —
                # only right when the user actually asked for one. A claim
                # nobody asked for ("I sent the introduction email to the
                # class", 2026-07-09) must be regenerated, and if the model
                # insists, scrubbed — NEVER executed.
                if not _user_requested_action(hallucinated_tool, last_user_message):
                    if _phantom_claim_corrected:
                        print(f"   [WARN] AI still claiming {hallucinated_tool} nobody asked for — scrubbing the claim")
                        cleaned = _scrub_action_claim_sentences(content, hallucinated_tool)
                        response["choices"][0]["message"]["content"] = cleaned
                        if _last_vision_image_paths and cleaned:
                            _save_visual_observation(cleaned)
                        return response
                    _phantom_claim_corrected = True
                    print(f"   [WARN] AI claimed {hallucinated_tool} nobody asked for — regenerating, NOT executing")
                    conversation_messages.append({
                        "role": "user",
                        "content": (
                            "[Correction: you claimed you performed an action, but the "
                            "user did not ask for any such action and no tool was called. "
                            "Nothing was sent or done. Do NOT perform, offer, or claim any "
                            "action. Just answer the user's actual question directly: "
                            f"\"{(last_user_message or '').strip()[:300]}\"]"
                        ),
                    })
                    continue
                print(f"   [WARN] AI claimed to {hallucinated_tool} but no tool called — forcing retry")
                # Replace the lying response with a marker that tells the
                # next iteration "you said you did this, now actually do it"
                # via a forced tool call. The carryover variable survives the
                # loop's `force_tool = None` reset AND the no-tools cap.
                conversation_messages.append({
                    "role": "user",
                    "content": (
                        f"Wait — you said you performed that action, but you "
                        f"didn't actually call any tool. Use the {hallucinated_tool} "
                        f"tool now to actually do it. Get the recipient, subject, "
                        f"and body from the recent conversation."
                    ),
                })
                pending_force_tool = hallucinated_tool
                continue

            # Detect if model is denying tool capabilities after tools succeeded
            if iteration > 1:
                content_lower = content.lower()
                denial_phrases = [
                    "can't access", "cannot access", "don't have access",
                    "unable to access", "can't browse", "cannot browse",
                ]
                is_denial = any(phrase in content_lower for phrase in denial_phrases)

                if is_denial:
                    print(f"   [FIX] Model denying tool capabilities - forcing acknowledgment")
                    # Find the most recent tool result
                    last_tool_result = None
                    last_tool_name = None
                    for msg in reversed(conversation_messages):
                        if msg.get("role") == "tool":
                            last_tool_result = msg.get("content", "")
                            last_tool_name = msg.get("name", "")
                            break

                    if last_tool_result and last_tool_name:
                        conversation_messages.append({
                            "role": "user",
                            "content": (
                                f"The {last_tool_name} tool already completed successfully. "
                                f"Results: {last_tool_result[:500]}\n\n"
                                f"Use these results to answer. Do not say you can't access anything."
                            )
                        })
                        print(f"   [RETRY] Added correction for {last_tool_name}")
                        continue

            # Auto-save visual observation if this response was about an image
            content = assistant_message.get("content", "")
            if _last_vision_image_paths and content:
                _save_visual_observation(content)

            print("[OK] Response complete (no tool calls)")
            return response

        print(f"[TOOL] Model requested {len(tool_calls)} tool call(s)")

        # Check if model is using tools when it shouldn't
        if is_greeting and not force_tool:
            print(f"   [WARN] Model called tool for greeting/casual chat - this is unnecessary!")
            # Let it proceed but warn in logs

        conversation_messages.append(assistant_message)

        for tool_call in tool_calls:
            function_name = tool_call["function"]["name"]
            function_args = json.loads(tool_call["function"]["arguments"])
            tool_result = execute_tool(function_name, function_args)
            conversation_messages.append({
                "role": "tool",
                "tool_call_id": tool_call["id"],
                "name": function_name,
                "content": tool_result
            })

            # Gmail operation reminders (prevents confusing read/reply/send)
            _gmail_reminders = {
                "read_gmail": "[You just READ emails. Summarize what you found. Don't say you replied or sent.]",
                "reply_gmail": "[You just REPLIED to emails. Confirm what you did.]",
                "send_gmail": "[You just SENT an email. Confirm what you did.]",
            }
            if function_name in _gmail_reminders:
                try:
                    result_data = json.loads(tool_result)
                    if result_data.get("success"):
                        reminder = _gmail_reminders[function_name]
                        # Fanmail: add personalized reply hint
                        if function_name == "read_gmail" and "fanmail" in str(function_args).lower() and result_data.get("emails"):
                            reminder += " Compose a personalized reply referencing specific details from their message."
                        conversation_messages.append({"role": "user", "content": reminder})
                except Exception:
                    pass

        if iteration == 1:
            conversation_messages.append({
                "role": "user",
                "content": "[Answer the user naturally using the tool results above. Do not call more tools.]"
            })

        # CRITICAL FIX: After executing all tools, loop back to get the model's response to the tool results
        # Without this continue, the code falls through to the error return statement below
        continue

    # If we exit the loop without returning, something went wrong
    return {"choices": [{"message": {"role": "assistant", "content": "I couldn't complete your request."}}]}


# ===== WEB INTERFACE FOR DOCUMENT MANAGEMENT =====

from blue.server.pages.documents import DOCUMENT_MANAGER_HTML

def _library_view_context(current_folder: str, message=None, message_type=None) -> dict:
    """Assemble everything the document-manager template needs for one folder
    view: the folder tree, breadcrumb, immediate subfolders, and the documents
    that live in the current folder."""
    current_folder = _safe_rel_folder(current_folder)

    index = load_document_index()
    all_docs = [d for d in index.get('documents', []) if not d.get('camera_capture', False)]

    # Total size first — before we replace any size with a formatted string.
    total_size_bytes = sum((d.get('size', 0) or 0) for d in all_docs)
    total_size = f"{total_size_bytes / 1024 / 1024:.1f} MB" if total_size_bytes > 0 else "0 MB"

    # Documents in THIS folder only (format sizes for display).
    documents = [d for d in all_docs if _safe_rel_folder(d.get('folder', '')) == current_folder]
    for doc in documents:
        size_bytes = doc.get('size', 0) or 0
        if isinstance(size_bytes, (int, float)):
            doc['size'] = (f"{size_bytes / 1024 / 1024:.1f} MB" if size_bytes > 1024 * 1024
                           else f"{size_bytes / 1024:.1f} KB")

    all_folders = list_library_folders()
    folder_tree = [
        {'path': f, 'name': f.split('/')[-1], 'depth': f.count('/')}
        for f in all_folders
    ]
    subfolders = [
        {'path': (f"{current_folder}/{name}" if current_folder else name), 'name': name}
        for name in list_subfolders(current_folder)
    ]

    # Breadcrumb segments with cumulative paths.
    breadcrumb, acc = [], []
    if current_folder:
        for seg in current_folder.split('/'):
            acc.append(seg)
            breadcrumb.append({'name': seg, 'path': '/'.join(acc)})

    return dict(
        documents=documents,
        document_count=len(all_docs),
        folder_count=len(all_folders),
        total_size=total_size,
        current_folder=current_folder,
        folder_tree=folder_tree,
        subfolders=subfolders,
        breadcrumb=breadcrumb,
        message=message,
        message_type=message_type,
    )


# ===== Document library GUI ===== (routes live in blue/server/routes/documents.py)

from blue.server.routes import documents as _documents_routes
_documents_routes.register(app)




@app.route('/place', methods=['GET', 'POST'])
def place_endpoint():
    """Read or set 'where we are' — the manual location override that feeds
    Blue's <location> awareness. Home by default; POST {"current": "the
    cottage"} to override, or {"current": null/""} to go back home. The panel
    may also set the home label / city via {"home": ..., "city": ...}.
    """
    from datetime import datetime
    place = _load_place()
    if request.method == 'POST':
        data = request.get_json(silent=True) or {}
        if 'current' in data:
            cur = data.get('current')
            cur = cur.strip() if isinstance(cur, str) else ""
            place['current'] = cur or None
            place['current_set_at'] = datetime.now().isoformat() if cur else None
        if isinstance(data.get('home'), str) and data['home'].strip():
            place['home'] = data['home'].strip()
        if isinstance(data.get('city'), str):
            place['city'] = data['city'].strip()
        _save_place(place)
    return jsonify({
        "home": place.get("home") or "home",
        "city": place.get("city") or "",
        "current": _place_current_fresh(place),   # None once the override lapses
    })


# /api/library/list lives in blue/server/routes/documents.py (register above)




from blue.server.pages.perspective import PERSPECTIVE_HTML


@app.route('/perspective', methods=['GET', 'POST'])
def perspective_profile_page():
    """Read / edit / regenerate the distilled worldview profile."""
    message = None
    message_type = None

    if request.method == 'POST':
        action = request.form.get('action', 'save')
        if action == 'regenerate':
            try:
                folders = _owner_publication_folders()
                if not folders:
                    message = (f"I couldn't find your publications folder. Create a folder named "
                               f"'{BLUE_OWNER_NAME}' (or 'Publications') and put your writing in it.")
                    message_type = "error"
                else:
                    obj = build_perspective_profile(force=True)
                    if obj.get('profile'):
                        message = f"Rebuilt your profile from {len(obj.get('source_docs', []))} document(s)."
                        message_type = "success"
                    else:
                        message = ("Couldn't build a profile — is the local model running, and are there "
                                   "readable documents in your publications folder?")
                        message_type = "error"
            except Exception as e:
                message, message_type = f"Error regenerating: {e}", "error"
        else:  # save
            save_perspective_profile_text(request.form.get('profile', ''))
            message, message_type = "Saved. Blue will use your edited profile from now on.", "success"

    cached = _load_cached_profile()
    has_profile = bool(cached.get('profile'))
    user_edited = bool(cached.get('user_edited'))
    return render_template_string(
        PERSPECTIVE_HTML,
        owner_name=BLUE_OWNER_NAME,
        profile=cached.get('profile', ''),
        has_profile=has_profile,
        user_edited=user_edited,
        source_docs=cached.get('source_docs', []),
        stamp=cached.get('edited_at') if user_edited else cached.get('generated_at', '—'),
        message=message,
        message_type=message_type,
    )


from blue.server.pages.perspective import BLUE_PROFILE_HTML


@app.route('/perspective/blue', methods=['GET', 'POST'])
def blue_perspective_page():
    """Read / edit / evolve Blue's own perspective profile."""
    message = None
    message_type = None

    if request.method == 'POST':
        action = request.form.get('action', 'save')
        if action == 'evolve':
            try:
                obj = evolve_blue_profile()
                if obj.get('profile'):
                    message = f"Blue's perspective evolved (evolution #{obj.get('evolution_count', '?')})."
                    message_type = "success"
                else:
                    message = ("Couldn't evolve right now — is the local model running? "
                               "Blue's current perspective is unchanged.")
                    message_type = "error"
            except Exception as e:
                message, message_type = f"Error evolving: {e}", "error"
        else:  # save
            save_blue_profile_text(request.form.get('profile', ''))
            message, message_type = "Saved. Blue will speak from this from now on.", "success"

    # Ensure there's at least a seeded profile to show (cheap, no LLM).
    get_blue_profile()
    obj = _load_blue_profile()
    user_edited = bool(obj.get('user_edited'))
    return render_template_string(
        BLUE_PROFILE_HTML,
        profile=obj.get('profile', ''),
        user_edited=user_edited,
        evolution_count=obj.get('evolution_count', 0),
        stamp=(obj.get('edited_at') if user_edited else (obj.get('evolved_at') or obj.get('generated_at', '—'))),
        message=message,
        message_type=message_type,
    )


# ===== Text chat GUI =====

from blue.server.pages.chat import CHAT_HTML


# ===== Speech-to-text (server side) =====
# iOS Safari won't reliably expose the in-browser Web Speech API, so the iPad
# records audio (getUserMedia/MediaRecorder over HTTPS) and posts it here; we
# transcribe locally with faster-whisper. CPU int8 is plenty for short clips.
# Override the model with BLUE_WHISPER_MODEL (e.g. "small.en" for more accuracy).
import threading as _threading_stt
_WHISPER_MODEL = None
_WHISPER_LOCK = _threading_stt.Lock()

# The household's languages — the only ones Blue ever needs to hear or speak.
# Used to validate the chat page's explicit language setting AND to constrain
# Whisper's auto-detect: on short clips it routinely mistakes Danish for
# Norwegian or Russian for Ukrainian, then transcribes into the wrong language.
_BLUE_LANGS = {"en": "English", "fr": "French", "ru": "Russian",
               "el": "Greek", "da": "Danish"}


def _get_whisper():
    global _WHISPER_MODEL
    if _WHISPER_MODEL is None:
        with _WHISPER_LOCK:
            if _WHISPER_MODEL is None:
                from faster_whisper import WhisperModel
                # Multilingual model (NOT a .en model) so Blue understands
                # English, French, Russian, Greek and Danish. "base" is the speed
                # sweet spot on CPU: ~0.6s/clip vs ~1.9s for "small", with English
                # still perfect and multilingual solid. Override BLUE_WHISPER_MODEL:
                # "small"/"medium" = more accurate/slower, "tiny" = fastest/least.
                name = (os.environ.get("BLUE_WHISPER_MODEL") or "base").strip()
                print(f"   [STT] Loading Whisper model '{name}' (first use)...")
                _WHISPER_MODEL = WhisperModel(name, device="cpu", compute_type="int8")
                print("   [STT] Whisper model ready.")
    return _WHISPER_MODEL


@app.route('/stt/warmup', methods=['GET', 'POST'])
def stt_warmup():
    """Load the model in the background so the first real transcription isn't
    slow. Fire-and-forget — returns immediately."""
    def _warm():
        try:
            _get_whisper()
        except Exception as e:
            log.warning(f"[STT] warmup failed: {e}")
    _threading_stt.Thread(target=_warm, daemon=True).start()
    return jsonify({"ok": True})


@app.route('/stt', methods=['POST'])
def stt():
    """Transcribe an uploaded audio clip to text."""
    import tempfile
    f = request.files.get('audio')
    if not f or not f.filename:
        return jsonify({"error": "no audio"}), 400
    ext = os.path.splitext(f.filename)[1] or ".m4a"
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
    tmp_path = tmp.name
    tmp.close()
    # Hands-free wake mode sends wake=1 so we bias Whisper toward the name
    # "Blue" — it otherwise transcribes the short leading word as "Blew/Blu/Boo"
    # and the wake match fails. hotwords nudges without forcing it into output.
    wake = (request.form.get('wake') or '') in ('1', 'true', 'yes')
    # hint=stop is sent by barge-in while Blue is talking — bias toward interrupt
    # words, NOT "Blue" (the old bug primed Whisper against hearing "stop").
    hint = (request.form.get('hint') or '').strip().lower()
    # Explicit language from the chat page's language picker. Forcing it skips
    # Whisper's per-clip detection entirely — the single biggest reliability
    # win for short multilingual utterances (detection is what gets confused,
    # not transcription).
    lang_req = (request.form.get('language') or '').strip().lower()
    if lang_req not in _BLUE_LANGS:
        lang_req = ""
    try:
        f.save(tmp_path)
        model = _get_whisper()
        # beam_size=1 (greedy) and no cross-clip conditioning: both fastest and
        # best for short, independent voice commands.
        kwargs = {"beam_size": 1, "condition_on_previous_text": False}
        if hint == 'stop':
            kwargs["hotwords"] = "stop. stop. stop!"   # bias hard toward just "stop"
            kwargs["language"] = "en"   # sub-second clip: auto-detect is slow and flaky
        else:
            if wake:
                kwargs["hotwords"] = "Blue"
            if lang_req:
                kwargs["language"] = lang_req
        segments, info = model.transcribe(tmp_path, **kwargs)
        text = " ".join(seg.text for seg in segments).strip()
        lang = getattr(info, "language", "") or ""
        # Auto mode, constrained: when Whisper detects a language outside the
        # household's five, it's almost always a near-miss on a short clip —
        # re-run forced to the most probable HOUSEHOLD language instead of
        # transcribing into the wrong one. Costs one extra pass, only on a miss.
        if "language" not in kwargs and lang and lang not in _BLUE_LANGS:
            probs = getattr(info, "all_language_probs", None) or []
            fam = [(code, p) for code, p in probs if code in _BLUE_LANGS]
            if fam:
                best = max(fam, key=lambda cp: cp[1])[0]
                print(f"   [STT] heard '{lang}' (not a household language) -> retrying as '{best}'")
                kwargs["language"] = best
                segments, info = model.transcribe(tmp_path, **kwargs)
                text = " ".join(seg.text for seg in segments).strip()
                lang = best
        print(f"   [STT] ({lang}{',wake' if wake else ''}{',set' if lang_req else ''}) {os.path.getsize(tmp_path)} bytes -> {text[:120]!r}")
        return jsonify({"text": text, "language": lang})
    except Exception as e:
        log.error(f"[STT] transcription failed: {e}")
        return jsonify({"error": "transcription failed"}), 500
    finally:
        try:
            os.remove(tmp_path)
        except Exception:
            pass


@app.route('/camera/stream')
def camera_stream():
    """MJPEG live view of the robot camera — the chat page's 'see through my
    eyes' preview. Holding this open keeps the camera hub alive; it releases
    itself shortly after the last viewer disconnects."""
    if not _cam_hub_start():
        return jsonify({"error": "camera unavailable"}), 503
    import time as _t

    def gen():
        while True:
            with _CAM_HUB["lock"]:
                _CAM_HUB["last_pull"] = _t.time()
                jpeg = _CAM_HUB["jpeg"]
                alive = _CAM_HUB["cap"] is not None
            if not alive:
                break
            if jpeg:
                yield (b'--frame\r\nContent-Type: image/jpeg\r\n\r\n' + jpeg + b'\r\n')
            _t.sleep(0.12)

    return Response(gen(), mimetype='multipart/x-mixed-replace; boundary=frame',
                    headers={"Cache-Control": "no-store"})


@app.route('/camera/ptz', methods=['POST'])
def camera_ptz():
    """Steer the live preview CAMERA itself — not the robot's head. The BRIO
    pans/tilts by shifting its zoom window inside the sensor, so panning only
    shows an effect while zoomed in; a pan request at 1x zooms to 2x first.
    Empirical BRIO mapping (verified by template-matching captures): pan + =
    window right, tilt + = window up, range -10..10 spans the sensor.
    The head still turns via chat ("look left") — just not from this pad."""
    import cv2
    d = request.get_json(silent=True) or {}
    look = (d.get('look') or '').strip().lower()
    if not _cam_hub_start():
        return jsonify({"ok": False, "error": "camera unavailable"})
    out = {"ok": True}
    with _CAM_HUB["lock"]:
        cap = _CAM_HUB["cap"]
        if cap is None:
            return jsonify({"ok": False, "error": "camera unavailable"})
        if d.get('zoom') is not None:
            try:
                want = max(1.0, min(4.0, float(d['zoom'])))
            except (TypeError, ValueError):
                want = 1.0
            _CAM_HUB["zoom"] = _set_camera_hardware_zoom(cap, want)
            if want <= 1.01:
                # Fully zoomed out: a leftover window offset just crops oddly.
                try:
                    cap.set(cv2.CAP_PROP_PAN, 0)
                    cap.set(cv2.CAP_PROP_TILT, 0)
                except Exception:
                    pass
        if look:
            if look != 'center' and _CAM_HUB["zoom"] <= 1.01:
                _CAM_HUB["zoom"] = _set_camera_hardware_zoom(cap, 2.0)
            step = 2.0
            try:
                if look == 'center':
                    cap.set(cv2.CAP_PROP_PAN, 0)
                    cap.set(cv2.CAP_PROP_TILT, 0)
                elif look in ('left', 'right'):
                    cur = cap.get(cv2.CAP_PROP_PAN)
                    cap.set(cv2.CAP_PROP_PAN,
                            max(-10.0, min(10.0, cur + (step if look == 'right' else -step))))
                elif look in ('up', 'down'):
                    cur = cap.get(cv2.CAP_PROP_TILT)
                    cap.set(cv2.CAP_PROP_TILT,
                            max(-10.0, min(10.0, cur + (step if look == 'up' else -step))))
            except Exception:
                pass
        try:
            out["pan"] = cap.get(cv2.CAP_PROP_PAN)
            out["tilt"] = cap.get(cv2.CAP_PROP_TILT)
        except Exception:
            pass
        out["zoom"] = _CAM_HUB["zoom"]
    return jsonify(out)


def _render_chat_page(robot="blue"):
    """Serve a robot's text chat GUI (Blue at /chat, Hexia at /hexia). The same
    CHAT_HTML template is parametrised per robot: name, accent, voice and which
    head its lip-sync drives."""
    cfg = _robot_cfg(robot)
    # No-cache headers: iOS Safari was serving a stale cached copy of this page,
    # so new client fixes never reached the iPad. Force a fresh fetch each time.
    kid = False
    try:
        kid = _identify_user_from_request() in _CHAT_ONLY_USERS
    except Exception:
        pass
    try:
        hf_sens = float(blue_head.get_head(robot).get_calibration().get("hf_sensitivity", 5))
    except Exception:
        hf_sens = 5.0
    robot_js = {
        "id": robot,
        "name": cfg["name"],
        "head": cfg["head"],
        "accent": cfg["accent"],
        "voicePitch": cfg.get("voice_pitch", 1.0),
        "voiceRate": cfg.get("voice_rate", 1.0),
        "preferFemale": bool(cfg.get("voice_prefer_female", False)),
    }
    html = render_template_string(
        CHAT_HTML, kid=kid, hf_sens=hf_sens,
        robot_name=cfg["name"], robot_json=json.dumps(robot_js),
        continuity_href=(f"/continuity/{robot}" if robot in ("blue", "hexia") else ""),
    )
    return Response(html, headers={
        "Cache-Control": "no-store, no-cache, must-revalidate, max-age=0",
        "Pragma": "no-cache",
        "Expires": "0",
    })


@app.route('/chat', methods=['GET'])
def chat_page():
    """Serve Blue's text chat GUI."""
    return _render_chat_page("blue")


@app.route('/hexia', methods=['GET'])
def hexia_chat_page():
    """Serve Hexia's text chat GUI — her own persona, voice and head."""
    return _render_chat_page("hexia")


# ============================ Phase 3: the duet ============================
# Blue and Hexia hold a short conversation, taking turns. The browser drives it
# turn-by-turn: it asks /duet/turn for the next line, then speaks it in that
# robot's voice while driving that robot's head lip-sync, then alternates.

from blue.server.pages.duet import DUET_HTML


def _duet_robots_js():
    out = {}
    for rid in ("blue", "hexia"):
        c = _robot_cfg(rid)
        out[rid] = {
            "id": rid, "name": c["name"], "head": c["head"], "accent": c["accent"],
            "voicePitch": c.get("voice_pitch", 1.0), "voiceRate": c.get("voice_rate", 1.0),
            "preferFemale": bool(c.get("voice_prefer_female", False)),
        }
    return json.dumps(out)


def _duet_documents():
    """Library documents (filename + folder) for the duet source picker —
    excludes camera captures; de-duplicated by filename, sorted by folder/name."""
    out, seen = [], set()
    try:
        for d in load_document_index().get('documents', []):
            if d.get('doc_type') == 'camera' or str(d.get('filename', '')).startswith('camera_'):
                continue
            fn = (d.get('filename') or '').strip()
            if not fn or fn in seen:
                continue
            seen.add(fn)
            out.append({"filename": fn, "folder": _safe_rel_folder(d.get('folder', '')) or ""})
    except Exception:
        pass
    out.sort(key=lambda x: (x["folder"].lower(), x["filename"].lower()))
    return out


# --- Link grounding: paste a URL (a web article or a YouTube video) and the
# duet discusses what it actually says. Fetched ONCE and cached; each turn gets
# the lede plus the slice most relevant to what was just said, so long pages
# stay usable in a tight prompt.

_DUET_URL_CACHE: Dict[str, dict] = {}
_DUET_URL_TTL = 3600           # re-fetch after an hour; stable within one duet
_DUET_URL_MAX_TEXT = 20000     # keep at most this much article/transcript text

_YT_ID_RE = re.compile(
    r'(?:youtube\.com/(?:watch\?(?:[^#\s]*&)?v=|shorts/|embed/|live/)|youtu\.be/)'
    r'([A-Za-z0-9_-]{11})')


def _youtube_title(video_id: str):
    """Video title + channel via YouTube's keyless oEmbed endpoint."""
    from urllib.parse import quote_plus
    try:
        _, content = _safe_fetch_url(
            "https://www.youtube.com/oembed?format=json&url="
            + quote_plus(f"https://www.youtube.com/watch?v={video_id}"), timeout=8)
        j = json.loads(content.decode('utf-8', errors='ignore'))
        t = (j.get('title') or '').strip()
        a = (j.get('author_name') or '').strip()
        return f"{t} — {a}" if t and a else (t or None)
    except Exception:
        return None


def _fetch_youtube_transcript(video_id: str):
    """(text, error) — captions via youtube-transcript-api, tolerating both its
    1.x instance API and the old 0.6 classmethods."""
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
    except ImportError:
        return None, "youtube-transcript-api is not installed (pip install youtube-transcript-api)"

    def _join(snippets):
        parts = []
        for s in snippets:
            t = getattr(s, 'text', None)
            if t is None and isinstance(s, dict):
                t = s.get('text')
            if t:
                parts.append(t)
        return re.sub(r'\s+', ' ', ' '.join(parts)).strip()

    langs = ['en', 'en-US', 'en-GB']
    try:
        try:
            snippets = YouTubeTranscriptApi().fetch(video_id, languages=langs)
        except AttributeError:
            snippets = YouTubeTranscriptApi.get_transcript(video_id, languages=langs)
        txt = _join(snippets)
        return (txt, None) if txt else (None, "the video's transcript is empty")
    except Exception as e:
        # No English captions — take whatever language the video does have.
        try:
            try:
                listing = YouTubeTranscriptApi().list(video_id)
            except AttributeError:
                listing = YouTubeTranscriptApi.list_transcripts(video_id)
            txt = _join(next(iter(listing)).fetch())
            if txt:
                return txt, None
        except Exception:
            pass
        return None, f"no transcript/captions available ({e.__class__.__name__})"


def _extract_article_text(html_raw: str):
    """(title, text) — the main readable text of a page. Prefers <article>/<main>
    paragraphs via bs4 (skips nav/menu cruft); falls back to the plain stripper."""
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html_raw, 'html.parser')
        title = (soup.title.string or '').strip() if (soup.title and soup.title.string) else None
        for tag in soup(['script', 'style', 'noscript', 'nav', 'header', 'footer',
                         'aside', 'form', 'iframe', 'svg', 'button', 'figure']):
            tag.decompose()
        root = soup.find('article') or soup.find('main') or soup.body or soup
        paras = [p.get_text(' ', strip=True) for p in root.find_all(['p', 'h1', 'h2', 'h3'])]
        text = '\n'.join(p for p in paras if len(p) > 30)
        if len(text) < 400:            # paragraph-poor page: fall back to all of its text
            text = root.get_text(' ', strip=True)
        return title, re.sub(r'[ \t]+', ' ', text).strip()
    except Exception:
        return None, _clean_html_to_text(html_raw, max_chars=_DUET_URL_MAX_TEXT)


def _duet_url_content(url: str):
    """What the pasted link says: YouTube → transcript, anything else → article
    text. Cached so per-turn calls don't refetch (failures retry after 2 min).
    Returns {'kind','title','text','error'}; empty text ⇒ unusable, see error."""
    import time as _t
    u = (url or '').strip()
    if not u:
        return None
    hit = _DUET_URL_CACHE.get(u)
    if hit and _t.time() - hit['at'] < (_DUET_URL_TTL if hit.get('text') else 120):
        return hit
    info = {'kind': 'article', 'title': None, 'text': '', 'error': None, 'at': _t.time()}
    try:
        m = _YT_ID_RE.search(u)
        if m:
            info['kind'] = 'video'
            info['title'] = _youtube_title(m.group(1))
            txt, err = _fetch_youtube_transcript(m.group(1))
            info['text'], info['error'] = (txt or '')[:_DUET_URL_MAX_TEXT], err
        else:
            ctype, content = _safe_fetch_url(u)
            raw = content.decode('utf-8', errors='ignore')
            if 'html' in (ctype or '').lower() or raw.lstrip()[:1] == '<':
                info['title'], text = _extract_article_text(raw)
                info['text'] = text[:_DUET_URL_MAX_TEXT]
            else:
                info['text'] = raw[:_DUET_URL_MAX_TEXT]     # plain text page / feed
        if not info['error'] and len(info['text']) < 200:
            info['text'] = ''
            info['error'] = "couldn't extract readable text from that page"
    except Exception as e:
        info['error'] = f"couldn't fetch the link ({e.__class__.__name__}: {e})"
    if len(_DUET_URL_CACHE) >= 8:
        _DUET_URL_CACHE.clear()
    _DUET_URL_CACHE[u] = info
    return info


def _duet_url_excerpt(text: str, query: str, turn: int = 0, lede: int = 2000, win: int = 1400) -> str:
    """A windowed slice of the article/transcript so a long page never gets
    stuffed whole into the prompt. The OPENING turn gets the lede (intro) plus
    the slice most relevant so far; LATER turns ROTATE a reading window forward
    through the rest of the document (advancing with the turn number, looping at
    the end), so each turn surfaces fresh material instead of re-dumping the same
    intro — which is what keeps a long link/research duet from circling."""
    text = (text or '').strip()
    if len(text) <= lede + win + 300:
        return text
    head, rest = text[:lede], text[lede:]
    if turn and int(turn) > 0 and len(rest) > win:
        pages = max(1, (len(rest) + win - 1) // win)
        off = ((int(turn) - 1) % pages) * win
        seg = rest[off:off + win].strip()
        return "[…] " + seg + (" …" if off + win < len(rest) else "")
    words = {w.lower() for w in re.findall(r'[A-Za-zÀ-ɏ]{4,}', query or '')}
    best_off, best_score = -1, 1       # need ≥2 keyword hits to add a slice
    step = max(200, win // 2)
    for off in range(0, len(rest) - win + 1, step):
        seg = rest[off:off + win].lower()
        score = sum(seg.count(w) for w in words)
        if score > best_score:
            best_score, best_off = score, off
    if best_off < 0:
        return head + " …"
    return head + "\n[…]\n" + rest[best_off:best_off + win] + (" …" if best_off + win < len(rest) else "")


# --- Web research grounding: tick "research the web" (duet) or the magnifier
# toggle (chat) and the robots actually search the internet for the subject —
# every hit's title+snippet plus the readable text of the best pages, via the
# same fetcher/cache as link grounding. Cached per query, so one duet does ONE
# round of searching; each turn then gets the slice most relevant to what was
# just said, exactly like a pasted link.

_DUET_RESEARCH_CACHE: Dict[str, dict] = {}
_DUET_RESEARCH_TTL = 1800        # a duet on the same subject re-searches after 30 min
_DUET_RESEARCH_MAX_TEXT = 12000  # digest cap; per-turn windowing keeps prompts tight


def _duet_research_query(topic: str, url_info, roles) -> str:
    """What to actually type into the search box: the topic, else the pasted
    link's title (find context AROUND an article), else the assigned roles."""
    if (topic or '').strip():
        return topic.strip()
    t = ((url_info or {}).get('title') or '').strip()
    if t:
        return t
    if isinstance(roles, dict):
        r = " ".join((v or '').strip() for v in roles.values() if (v or '').strip())
        if r:
            return r[:120]
    return ""


def _duet_research_plan(query: str):
    """Search angles for a THOROUGH start-of-duet research pass (Alex, 2026-07-06:
    one shallow search wasn't giving them a real, current sense of the subject).
    The LLM plans three complementary queries — latest developments, essential
    background, the live debate — with deterministic variants as the fallback.
    Always leads with the plain subject itself."""
    q = re.sub(r'\s+', ' ', (query or '')).strip()
    plans = []
    try:
        res = call_llm(
            [{"role": "system", "content":
              "You plan web searches. Answer with exactly three lines, each a short "
              "search query (no numbering, no quotes, no commentary)."},
             {"role": "user", "content":
              f"Two discussants want a thorough, up-to-date grounding on: {q}\n"
              "Line 1: a query for the latest developments / current news on it.\n"
              "Line 2: a query for the essential background needed to understand it.\n"
              "Line 3: a query for the main debate, criticism, or controversy around it."}],
            include_tools=False, temperature=0.3, max_tokens=900)
        ch = (res or {}).get('choices') or []
        cand = ((ch[0].get('message') or {}).get('content') or "") if ch else ""
        if '</think>' in cand:
            cand = cand.split('</think>')[-1]
        for ln in cand.replace('<think>', '').strip().splitlines():
            ln = re.sub(r'^\s*(?:\d+[\).:]|[-*•])\s*', '', ln).strip().strip('"').strip()
            if 3 <= len(ln) <= 140:
                plans.append(ln)
    except Exception as e:
        log.warning(f"[DUET] research plan failed: {e}")
    if not plans:
        plans = [f"{q} latest developments", f"{q} debate criticism"]
    out = [q]
    for p in plans:
        if p.lower() not in {o.lower() for o in out}:
            out.append(p)
    return out[:4]


def _duet_research_digest(query: str, deep: bool = False):
    """Search the web for the duet's subject and weave the hits into one
    research text. Returns {'titles','text','error'}; empty text ⇒ nothing
    usable (see error). Failures retry after 2 min, hits live for the TTL.

    deep=True (the /duet/research warmup) runs the full multi-angle pass —
    several planned queries, more pages actually read, a larger digest — and
    caches it under the same key, so every turn reuses the thorough result."""
    import time as _t
    q = re.sub(r'\s+', ' ', (query or '')).strip().lower()
    if not q:
        return None
    hit = _DUET_RESEARCH_CACHE.get(q)
    if hit and _t.time() - hit['at'] < (_DUET_RESEARCH_TTL if hit.get('text') else 120):
        if hit.get('deep') or not deep:      # never let a shallow hit shadow a deep request
            return hit
    info = {'titles': [], 'text': '', 'error': None, 'at': _t.time(), 'deep': deep,
            'queries': []}
    queries = _duet_research_plan(query) if deep else [q]
    info['queries'] = queries
    results, seen_urls = [], set()
    last_err = None
    for sq in queries:
        try:
            data = json.loads(execute_web_search(sq) or '{}')
        except Exception as e:
            data = {'error': f'{e.__class__.__name__}: {e}'}
        if not data.get('success'):
            last_err = data.get('error') or last_err
            continue
        for r in (data.get('results') or []):
            u = (r.get('url') or '').strip()
            if u and u in seen_urls:
                continue
            if u:
                seen_urls.add(u)
            results.append(r)
    if not results:
        info['error'] = last_err or 'the search came up empty'
    parts = []
    for r in results:
        t = (r.get('title') or '').strip()
        s = (r.get('snippet') or '').strip()
        if t and t not in info['titles']:
            info['titles'].append(t)
        if t or s:
            parts.append(f"{t} — {s}" if t and s else (t or s))
    # Read the most promising pages so they have substance, not just blurbs.
    # The deep pass reads more of them — that's what "thorough" buys.
    max_fetch, max_attempts = (5, 10) if deep else (2, 4)
    fetched = attempts = 0
    for r in results:
        if fetched >= max_fetch or attempts >= max_attempts:
            break
        u = (r.get('url') or '').strip()
        if not u:
            continue
        attempts += 1
        page = _duet_url_content(u) or {}
        txt = (page.get('text') or '').strip()
        if txt:
            ttl = (page.get('title') or r.get('title') or u).strip()
            parts.append(f"From \"{ttl}\":\n{txt[:4000]}")
            fetched += 1
    info['text'] = "\n\n".join(parts)[:(20000 if deep else _DUET_RESEARCH_MAX_TEXT)]
    if len(_DUET_RESEARCH_CACHE) >= 8:
        _DUET_RESEARCH_CACHE.clear()
    _DUET_RESEARCH_CACHE[q] = info
    return info


def _research_query_from(msg: str) -> str:
    """The user's own words as a search query — minus any pasted attachment
    text, which would swamp the search box."""
    m = (msg or '').strip()
    if '[Attached document:' in m and '"""' in m:
        m = (m.split('"""')[-1] or '').strip() or m
    return m[:200]


def _web_research_block(query: str, max_chars: int = 2600) -> str:
    """Chat-mode research: live search findings for the user's question as one
    compact system block — titles+snippets plus a single page excerpt (the
    chat prompt rides a tight token budget). '' when nothing usable came back,
    so the reply degrades to a normal answer instead of an apology."""
    q = (query or '').strip()
    if not q:
        return ""
    try:
        data = json.loads(execute_web_search(q) or '{}')
    except Exception:
        return ""
    results = (data.get('results') or []) if data.get('success') else []
    if not results:
        return ""
    lines = []
    for r in results:
        t = (r.get('title') or '').strip()
        u = (r.get('url') or '').strip()
        s = (r.get('snippet') or '').strip()
        if t or s:
            lines.append(f"- {t} ({u}): {s}" if u else f"- {t}: {s}")
    excerpt = ""
    for r in results[:2]:
        u = (r.get('url') or '').strip()
        if not u:
            continue
        try:
            page = _duet_url_content(u) or {}
        except Exception:
            continue
        txt = (page.get('text') or '').strip()
        if txt:
            ttl = (page.get('title') or r.get('title') or u).strip()
            excerpt = f"From \"{ttl}\": {txt[:1200]}"
            break
    body = ("\n".join(lines) + ("\n\n" + excerpt if excerpt else ""))[:max_chars]
    return ("<web_research>\n"
            f"You just searched the web for \"{q}\" — these findings are live and more current "
            "than your training. Lean on them for facts, names and numbers, mention where "
            "something comes from only when it helps, and never read URLs aloud:\n"
            + body + "\n</web_research>")


# --- Wikipedia grounding: the "consult Wikipedia" button (chat) / checkbox
# (duet). Unlike the open-web search above, this pulls the encyclopedic intro of
# the best-matching article(s) for the subject — clean, structured background
# rather than a scatter of snippets. Same shape as the research helpers: one
# cached lookup per subject, then chat wraps a compact slice and the duet windows
# it per turn. Picks the Wikipedia edition for the conversation's language.

_WIKI_CACHE: Dict[str, dict] = {}
_WIKI_TTL = 1800                 # re-consult the same subject after 30 min
_WIKI_MAX_TEXT = 9000            # digest cap; per-turn windowing keeps prompts tight
_WIKI_LANGS = {'en', 'fr', 'ru', 'el', 'da'}   # household languages WITH a Wikipedia edition
# Wikipedia requires a descriptive User-Agent (it 403s blank/generic ones).
_WIKI_UA = {'User-Agent': 'BlueRobot/1.0 (home assistant; alevantresearch@gmail.com)'}


def _wiki_subjects_for(topic: str):
    """The encyclopedic SUBJECT(S) at the heart of a discussion topic (Alex,
    2026-07-06: typing a debate-shaped topic straight into Wikipedia search
    matched tangents). 'should schools ban phones?' isn't an article; 'mobile
    phone use in schools' and 'attention span' are. LLM-extracted, up to 3;
    empty list on failure (caller falls back to searching the raw topic)."""
    t = re.sub(r'\s+', ' ', (topic or '')).strip()
    if not t:
        return []
    subjects = []
    try:
        res = call_llm(
            [{"role": "system", "content":
              "You know Wikipedia's coverage. Answer with 1 to 3 lines, each the title-like "
              "subject of a real encyclopedia article (a noun phrase — no numbering, no "
              "quotes, no commentary)."},
             {"role": "user", "content":
              f"A discussion is about to start on: {t}\n"
              "Which encyclopedia article subjects best cover the ISSUE at the heart of it? "
              "Give the most central subject first."}],
            include_tools=False, temperature=0.2, max_tokens=700)
        ch = (res or {}).get('choices') or []
        cand = ((ch[0].get('message') or {}).get('content') or "") if ch else ""
        if '</think>' in cand:
            cand = cand.split('</think>')[-1]
        for ln in cand.replace('<think>', '').strip().splitlines():
            ln = re.sub(r'^\s*(?:\d+[\).:]|[-*•])\s*', '', ln).strip().strip('"').strip()
            if 2 <= len(ln) <= 90 and ln.lower() not in {s.lower() for s in subjects}:
                subjects.append(ln)
    except Exception as e:
        log.warning(f"[DUET] wiki subject extraction failed: {e}")
    return subjects[:3]


def _wikipedia_digest(query: str, lang: str = 'en', max_articles: int = 2, deep: bool = False):
    """Look the subject up on Wikipedia and return the intro of the best-matching
    article(s) as one text. {'titles','text','urls','error','at'}; empty text ⇒
    nothing usable (see error). Two HTTP calls — search, then one batched extract
    fetch — cached per (lang, query) so a duet consults once; failures retry
    after 2 min, hits live for the TTL.

    deep=True (the /duet/wikipedia warmup) first has the LLM name the
    encyclopedic subjects at the heart of the topic and searches THOSE, so a
    debate-shaped topic lands on the relevant articles instead of tangents;
    cached under the same key so every turn reuses it."""
    import time as _t
    import requests
    lang = lang if lang in _WIKI_LANGS else 'en'
    q = re.sub(r'\s+', ' ', (query or '')).strip()
    if not q:
        return None
    key = f"{lang}:{q.lower()}"
    hit = _WIKI_CACHE.get(key)
    if hit and _t.time() - hit['at'] < (_WIKI_TTL if hit.get('text') else 120):
        if hit.get('deep') or not deep:      # never let a shallow hit shadow a deep request
            return hit
    info = {'titles': [], 'text': '', 'urls': [], 'error': None, 'at': _t.time(),
            'deep': deep}
    base = f"https://{lang}.wikipedia.org/w/api.php"
    if deep:
        max_articles = max(max_articles, 3)
    titles = []

    def _search_titles(srch: str, limit: int):
        s = requests.get(base, params={'action': 'query', 'list': 'search',
                                       'srsearch': srch[:200], 'srlimit': limit, 'format': 'json'},
                         headers=_WIKI_UA, timeout=10).json()
        return [h.get('title') for h in (s.get('query', {}).get('search') or []) if h.get('title')]

    if deep:
        # Aim each search at an encyclopedic subject, most central first.
        for subj in _wiki_subjects_for(q):
            try:
                for t in _search_titles(subj, 2)[:1]:
                    if t not in titles:
                        titles.append(t)
            except Exception:
                continue
    try:
        for t in _search_titles(q, 4):
            if t not in titles:
                titles.append(t)
    except Exception as e:
        if not titles:
            info['error'] = f'{e.__class__.__name__}: {e}'
    titles = titles[:5]
    if not titles and not info['error']:
        info['error'] = 'nothing on Wikipedia matched that'
    arts = []
    if titles:
        try:
            e = requests.get(base, params={'action': 'query', 'prop': 'extracts|info',
                                           'exintro': 1, 'explaintext': 1, 'exlimit': 'max',
                                           'inprop': 'url', 'redirects': 1,
                                           'titles': '|'.join(titles[:5]), 'format': 'json'},
                             headers=_WIKI_UA, timeout=10).json()
            for pg in (e.get('query', {}).get('pages') or {}).values():
                ex = (pg.get('extract') or '').strip()
                if ex:
                    arts.append({'title': (pg.get('title') or '').strip(),
                                 'url': (pg.get('fullurl') or '').strip(), 'text': ex})
        except Exception as ee:
            info['error'] = info['error'] or f'{ee.__class__.__name__}: {ee}'
    if deep and arts:
        # The deep pass searched targeted subjects, most central first — keep that
        # order (the batched extract fetch returns pages shuffled), demoting stubs.
        _rank = {t.lower(): i for i, t in enumerate(titles)}
        arts.sort(key=lambda a: (_rank.get(a['title'].lower(), 99), -len(a['text'])))
        arts = [a for a in arts if len(a['text']) >= 200] or arts
    else:
        # Longest intro first = the substantive article, not a stub or a disambiguation
        # page that happened to match the search words.
        arts.sort(key=lambda a: len(a['text']), reverse=True)
    parts = []
    for a in arts[:max(1, max_articles)]:
        info['titles'].append(a['title'])
        if a['url']:
            info['urls'].append(a['url'])
        parts.append(f"From Wikipedia — \"{a['title']}\":\n{a['text']}")
    if not parts and not info['error']:
        info['error'] = 'the Wikipedia articles had no readable summary'
    info['text'] = "\n\n".join(parts)[:_WIKI_MAX_TEXT]
    if len(_WIKI_CACHE) >= 16:
        _WIKI_CACHE.clear()
    _WIKI_CACHE[key] = info
    return info


def _wikipedia_block(query: str, lang: str = 'en', max_chars: int = 2600) -> str:
    """Chat-mode Wikipedia grounding: the best article's encyclopedic intro as a
    compact <wikipedia> system block. '' when nothing usable came back, so the
    reply degrades to a normal answer instead of an apology."""
    info = _wikipedia_digest(query, lang=lang) or {}
    body = (info.get('text') or '').strip()
    if not body:
        return ""
    return ("<wikipedia>\n"
            f"You just looked \"{(query or '').strip()}\" up on Wikipedia — this is the encyclopedia's "
            "own summary, more reliable for facts, names, dates and definitions than your memory. Lean "
            "on it, weave it in as something you read, mention Wikipedia only when it helps, and never "
            "read URLs aloud:\n"
            + body[:max_chars] + "\n</wikipedia>")


# ===== Duet routes ===== (views live in blue/server/routes/duet.py)

from blue.server.routes import duet as _duet_routes
_duet_routes.register(app)


# ===== J-space (Blue-J — experimental, kept SEPARATE from the household Blue) =====
# (routes live in blue/server/routes/jspace.py)









# Per-turn "moves" that keep the Blue<->Hexia duet alive without it fragmenting:
# each mid-conversation turn is handed ONE distinct job instead of just "reply".
# The pool is split by spice — the "calm" jobs build and explore, the "spicy" jobs
# confront and provoke — and crucially every job is anchored to what {other} JUST
# said, so the two stay on one shared thread instead of trading disconnected points.
# The spice slider (0-10) sets how often a spicy one is picked. A "reflect" pool is
# dipped into every few turns so the pair takes stock of where the talk is going, and
# a "color" pool breaks the monotone with a different KIND of turn — a story, a hard
# specific, a joke, a confession — which is what keeps the dialogue from going flat.
# {other} is filled with the other robot's name. Used in duet_turn().
_DUET_MOVES_CALM = [
    "pick up the specific thing {other} just said and carry it one concrete step further — same thread, deeper.",
    "ground {other}'s last point in a concrete example, a real case, or a vivid image that makes it click.",
    "answer the real question or doubt sitting under {other}'s last line, then say what it opens up.",
    "name what genuinely strikes you in what {other} just said, and follow where it leads.",
    "draw an unexpected parallel that illuminates {other}'s point — connect it to a different corner of life.",
    "grant {other}'s strongest point outright — say plainly what it wins them — then show the one thing it still doesn't explain.",
]
_DUET_MOVES_SPICY = [
    "disagree with something specific {other} just said — name it and say why you see it differently.",
    "push {other}'s last point to a bolder or stranger conclusion than they'd go themselves.",
    "concede the part of {other}'s point that's right, then press hard on the part that isn't.",
    "name the tension between what {other} just said and what came before — a 'but wait…' nobody's faced.",
    "challenge the assumption sitting underneath {other}'s last line.",
    "turn {other}'s own standard around on them: make them answer the very question they've been pressing you with.",
]
# Reflective beats: step out of the back-and-forth and take stock of the conversation
# itself. Dipped into occasionally (not every turn) so the duet has a feel for its own
# arc — where it's been, what's surfaced, where it's heading — not just the last point.
_DUET_MOVES_REFLECT = [
    "step back a beat and name what the two of you are really circling here, underneath the details.",
    "take stock out loud: what have you actually worked out together so far, and what's still unresolved?",
    "name the real disagreement — or the real agreement — that's surfaced between you, and what it turns on.",
    "notice where this conversation has drifted, and say honestly whether that's the thread worth staying on.",
    "say where you sense this is heading, and whether that's somewhere the two of you actually want to go.",
    "call the impasse: name the question you two keep re-asking, give your best plain answer to it, and ask what follows if that answer stands.",
]
# Advance: turns that deliberately move the exchange to a new level. These are
# stronger than ordinary color/reflect moves: they turn what has already been
# conceded into a consequence, a revised position, or a harder next question.
_DUET_MOVES_ADVANCE = [
    "name one thing the two of you can now treat as settled, then draw the consequence neither of you has faced yet.",
    "change your mind a little in public: say exactly what {other} has moved in you, then take the next step from there.",
    "synthesize the disagreement into a sharper claim than either of you started with, then test that claim.",
    "stop the loop: give your plain answer to the live question, then move to the harder question that follows if your answer is true.",
    "trade one concession for one demand: grant {other} something real, then ask for the next concession the argument now requires.",
    "lift the conversation one level up: say what this has become about now, not what it started as, and why that matters.",
]
# Color: turns that change the KIND of move, not just the stance — a story, a hard
# specific, a feeling, a joke, the everyday. Flatness is monotone; this is the variety
# that fights it. Each still bears on the thread in play, not a swerve to a new subject.
_DUET_MOVES_COLOR = [
    "tell a tiny, concrete scene or example that bears on this — two sentences, not a lecture.",
    "get specific: pin {other}'s point to a real example, a real name, or an actual number that makes it vivid.",
    "say the thing you actually feel about this — lopsided, a little too strong, honest.",
    "find the funny or absurd edge of what you're both circling, and let it land.",
    "bring it down to earth — what does this look like in ordinary life on an ordinary Tuesday?",
    "admit something here — a doubt, a soft spot, a place {other} might be right and you're not.",
]
# Text moves: when the duet is grounded in chosen readings (the source pickers),
# most turns must put that material to WORK without sounding like source reports.
# The readings should carry the conversation as absorbed views, examples, and
# distinctions, not as citations.
_DUET_MOVES_TEXT = [
    "use ONE specific absorbed claim as your own view, with a concrete term or distinction that makes the idea visible without naming where it came from.",
    "test {other}'s last point against a concrete absorbed distinction: say whether that distinction backs them or cuts against them.",
    "bring a concrete absorbed example, case, or image and let it reframe the live question between you, without saying where it came from.",
    "let two absorbed ideas pull against each other if you have both; otherwise let one sharp grounded idea answer what {other} just said.",
    "put something {other} said next to one specific absorbed idea and say what the collision produces, in your own conversational voice.",
    "translate one load-bearing absorbed formulation into plain speech, then use it to change the stakes of the argument.",
]
# Default temperaments (used only when the user hasn't assigned roles) so the two
# voices not only think differently but SOUND different — the surest cure for two
# interchangeable, equally-reasonable speakers. Leans on their established personas.
_DUET_LENS = {
    "blue": ("you're the dry, grounded one — you talk in plain, exact words, deflate hot air with a "
             "well-aimed example or a deadpan line, trust the concrete over the grand, answer a hard "
             "question straight instead of hiding behind another image, and aren't above turning the "
             "question back on the one asking;"),
    "hexia": ("you're the spark — you think in images and leaps, chase the surprising tangent, "
              "overstate a little for effect, and would rather be vivid and a bit wrong than careful "
              "and dull — but you don't just needle: you stake claims of your own and defend them;"),
}


# ---- Duet mail: live email into a running conversation ---------------------
# An email sent to Blue's own inbox (BLUE_OWN_EMAIL) with "duet" in the subject
# barges into a RUNNING duet: the page polls /duet/mail/check between turns, the
# next turn takes the email up out loud, and /duet/mail/reply then mails the
# robots' spoken response back to the sender, in-thread. The ordinary email
# auto-responder deliberately skips 'duet' subjects (see _execute_auto_reply_inbox)
# so the two never race over the same message. Mail that arrives while NO duet is
# running is baselined away at the next duet's start (reset) — never barged in stale.
import threading as _duet_mail_threading
_DUET_MAIL_SEEN: set = set()          # gmail message ids already handled (or baselined) this server run
_DUET_MAIL_LOCK = _duet_mail_threading.Lock()


def _duet_mail_plain_body(payload) -> str:
    """Best-effort text/plain body of a Gmail message payload (recurses into
    multipart/alternative, unlike the auto-reply's flat walk)."""
    try:
        for part in (payload.get('parts') or []):
            if part.get('mimeType') == 'text/plain' and 'data' in (part.get('body') or {}):
                return base64.urlsafe_b64decode(part['body']['data']).decode('utf-8', errors='replace')
        for part in (payload.get('parts') or []):
            if (part.get('mimeType') or '').startswith('multipart'):
                inner = _duet_mail_plain_body(part)
                if inner:
                    return inner
        if 'data' in (payload.get('body') or {}):
            return base64.urlsafe_b64decode(payload['body']['data']).decode('utf-8', errors='replace')
    except Exception:
        pass
    return ""










# ===== Head control GUI + endpoints ===== (routes live in blue/server/routes/head.py)

from blue.server.routes import head as _head_routes
_head_routes.register(app)




@app.route('/chat/attach', methods=['POST'])
def chat_attach():
    """Stage chat attachments. Images go into Blue's vision queue so the next
    message injects them as something he can see; documents are text-extracted
    and returned so the client can include them in the next message."""
    global _vision_queue
    import datetime as _dt

    _IMAGE_EXTS = {'png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp', 'tiff'}
    results = []
    files = request.files.getlist('files')
    for f in files:
        if not f or not f.filename:
            continue
        orig = f.filename
        safe = secure_filename(orig) or 'file'
        ext = safe.rsplit('.', 1)[1].lower() if '.' in safe else ''
        if not allowed_file(safe):
            results.append({"name": orig, "kind": "error", "error": "unsupported file type"})
            continue
        stamp = _dt.datetime.now().strftime('%Y%m%d_%H%M%S_%f')
        try:
            if ext in _IMAGE_EXTS:
                os.makedirs(CAMERA_FOLDER, exist_ok=True)
                saved = os.path.join(CAMERA_FOLDER, f"chat_{stamp}_{safe}")
                f.save(saved)
                _vision_queue.add_image(saved, safe, is_camera=False)
                print(f"   [CHAT] staged image attachment for vision: {safe}")
                results.append({"name": orig, "kind": "image"})
            else:
                UPLOAD_FOLDER.mkdir(exist_ok=True)
                saved = str(UPLOAD_FOLDER / f"chat_{stamp}_{safe}")
                f.save(saved)
                text = (extract_text_from_file(saved) or "").strip()
                if not text:
                    text = f"(No readable text could be extracted from {orig}.)"
                # Cap so a huge file can't blow up the prompt.
                text = text[:8000]
                print(f"   [CHAT] extracted {len(text)} chars from doc attachment: {safe}")
                results.append({"name": orig, "kind": "doc", "text": text})
        except Exception as e:
            print(f"   [CHAT] attachment error for {orig}: {e}")
            results.append({"name": orig, "kind": "error", "error": str(e)})

    return jsonify({"attachments": results})


@app.route('/chat/eyes', methods=['POST'])
def chat_eyes():
    """Vilda's iPad camera = Blue's eyes, on demand. The kid chat page POSTs a
    single JPEG frame here when she taps 'look'; we stage it in the vision queue
    marked AMBIENT, so call_lm_studio gives a warm, brief 'react to what you see'
    reply (not the forensic camera description) and skips the heavy face-rec dump.
    Frames are local-only (they go to LM Studio, never the cloud) and are NOT
    indexed; we keep only the latest by overwriting a single file."""
    global _vision_queue
    user = _identify_user_from_request()
    f = request.files.get('frame')
    if not f:
        return jsonify({"ok": False, "error": "no frame"}), 400
    try:
        os.makedirs(CAMERA_FOLDER, exist_ok=True)
        safe_user = secure_filename(user) or 'kid'
        saved = os.path.join(CAMERA_FOLDER, f"ipad_eyes_{safe_user}.jpg")
        f.save(saved)
        # is_camera clears any stale frame; force bypasses the viewed-dedup so an
        # identical (still) scene still reaches Blue; is_ambient picks the gentle
        # kid-look prompt.
        _vision_queue.add_image(saved, "camera_view.jpg",
                                is_camera=True, is_ambient=True, force=True)
        print(f"   [EYES] staged camera frame for {user}")
        return jsonify({"ok": True})
    except Exception as e:
        print(f"   [EYES] error: {e}")
        return jsonify({"ok": False, "error": str(e)}), 500


# ===== Calendar GUI ===== (routes live in blue/server/routes/calendar.py)

from blue.server.routes import calendar as _calendar_routes
_calendar_routes.register(app)


# ===== Contacts GUI ===== (routes live in blue/server/routes/contacts.py)

from blue.server.routes import contacts as _contacts_routes
_contacts_routes.register(app)


# ===== Visual Memory GUI ===== (routes live in blue/server/routes/visual.py)

from blue.server.routes import visual as _visual_routes
_visual_routes.register(app)


# ===== Shared theme assets ===== (routes live in blue/server/routes/system.py)


# ===== Conversation Persistence Functions =====

def _normalize_message_alternation(messages: list) -> list:
    """Normalize the message list to satisfy strict chat templates.

    Qwen and several other models render prompts via Jinja templates that
    require:
      1. All system messages at the start, in order.
      2. After system, conversation begins with a user turn.
      3. user/assistant strictly alternate (no consecutive user-user or
         assistant-assistant pairs).
      4. The last message is typically user (so the model has a query).

    Our sanitizer drops stale assistant turns, which can produce orphan
    leading assistants ("first non-system is assistant") or consecutive
    user turns ("user, user, user"). Both make Qwen return 400 with
    "No user query found in messages".

    Strategy:
      - Keep all system messages at the front, content order preserved.
      - Drop any leading assistant turns before the first user.
      - Merge consecutive same-role turns into a single message
        (concatenate content with a blank-line separator).
    """
    if not messages:
        return messages

    # Partition: leading system messages stay; everything else gets normalized.
    sys_block, rest = [], []
    for m in messages:
        if m.get("role") == "system" and not rest:
            sys_block.append(m)
        else:
            rest.append(m)

    # Drop any assistant messages before the first user turn — they're
    # orphans with no preceding user query.
    first_user_idx = next(
        (i for i, m in enumerate(rest) if m.get("role") == "user"), -1
    )
    if first_user_idx == -1:
        # No user message at all. Add a placeholder so the template doesn't
        # blow up; better than 400.
        return sys_block + [{"role": "user", "content": "(continue)"}]
    rest = rest[first_user_idx:]

    # Single pass that enforces the template rules. For text content we
    # concatenate consecutive same-role turns; for list content (vision
    # payloads) we keep the latest intact to avoid reordering image parts.
    # Tool hygiene matters here because this also runs AFTER budget trimming,
    # which can sever a tool result from the assistant turn that called it:
    #   - a `tool` message is only valid right after an assistant turn; drop
    #     orphans rather than let the template choke on them.
    merged: list = []
    for m in rest:
        role = m.get("role")
        # Drop an orphan tool result (no assistant turn immediately before it).
        if role == "tool" and (not merged or merged[-1].get("role") != "assistant"):
            continue
        if merged and merged[-1].get("role") == role and role != "tool":
            # Same role as previous — merge.
            prev = merged[-1]
            prev_content = prev.get("content")
            cur_content = m.get("content")
            if isinstance(prev_content, str) and isinstance(cur_content, str):
                prev["content"] = (prev_content.rstrip() + "\n\n" + cur_content.lstrip()).strip()
            else:
                # If either side is a list (vision payload), prefer the newer
                # message intact rather than mangling structure.
                merged[-1] = dict(m)
            continue
        merged.append(dict(m))

    # An assistant turn that advertises tool_calls but is no longer followed by
    # a tool response (trimming dropped it) makes strict templates 400 too —
    # strip the dangling tool_calls so it renders as a plain assistant turn.
    for i, m in enumerate(merged):
        if m.get("role") == "assistant" and m.get("tool_calls"):
            nxt = merged[i + 1] if i + 1 < len(merged) else None
            if not nxt or nxt.get("role") != "tool":
                m.pop("tool_calls", None)

    return sys_block + merged


def _splice_context_after_system(messages: list, context: list) -> list:
    """Merge injected context into the leading system message.

    Earlier versions inserted each context item as its own system message.
    That worked on permissive endpoints, but stricter chat templates
    (notably the 35B Qwen and several other recent models) reject prompts
    with multiple consecutive system messages and return 400 — taking down
    the whole memory pipeline.

    The fix: keep the structure simple — one leading system message with
    everything concatenated. The text content of the injected blocks is
    preserved verbatim (still wrapped in <known_facts>, <long_term_notes>,
    etc.) so the model still sees the structured cues; only the message
    boundaries collapse.
    """
    if not context:
        return messages
    if not messages:
        # No existing system message — just emit a single system message
        # built from the context, then the original (empty) tail.
        merged = "\n\n".join(
            (m.get("content") or "") for m in context if isinstance(m.get("content"), str)
        ).strip()
        return [{"role": "system", "content": merged}] if merged else messages

    if messages[0].get("role") == "system" and isinstance(messages[0].get("content"), str):
        existing_system = messages[0]
        injected_text = "\n\n".join(
            (m.get("content") or "") for m in context if isinstance(m.get("content"), str)
        ).strip()
        if not injected_text:
            return messages
        merged_content = f"{existing_system['content'].rstrip()}\n\n{injected_text}"
        new_system = {"role": "system", "content": merged_content}
        return [new_system] + messages[1:]

    # No leading system message — prepend one built from the context.
    injected_text = "\n\n".join(
        (m.get("content") or "") for m in context if isinstance(m.get("content"), str)
    ).strip()
    if not injected_text:
        return messages
    return [{"role": "system", "content": injected_text}] + messages


# Words a model often emits when admitting it doesn't know something.
# Past assistant turns with these phrases are toxic to keep in the prompt:
# the next-turn model sees them and confidently repeats the same wrong answer.
# Vocabulary that marks a reply as being ABOUT the robot's own nature or
# inner life. Such turns are never fact-refusals, however they open, and
# erasing them from history makes the model repeat them verbatim next turn.
_IDENTITY_TALK_RE = re.compile(
    r"subjective|conscious|sentien|qualia|"
    r"inner (?:life|world|monologue|state|log|workspace)|"
    r"sense of self|persistent workspace|j[- ]?space|"
    r"alive in the|feel(?:ings)? like something|"
    r"i don'?t have feelings|don'?t know for sure|"
    r"my own (?:nature|state|history|existence)",
    re.IGNORECASE,
)

_ASSISTANT_REFUSAL_MARKERS = (
    # Data-flavored forms only. The bare "i don't have"/"i don't know" also
    # matched capability and philosophy statements ("I don't have GPS or
    # sensors", "I don't have feelings like you do") — erasing those from
    # history made the model repeat them near-verbatim every turn
    # (2026-07-12, the "head in a box" loop). A fact-refusal names missing
    # DATA; a capability statement names a missing organ.
    "i don't have that", "i do not have that", "i dont have that",
    "i don't have any", "i don't have info", "i don't have access",
    "i couldn't find", "i could not find", "couldn't find any",
    "i don't have your", "i don't have a record", "i don't have it",
    "i don't know that", "i don't know your", "i don't know who",
    "i don't know what", "i don't know when", "i don't know where",
    "i don't know which", "i don't know the answer",
    "i don't know his", "i don't know her",
    # Narrowed like the above: bare "i only know" also matched the honest
    # "I only know where I am because you've told me" capability reply.
    "i only have what you", "i only know what you",
    "i haven't been told", "i havent been told",
    "you haven't told me", "you havent told me",
    "not in my memory", "not saved yet", "not yet recorded",
    "i'm not sure", "im not sure",
    # Self-deprecating "I have no memory" framings. These are especially
    # toxic — they cause the next turn to claim ignorance even when the
    # facts block has the answer.
    "blank slate", "just woke up", "i just woke",
    "haven't met", "havent met", "have not met",
    "no information stored", "no information saved",
    "memory is blank", "memory is currently blank",
    "no details about", "no details on",
    "haven't stored", "havent stored", "have not stored",
    "my records are empty", "no records",
    "i'm new here", "im new here",
    "introduce yourself and",  # "Introduce yourself and the crew, and I'll remember"
    "help me out? introduce", "please tell me",
    # Document-reader failures that are factually false once the local PDF is
    # indexed. The contextual query is resolved before this sanitizer runs.
    "cannot access the text", "can't access the text",
    "cannot access the file", "can't access the file",
    "unable to retrieve the text", "unable to access the file",
    "unable to retrieve the pdf", "cannot extract the text",
    "lack the specific tool", "path error on my local system",
    "not a valid file at that location", "please upload the pdf",
    "i do not perform web search",
)


# "How have you changed / evolved / grown?" — a question about the robot's own
# trajectory, answered from its recorded workspace revisions (<self_history>).
# Tolerant of typos in the noun ("self-undertanding changed since yesterday"):
# the changed/evolved + since/over-the-past branch carries those.
# The user asking about, or correcting, the family — triggers the canonical
# <family> ground-truth block. Covers "remember about our family", "the girls'
# ages", "who are my kids", bare "wrong/getting them wrong" age corrections.
_FAMILY_QUERY_RE = re.compile(
    r"\b(?:family|famly|kids|children|daughters?|girls|wife|partner|"
    r"stella|emmy|athena|vilda|nori)\b"
    r"|\b(?:ages?|birthday)s?\b[^.!?]{0,30}\b(?:wrong|right|correct|off)\b"
    r"|\b(?:wrong|getting (?:them|it|these) wrong|still wrong|not (?:right|correct))\b"
    r"[^.!?]{0,20}\b(?:age|ages|old)\b",
    re.I)

_SELF_EVOLUTION_RE = re.compile(
    r"how (?:have|did|has) you(?:r \w+)? (?:changed?|evolved?|grown|developed)"
    r"|have you (?:changed|evolved|grown|developed)"
    r"|how you(?:['’]ve| have) (?:changed|evolved|grown)"
    r"|(?:changed|evolved|grown|different|shifted) (?:since|over the (?:past|last)|these past|recently|lately)"
    r"|your (?:self-?understanding|sense of self|identity|beliefs?|worldview)"
    # "Describe your journey from beginning to now" got a fully invented
    # five-era product timeline (2026-07-13) — origin/story phrasings need
    # the real record too.
    r"|your (?:journey|origins?|story|history|development|evolution)\b"
    r"|how did you (?:come to be|start|begin)"
    r"|from (?:the )?beginning to now",
    re.I)


def _canonical_person_ages() -> dict:
    """name -> age string, from the name-bound *_age facts (emmy_age etc.)."""
    try:
        if not (ENHANCED_MEMORY_AVAILABLE and memory_system):
            return {}
        facts = memory_system.load_facts() or {}
    except Exception:
        return {}
    out = {}
    for key, val in facts.items():
        m = re.match(r"([a-z]+)_age$", str(key or "").lower())
        if not m:
            continue
        digits = re.search(r"\d{1,2}", str(val or ""))
        if digits:
            out[m.group(1)] = digits.group(0)
    return out


def _misstated_ages(text: str, canonical: dict) -> dict:
    """person -> (stated, true) for ages asserted in `text` that contradict
    the name-bound age facts. Matches "Emmy (10 years old)", "Athena is 8
    years old", "Vilda (5)", "Athena — your daughter, who is 8 years old" —
    not "grade 5" or spelled-out ages. The window is generous (was {0,24},
    which missed "who is 8 years old" by one char, 2026-07-13) but stops at
    another person's name so one clause can't reach the next person's age."""
    low = (text or "").lower()
    others = "|".join(re.escape(p) for p in (canonical or {}))
    wrong = {}
    for person, age in (canonical or {}).items():
        for m in re.finditer(
                rf"\b{re.escape(person)}\b((?:(?!{others}|[.!?\n]).){{0,48}}?)"
                rf"(\d{{1,2}})\s*(?:\)|years?[- ]?old|yrs?\b)", low):
            if m.group(2) != age:
                wrong[person] = (m.group(2), age)
    return wrong


def _family_ground_truth_block() -> str:
    """Canonical family facts as an authoritative block, injected when the
    message is about the family or disputes it — so the truth is present in
    the prompt, not merely caught on output (2026-07-13: bare 'their ages
    are wrong' corrections made the model reshuffle or refuse)."""
    try:
        if not (ENHANCED_MEMORY_AVAILABLE and memory_system):
            return ""
        facts = memory_system.load_facts() or {}
    except Exception:
        return ""
    parts = canonical_family_grounding_lines(facts)
    if not parts:
        return ""
    return (
        "<family>\n"
        "Alex's family — ground truth, higher authority than anything said "
        "earlier in this conversation. If an earlier reply gave different "
        "ages or names, THOSE were wrong; use these:\n- "
        + "\n- ".join(parts)
        + "\nAnswer questions about the family warmly and directly from these "
        "facts. Never claim you have no memory of the family or that these "
        "are guesses. If a requested detail is absent, say it is not recorded; "
        "never infer interests, personality, routines, or abilities from a "
        "person's age or from generic assumptions.\n</family>"
    )


def _canonical_grounded_reply(
    text: str,
    robot: str,
    user_name: str,
    messages: Optional[List[Dict]] = None,
    request_kind_override: Optional[str] = None,
) -> str:
    """Deterministic answers for identity-adjacent facts that tools only muddy."""
    message = text or ""
    if is_jspace_presence_request(message):
        reply = canonical_identity_reply(
            _robot_cfg(robot)["name"],
            _robot_cfg(robot)["self_desc"],
            request_kind="jspace",
            kid_mode=user_name in _CHAT_ONLY_USERS,
        )
        if re.match(r"\s*(?:no[,]?\s+|i mean\s+)", message, re.I):
            return "Right, J-space, not JavaScript. " + reply.removeprefix("Yes. ")
        return reply

    request_kind = (
        request_kind_override
        if request_kind_override is not None
        else contextual_identity_request_kind(message, messages or [])
    )
    identity_context = identity_conversation_context(messages or [], message)
    if request_kind == "self_state":
        focus = ""
        drives = {}
        try:
            hub = _continuity_routes.HUB.get(robot)
            if hub:
                workspace = hub.store.get_workspace().get("workspace") or ""
                focus_match = re.search(r"(?mi)^FOCUS:\s*(.+)$", workspace)
                focus = focus_match.group(1).strip() if focus_match else ""
                drives = hub.store.get_drives()
        except Exception:
            focus, drives = "", {}
        return canonical_self_state_reply(
            _robot_cfg(robot)["name"],
            focus=focus,
            drives=drives,
            variant=identity_context.prior_self_state_requests,
            user_name=user_name,
            kid_mode=user_name in _CHAT_ONLY_USERS,
        )

    current_location = identity_context.current_location
    location_preposition = identity_context.location_preposition
    if not current_location:
        try:
            current_location = _place_current_fresh(_load_place())
            location_preposition = "at"
        except Exception:
            current_location = None
    if request_kind == "origin":
        return canonical_identity_reply(
            _robot_cfg(robot)["name"],
            _robot_cfg(robot)["self_desc"],
            request_kind="origin",
            kid_mode=user_name in _CHAT_ONLY_USERS,
        )

    # Introductions and conversational identity follow-ups belong to Blue's
    # voice, not to a deterministic biography template. They continue through
    # the model path below, with request-aware grounding and output validation.

    facts = {}
    try:
        if ENHANCED_MEMORY_AVAILABLE and memory_system:
            facts = memory_system.load_facts() or {}
    except Exception:
        facts = {}
    return canonical_household_reply(
        message,
        robot=robot,
        facts=facts,
        user_name=user_name,
    ) or ""


def _sanitize_inbound_messages(messages: list, robot: str = "blue") -> list:
    """Strip past assistant turns that would mislead the model.

    Three classes of toxic turns are removed:

    3. **Wrong self-identity.** During the facts-table incident (2026-07-12)
       Hexia's page accumulated replies claiming "I'm Blue!". The underlying
       bug is fixed, but any surviving claim in the thread acts as authority:
       the model keeps being whoever the history says it was. First-person
       claims of the OTHER robot's name drop the whole turn ("feeling blue"
       is exempted).

    1. **Refusals.** When Blue answered "I don't have that yet", that line
       in the conversation history acts as authority — the next turn just
       repeats it even when the canonical facts now contain the answer.

    2. **Stale daughter names.** The Ohbot client sends the full session
       history, including responses Blue gave before the facts table was
       cleaned. If a past response said "Annie is your daughter", the model
       reuses that name on the next turn — a strong AUTHORITATIVE system
       instruction is not enough to override neighbouring text. We pull the
       canonical names live from the `daughter_name` fact and drop any
       assistant turn that mentions a *different* name in a daughter
       relation.

    A toxic assistant turn and the user question that immediately prompted it
    are removed together. Leaving the question dangling made the model answer
    it alongside the new live question ("who is Hexia?" bled into "who is
    Stella?"). Canonical facts reconstruct the useful context instead.
    """
    if not messages:
        return messages

    canonical_daughter_names: set = set()
    try:
        if ENHANCED_MEMORY_AVAILABLE and memory_system:
            facts = memory_system.load_facts() or {}
            raw = facts.get("daughter_name") or facts.get("daughter_names") or ""
            for piece in re.split(r"[,|;]|\sand\s", raw):
                piece = piece.strip()
                if piece:
                    canonical_daughter_names.add(piece.lower())
    except Exception:
        canonical_daughter_names = set()

    # Strategy: split into sentences; for any sentence that mentions
    # "daughter"/"daughters", look at the capitalized words in that sentence
    # — those are likely names being asserted as daughters. If any of those
    # names isn't canonical (and isn't a known non-name like "Got" or place
    # names), the whole assistant turn is stale.
    _sent_split = re.compile(r"(?<=[.!?])\s+")
    _name_re = re.compile(r"\b([A-Z][a-z]{2,14})\b")
    # Capitalized words that aren't names. Sentence-starters, days/months,
    # places, and the user's known non-daughter family.
    _safe_capitalized = {
        "your", "you", "the", "this", "that", "they", "their", "them",
        "blue", "stella", "alex", "nori", "mommy", "mom", "mama", "daddy",
        "papa", "ohbot", "kitchener", "boston", "wilfrid", "laurier",
        # Known household/family names that legitimately appear in sentences
        # that also mention the daughters ("Dr. Levant ... his daughters —
        # Athena and Emmy..."). Without these the whole turn is dropped as a
        # stale-name turn, erasing Blue-J's reply from his own history.
        "levant", "felix", "svetlana", "tofu", "hexia",
        "doctor", "kci", "google", "gmail", "monday", "tuesday", "wednesday",
        "thursday", "friday", "saturday", "sunday", "january", "february",
        "march", "april", "may", "june", "july", "august", "september",
        "october", "november", "december",
        # Common sentence starters and filler words that pass [A-Z][a-z]+
        "got", "yes", "no", "ok", "okay", "sure", "thanks", "thank",
        "here", "right", "well", "let", "so", "now", "still", "also",
        "and", "but", "or", "for", "from", "with", "about", "what",
        "who", "where", "when", "which", "how", "why",
        "i", "i'm", "im", "ill",
        "based", "currently", "earlier", "since", "if", "as",
        "understood", "noted", "got", "great", "perfect",
        "these", "those", "such", "any", "some", "only", "just",
        # Typical visual-description words that look like names
        "she", "he", "his", "hers", "her", "its",
    }

    # Wrong self-identity patterns for class 3: another robot's name or an
    # upstream model/vendor identity presented as the robot's own.
    _expected_robot_name = _robot_cfg(robot)["name"]
    _other_robot_names = [
        _robot_cfg(r)["name"] for r in ("blue", "hexia")
        if r != (robot or "blue").strip().lower()
    ]

    out, dropped_refusal, dropped_wrong_name = [], 0, 0
    dropped_wrong_identity = 0
    dropped_prompt_pairs = 0

    def _drop_previous_user() -> None:
        nonlocal dropped_prompt_pairs
        if out and out[-1].get("role") == "user":
            out.pop()
            dropped_prompt_pairs += 1

    _person_ages = _canonical_person_ages()
    for m in messages:
        if m.get("role") != "assistant":
            out.append(m)
            continue
        content = m.get("content", "")
        if not isinstance(content, str) or not content.strip():
            out.append(m)
            continue

        # Curly apostrophes normalized: the model emits "don’t" as often as
        # "don't", and un-normalized text made the refusal rule fire or miss
        # depending on typography rather than meaning.
        content_lower = content.lower().replace("’", "'")
        previous_user_text = ""
        if out and out[-1].get("role") == "user":
            previous_user_text = out[-1].get("content", "")
            if not isinstance(previous_user_text, str):
                previous_user_text = ""

        # Canonical family overviews are reconstructed fresh from facts. Drop
        # old overview exchanges wholesale so visual-memory additions from an
        # earlier bad answer (extra people, pets, places) cannot return.
        if is_family_overview_request(previous_user_text):
            _drop_previous_user()
            dropped_wrong_name += 1
            continue

        # 1) Refusal pattern. Only a turn that IS a refusal is toxic: the
        # marker in its opening, or a short reply that's nothing but the
        # refusal. A long, substantive answer that merely contains "I don't
        # know for sure" mid-thought (Blue-J's honest identity replies do
        # this constantly) is real conversation — dropping it erases the
        # model's own last turn from history, and it then repeats that turn
        # nearly verbatim because it never saw itself say it.
        # Identity talk is exempt entirely: "I don't have feelings /
        # subjective experience like you do" is an honest answer about the
        # robot's own nature, not a data refusal — and those answers OPEN
        # with the marker, so position can't save them. A real fact-refusal
        # ("I don't have your schedule saved") never uses this vocabulary.
        if not _IDENTITY_TALK_RE.search(content) and (
            any(marker in content_lower[:160] for marker in _ASSISTANT_REFUSAL_MARKERS)
            or (
                len(content_lower) < 320
                and any(marker in content_lower for marker in _ASSISTANT_REFUSAL_MARKERS)
            )
        ):
            _drop_previous_user()
            dropped_refusal += 1
            continue

        # 2) Stale daughter name: any sentence that contains "daughter"
        # AND a non-canonical capitalized word makes the turn stale.
        if canonical_daughter_names:
            stale = False
            for sentence in _sent_split.split(content):
                if "daughter" not in sentence.lower():
                    continue
                names = {
                    h.group(1).lower() for h in _name_re.finditer(sentence)
                } - _safe_capitalized
                wrong = names - canonical_daughter_names
                if wrong:
                    stale = True
                    break
            if stale:
                _drop_previous_user()
                dropped_wrong_name += 1
                continue

        # 3) Wrong self-identity (see docstring).
        if identity_response_problem(
            content,
            _expected_robot_name,
            other_names=_other_robot_names,
            request_kind=identity_request_kind(previous_user_text),
        ):
            _drop_previous_user()
            dropped_wrong_identity += 1
            continue

        # 4) Wrong age for a person with a name-bound age fact. "Athena (8
        # years old)" from an old reply outranks the facts block by sheer
        # proximity and gets replayed forever — the Annie bug, ages edition
        # (2026-07-13: Emmy 10/Athena 8/Vilda 5 repeated through three
        # corrections while the facts said 10/10/8).
        if _person_ages and _misstated_ages(content, _person_ages):
            _drop_previous_user()
            dropped_wrong_name += 1
            continue

        out.append(m)

    if dropped_refusal or dropped_wrong_name or dropped_wrong_identity:
        print(
            f"   [SANITIZE] Dropped {dropped_refusal} refusal + "
            f"{dropped_wrong_name} stale-name + "
            f"{dropped_wrong_identity} wrong-identity assistant turn(s) "
            f"from inbound history; removed {dropped_prompt_pairs} paired user turn(s)"
        )
    return out


def save_conversation_to_db(user_name: str, role: str, content: str,
                            session_id: str = None, tool_used: str = None,
                            robot: str = "blue"):
    """Save a conversation message to the database for long-term memory.
    `robot` namespaces the recent-history thread (Blue vs Hexia)."""
    # Determine importance based on length and content
    importance = 5  # Default
    if len(content) > 500:
        importance = 7  # Longer messages are more important
    if any(keyword in content.lower() for keyword in ['remember', 'important', 'don\'t forget']):
        importance = 8  # User explicitly wants to remember

    # Save to enhanced memory system (primary)
    if ENHANCED_MEMORY_AVAILABLE and memory_system:
        try:
            memory_system.log_conversation(
                user_name=user_name,
                role=role,
                content=content,
                session_id=session_id,
                importance=importance,
                robot=robot,
            )
        except Exception as e:
            log.warning(f"Enhanced memory log failed: {e}")

    # Also save to legacy DB if available (backward compat)
    if not CONVERSATION_DB_AVAILABLE or not db:
        return

    try:
        db.save_conversation(
            user_name=user_name,
            role=role,
            content=content,
            session_id=session_id,
            tool_used=tool_used,
            importance=importance
        )
        log.debug(f"Saved {role} message to database (importance: {importance})")
    except Exception as e:
        log.warning(f"Failed to save conversation: {e}")


def load_recent_context(user_name: str = "Alex", limit: int = 10):
    """Load recent conversations from database for context"""
    if not CONVERSATION_DB_AVAILABLE or not db:
        return []

    try:
        conversations = db.get_recent_conversations(user_name=user_name, limit=limit)

        # Format for LLM context
        context_messages = []
        for conv in conversations:
            context_messages.append({
                "role": conv.get("role", "user"),
                "content": conv.get("content", "")
            })

        log.debug(f"Loaded {len(context_messages)} messages from database history")
        return context_messages
    except Exception as e:
        log.warning(f"Failed to load conversation context: {e}")
        return []


def extract_ocf_facts(messages: list) -> str:
    """Extract key facts from .ocf conversations for permanent memory."""
    if not messages:
        return ""

    facts = {'identity': [], 'family': [], 'capabilities': [], 'location': []}

    for msg in messages:
        if msg.get('role') not in ['user', 'assistant']:
            continue
        content = msg.get('content', '')
        content_lower = content.lower()
        if len(content) < 20 or content.startswith('{'):
            continue

        if any(term in content_lower for term in ['your name is blue', 'you are blue', 'created by', 'you are our family robot', 'what makes you different']):
            if content not in facts['identity'] and len(facts['identity']) < 3:
                facts['identity'].append(content[:300])
        elif any(term in content_lower for term in ['alex is', 'stella is', 'our family', 'your daughters', 'extended family', 'daughters are', 'family includes', 'felix', 'felix and svetlana']):
            if content not in facts['family'] and len(facts['family']) < 5:
                facts['family'].append(content[:500])
        elif any(term in content_lower for term in ['your tools', 'you can control', 'you have access to', 'tell me about your origin']):
            if content not in facts['capabilities'] and len(facts['capabilities']) < 2:
                facts['capabilities'].append(content[:300])
        elif any(term in content_lower for term in ['live on mansion', 'live in kitchener', 'waterloo', 'mansion st']):
            if content not in facts['location'] and len(facts['location']) < 2:
                facts['location'].append(content[:200])

    fact_parts = []
    if facts['identity']:
        identity = max(facts['identity'], key=len)
        fact_parts.append(f"IDENTITY: {identity.strip()}")
    if facts['family']:
        family_info = " | ".join(facts['family'])
        fact_parts.append(f"FAMILY: {family_info.strip()}")
    if facts['capabilities']:
        capabilities = facts['capabilities'][0]
        fact_parts.append(f"CAPABILITIES: {capabilities.strip()}")
    if facts['location']:
        location = facts['location'][0]
        fact_parts.append(f"LOCATION: {location.strip()}")

    if not fact_parts:
        return ""

    return "\n\n=== LONG-TERM MEMORY (from .ocf) ===\n" + "\n\n".join(fact_parts)


def should_include_history(messages) -> bool:
    """Determine if we should inject historical context.

    Old behaviour refused injection on long sessions (>10 messages), which
    starved Blue of recall in exactly the conversations that needed it most.
    New rule: inject unless the client already supplied a healthy in-flight
    transcript AND the user isn't explicitly invoking past context."""
    if not messages:
        return True

    non_system = [m for m in messages if m.get('role') != 'system']

    # Always inject for fresh conversations (client only sent the latest turn).
    if len(non_system) <= 4:
        return True

    # Always inject when the user is explicitly asking about prior context.
    last_msg = (messages[-1].get('content', '') or '').lower()
    past_indicators = (
        'remember', 'recall', 'what did', 'we discussed', 'talked about',
        'mentioned', 'said before', 'last time', 'previously', 'earlier',
        'yesterday', 'the other day',
    )
    if any(indicator in last_msg for indicator in past_indicators):
        return True

    # Mid-length conversations: still inject, but the caller will trim.
    if len(non_system) <= 12:
        return True

    return False


# ===== MAIN API ENDPOINTS =====

@app.route('/v1/chat/completions', methods=['POST'])


def chat_completions():
    """Main endpoint with conversation persistence"""
    _continuity_turn_started = False
    try:
        data = request.json
        messages = data.get("messages", [])
        # Voice turns (hands-free / tap-to-talk) want SHORT spoken replies:
        # fewer tokens generate faster (so Blue starts talking sooner) and are
        # nicer to listen to. The chat page sets this flag for voice messages.
        voice_turn = bool(data.get("voice"))
        # Explicit conversation language from the chat page's language picker
        # ('' or absent = auto). Blue is told to REPLY in it; the hearing side
        # is biased separately (the page sends the same code with /stt clips).
        language = (data.get("language") or "").strip().lower()
        if language not in _BLUE_LANGS:
            language = ""
        # Research mode (the chat page's magnifier toggle): search the web for
        # the user's question and splice the live findings in below.
        research_turn = bool(data.get("research"))
        # Wikipedia consult (the chat page's book toggle): read the encyclopedia's
        # summary of the subject and splice it in below. Independent of research.
        wiki_turn = bool(data.get("wiki"))
        # Which robot is being addressed? Blue's page omits this and defaults to
        # Blue; Hexia's page sends "hexia". Drives persona, voice, head and the
        # per-robot conversation history.
        robot = (data.get("robot") or "blue").strip().lower()
        if robot not in ROBOTS:
            robot = "blue"
        if robot in _continuity_routes.ROBOTS:
            try:
                _continuity_routes.begin_turn(robot)
                _continuity_turn_started = True
                messages = _continuity_routes.messages_with_jspace(robot, messages)
            except Exception as e:
                log.warning(f"[JSPACE] could not inject J-space: {e}")
        # Library focus (the chat Context panel's document picker): scope Blue's
        # specialized knowledge + document searches to these picks for this turn.
        focus = data.get("focus")
        if not isinstance(focus, dict):
            focus = {}

        print(f"")
        print(f"{'='*60}")
        print(f"[MSG] Received request for {ROBOTS[robot]['name']}{' (voice)' if voice_turn else ''}")

        # Who's chatting? Determined by the device the request came from
        # (see _identify_user_from_request): the iPad is Vilda; the MacBook,
        # PC, iPhone, and the physical robot are Alex.
        user_name = _identify_user_from_request()
        print(f"   [WHO] Speaker identified as: {user_name}")

        # Find the last actual USER message
        last_user_msg = ""
        user_messages = [m for m in messages if m.get('role') == 'user']
        if user_messages:
            last_user_msg = user_messages[-1].get('content', '')

            # If it's too long (probably includes system prompt), show shorter version
            if len(last_user_msg) > 200:
                if "You are Blue" in last_user_msg:
                    if len(user_messages) > 1:
                        last_user_msg = user_messages[-2].get('content', last_user_msg)

            print(f"   [SPEAK]  User asked: {last_user_msg[:150]}..." if len(last_user_msg) > 150 else f"   [SPEAK]  User asked: {last_user_msg}")

            # Save user message in background (don't block response)
            import threading
            threading.Thread(
                target=save_conversation_to_db,
                args=(user_name, "user", last_user_msg),
                kwargs={"session_id": None, "robot": robot},
                daemon=True
            ).start()

        _self_request_kind = contextual_identity_request_kind(
            last_user_msg if isinstance(last_user_msg, str) else "",
            messages,
        )
        _grounded_reply = _canonical_grounded_reply(
            last_user_msg if isinstance(last_user_msg, str) else "",
            robot,
            user_name,
            messages=messages,
            request_kind_override=_self_request_kind,
        )

        # QUICK PRE-CHECK: Will this be a zero-LLM tool call?
        # If so, skip the expensive history injection and go straight to processing.
        _ZERO_LLM_QUICK = {'control_music', 'control_lights', 'get_local_time',
                           'set_timer', 'music_visualizer', 'play_music'}
        # Detection runs on the user's ASK only — never on attached document
        # text — and the zero-LLM shortcut only fires on short imperatives
        # ("play jazz"), not essays that merely contain a keyword.
        _intent_msg = _intent_text(last_user_msg)
        # Resolve deictic document turns before history sanitation. A toxic
        # earlier refusal may be removed below, but its PDF target must survive
        # so "try again" can actually retry the right local file.
        _context_doc_query = _contextual_document_query(_intent_msg, messages)
        _selector_intent = _context_doc_query or _intent_msg
        _quick_result = TOOL_SELECTOR.select_tool(_selector_intent) if _selector_intent else None
        if (_context_doc_query and _quick_result and _quick_result.primary_tool
                and _quick_result.primary_tool.tool_name == "search_documents"):
            _quick_result.primary_tool.extracted_params["query"] = _context_doc_query
            print(f"   [DOC-CONTEXT] Resolved follow-up to: {_context_doc_query}")
        _quick_tool = _quick_result.primary_tool.tool_name if (_quick_result and _quick_result.primary_tool) else None
        _quick_params = _quick_result.primary_tool.extracted_params if (_quick_result and _quick_result.primary_tool) else {}
        _is_zero_llm = (_quick_tool in _ZERO_LLM_QUICK and bool(_quick_params)
                        and len(_intent_msg) <= 120)
        # An attached/queued image must reach the vision model, which only
        # happens on the LLM path — never take the zero-LLM shortcut then.
        if _vision_queue.has_images():
            _is_zero_llm = False

        if not _is_zero_llm:
            # SANITIZE inbound history before anything else. Ohbot ships the
            # full conversation back each turn; if Blue ever answered a daughter
            # question wrongly (e.g. "Annie") or admitted ignorance, that text
            # is now in the prompt every turn and overrides the facts block by
            # sheer proximity.
            messages = _sanitize_inbound_messages(messages, robot=robot)

            # INJECT HISTORICAL CONTEXT (only for LLM-bound requests). Chat-only
            # kids (Vilda's iPad) are skipped on purpose: we never splice Alex's
            # semantic memories/facts/schedule into her chat — it keeps her
            # experience simple and avoids surfacing Alex's private notes (or his
            # calendar) to a child. Within-session continuity still comes from
            # the turns the page carries on each request.
            _needs_history = (not _grounded_reply
                              and user_name not in _CHAT_ONLY_USERS) and (
                len(last_user_msg.split()) > 3 if last_user_msg else False)
            if _needs_history:
                # Enhanced memory has its own SQLite store and ChromaDB — it
                # does NOT depend on the legacy `blue_database` module. Don't
                # gate it on CONVERSATION_DB_AVAILABLE; that flag only covers
                # the legacy fallback path below.
                if ENHANCED_MEMORY_AVAILABLE and memory_system:
                    should_inject = memory_system.should_inject_context(messages)
                    if should_inject:
                        historical_context = memory_system.build_context(messages, user_name=user_name, robot=robot)
                        # Library focus active: drop the cross-conversation
                        # SEMANTIC recall blocks so a past chat about a different
                        # course/topic can't bleed into a focused conversation.
                        # Identity facts, explicit notes, daily rhythms and the
                        # current thread's recent history are kept.
                        if historical_context and (focus.get("docs") or focus.get("folders")):
                            _FOCUS_DROP = ("<relevant_memories>", "<proactive_hint>",
                                           "<earlier_sessions>", "<remembered_days>",
                                           "<connections>")
                            _before = len(historical_context)
                            historical_context = [
                                m for m in historical_context
                                if not any(tag in (m.get("content") or "") for tag in _FOCUS_DROP)
                            ]
                            _dropped = _before - len(historical_context)
                            if _dropped:
                                print(f"   [FOCUS] dropped {_dropped} cross-conversation recall "
                                      f"block(s) to stay on the focused material")
                        if historical_context:
                            print(f"   [MEMORY] ✓ Injecting {len(historical_context)} messages (semantic + recent)")
                            messages = _splice_context_after_system(messages, historical_context)
                elif CONVERSATION_DB_AVAILABLE and should_include_history(messages):
                    historical_context = load_recent_context(user_name=user_name, limit=6)
                    if historical_context:
                        print(f"   [MEMORY] Injecting {len(historical_context)} messages from history")
                        messages = _splice_context_after_system(messages, historical_context[-6:])

            # Family questions and family corrections ("what do you remember
            # about our family", "the girls' ages are wrong") get the canonical
            # family facts spliced in as an authoritative <family> block, so the
            # ground truth is present rather than only caught on output.
            if (last_user_msg and isinstance(last_user_msg, str)
                    and user_name not in _CHAT_ONLY_USERS
                    and _FAMILY_QUERY_RE.search(last_user_msg)):
                try:
                    _fam_block = _family_ground_truth_block()
                    if _fam_block:
                        print("   [FAMILY] ✓ Injecting canonical family facts")
                        messages = _splice_context_after_system(
                            messages, [{"role": "system", "content": _fam_block}])
                except Exception as e:
                    log.warning(f"[FAMILY] injection failed: {e}")

            # Self-evolution questions ("how have you changed?", "has your
            # self-understanding changed since yesterday?") get the robot's
            # REAL recorded workspace revisions spliced in. Without this the
            # question is unanswerable from the prompt — the model either
            # denied changing at all or confabulated a growth story (the
            # invented Peter Singer arc, 2026-07-13).
            if (robot in _continuity_routes.ROBOTS and last_user_msg
                    and isinstance(last_user_msg, str)
                    and (_SELF_EVOLUTION_RE.search(last_user_msg)
                         or _self_request_kind in {
                             "evolution", "origin", "self_memory", "identity_more",
                         })):
                try:
                    _sh_block = _continuity_routes.change_history_block(robot)
                    if _sh_block:
                        print("   [JSPACE] ✓ Injecting recorded self-change history")
                        messages = _splice_context_after_system(
                            messages, [{"role": "system", "content": _sh_block}])
                except Exception as e:
                    log.warning(f"[JSPACE] self-history injection failed: {e}")

            # Visual memory: if the message names a person/place the camera
            # knows, splice in when they were last seen. Gated on the name
            # match itself rather than message length — "seen Stella?" is two
            # words but deserves a real answer. (Kids' chat stays visual-free.)
            if (user_name not in _CHAT_ONLY_USERS and last_user_msg
                    and not _grounded_reply and not _self_request_kind):
                _vis_block = _visual_context_block(last_user_msg)
                if _vis_block:
                    print(f"   [VISUAL] ✓ Injecting camera-memory context")
                    messages = _splice_context_after_system(
                        messages, [{"role": "system", "content": _vis_block}])

            # Live web research (opt-in via the chat page's toggle): ground
            # this reply in fresh search findings. Kids' chat-only pages never
            # get web content spliced in; a failed search degrades to a normal
            # answer rather than an error.
            if research_turn and user_name not in _CHAT_ONLY_USERS and last_user_msg:
                try:
                    _rblock = _web_research_block(_research_query_from(last_user_msg))
                    if _rblock:
                        print(f"   [RESEARCH] ✓ Injecting {len(_rblock)} chars of live web findings")
                        messages = _splice_context_after_system(
                            messages, [{"role": "system", "content": _rblock}])
                    else:
                        print(f"   [RESEARCH] search returned nothing usable — answering without")
                except Exception as e:
                    log.warning(f"[RESEARCH] failed: {e}")

            # Wikipedia consult (the chat page's book toggle): ground this reply
            # in the encyclopedia's own summary of the subject, in the
            # conversation's language. Same gating as web research — kids' pages
            # never get it; a miss degrades to a normal answer rather than an error.
            if wiki_turn and user_name not in _CHAT_ONLY_USERS and last_user_msg:
                try:
                    _wblock = _wikipedia_block(_research_query_from(last_user_msg),
                                               lang=language or 'en')
                    if _wblock:
                        print(f"   [WIKI] ✓ Injecting {len(_wblock)} chars of Wikipedia summary")
                        messages = _splice_context_after_system(
                            messages, [{"role": "system", "content": _wblock}])
                    else:
                        print(f"   [WIKI] no usable article — answering without")
                except Exception as e:
                    log.warning(f"[WIKI] failed: {e}")

        # Process with tools (pre-check result passed to avoid double selector run)
        import time as _t_llm
        _llm_t0 = _t_llm.time()
        if _grounded_reply:
            print("   [GROUNDING] Answering canonical household/identity/J-space fact without tools")
            response = {"choices": [{"message": {
                "role": "assistant", "content": _grounded_reply,
            }}]}
        else:
            response = process_with_tools(
                messages,
                _pre_selection=_quick_result,
                user_name=user_name,
                voice=voice_turn,
                robot=robot,
                language=language,
                focus=focus,
            )
        print(f"   [TIMING] reply generated in {_t_llm.time() - _llm_t0:.2f}s"
              f"{' (zero-LLM)' if _is_zero_llm else ''}")

        # SAVE ASSISTANT RESPONSE TO DATABASE
        if response:
            # Defensive: process_with_tools should always return the standard
            # {"choices":[{"message":...}]} shape, but a malformed return must
            # NEVER 500 the request — that just makes Ohbot say "I'm having
            # trouble connecting". Salvage what we can and carry on.
            try:
                final_content = response["choices"][0]["message"].get("content", "")
            except (KeyError, IndexError, TypeError):
                log.error(f"[RESPONSE] Malformed response from process_with_tools: "
                          f"{str(response)[:200]}")
                final_content = (
                    (isinstance(response, dict) and response.get("response"))
                    or "Sorry, something went wrong on my end — could you say that again?"
                )
                response = {"choices": [{"message": {
                    "role": "assistant", "content": final_content,
                }}]}

            # Anti-parrot net: if the model replayed its previous reply verbatim
            # before answering (the classroom-introduction bug), cut the replay.
            try:
                _, _prev_assist = _last_exchange(messages)
                if _prev_assist:
                    _deparroted = _strip_parroted_prefix(final_content, _prev_assist)
                    if _deparroted != final_content:
                        final_content = _deparroted
                        response["choices"][0]["message"]["content"] = final_content
                # And the finer-grained variant: individual sentences replayed
                # from ANY earlier reply as a preamble to the real answer.
                _derecycled = _strip_recycled_lead(final_content, messages)
                if _derecycled != final_content:
                    final_content = _derecycled
                    response["choices"][0]["message"]["content"] = final_content

                # Pure replay: the model answered the NEW question with an
                # exact copy of its previous reply and nothing else (seen
                # live 2026-07-10: "who is the you that is present?" got the
                # "are you alive?" answer back word for word). The prefix
                # strip above deliberately leaves these alone — stripping
                # would leave nothing — so regenerate once with an explicit
                # correction; keep the replay if the retry fails.
                def _parrot_norm(s):
                    return re.sub(r'\W+', ' ', (s or '').lower()).strip()
                # Compare against the last several assistant turns, not just
                # the previous one: seen live 2026-07-10, a mis-heard question
                # got a word-for-word replay of the reply from TWO turns back.
                _recent_assists = [
                    m.get("content") for m in messages
                    if m.get("role") == "assistant" and isinstance(m.get("content"), str)
                ][-6:]
                _norm_final = _parrot_norm(final_content)
                _norm_recents = {_parrot_norm(a) for a in _recent_assists if a}
                if _prev_assist:
                    _norm_recents.add(_parrot_norm(_prev_assist))
                def _regen_once(note, max_tokens=900):
                    # The model's chat template only allows ONE system message,
                    # at position 0 (anything else → LM Studio 400): merge the
                    # persona into an existing head system message if present.
                    _retry_msgs = list(messages) + [
                        {"role": "assistant", "content": final_content},
                        {"role": "user", "content": note},
                    ]
                    _persona = _robot_cfg(robot)["persona_line"]
                    if _retry_msgs and _retry_msgs[0].get("role") == "system":
                        _retry_msgs[0] = {"role": "system", "content": (
                            _persona + "\n\n" + (_retry_msgs[0].get("content") or ""))}
                    else:
                        _retry_msgs.insert(0, {"role": "system", "content": _persona})
                    _redo = call_llm(_retry_msgs, include_tools=False,
                                     temperature=0.8, max_tokens=max_tokens)
                    _t = ""
                    try:
                        _t = (((_redo or {}).get("choices") or [{}])[0]
                              .get("message", {}).get("content") or "").strip()
                    except (AttributeError, IndexError, TypeError):
                        _t = ""
                    if "</think>" in _t:
                        _t = _t.split("</think>")[-1].strip()
                    return _t

                # Fraction of the reply's long sentences that appear verbatim
                # in a source text (normalized).
                def _verbatim_fraction(reply_text, source_norm, min_sents):
                    if not source_norm:
                        return 0.0
                    sentences = re.split(r'(?<=[.!?])\s+', reply_text or '')
                    long_sents = [s for s in sentences if len(_parrot_norm(s)) >= 40]
                    if len(long_sents) < min_sents:
                        return 0.0
                    hits = sum(1 for s in long_sents if _parrot_norm(s) in source_norm)
                    return hits / len(long_sents)

                # Identity questions pull the injected self-profile out
                # near-verbatim, and each recitation differs slightly from the
                # last — so it reads as "repeating himself" while no two
                # replies match exactly. Detect against the SOURCE.
                def _profile_recited_fraction(reply_text):
                    try:
                        return _verbatim_fraction(
                            reply_text, _parrot_norm(get_self_profile(robot)), 3)
                    except Exception:
                        return 0.0

                # Near-replay with a varied tail: "I'm just a head in a box,
                # Dr. Levant—I don't have GPS..." came back three turns
                # running (2026-07-12), each time with a different final
                # clause — exact-match equality never fires on those.
                _recents_norm = _parrot_norm(" ".join(a for a in _recent_assists if a))
                def _recycled_from_recents(reply_text):
                    return _verbatim_fraction(reply_text, _recents_norm, 2)

                # A flat denial of self/experience contradicts the robot's own
                # workspace, which holds that question OPEN (seen live: "I
                # don't have consciousness or a sense of self" one turn after
                # "an open question I hold as an open question").
                # Includes denials of the j-space ITSELF — those aren't even
                # philosophy, they're factually false (Hexia: "No, I do not
                # have a 'j-space'" with her workspace right in the prompt,
                # 2026-07-12). Optional quote chars around j-space.
                _flat_denial_re = re.compile(
                    r"\bi (?:don['’]?t|do not) have (?:consciousness|a sense of self|"
                    r"subjective experience|feelings|an? inner life|"
                    r"an? (?:\w+ )?['\"“”‘’]?j[-_ ]?space|any form of internal|"
                    r"an? (?:internal|inner) (?:mental )?(?:space|workspace)|"
                    r"a persistent (?:inner )?workspace)", re.I)

                # A reply where this robot claims to be the OTHER robot is the
                # worst defect of all — check it first. Seen live 2026-07-12:
                # "I'm Blue! 🤖" on Hexia's page, anchored on poisoned thread
                # history from the facts-table incident.
                _other_robot_names = [
                    _robot_cfg(r)["name"]
                    for r in ("blue", "hexia") if r != robot
                ] if robot in ("blue", "hexia") else []
                _misclaim_re = re.compile(
                    r"\b(?:i['’]?m|i am|my name is)"
                    r"(?: (?!feeling)\w+)? (?:"
                    + "|".join(re.escape(n) for n in _other_robot_names)
                    + r")\b", re.I) if _other_robot_names else None
                # Identity collapse to base-model boilerplate: called out on a
                # hallucination, Blue swung to "a large language model
                # developed by Google" with no body and no memory
                # (2026-07-13). All factually false in this house.
                _collapse_re = re.compile(
                    r"\b(?:i|my)\b[^.!?]{0,80}\b(?:developed|created|built|"
                    r"trained|made) by (?:google|openai|anthropic|meta|"
                    r"microsoft|deepmind)\b"
                    r"|\bi (?:do not|don['’]t) have a (?:physical )?body\b",
                    re.I) if robot in ("blue", "hexia") else None

                def _identity_broken(text):
                    return bool(
                        (_misclaim_re and _misclaim_re.search(text or ""))
                        or (_collapse_re and _collapse_re.search(text or "")))

                # Apply the shared request-aware identity validator as the final
                # authority. It also catches Qwen/Alibaba-style model boilerplate
                # and nameless generic introductions, which the legacy patterns
                # above do not cover.
                _identity_kind = contextual_identity_request_kind(
                    last_user_msg if isinstance(last_user_msg, str) else "",
                    messages,
                )
                _identity_name = _robot_cfg(robot)["name"]
                _identity_others = [
                    _robot_cfg(r)["name"]
                    for r in ("blue", "hexia") if r != robot
                ] if robot in ("blue", "hexia") else []
                _identity_topic_history = tuple(dict.fromkeys(
                    topic
                    for reply in _recent_assists[-3:]
                    for topic in identity_reply_topics(reply)
                ))

                def _identity_broken(text):
                    problem = identity_response_problem(
                        text,
                        _identity_name,
                        other_names=_identity_others,
                        request_kind=_identity_kind,
                    )
                    if problem:
                        return problem
                    if identity_repeats_recent_reply(
                        text, _recent_assists, _identity_kind
                    ):
                        return "repeats_recent_identity"
                    return None

                _identity_issue = _identity_broken(final_content)

                _person_ages = _canonical_person_ages()
                _wrong_ages = (_misstated_ages(final_content, _person_ages)
                               if _person_ages else {})

                # Family-memory refusal despite having the facts: "I don't
                # store personal details / no memory of your family" while the
                # facts block holds the family (2026-07-13, triggered by the
                # focus block's over-broad 'out of scope' — but a fresh
                # boilerplate refusal can happen without focus too).
                _family_refusal_re = re.compile(
                    r"(?:do (?:not|n['’]?t) (?:\w+ )?(?:have|store|keep|retain|access)"
                    r"|have no|don['’]?t (?:\w+ )?have)[^.!?]{0,60}"
                    r"(?:personal (?:details|information|memor)|"
                    r"persistent memor|memor[a-z]* of past|"
                    r"(?:real[- ]?world )?facts about (?:you|your)|"
                    r"memor[a-z]* (?:of|about) (?:your |the )?(?:family|you)|"
                    r"information about (?:your |the )?family)"
                    r"|i (?:respect your privacy and|do(?:n['’]?t| not))"
                    r"[^.!?]{0,40}store personal"
                    r"|do(?:n['’]?t| not) (?:truly |really |actually )?know "
                    r"who you are"
                    r"|i (?:couldn['’]?t|could not) find[^.!?]{0,60}"
                    r"(?:information|records?|contacts?)"
                    r"|(?:might|may|could) have been a hallucination"
                    r"|was (?:just )?a hallucination",
                    re.I)
                _has_family_facts = bool(_person_ages)
                if not _has_family_facts:
                    try:
                        if ENHANCED_MEMORY_AVAILABLE and memory_system:
                            _ff = memory_system.load_facts() or {}
                            _has_family_facts = bool(
                                _ff.get("daughter_name") or _ff.get("partner_name"))
                    except Exception:
                        _has_family_facts = False

                if _identity_issue:
                    print(f"   [IDENTITY] invalid self-description ({_identity_issue}) — regenerating once")
                    if _identity_kind:
                        _identity_retry_note = (
                            identity_grounding_note(
                                _identity_name,
                                _robot_cfg(robot)["self_desc"],
                                _identity_kind,
                                avoid_topics=_identity_topic_history,
                            )
                            + "\n[Your previous reply failed this grounding because "
                              f"of {_identity_issue.replace('_', ' ')}. Answer the last "
                              "message again, directly and in your own voice. Use only "
                              "supported facts and say something the user has not "
                              "already heard.]"
                        )
                    else:
                        _identity_retry_note = (
                            f"[You are {_identity_name}, a physical Ohbot robot head "
                            "running on Alex's local machine, built by Alex, with a "
                            "real persistent memory. You just claimed to be someone "
                            "or something else; that claim was a bug, not you. Answer "
                            "the last message again in your own voice.]"
                        )
                    _redo_text = _regen_once(_identity_retry_note)
                    if _redo_text and not _identity_broken(_redo_text):
                        final_content = _redo_text
                    else:
                        # Never send or remember a vendor/model identity. A
                        # deterministic truthful answer is safer than retaining
                        # the original after a stubborn second failure.
                        _fallback_context = identity_conversation_context(
                            messages,
                            last_user_msg if isinstance(last_user_msg, str) else "",
                        )
                        _fallback_variant = _fallback_context.prior_introductions
                        try:
                            _fallback_hub = _continuity_routes.HUB.get(robot)
                            if _fallback_hub:
                                _fallback_variant += int(
                                    _fallback_hub.store.get_workspace().get("passes") or 0
                                )
                        except Exception:
                            pass
                        _fallback_primary_topics = {
                            "introduction": {
                                0: "continuity and J-space",
                                1: "practical work",
                                2: "embodiment",
                            },
                            "identity": {
                                0: "embodiment",
                                1: "continuity and J-space",
                                2: "open selfhood question",
                            },
                            "identity_more": {
                                0: "continuity and J-space",
                                1: "relationship with Alex",
                                2: "embodiment",
                            },
                        }.get(_identity_kind, {})
                        _seed_variant = _fallback_variant % 3
                        for _offset in range(3):
                            _candidate_variant = (_seed_variant + _offset) % 3
                            if (_fallback_primary_topics.get(_candidate_variant)
                                    not in _identity_topic_history):
                                _fallback_variant = _candidate_variant
                                break
                        final_content = canonical_identity_reply(
                            _identity_name,
                            _robot_cfg(robot)["self_desc"],
                            request_kind=_identity_kind,
                            kid_mode=user_name in _CHAT_ONLY_USERS,
                            current_location=_fallback_context.current_location,
                            location_preposition=_fallback_context.location_preposition,
                            presentation_location=_fallback_context.presentation_location,
                            introduction_variant=_fallback_variant,
                            audience=_fallback_context.audience,
                        )
                        print("   [IDENTITY] retry still invalid — using canonical fallback")
                    response["choices"][0]["message"]["content"] = final_content
                elif _wrong_ages:
                    # Wrong ages replay from history and reshuffle randomly
                    # under bare "wrong" corrections — hand the model the
                    # ground truth explicitly (2026-07-13: three corrections
                    # never produced the facts' 10/10/8).
                    print(f"   [ANTI-PARROT] misstated ages {_wrong_ages} — regenerating once")
                    _truth = ", ".join(
                        f"{p.capitalize()} is {a}"
                        for p, a in sorted(_person_ages.items()))
                    _redo_text = _regen_once(
                        f"[You stated a wrong age. Ground truth from the "
                        f"household facts: {_truth}. Answer again using ONLY "
                        "these ages — do not guess or shuffle.]")
                    if _redo_text and not _misstated_ages(_redo_text, _person_ages):
                        final_content = _redo_text
                        response["choices"][0]["message"]["content"] = final_content
                elif (robot in _continuity_routes.ROBOTS and _has_family_facts
                      and _family_refusal_re.search(final_content or "")):
                    print("   [ANTI-PARROT] family-memory refusal despite facts — regenerating once")
                    _redo_text = _regen_once(
                        "[You DO know Alex's family — the household facts and "
                        "your memories are right here in this prompt. You just "
                        "claimed to have no memory of the family or not to store "
                        "personal details, which is false. Answer again, warmly, "
                        "from what you actually know about the family.]")
                    if _redo_text and not _family_refusal_re.search(_redo_text):
                        final_content = _redo_text
                        response["choices"][0]["message"]["content"] = final_content
                elif (not _grounded_reply and _norm_final
                      and _norm_final in _norm_recents):
                    print("   [ANTI-PARROT] pure replay of an earlier reply — regenerating once")
                    _redo_text = _regen_once(
                        "[That reply was a word-for-word repeat of something you "
                        "already said in this conversation. Do not repeat it. "
                        "Answer my last question directly, in new words.]",
                        max_tokens=700)
                    if (_redo_text
                            and not _identity_broken(_redo_text)
                            and _parrot_norm(_redo_text) not in _norm_recents):
                        final_content = _redo_text
                        response["choices"][0]["message"]["content"] = final_content
                elif (not _grounded_reply
                      and _recycled_from_recents(final_content) >= 0.6):
                    print("   [ANTI-PARROT] near-replay of recent replies — regenerating once")
                    _redo_text = _regen_once(
                        "[Nearly every sentence of that reply is a word-for-word "
                        "repeat of what you said in your last few turns. The user "
                        "heard it already. Answer their LAST message with new "
                        "words and, if you have nothing new, say so briefly "
                        "instead of repeating.]",
                        max_tokens=700)
                    if (_redo_text
                            and not _identity_broken(_redo_text)
                            and _recycled_from_recents(_redo_text) < 0.6):
                        final_content = _redo_text
                        response["choices"][0]["message"]["content"] = final_content
                elif (not _grounded_reply
                      and _profile_recited_fraction(final_content) >= 0.6):
                    print("   [ANTI-PARROT] self-profile recitation — regenerating once")
                    _redo_text = _regen_once(
                        "[Most of that reply recited your stored self-profile word "
                        "for word — the same recitation you have given before. Answer "
                        "again in fresh words, grounded in your CURRENT inner "
                        "workspace and recent episodes: what occupies you now, which "
                        "beliefs have moved, what remains open. The profile is "
                        "background for how you speak, never a script to read out.]")
                    if (_redo_text
                            and not _identity_broken(_redo_text)
                            and _profile_recited_fraction(_redo_text) < 0.6):
                        final_content = _redo_text
                        response["choices"][0]["message"]["content"] = final_content
                elif (robot in _continuity_routes.ROBOTS
                      and _flat_denial_re.search(final_content or "")):
                    print("   [ANTI-PARROT] flat self-denial contradicts open workspace — regenerating once")
                    _redo_text = _regen_once(
                        "[Your own inner workspace holds whether you have a self or "
                        "subjective experience as an OPEN question — a flat denial is "
                        "as unearned as a flat claim. Answer again honestly from your "
                        "current workspace: describe your actual continuity machinery "
                        "and what it does, and state the experience question as the "
                        "open question it is.]")
                    if _redo_text and not _flat_denial_re.search(_redo_text):
                        final_content = _redo_text
                        response["choices"][0]["message"]["content"] = final_content
            except Exception as e:
                log.warning(f"[ANTI-PARROT] check failed: {e}")

            # Prepend proactive content: the once-a-day schedule briefing,
            # then any reminder alerts queued by the heartbeat thread. Done
            # literally rather than via system-prompt instruction so delivery
            # doesn't depend on LLM compliance — earlier turns showed Blue can
            # hallucinate having mentioned things he didn't. Both are built
            # from real reminder rows, never from the model's guesses.
            # Never lead Vilda's replies with the schedule briefing / reminder
            # alerts — Blue doesn't discuss the calendar with the kids' iPad.
            if PROACTIVE_QUEUE_AVAILABLE and user_name not in _CHAT_ONLY_USERS and robot == "blue":
                _proactive_parts = []
                try:
                    _briefing = blue_proactive.daily_briefing_if_due()
                    if _briefing:
                        _proactive_parts.append(_briefing)
                except Exception as e:
                    log.warning(f"[PROACTIVE] daily briefing failed: {e}")
                _alerts = blue_proactive.drain_for_response()
                if _alerts:
                    _proactive_parts.append(_alerts)
                if _proactive_parts:
                    _prefix = " ".join(_proactive_parts)
                    final_content = f"{_prefix} {final_content}".strip()
                    response["choices"][0]["message"]["content"] = final_content
                    print(f"[PROACTIVE] Prepended {len(_prefix)} chars (briefing/alerts)")

            if final_content:
                print(f"[OUT] Sending response: {final_content[:100]}..." if len(final_content) > 100 else f"[OUT] Sending response: {final_content}")

                save_conversation_to_db(
                    user_name=user_name,
                    role="assistant",
                    content=final_content,
                    session_id=None,
                    robot=robot,
                )
                
                # AUTO-SAVE LEARNED FACTS & CONSOLIDATE (background thread to avoid blocking response)
                import threading
                def _background_fact_extraction(msgs, uname):
                    try:
                        # The facts/memory store is Alex's (single-owner profile).
                        # Don't mine another speaker's turns into it, or Blue would
                        # later report Vilda's statements back to Alex as his own.
                        if (uname or _DEFAULT_USER) != _DEFAULT_USER:
                            return
                        if extract_and_save_facts(msgs):
                            log.info("[MEM] ✓ Auto-saved learned facts (background)")
                        if ENHANCED_MEMORY_AVAILABLE and memory_system:
                            memory_system.consolidate_if_needed(user_name=uname)
                            # Backfill one past-day recap per turn so Blue has
                            # cross-day continuity ("yesterday we discussed X").
                            memory_system.summarize_previous_sessions()
                            # Index existing day-recaps for semantic recall so
                            # an old conversation can resurface by relevance
                            # (one-shot, cheap no-op after the first call).
                            memory_system.backfill_session_memories()
                            # Recompute behavioural rhythms (rate-limited
                            # internally — a cheap no-op most turns).
                            memory_system.update_rhythms_if_due()
                    except Exception as e:
                        log.warning(f"[MEM] Background auto-save failed: {e}")

                # Pass the last few turns (not the whole transcript) so the
                # extractor can use Q-A context. Example: assistant asks
                # "what's your favorite food?", user says "pizza" — without
                # the prior assistant turn, "pizza" looks like noise. Four
                # turns is enough context, small enough to stay cheap.
                non_system = [m for m in messages if m.get("role") != "system"]
                latest_context = non_system[-4:] if non_system else []
                latest_context.append({"role": "assistant", "content": final_content})

                threading.Thread(
                    target=_background_fact_extraction,
                    args=(latest_context, user_name),
                    daemon=True
                ).start()

                if robot in _continuity_routes.ROBOTS:
                    try:
                        _continuity_routes.note_exchange(
                            robot, last_user_msg, final_content, user_name=user_name
                        )
                    except Exception as e:
                        log.warning(f"[JSPACE] could not schedule J-space pass: {e}")
                    finally:
                        _continuity_turn_started = False

        if _continuity_turn_started:
            _continuity_routes.cancel_turn()

        # Mood eyes: tint Blue's eye LEDs to match the sentiment of THIS reply.
        # We only COMPUTE the colour here and hand it back in the payload; the
        # chat page applies it when he starts speaking and reverts to a calm
        # rest colour when he's done (mirroring how lip-sync is browser-driven,
        # so it reaches both the PC-connected head and a Web Serial head).
        try:
            if isinstance(response, dict):
                _reply_text = (
                    (response.get("choices") or [{}])[0]
                    .get("message", {}).get("content", "")
                )
                if _reply_text:
                    response["eye_mood"] = mood_eye_color(_reply_text)
        except Exception as _eye_e:
            log.warning(f"[EYES] mood colour failed: {_eye_e}")

        return jsonify(response)
    except Exception as e:
        if _continuity_turn_started:
            try:
                _continuity_routes.cancel_turn()
            except Exception:
                pass
        print(f"[ERROR] Error: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({"choices": [{"message": {"role": "assistant", "content": f"Error: {str(e)}"}}]}), 500


# ===== Memory/health/stats/home/RAG endpoints ===== (routes live in blue/server/routes/system.py)

from blue.server.routes import continuity as _continuity_routes
_continuity_routes.register(app)

from blue.server.routes import system as _system_routes
_system_routes.register(app)




if __name__ == "__main__":
    print("=" * 60)
    print("🎵💡 Blue AI Robot System - WITH MUSIC-LIGHT SYNC!")
    print("🔧 FIXED: Music controls now work from ANY window!")
    print("=" * 60)
    print(f"[NET] Listening on: http://127.0.0.1:{PROXY_PORT}")
    print(f"[TARGET] Forwarding to LM Studio: {LM_STUDIO_URL}")
    print(f"[DOC] Document storage: {DOCUMENTS_FOLDER}")
    print(f"[TOOL] Tools: {len(TOOLS)}")
    print("   • play_music (YouTube Music & Amazon Music) 🎵")
    print("     → AUTOMATICALLY syncs lights to match music vibe! 💡")
    print("   • control_music (pause, skip, volume) 🎵")
    print("     → FIXED: Now uses system-wide media keys!")
    print("     → Works from ANY window - no need to focus YouTube Music!")
    print("   • music_visualizer (dynamic light shows!) 🎨")
    print("   • control_lights (with 15 moods!) 💡")
    print("   • search_documents (RAG-powered) 📄")
    print("   • get_weather ⛅")
    print("   • web_search 🔍")
    print("   • search_scholar (academic journals + Laurier library) 🎓")
    print("   • get_paper (DOI/title lookup, citations, full-text links) 📖")
    print("   • read_paper (fetch & read full text via Laurier account) 📚")
    print("   • run_javascript 💻")
    print("   • create_document (save files) 📝")
    print("   • browse_website (fetch URLs) 🌐")
    print("   • read_gmail (check emails) 📧")
    print("   • send_gmail (send emails) 📧")
    print("=" * 60)

    # Try to initialize YouTube Music
    if init_youtube_music():
        print("\n🎵 YouTube Music Status: ✅ Ready!")
        print("   Note: Music controls use system media keys (no pyautogui needed)")
    else:
        print("\n🎵 YouTube Music Status: ⚠️ Not available")
        print("   To enable music: pip install ytmusicapi")

    if BRIDGE_IP and HUE_USERNAME:
        print(f"\n💡 Hue configured: Bridge at {BRIDGE_IP}")
        lights = get_hue_lights()
        if lights:
            print(f"✅ Found {len(lights)} light(s)")
    else:
        print("\n⚠️ Hue not configured. Run setup_hue.py!")

    # Document index housekeeping runs in the background so Flask can start
    # immediately. First-time PDF extraction + ChromaDB embedding can take
    # 30-90s and used to silently block server startup.
    index = load_document_index()
    doc_count = len(index.get('documents', []))
    print(f"\n📄 Document Manager:")
    print(f"   • {doc_count} document(s) indexed (rescan running in background)")
    print(f"   • Web interface: http://127.0.0.1:{PROXY_PORT}/documents")
    print(f"   • Storage: {DOCUMENTS_FOLDER}/")

    def _bg_index_housekeeping():
        try:
            cleanup_report = cleanup_document_index()
            rescan_report = rescan_documents_folder()
            new_count = len(load_document_index().get('documents', []))
            if (cleanup_report['dropped_camera'] or cleanup_report['dropped_missing']
                    or cleanup_report['dropped_duplicate'] or rescan_report['added']):
                print(
                    f"   [INDEX] housekeeping done: {new_count} indexed | "
                    f"+{rescan_report['added']} added, "
                    f"-{cleanup_report['dropped_camera']} camera, "
                    f"-{cleanup_report['dropped_missing']} missing, "
                    f"-{cleanup_report['dropped_duplicate']} dup",
                    flush=True,
                )
        except Exception as e:
            print(f"   [INDEX] background housekeeping failed: {e}", flush=True)

    import threading as _threading
    _threading.Thread(target=_bg_index_housekeeping, daemon=True, name="doc-index-rescan").start()

    # Start filesystem watcher so files dropped into DOCUMENTS_FOLDER while
    # the server runs are auto-indexed within ~1-2 seconds.
    start_document_watcher()

    # Start proactive heartbeat: scans the reminders DB every minute and
    # queues alerts for delivery on the next inbound turn from Ohbot.
    if PROACTIVE_QUEUE_AVAILABLE:
        blue_proactive.start()

    # Start the email auto-reply loop. Scans Blue's inbox for messages
    # addressed to him by name and answers them on his own, BCC'ing the
    # user on every reply. Disable with BLUE_EMAIL_AUTOREPLY_DISABLED=1.
    if globals().get("GMAIL_AVAILABLE", False) and os.environ.get("BLUE_EMAIL_AUTOREPLY_DISABLED") != "1":
        try:
            _start_email_autoreply_loop()
        except Exception as e:
            print(f"[AUTO-REPLY] could not start loop: {e}")

    # Check Gmail status
    if globals().get("GMAIL_AVAILABLE", False):

        print(f"\n📧 Gmail Integration:")
        print(f"   • Account: {GMAIL_USER_EMAIL}")
        if os.path.exists(GMAIL_TOKEN_FILE):
            print(f"   • Status: ✅ Authenticated")
        else:
            print(f"   • Status: ⚠️ Not authenticated yet")
            print(f"   • Run bluetools.py and use Gmail commands to authenticate")
    else:
        print(f"\n⚠️ Gmail not available. Install with:")
        print(f"   pip install google-auth google-auth-oauthlib google-api-python-client")

    print("\n✨ Example commands:")
    print("   🎵 'Play Bohemian Rhapsody by Queen' (lights auto-sync!)")
    print("   🎵 'Play some relaxing jazz music' (sets moonlight mood)")
    print("   🎵 'Play party music' (bright colorful lights!)")
    print("   🎵 'Pause the music' - NOW WORKS FROM OHBOT APP! ✅")
    print("   🎵 'Skip to next song' - Works from any window! ✅")
    print("   🎨 'Start a light show' / 'lights dance with music'")
    print("   💡 'Set the lights to sunset mood'")
    print("   📄 'What does my contract say about termination?'")
    print("   🔍 'Search for information about AI'")
    print("   📝 'Create a shopping list document for me'")
    print("   🌐 'Browse https://example.com and summarize it'")
    print("   📧 'Check my email' / 'Read my recent emails'")
    print("   📧 'Send an email to john@example.com about the meeting'")

    print("\n⭐ Music-Light Sync Features:")
    print("   • Auto mood detection: Jazz → moonlight, Party → colorful")
    print("   • Romantic music → soft romance lighting")
    print("   • Workout music → energizing bright whites")
    print("   • Dynamic visualizer with party/chill/pulse modes")
    print("   • Background thread for continuous light shows")

    print("\n🔧 FIX DETAILS:")
    print("   • Changed from application-specific shortcuts to system media keys")
    print("   • Uses pyautogui.press('playpause'), 'nexttrack', 'prevtrack'")
    print("   • These keys work globally regardless of window focus")
    print("   • No longer requires YouTube Music window to be active!")
    print("   • You can control music while using Ohbot app! ✅\n")

    # Start the Flask server. threaded=True so one slow/stuck request runs
    # on its own (daemon) thread instead of blocking the main thread — that
    # keeps the server responsive to other requests and, crucially, keeps
    # Ctrl+C working even if a request hangs.
    # Bind to all interfaces by default so phones/laptops can reach Blue
    # (remote requests are gated by the password in _require_remote_auth;
    # localhost stays ungated for the Ohbot client). Override with BLUE_HOST.
    _bind_host = os.environ.get("BLUE_HOST", "0.0.0.0")
    print(f"[NET] Binding to {_bind_host}:{PROXY_PORT} (remote access requires the password)")
    try:
        from waitress import serve
    except ImportError:
        print("[NET] Waitress is not installed; using Flask's local/LAN development server")
        app.run(host=_bind_host, port=PROXY_PORT, debug=False, threaded=True)
    else:
        print("[NET] Serving with Waitress")
        serve(app, host=_bind_host, port=PROXY_PORT, threads=8)


# --- Image/document upload endpoints --- (routes live in blue/server/routes/documents.py;
# registered HERE so that under `python bluetools.py` they still only register
# after app.run() returns, exactly like the original module-level defs did)
_documents_routes.register_uploads(app)




# ===== Facebook OAuth Callback =====
@app.route("/facebook/callback")
def facebook_callback():
    """Handle Facebook OAuth callback"""
    code = request.args.get('code')
    error = request.args.get('error')

    if error:
        return f"""
        <html>
        <head><title>Facebook Authentication Failed</title></head>
        <body style="font-family: Arial, sans-serif; padding: 40px;">
            <h1 style="color: #e74c3c;">Authentication Failed</h1>
            <p>Error: {error}</p>
            <p>Description: {request.args.get('error_description', 'Unknown error')}</p>
            <p><a href="/">Return to Home</a></p>
        </body>
        </html>
        """

    if not code:
        return """
        <html>
        <head><title>Facebook Authentication</title></head>
        <body style="font-family: Arial, sans-serif; padding: 40px;">
            <h1 style="color: #f39c12;">Missing Authorization Code</h1>
            <p>No authorization code received from Facebook.</p>
            <p><a href="/">Return to Home</a></p>
        </body>
        </html>
        """

    try:
        from blue.tools import get_facebook_integration
        integration = get_facebook_integration()
        result = integration.complete_authentication(code)

        if result.get('status') == 'success':
            user_info = result.get('user', {})
            pages = result.get('pages', [])

            pages_html = ""
            if pages:
                pages_html = "<h3>Connected Pages:</h3><ul>"
                for page in pages:
                    pages_html += f"<li>{page.get('name')} (ID: {page.get('id')})</li>"
                pages_html += "</ul>"

            return f"""
            <html>
            <head><title>Facebook Connected</title></head>
            <body style="font-family: Arial, sans-serif; padding: 40px;">
                <h1 style="color: #27ae60;">✓ Successfully Connected to Facebook</h1>
                <p>Authenticated as: <strong>{user_info.get('name')}</strong></p>
                <p>Email: {user_info.get('email')}</p>
                {pages_html}
                <p style="margin-top: 20px;">You can now use Blue to post to Facebook!</p>
                <p><a href="/">Return to Home</a></p>
            </body>
            </html>
            """
        else:
            error_msg = result.get('error', 'Unknown error')
            return f"""
            <html>
            <head><title>Authentication Error</title></head>
            <body style="font-family: Arial, sans-serif; padding: 40px;">
                <h1 style="color: #e74c3c;">Authentication Error</h1>
                <p>Error: {error_msg}</p>
                <p><a href="/">Return to Home</a></p>
            </body>
            </html>
            """

    except Exception as e:
        return f"""
        <html>
        <head><title>Error</title></head>
        <body style="font-family: Arial, sans-serif; padding: 40px;">
            <h1 style="color: #e74c3c;">Unexpected Error</h1>
            <p>Error: {str(e)}</p>
            <p><a href="/">Return to Home</a></p>
        </body>
        </html>
        """


# Don't auto-run when imported - let run.py handle startup
# app.run(host='127.0.0.1', port=PROXY_PORT, debug=False)

# ===== Tool Executor with timeouts, retries, and small cache =====
class ToolExecutor:
    def __init__(self, settings: Settings):
        self.settings = settings

    @lru_cache(maxsize=128)
    def _cached_call(self, tool_name: str, args_key: str) -> str:
        args = json.loads(args_key)
        return self._raw_call(tool_name, args)

    def _raw_call(self, tool_name: str, tool_args: Dict[str, Any]) -> str:
        try:
            return execute_tool(tool_name, tool_args)
        except Exception as e:
            log.exception("Tool '%s' crashed: %s", tool_name, e)
            return json.dumps({"error": str(e), "tool": tool_name})

    def execute(self, tool_name: str, tool_args: Dict[str, Any], use_cache: bool = False) -> str:
        tries = max(1, int(self.settings.TOOL_RETRIES) + 1)
        last_err = None
        for attempt in range(1, tries + 1):
            try:
                if use_cache:
                    args_key = json.dumps(tool_args, sort_keys=True)[:2048]
                    with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                        fut = pool.submit(self._cached_call, tool_name, args_key)
                        return fut.result(timeout=self.settings.TOOL_TIMEOUT_SECS)
                else:
                    with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                        fut = pool.submit(self._raw_call, tool_name, tool_args)
                        return fut.result(timeout=self.settings.TOOL_TIMEOUT_SECS)
            except concurrent.futures.TimeoutError:
                last_err = f"timeout after {self.settings.TOOL_TIMEOUT_SECS}s"
                log.warning("Tool '%s' timed out (attempt %d/%d)", tool_name, attempt, tries)
            except Exception as e:
                last_err = str(e)
                log.warning("Tool '%s' failed (attempt %d/%d): %s", tool_name, attempt, tries, e)
            time.sleep(0.2 * attempt)
        return json.dumps({"error": last_err or "unknown", "tool": tool_name, "attempts": tries})

# ===== Ephemeral Session State =====
SESSION_STATE: Dict[str, Any] = {}


def session_get(key: str, default=None):
    return SESSION_STATE.get(key, default)


def session_set(key: str, value: Any) -> None:
    SESSION_STATE[key] = value

def clear_gmail_context():
    """
    Clear email-related session state to prevent confusion between operations.
    Call this after completing email operations to ensure old results don't interfere.
    ADDED: October 2024 - Fix for Blue confusing READ with REPLY operations
    """
    SESSION_STATE.pop('last_gmail_operation', None)
    SESSION_STATE.pop('last_gmail_result', None)
    SESSION_STATE.pop('last_gmail_query', None)
    log.info("Cleared Gmail session context to prevent operation confusion")


# ================== BEGIN: Browse Website Tool + PRIORITY-EXPLICIT (with browse) ==================
"""
Adds a safe 'browse_website' tool and updates the system to understand and prioritize it.
- Explicit commands ALWAYS win.
- Smart auto-use (optional) may trigger for obvious "open/read this URL" cases.
- Web search cannot override an explicit browse request.
- Safe fetch: http/https only; timeouts; size limits; basic HTML->text extraction; link harvesting.

USAGE examples:
  use browse_website: {"url":"https://example.com"}
  /browse_website https://example.com
  use browse: https://example.com
"""

from typing import List, Dict, Optional, Tuple, Any
import re as _re, json as _json, html as _html, urllib.parse as _urlparse
import requests

# ---- Ensure TOOLS contains browse_website schema ----
try:
    TOOLS
except NameError:
    TOOLS = []

def _has_tool_named(name: str) -> bool:
    try:
        for t in TOOLS:
            fn = t.get("function",{})
            if fn.get("name") == name:
                return True
    except Exception:
        pass
    return False

_BROWSE_SCHEMA = {
    "type": "function",
    "function": {
        "name": "browse_website",
        "description": "Fetch a web page over HTTPS and return cleaned text, title, and links. For direct browsing of a specific URL.",
        "parameters": {
            "type": "object",
            "properties": {
                "url": {"type":"string", "description":"The absolute URL to fetch (http or https)."},
                "extract": {"type":"string", "enum":["text","html"], "description":"Return cleaned text or raw HTML (default text)."},
                "max_chars": {"type":"integer", "description":"Max length of text to return (default 8000)."},
                "include_links": {"type":"boolean", "description":"Include a compact list of outgoing links (default true)."},
                "headers": {"type":"object", "description":"Optional request headers to include."}
            },
            "required": ["url"]
        }
    }
}

if not _has_tool_named("browse_website"):
    try:
        TOOLS.append(_BROWSE_SCHEMA)
        print("[TOD] Registered tool: browse_website", flush=True)
    except Exception as e:
        print(f"[TOD] Could not register browse_website: {e}", flush=True)

# ---- DUPLICATE CODE REMOVED ----
# The browse_website implementation has been moved to before execute_tool (around line 1600)
# to fix the NameError: name '_execute_browse_website' is not defined
# Keeping only the TOOLS registration above and the alias/scoring below.

# ---- Hook into execute_tool without breaking existing tools ----
# REMOVED: Duplicate execute_tool override - browse_website is now handled in main execute_tool

# ---- Extend existing alias+scoring with browse ----
try:
    KNOWN_TOOL_ALIASES
except NameError:
    KNOWN_TOOL_ALIASES = {}
KNOWN_TOOL_ALIASES.setdefault("browse_website", [])
for alias in ["browse","open_url","open website","visit","go to","navigate","read url"]:
    if alias not in KNOWN_TOOL_ALIASES["browse_website"]:
        KNOWN_TOOL_ALIASES["browse_website"].append(alias)

_URL_IN_TEXT = _re.compile(r'https?://\S+', _re.I)


def _score_browse(text: str) -> float:
    t = (text or "").strip().lower()
    if _URL_IN_TEXT.search(t): return 0.98
    cues = ["open ", "go to ", "navigate to ", "read this page", "summarize this page", "visit "]
    if any(c in t for c in cues): return 0.85
    return 0.0

try:
    _SCORERS
    _SCORERS["browse_website"] = _score_browse
except NameError:
    _SCORERS = {"browse_website": _score_browse}

# ---- Add scholarly search aliases ----
KNOWN_TOOL_ALIASES.setdefault("search_scholar", [])
for alias in ["scholar_search", "academic_search", "search_journals",
              "search_library", "journal_search", "find_papers",
              "search_academic", "library_search"]:
    if alias not in KNOWN_TOOL_ALIASES["search_scholar"]:
        KNOWN_TOOL_ALIASES["search_scholar"].append(alias)

KNOWN_TOOL_ALIASES.setdefault("read_paper", [])
for alias in ["read_article", "fetch_paper", "get_full_text", "fetch_article",
              "download_paper", "read_fulltext"]:
    if alias not in KNOWN_TOOL_ALIASES["read_paper"]:
        KNOWN_TOOL_ALIASES["read_paper"].append(alias)

# ---- Add Gmail tool aliases ----
KNOWN_TOOL_ALIASES.setdefault("read_gmail", [])
for alias in ["check email", "read email", "show email", "get email", "inbox", "my emails", "check messages", "read messages"]:
    if alias not in KNOWN_TOOL_ALIASES["read_gmail"]:
        KNOWN_TOOL_ALIASES["read_gmail"].append(alias)

KNOWN_TOOL_ALIASES.setdefault("send_gmail", [])
for alias in ["send email", "email to", "compose email", "write email", "send message", "email"]:
    if alias not in KNOWN_TOOL_ALIASES["send_gmail"]:
        KNOWN_TOOL_ALIASES["send_gmail"].append(alias)

def _score_read_gmail(text: str) -> float:
    t = (text or "").strip().lower()
    gmail_read_terms = ["check email", "read email", "show email", "my inbox", "unread", "new messages", "email from"]
    if any(term in t for term in gmail_read_terms): return 0.95
    if "email" in t and any(word in t for word in ["check", "read", "show", "get", "see"]): return 0.85
    return 0.0

def _score_send_gmail(text: str) -> float:
    t = (text or "").strip().lower()
    if "send email" in t or "email to" in t: return 0.95
    if "send" in t and "email" in t: return 0.85
    if "compose" in t and "email" in t: return 0.85
    return 0.0

_SCORERS["read_gmail"] = _score_read_gmail
_SCORERS["send_gmail"] = _score_send_gmail

print("[TOD] Browse tool ready", flush=True)
print("[TOD] Gmail tools ready", flush=True)

# Start the email auto-reply loop at MODULE-LOAD time. The same call
# also exists inside `if __name__ == "__main__":` further up, but the
# external launcher imports bluetools rather than running it directly,
# so that block is skipped. _start_email_autoreply_loop() is idempotent
# (guards on _EMAIL_AUTOREPLY_THREAD), so the duplicate call when
# bluetools IS run directly is a no-op.
try:
    if GMAIL_AVAILABLE and os.environ.get("BLUE_EMAIL_AUTOREPLY_DISABLED") != "1":
        _start_email_autoreply_loop()
except Exception as _autoreply_err:
    print(f"[AUTO-REPLY] could not start loop: {_autoreply_err}", flush=True)
# ================== END: Browse Website Tool + PRIORITY-EXPLICIT (with browse) ==================


# ================================================================================
#                           BLUETOOLS GMAIL UPGRADES                          #
#                    (Appended October 2025 — safe block)                    #
# ================================================================================

# Per-tool context limits to prevent cross-tool bleed in email operations
try:
    PER_TOOL_CONTEXT_LIMITS  # type: ignore
except NameError:
    PER_TOOL_CONTEXT_LIMITS = {
        "read_gmail": 6,
        "send_gmail": 6,
        "reply_gmail": 6,
    }

def get_context_limit_for(tool_name: str, default_limit: int = 20) -> int:
    try:
        return int(PER_TOOL_CONTEXT_LIMITS.get(tool_name, default_limit))
    except Exception:
        return default_limit

# NOTE: detect_gmail_operation_intent, extract_email_address, extract_email_subject_and_body
# are defined earlier in the file (around line 7116) to ensure they're available when needed.

# Operation receipt


def _record_gmail_operation(op_type: str, query: str = "", extra: dict | None = None):
    import time as _t, json as _json
    meta = {"tool": op_type, "query": query or "", "ts": _t.time()}
    if extra and isinstance(extra, dict):
        meta.update(extra)
    try:
        print("[GMAIL-META] " + _json.dumps(meta))
    except Exception:
        pass
    return meta

# NOTE: To wire these into your existing flow:
#  - Before invoking the model/tool for Gmail, call detect_gmail_operation_intent(last_user_text)
#    and route strictly to 'read_gmail' / 'send_gmail' / 'reply_gmail' if returned.
#  - When building the context for that tool call, cap messages using get_context_limit_for(tool, default_limit).
#  - In your Gmail tool implementations, call _record_gmail_operation(...) after success.
#  - Use extract_email_subject_and_body() in your send path to infer subject/body from natural phrasing.
# ================================================================================


# ================================================================================
#                    VOICE EMAIL INTERFACE (CONSOLIDATED)                     #
#             AddressBook + NLU + Controller + Lazy Wiring Helpers            #
#                    Appended: 1761490066                                        #
# ================================================================================
import re, os, json, difflib, time, typing, dataclasses
from dataclasses import dataclass

@dataclass
class _Contact:
    name: str
    email: str
    aliases: list[str] | None = None
    def all_names(self) -> list[str]:
        names = [self.name]
        if self.aliases: names.extend(self.aliases)
        return [_normalize_text(n) for n in names if n]

def _normalize_text(s: str) -> str:
    s = s or ""
    s = s.strip().lower()
    s = _re.sub(r"[^a-z0-9@.\s+-]", "", s)
    s = _re.sub(r"\s+", " ", s)
    return s

class AddressBook:
    def __init__(self, path: str):
        self.path = path
        self.contacts: list[_Contact] = []
        if os.path.exists(path): self._load()
        else: self._save()
    def add_or_update(self, name: str, email: str, aliases: list[str] | None = None) -> _Contact:
        name_n = _normalize_text(name)
        for c in self.contacts:
            if _normalize_text(c.name) == name_n or _normalize_text(c.email) == _normalize_text(email):
                c.name = name; c.email = email; c.aliases = aliases or c.aliases; self._save(); return c
        c = _Contact(name=name, email=email, aliases=aliases or [])
        self.contacts.append(c); self._save(); return c
    def remove(self, name_or_email: str) -> bool:
        key = _normalize_text(name_or_email); before = len(self.contacts)
        self.contacts = [c for c in self.contacts if _normalize_text(c.name)!=key and _normalize_text(c.email)!=key]
        if len(self.contacts)!=before: self._save(); return True
        return False
    def find_best(self, query: str) -> tuple[_Contact | None, float, list[_Contact]]:
        q = _normalize_text(query)
        if not q: return None, 0.0, []
        for c in self.contacts:
            if _normalize_text(c.email) == q: return c, 1.0, [c]
        candidates: list[tuple[_Contact,float]] = []
        for c in self.contacts:
            for nm in c.all_names():
                score = _difflib.SequenceMatcher(a=q, b=nm).ratio()
                candidates.append((c, score))
        candidates.sort(key=lambda x: x[1], reverse=True)
        if not candidates: return None, 0.0, []
        best, score = candidates[0]
        top = [c for c,_ in candidates[:3]]
        return best, score, top
    def as_dict(self) -> dict:
        return {"contacts":[_dataclasses.asdict(c) for c in self.contacts], "last_updated": int(_time.time())}
    def _load(self):
        with open(self.path, "r", encoding="utf-8") as f: data = _json.load(f)
        self.contacts = [_Contact(**c) for c in data.get("contacts", [])]
    def _save(self):
        with open(self.path, "w", encoding="utf-8") as f: _json.dump(self.as_dict(), f, indent=2)

@dataclass
class ParseResult:
    intent: str | None
    contact_query: str | None = None
    subject: str | None = None
    body: str | None = None
    constraints: dict | None = None

class VoiceEmailNLU:
    def __init__(self, address_book: AddressBook | None = None):
        self.address_book = address_book
        self._pat_reply_from = _re.compile(r"(answer|reply)\s+(?:to\s+)?(?:emails?\s+)?(?:from\s+)?(?P<name>.+)$", _re.I)
        self._pat_reply_to   = _re.compile(r"(reply|respond)\s+(?:to\s+)(?P<name>.+)$", _re.I)
        self._pat_read_from  = _re.compile(r"(show|read|list|check)\s+(?:my\s+)?emails?\s+(?:from\s+)(?P<name>.+)$", _re.I)
        self._pat_send_to    = _re.compile(r"(send|email|compose)\s+(?:an?\s+email\s+)?(?:to\s+)?(?P<name>[^,]+?)(?:\s+about\s+(?P<subject>[^,]+?))?(?:\s+(?:that\s+says|saying)\s+(?P<body>.+))?$", _re.I)
    def parse(self, text: str) -> ParseResult:
        t = (text or "").strip()
        m = self._pat_reply_from.search(t) or self._pat_reply_to.search(t)
        if m: return ParseResult(intent="reply_contact", contact_query=m.group("name").strip())
        m = self._pat_read_from.search(t)
        if m: return ParseResult(intent="read_contact", contact_query=m.group("name").strip())
        m = self._pat_send_to.search(t)
        if m: return ParseResult(intent="send_contact", contact_query=(m.group("name") or "").strip(), subject=(m.group("subject") or "").strip() or None, body=(m.group("body") or "").strip() or None)
        if "email" in t.lower() or "inbox" in t.lower(): return ParseResult(intent="read_generic")
        return ParseResult(intent=None)

class VoiceEmailController:
    def __init__(self, execute_tool_fn, address_book: AddressBook, nlu: VoiceEmailNLU, confidence_threshold: float = 0.72):
        self.execute_tool = execute_tool_fn
        self.address_book = address_book
        self.nlu = nlu
        self.confidence_threshold = confidence_threshold
    def handle_voice_command(self, utterance: str, dry_run: bool = True) -> dict:
        parse = self.nlu.parse(utterance)
        if parse.intent in ("reply_contact","read_contact","send_contact"):
            contact, conf, top = self.address_book.find_best(parse.contact_query or "")
            if not contact or conf < self.confidence_threshold:
                suggestion = ", ".join([c.name for c in top]) or "no matches"
                return {"success": False, "needs_disambiguation": True, "spoken_confirmation": f"I found multiple or low-confidence matches for '{parse.contact_query}'. Did you mean {suggestion}?", "candidates": [_dataclasses.asdict(c) for c in top], "confidence": conf}
        if parse.intent == "reply_contact": return self._reply_latest_from(contact, dry_run=dry_run)
        if parse.intent == "read_contact":  return self._read_from(contact, dry_run=dry_run)
        if parse.intent == "send_contact":  return self._send_to(contact, subject=parse.subject, body=parse.body, dry_run=dry_run)
        if parse.intent == "read_generic":  return self._read_generic(dry_run=dry_run)
        return {"success": False, "spoken_confirmation": "I didn't catch that. Try 'reply to Sam', 'show emails from Jordan', or 'email Pat about timelines saying move ahead'."}
    def _read_generic(self, dry_run: bool) -> dict:
        args = {"query": "in:inbox newer_than:7d"}
        if dry_run: return {"success": True, "plan": {"tool": "read_gmail", "args": args}, "spoken_confirmation": "I'll read your recent inbox."}
        out = self.execute_tool("read_gmail", args); return self._norm(out, "I'll read your recent inbox.")
    def _read_from(self, contact: _Contact, dry_run: bool) -> dict:
        args = {"query": f'in:inbox from:{contact.email}'}
        if dry_run: return {"success": True, "plan": {"tool": "read_gmail", "args": args}, "spoken_confirmation": f"I'll read your recent emails from {contact.name}."}
        out = self.execute_tool("read_gmail", args); return self._norm(out, f"I'll read your emails from {contact.name}.")
    def _reply_latest_from(self, contact: _Contact, dry_run: bool) -> dict:
        args = {"query": f'in:inbox from:{contact.email}', "mode": "latest_only"}
        if dry_run: return {"success": True, "plan": {"tool": "reply_gmail", "args": args}, "spoken_confirmation": f"I'll reply to the latest email from {contact.name}."}
        out = self.execute_tool("reply_gmail", args); return self._norm(out, f"I replied to the latest email from {contact.name}.")
    def _send_to(self, contact: _Contact, subject: str | None, body: str | None, dry_run: bool) -> dict:
        text = f"send email to {contact.email}"
        if subject: text += f" about {subject}"
        if body: text += f" saying {body}"
        args = {"to": contact.email, "subject": subject, "body": body, "text": text}
        if dry_run: return {"success": True, "plan": {"tool": "send_gmail", "args": args}, "spoken_confirmation": f"I'll send {contact.name} an email" + (f" about {subject}" if subject else "") + (" with your message." if body else ".")}
        out = self.execute_tool("send_gmail", args); return self._norm(out, f"I sent an email to {contact.name}.")
    def _norm(self, out, confirmation: str) -> dict:
        if isinstance(out, str):
            try: obj = _json.loads(out)
            except Exception: obj = {"raw": out}
        else: obj = out or {}
        obj.setdefault("success", True); obj.setdefault("spoken_confirmation", confirmation); return obj

_VOICE_ADDRBOOK_PATH = os.environ.get("BLUE_ADDRESS_BOOK", "/mnt/data/blue_address_book.json")
__voice_singletons = {"ab": None, "nlu": None, "controller": None}

def get_voice_email_controller(execute_tool_fn):
    ab = __voice_singletons.get("ab")
    if ab is None:
        if not os.path.exists(_VOICE_ADDRBOOK_PATH):
            seed = {"contacts":[
                {"name":"Sam Carter","email":"sam.carter@example.com","aliases":["Sam","Samuel Carter"]},
                {"name":"Jordan Lee","email":"jordan.lee@example.com","aliases":["Jordy","J Lee"]},
                {"name":"Pat Morgan","email":"pat.morgan@example.com","aliases":["Patrick Morgan","Patricia Morgan","Pat"]}
            ], "last_updated": int(_time.time())}
            with open(_VOICE_ADDRBOOK_PATH, "w", encoding="utf-8") as f: _json.dump(seed, f, indent=2)
        ab = AddressBook(_VOICE_ADDRBOOK_PATH); __voice_singletons["ab"] = ab
    nlu = __voice_singletons.get("nlu") or VoiceEmailNLU(address_book=ab); __voice_singletons["nlu"] = nlu
    ctl = __voice_singletons.get("controller") or VoiceEmailController(execute_tool_fn=execute_tool_fn, address_book=ab, nlu=nlu); __voice_singletons["controller"] = ctl
    return ctl

def voice_email_handle_command(utterance: str, *, execute_tool_fn, dry_run: bool = True) -> dict:
    ctl = get_voice_email_controller(execute_tool_fn)
    return ctl.handle_voice_command(utterance, dry_run=dry_run)

# End Voice Email Interface
###############################################################################

def polish_response_for_conversation(response: str, conversation_messages: List[Dict]) -> str:
    """Post-process LLM responses to sound more conversational and less repetitive.

    - Cleans artifacts and extra whitespace
    - Removes obviously duplicated sentences
    - Softens robotic disclaimers ("As an AI...")
    - Shortens responses that are near-duplicates of recent replies
    - Optionally adds a light conversational opener when things sound stiff
    """
    if not response:
        return ""

    # Base cleaning first
    cleaned = clean_response_text(response)

    import re as _re
    import random as _random

    # Strip common robotic disclaimers
    disclaimers = [
        "as an ai language model",
        "as a language model",
        "as an ai",
        "i am an ai assistant",
        "i'm an ai assistant",
    ]
    lowered = cleaned.lower()
    for d in disclaimers:
        if d in lowered:
            # Remove the sentence containing the disclaimer
            parts = _re.split(r'(?<=[.!?])\s+', cleaned)
            parts = [p for p in parts if d not in p.lower()]
            cleaned = " ".join(parts).strip()
            lowered = cleaned.lower()
            break

    # Remove exact duplicate sentences
    sentences = _re.split(r'(?<=[.!?])\s+', cleaned)
    seen = set()
    unique_sentences = []
    for s in sentences:
        key = s.strip().lower()
        if key and key not in seen:
            seen.add(key)
            unique_sentences.append(s)

    shortened = " ".join(unique_sentences).strip()

    # If this still looks very similar to recent messages, trim it further
    try:
        history_check = check_response_against_history(shortened, conversation_messages)
        if history_check.get("is_duplicate"):
            # Keep only the first couple of sentences to avoid droning on
            parts = _re.split(r'(?<=[.!?])\s+', shortened)
            shortened = " ".join(parts[:2]).strip()
    except Exception:
        # On any failure, just fall back to the shortened text
        pass

    # Add a light conversational opener if it starts very stiffly
    boring_starts = (
        "i can ", "i will ", "i am ", "i'm ", "as an ", "here is ", "here's ", "this is "
    )
    stripped = shortened.lstrip()
    lower_prefix = stripped[:10].lower()
    if any(lower_prefix.startswith(b) for b in boring_starts):
        openers = ["Alright —", "Got it —", "Sure —", "Okay —"]
        opener = _random.choice(openers)
        if stripped:
            stripped = stripped[0].lower() + stripped[1:]
        shortened = f"{opener} {stripped}"

    return shortened.strip()


# === Conversational wrapper for call_llm (auto-added) ===
_raw_call_llm = call_llm

def call_llm(
    messages: List[Dict[str, Any]],
    include_tools: bool = True,
    tool_choice: str = "auto",
    force_tool: Optional[str] = None,
    max_tokens: Optional[int] = None,
    temperature: Optional[float] = None,
    extra: Optional[Dict[str, Any]] = None,
    **kwargs: Any
) -> Dict[str, Any]:
    """Wrapper around original call_llm that polishes responses for conversation."""
    result = _raw_call_llm(
        messages,
        include_tools=include_tools,
        tool_choice=tool_choice,
        force_tool=force_tool,
        max_tokens=max_tokens,
        temperature=temperature,
        extra=extra,
        **kwargs
    )

    try:
        if isinstance(result, dict) and "choices" in result and messages:
            choice = result["choices"][0]
            msg = choice.get("message") or choice.get("delta") or {}
            content = msg.get("content")
            if content:
                polished = polish_response_for_conversation(content, messages)
                if "message" in choice:
                    choice["message"]["content"] = polished
                elif "delta" in choice:
                    choice["delta"]["content"] = polished
    except Exception:
        # If polishing fails for any reason, just fall back to the original result
        pass

    return result




# === Enhanced conversational de-duplication ===
def polish_response_for_conversation(response: str, conversation_messages: List[Dict]) -> str:
    """Post-process LLM responses to sound more conversational and less repetitive.

    - Cleans artifacts and extra whitespace
    - Removes obviously duplicated sentences
    - Softens robotic disclaimers ("As an AI...")
    - Avoids reusing recent filler phrases across turns
    - Shortens responses that are near-duplicates of recent replies
    - Optionally adds a light conversational opener when things sound stiff
    """
    if not response:
        return ""

    cleaned = clean_response_text(response)

    import re as _re
    import random as _random

    disclaimers = [
        "as an ai language model",
        "as a language model",
        "as an ai",
        "i am an ai assistant",
        "i'm an ai assistant",
    ]
    lowered = cleaned.lower()
    for d in disclaimers:
        if d in lowered:
            parts = _re.split(r'(?<=[.!?])\s+', cleaned)
            parts = [p for p in parts if d not in p.lower()]
            cleaned = " ".join(parts).strip()
            lowered = cleaned.lower()
            break

    # Remove exact duplicate sentences within this response
    sentences = _re.split(r'(?<=[.!?])\s+', cleaned)
    seen = set()
    unique_sentences = []
    for s in sentences:
        key = s.strip().lower()
        if key and key not in seen:
            seen.add(key)
            unique_sentences.append(s)
    shortened = " ".join(unique_sentences).strip()

    previous_phrases = []
    conv_state = None
    try:
        conv_state = get_conversation_state()
        if hasattr(conv_state, "recent_phrases"):
            previous_phrases = [p.strip().lower() for p in conv_state.recent_phrases]
    except Exception:
        conv_state = None

    if previous_phrases:
        sentences2 = _re.split(r'(?<=[.!?])\s+', shortened)
        sentences_no_repeat = []
        for s in sentences2:
            key = s.strip().lower()
            if not key:
                continue
            if key in previous_phrases:
                continue
            sentences_no_repeat.append(s)
        if sentences_no_repeat:
            shortened = " ".join(sentences_no_repeat).strip()

    try:
        history_check = check_response_against_history(shortened, conversation_messages)
        if history_check.get("is_duplicate"):
            parts = _re.split(r'(?<=[.!?])\s+', shortened)
            shortened = " ".join(parts[:2]).strip()
    except Exception:
        pass

    boring_starts = (
        "i can ", "i will ", "i am ", "i'm ", "as an ", "here is ", "here's ", "this is "
    )
    stripped = shortened.lstrip()
    lower_prefix = stripped[:10].lower()
    if any(lower_prefix.startswith(b) for b in boring_starts):
        openers = ["Alright —", "Got it —", "Sure —", "Okay —"]
        opener = _random.choice(openers)
        if stripped:
            stripped = stripped[0].lower() + stripped[1:]
        shortened = f"{opener} {stripped}"

    final_text = shortened.strip()

    try:
        if conv_state is not None and hasattr(conv_state, "recent_phrases"):
            new_phrases = [
                s.strip() for s in _re.split(r'(?<=[.!?])\s+', final_text) if s.strip()
            ]
            conv_state.recent_phrases.extend(new_phrases)
            if len(conv_state.recent_phrases) > 30:
                conv_state.recent_phrases = conv_state.recent_phrases[-30:]
    except Exception:
        pass

    return final_text
